"""⚡ Fusion.AI v14 — FastAPI Edition
Requirements: fastapi uvicorn python-pptx openpyxl python-docx requests cryptography
"""
import subprocess, sys, os, json, uuid, base64, hashlib, sqlite3, time as _time
from datetime import datetime
from functools import wraps

# deps installed via requirements.txt
from fastapi import FastAPI, Request, Header, HTTPException
from fastapi.responses import JSONResponse, HTMLResponse, StreamingResponse
from cryptography.fernet import Fernet
import requests as req
import io, re as _re, base64 as _b64

app = FastAPI()

_raw_secret    = os.environ.get("SECRET_KEY","")
SECRET         = _raw_secret.encode() if _raw_secret else os.urandom(32)
GROQ_KEY       = os.environ.get("GROQ_KEY","")
OPENROUTER_KEY = os.environ.get("OPENROUTER_KEY","")
HF_TOKEN       = os.environ.get("HF_TOKEN","")
GITHUB_TOKEN   = os.environ.get("GITHUB_TOKEN","")
STABILITY_KEY  = os.environ.get("STABILITY_KEY","")
WORKER_URL     = os.environ.get("WORKER_URL","https://fusionai.pantathagat.workers.dev/")
WORKER_KEY     = os.environ.get("WORKER_KEY","")
CF_ACCOUNT_ID  = os.environ.get("CF_ACCOUNT_ID","")
CF_KEY         = os.environ.get("CF_KEY","")
CF_ACCOUNT_ID2 = os.environ.get("CF_ACCOUNT_ID2","")
CF_KEY2        = os.environ.get("CF_KEY2","")
TP_USERNAME    = os.environ.get("TP_USERNAME","")
DEV_PASSWORD   = os.environ.get("DEV_PASSWORD","")
GOOGLE_CLIENT_ID  = os.environ.get("GOOGLE_CLIENT_ID","")
GOOGLE_CLIENT_SECRET = os.environ.get("GOOGLE_CLIENT_SECRET","")
OPENWEATHER_KEY= os.environ.get("OPENWEATHER_KEY","")
EXTRA_API_KEY  = os.environ.get("EXTRA_API_KEY","")
EXTRA_API_URL  = os.environ.get("EXTRA_API_URL","")
HCAPTCHA_SECRET  = os.environ.get("HCAPTCHA_SECRET","")
HCAPTCHA_SITE_KEY= os.environ.get("HCAPTCHA_SITE_KEY","10000000-ffff-ffff-ffff-000000000001")
# ── Direct provider API keys ───────────────────────────────────────────────────
OPENAI_KEY     = os.environ.get("OPENAI_KEY","")        # OpenAI direct
ANTHROPIC_KEY  = os.environ.get("ANTHROPIC_KEY","")     # Claude / Anthropic
GEMINI_KEY     = os.environ.get("GEMINI_KEY","")        # Google Gemini
DEEPSEEK_KEY   = os.environ.get("DEEPSEEK_KEY","")      # DeepSeek direct
MOONSHOT_KEY   = os.environ.get("MOONSHOT_KEY","")      # Kimi / Moonshot AI
MISTRAL_KEY    = os.environ.get("MISTRAL_KEY","")       # Mistral AI direct
XAI_KEY        = os.environ.get("XAI_KEY","")           # xAI Grok
COHERE_KEY     = os.environ.get("COHERE_KEY","")        # Cohere Command R
TOGETHER_KEY   = os.environ.get("TOGETHER_KEY","")      # Together AI
PERPLEXITY_KEY = os.environ.get("PERPLEXITY_KEY","")    # Perplexity Sonar
# ── CF Worker 2 (gamesdohas — same WORKER_KEY) ────────────────────────────────
WORKER_URL2    = os.environ.get("WORKER_URL2","https://fusionai.gamesdohas.workers.dev/")
GITHUB_ENDPOINT= "https://models.github.ai/inference/chat/completions"
OVERSEER_MODEL = "llama-3.3-70b-versatile"
OVERSEER_ENHANCE_MODEL = "llama-3.1-8b-instant"  # cheapest for prompt enhancement
def _get_data_dir():
    for d in ["/data", "/tmp", "."]:
        if os.path.isdir(d):
            try:
                t=os.path.join(d,".wtest"); open(t,"w").write("1"); os.remove(t); return d
            except: continue
    return "."
_data_dir = _get_data_dir()
DB             = os.path.join(_data_dir,"fusionai.db")

def cf_ep(model,acc=None): return f"https://api.cloudflare.com/client/v4/accounts/{acc or CF_ACCOUNT_ID}/ai/run/{model}"
def cf_key_for(model,explicit_key=None):
    if explicit_key: return explicit_key
    # Round-robin: if CF2 set and model hash picks it, use CF2
    if CF_ACCOUNT_ID2.strip() and CF_KEY2.strip() and hash(model)%2==0: return CF_KEY2
    return CF_KEY
def cf_acc_for(model):
    if CF_ACCOUNT_ID2.strip() and CF_KEY2.strip() and hash(model)%2==0: return CF_ACCOUNT_ID2
    return CF_ACCOUNT_ID

def _fernet(salt):
    key = base64.urlsafe_b64encode(hashlib.sha256(SECRET + salt.encode()).digest())
    return Fernet(key)
def encrypt(t,s): return _fernet(s).encrypt(t.encode()).decode()
def decrypt(t,s): return _fernet(s).decrypt(t.encode()).decode()

def db():
    c = sqlite3.connect(DB, timeout=5, check_same_thread=False)
    c.row_factory = sqlite3.Row
    c.execute("PRAGMA journal_mode=WAL")
    c.execute("PRAGMA synchronous=NORMAL")
    return c

def init_db():
    c=db()
    c.executescript("""
    CREATE TABLE IF NOT EXISTS users(id INTEGER PRIMARY KEY AUTOINCREMENT,username TEXT UNIQUE NOT NULL,password TEXT NOT NULL,salt TEXT NOT NULL,created TEXT NOT NULL,google_id TEXT,avatar_url TEXT,avatar_emoji TEXT DEFAULT NULL,guest_id TEXT UNIQUE,display_name TEXT);
    CREATE TABLE IF NOT EXISTS oauth_state(state TEXT PRIMARY KEY,created TEXT NOT NULL);
    CREATE TABLE IF NOT EXISTS tokens(token TEXT PRIMARY KEY,user_id INTEGER NOT NULL,created TEXT NOT NULL);
    CREATE TABLE IF NOT EXISTS api_keys(user_id INTEGER NOT NULL,provider TEXT NOT NULL,key_enc TEXT NOT NULL,PRIMARY KEY(user_id,provider));
    CREATE TABLE IF NOT EXISTS messages(id INTEGER PRIMARY KEY AUTOINCREMENT,user_id INTEGER NOT NULL,role TEXT NOT NULL,content TEXT NOT NULL,model TEXT,ts TEXT NOT NULL);
    CREATE TABLE IF NOT EXISTS preferences(user_id INTEGER PRIMARY KEY,theme TEXT DEFAULT 'system',model_key TEXT DEFAULT 'auto');
    CREATE TABLE IF NOT EXISTS saved_items(id INTEGER PRIMARY KEY AUTOINCREMENT,user_id INTEGER NOT NULL,title TEXT NOT NULL,content TEXT NOT NULL,ts TEXT NOT NULL);
    CREATE TABLE IF NOT EXISTS memory(id INTEGER PRIMARY KEY AUTOINCREMENT,user_id INTEGER NOT NULL,key TEXT NOT NULL,value TEXT NOT NULL,ts TEXT NOT NULL,UNIQUE(user_id,key));
    CREATE TABLE IF NOT EXISTS visitor_log(id INTEGER PRIMARY KEY AUTOINCREMENT,user_id INTEGER,username TEXT,event TEXT NOT NULL,ip TEXT,ua TEXT,ts TEXT NOT NULL);
    CREATE TABLE IF NOT EXISTS conversations(id INTEGER PRIMARY KEY AUTOINCREMENT,user_id INTEGER NOT NULL,title TEXT NOT NULL,created TEXT NOT NULL,updated TEXT NOT NULL);
    CREATE TABLE IF NOT EXISTS conv_messages(id INTEGER PRIMARY KEY AUTOINCREMENT,conv_id INTEGER NOT NULL,user_id INTEGER NOT NULL,role TEXT NOT NULL,content TEXT NOT NULL,model TEXT,ts TEXT NOT NULL);
    """)
    c.commit(); c.close()

def _safe_init_db():
    global DB, _data_dir
    for attempt_dir in [_data_dir, "/tmp", "."]:
        try:
            DB = os.path.join(attempt_dir, "fusionai.db")
            init_db()
            print(f"✅ DB: {DB}", flush=True)
            return
        except Exception as e:
            print(f"⚠️  DB failed at {attempt_dir}: {e}", flush=True)
    print("❌ All DB locations failed", flush=True)
_safe_init_db()

def hash_pw(pw,salt): return hashlib.pbkdf2_hmac("sha256",pw.encode(),salt.encode(),100000).hex()
def make_token(uid):
    tok = uuid.uuid4().hex+uuid.uuid4().hex
    with db() as c: c.execute("INSERT OR REPLACE INTO tokens(token,user_id,created)VALUES(?,?,?)",(tok,uid,datetime.now().isoformat()))
    return tok

def _now(): return datetime.now().isoformat()

def log_event(uid,uname,event,request=None):
    try:
        ip = (request.headers.get("x-forwarded-for","") or (request.client.host if request and request.client else ""))[:64]
        ua = request.headers.get("user-agent","")[:200] if request else ""
        with db() as c: c.execute("INSERT INTO visitor_log(user_id,username,event,ip,ua,ts)VALUES(?,?,?,?,?,?)",(uid,uname,event,ip,ua,_now()))
    except: pass

def get_user_by_token(token:str):
    if not token: return None
    with db() as c: row = c.execute("SELECT user_id FROM tokens WHERE token=?",(token,)).fetchone()
    if not row: return None
    with db() as c: return c.execute("SELECT * FROM users WHERE id=?",(row["user_id"],)).fetchone()

def auth_user(request:Request):
    tok = request.headers.get("x-auth-token","").strip()
    u = get_user_by_token(tok)
    if not u: raise HTTPException(401,{"error":"Not authenticated"})
    return u

def dev_check(request:Request):
    u = auth_user(request)
    if not TP_USERNAME.strip() or u["username"].lower() != TP_USERNAME.strip().lower(): raise HTTPException(403,{"error":"Dev access denied"})
    return u

def is_dev(u): return bool(TP_USERNAME.strip()) and (u["username"].lower()==TP_USERNAME.strip().lower())

def get_available(user_id=None,salt=None):
    a = {}
    if GITHUB_TOKEN.strip():  a["github"]     = GITHUB_TOKEN.strip()
    if CF_ACCOUNT_ID.strip() and CF_KEY.strip(): a["cloudflare"] = CF_KEY.strip()
    if CF_ACCOUNT_ID2.strip() and CF_KEY2.strip(): a["cloudflare2"] = CF_KEY2.strip()
    if GROQ_KEY.strip():      a["groq"]       = GROQ_KEY.strip()
    if OPENROUTER_KEY.strip():a["openrouter"] = OPENROUTER_KEY.strip()
    if HF_TOKEN.strip():      a["huggingface"]= HF_TOKEN.strip()
    if STABILITY_KEY.strip(): a["stability"]  = STABILITY_KEY.strip()
    if OPENAI_KEY.strip():     a["openai"]     = OPENAI_KEY.strip()
    if ANTHROPIC_KEY.strip():  a["anthropic"]  = ANTHROPIC_KEY.strip()
    if GEMINI_KEY.strip():     a["gemini"]     = GEMINI_KEY.strip()
    if DEEPSEEK_KEY.strip():   a["deepseek"]   = DEEPSEEK_KEY.strip()
    if MOONSHOT_KEY.strip():   a["moonshot"]   = MOONSHOT_KEY.strip()
    if MISTRAL_KEY.strip():    a["mistral"]    = MISTRAL_KEY.strip()
    if XAI_KEY.strip():        a["xai"]        = XAI_KEY.strip()
    if COHERE_KEY.strip():     a["cohere"]     = COHERE_KEY.strip()
    if TOGETHER_KEY.strip():   a["together"]   = TOGETHER_KEY.strip()
    if PERPLEXITY_KEY.strip(): a["perplexity"] = PERPLEXITY_KEY.strip()
    if user_id:
        with db() as c: rows = c.execute("SELECT provider,key_enc FROM api_keys WHERE user_id=?",(user_id,)).fetchall()
        for r in rows:
            try:
                v = decrypt(r["key_enc"],salt)
                if v: a[r["provider"]] = v
            except: pass
    return a

def J(data,status=200): return JSONResponse(data,status_code=status)
def err(msg,code=400): return JSONResponse({"error":msg},status_code=code)

MODELS = {
    "gh_gpt4o": {"provider":"github","model":"openai/gpt-4o","label":"GPT-4o","company":"OpenAI","emoji":"🤖","desc":"GPT-4o — vision + fast","thinking":False,"type":"chat","vision":True},
    "gh_gpt4o_mini": {"provider":"github","model":"openai/gpt-4o-mini","label":"GPT-4o Mini","company":"OpenAI","emoji":"🤖","desc":"GPT-4o Mini — cheapest OpenAI","thinking":False,"type":"chat"},
    "gh_o1": {"provider":"github","model":"openai/o1","label":"OpenAI o1","company":"OpenAI","emoji":"🧮","desc":"o1 — deep reasoning, science & math","thinking":True,"type":"chat"},
    "gh_o1_mini": {"provider":"github","model":"openai/o1-mini","label":"OpenAI o1 Mini","company":"OpenAI","emoji":"🧮","desc":"o1 Mini — fast reasoning","thinking":True,"type":"chat"},
    "gh_o3": {"provider":"github","model":"openai/o3","label":"OpenAI o3","company":"OpenAI","emoji":"🧮","desc":"o3 — strongest reasoning model","thinking":True,"type":"chat"},
    "gh_o3_mini": {"provider":"github","model":"openai/o3-mini","label":"OpenAI o3 Mini","company":"OpenAI","emoji":"🧮","desc":"o3 Mini — fast powerful reasoning","thinking":True,"type":"chat"},
    "gh_o4_mini": {"provider":"github","model":"openai/o4-mini","label":"OpenAI o4 Mini","company":"OpenAI","emoji":"🧮","desc":"o4 Mini — latest fast reasoning","thinking":True,"type":"chat"},
    "gh_deepseek_r1": {"provider":"github","model":"deepseek/DeepSeek-R1","label":"DeepSeek R1","company":"DeepSeek","emoji":"🔵","desc":"DeepSeek R1 — open reasoning","thinking":True,"type":"chat"},
    "gh_deepseek_v3": {"provider":"github","model":"deepseek/DeepSeek-V3-0324","label":"DeepSeek V3","company":"DeepSeek","emoji":"🔵","desc":"DeepSeek V3 — fast & capable","thinking":False,"type":"chat"},
    "gh_llama4_scout":{"provider":"github","model":"meta/Llama-4-Scout-17B-16E-Instruct","label":"Llama 4 Scout","company":"Meta","emoji":"🦙","desc":"Llama 4 Scout — vision + multimodal","thinking":False,"type":"chat","vision":True},
    "groq_compound": {"provider":"groq","model":"compound-beta","label":"Groq · Compound Beta","company":"Groq","emoji":"⚡","desc":"Web search + agentic reasoning","thinking":False,"type":"chat"},
    "groq_compound_mini": {"provider":"groq","model":"compound-beta-mini","label":"Groq · Compound Mini","company":"Groq","emoji":"⚡","desc":"Fast agentic, web search","thinking":False,"type":"chat"},
    "groq_llama33_70b": {"provider":"groq","model":"llama-3.3-70b-versatile","label":"Llama 3.3 70B (Groq)","company":"Meta","emoji":"🦙","desc":"Groq ultra-fast Llama","thinking":False,"type":"chat"},
    "groq_llama31_8b": {"provider":"groq","model":"llama-3.1-8b-instant","label":"Llama 3.1 8B (Groq)","company":"Meta","emoji":"🦙","desc":"Groq fastest — instant","thinking":False,"type":"chat"},
    "groq_llama4_scout": {"provider":"groq","model":"meta-llama/llama-4-scout-17b-16e-instruct","label":"Llama 4 Scout (Groq)","company":"Meta","emoji":"🦙","desc":"Llama 4 Scout on Groq","thinking":False,"type":"chat","vision":True},
    "groq_qwen3_32b": {"provider":"groq","model":"qwen/qwen3-32b","label":"Qwen3 32B (Groq)","company":"Alibaba","emoji":"🧠","desc":"Qwen3 32B thinking on Groq","thinking":True,"type":"chat"},
    "groq_gpt_oss_120b": {"provider":"groq","model":"openai/gpt-oss-120b","label":"GPT-OSS 120B (Groq)","company":"OpenAI","emoji":"🤖","desc":"GPT-OSS 120B on Groq","thinking":True,"type":"chat"},
    "or_nemotron_super": {"provider":"openrouter","model":"nvidia/llama-3.1-nemotron-ultra-253b-v1:free","label":"Nemotron Ultra 253B","company":"NVIDIA","emoji":"🟢","desc":"253B ultra reasoning","thinking":True,"type":"chat"},
    "or_nemotron_70b": {"provider":"openrouter","model":"nvidia/llama-3.3-nemotron-super-49b-v1:free","label":"Nemotron Super 49B","company":"NVIDIA","emoji":"🟢","desc":"49B fast + thinking","thinking":True,"type":"chat"},
    "or_nemotron_12b_vl": {"provider":"openrouter","model":"nvidia/nemotron-nano-12b-v2-vl:free","label":"Nemotron 12B Vision","company":"NVIDIA","emoji":"🟢","desc":"12B vision + tools","thinking":False,"type":"chat","vision":True},
    "or_qwen3_235b": {"provider":"openrouter","model":"qwen/qwen3-235b-a22b:free","label":"Qwen3 235B MoE","company":"Alibaba","emoji":"💻","desc":"235B MoE reasoning","thinking":True,"type":"chat"},
    "or_qwen3_30b": {"provider":"openrouter","model":"qwen/qwen3-30b-a3b:free","label":"Qwen3 30B MoE","company":"Alibaba","emoji":"💻","desc":"30B MoE thinking","thinking":True,"type":"chat"},
    "or_deepseek_r1": {"provider":"openrouter","model":"deepseek/deepseek-r1-0528:free","label":"DeepSeek R1 (OR)","company":"DeepSeek","emoji":"🔵","desc":"DeepSeek R1 via OpenRouter","thinking":True,"type":"chat"},
    "or_deepseek_v3": {"provider":"openrouter","model":"deepseek/deepseek-chat-v3-0324:free","label":"DeepSeek V3 (OR)","company":"DeepSeek","emoji":"🔵","desc":"DeepSeek V3 via OpenRouter","thinking":False,"type":"chat"},
    "or_mistral_small": {"provider":"openrouter","model":"mistralai/mistral-small-3.2-24b-instruct:free","label":"Mistral Small 3.2","company":"Mistral AI","emoji":"🌀","desc":"Vision + tools","thinking":False,"type":"chat","vision":True},
    "or_glm_45_air": {"provider":"openrouter","model":"z-ai/glm-4.5-air:free","label":"GLM 4.5 Air","company":"Z-AI","emoji":"🌐","desc":"GLM 4.5 Air — free","thinking":False,"type":"chat"},
    "or_lfm_instruct": {"provider":"openrouter","model":"liquid/lfm-2.5-1.2b-instruct:free","label":"Liquid LFM 1.2B","company":"Liquid AI","emoji":"💧","desc":"Tiny & fast","thinking":False,"type":"chat"},
    "gh_phi4": {"provider":"github","model":"microsoft/Phi-4","label":"Phi-4","company":"Microsoft","emoji":"🔷","desc":"Phi-4 — compact & powerful","thinking":False,"type":"chat"},
    "gh_phi4_mini": {"provider":"github","model":"microsoft/Phi-4-mini","label":"Phi-4 Mini","company":"Microsoft","emoji":"🔷","desc":"Phi-4 Mini — ultra-fast","thinking":False,"type":"chat"},
    "gh_phi4_mm": {"provider":"github","model":"microsoft/Phi-4-multimodal-instruct","label":"Phi-4 Multimodal","company":"Microsoft","emoji":"🔷","desc":"Phi-4 vision + audio","thinking":False,"type":"chat","vision":True},
    "gh_phi35_moe": {"provider":"github","model":"microsoft/Phi-3.5-MoE-instruct","label":"Phi-3.5 MoE","company":"Microsoft","emoji":"🔷","desc":"Phi-3.5 MoE — fast","thinking":False,"type":"chat"},
    "cf_kimi_k2": {"provider":"cloudflare","model":"@cf/moonshotai/kimi-k2.5",                         "label":"Kimi K2.5",          "company":"Moonshot","emoji":"🌙","desc":"Kimi K2.5 — powerful reasoning","thinking":False,"type":"chat"},
    "cf_glm47": {"provider":"cloudflare","model":"@cf/zai-org/glm-4.7-flash",                        "label":"GLM 4.7 Flash",      "company":"Z-AI",    "emoji":"🌐","desc":"GLM 4.7 Flash — fast","thinking":False,"type":"chat"},
    "cf_gpt_oss_120b":{"provider":"cloudflare","model":"@cf/openai/gpt-oss-120b",                          "label":"GPT-OSS 120B (CF)",  "company":"OpenAI",  "emoji":"🤖","desc":"GPT-OSS 120B on Cloudflare","thinking":True,"type":"chat"},
    "cf_gpt_oss_20b": {"provider":"cloudflare","model":"@cf/openai/gpt-oss-20b",                           "label":"GPT-OSS 20B (CF)",   "company":"OpenAI",  "emoji":"🤖","desc":"GPT-OSS 20B — fast reasoning","thinking":False,"type":"chat"},
    "cf_llama4_scout":{"provider":"cloudflare","model":"@cf/meta/llama-4-scout-17b-16e-instruct",          "label":"Llama 4 Scout (CF)", "company":"Meta",    "emoji":"🦙","desc":"Llama 4 Scout on Cloudflare","thinking":False,"type":"chat","vision":True},
    "cf_nemotron": {"provider":"cloudflare","model":"@cf/nvidia/nemotron-3-120b-a12b",                  "label":"Nemotron 120B (CF)", "company":"NVIDIA",  "emoji":"🟢","desc":"Nemotron 120B on Cloudflare","thinking":True,"type":"chat"},
    "cf_qwq_32b": {"provider":"cloudflare","model":"@cf/qwen/qwq-32b",                                 "label":"QwQ 32B (CF)",       "company":"Alibaba", "emoji":"🧠","desc":"QwQ 32B deep thinking","thinking":True,"type":"chat"},
    "cf_qwen3_30b": {"provider":"cloudflare","model":"@cf/qwen/qwen3-30b-a3b-fp8",                       "label":"Qwen3 30B (CF)",     "company":"Alibaba", "emoji":"🧠","desc":"Qwen3 30B MoE on CF","thinking":True,"type":"chat"},
    "cf_qwen25_coder":{"provider":"cloudflare","model":"@cf/qwen/qwen2.5-coder-32b-instruct",              "label":"Qwen2.5 Coder 32B",  "company":"Alibaba", "emoji":"💻","desc":"Qwen 2.5 Coder — code specialist","thinking":False,"type":"chat"},
    "cf_gemma3_12b": {"provider":"cloudflare","model":"@cf/google/gemma-3-12b-it",                        "label":"Gemma 3 12B (CF)",   "company":"Google",  "emoji":"🌟","desc":"Gemma 3 12B on Cloudflare","thinking":False,"type":"chat"},
    "cf_mistral_24b": {"provider":"cloudflare","model":"@cf/mistralai/mistral-small-3.1-24b-instruct",     "label":"Mistral Small 3.1",  "company":"Mistral", "emoji":"🌀","desc":"Mistral Small 3.1 on CF","thinking":False,"type":"chat"},
    "cf_deepseek_r1": {"provider":"cloudflare","model":"@cf/deepseek-ai/deepseek-r1-distill-qwen-32b",     "label":"DeepSeek R1 (CF)",   "company":"DeepSeek","emoji":"🔵","desc":"DeepSeek R1 distill on CF","thinking":True,"type":"chat"},
    "cf_sea_lion": {"provider":"cloudflare","model":"@cf/aisingapore/gemma-sea-lion-v4-27b-it",         "label":"SEA-LION 27B",       "company":"AI SG",   "emoji":"🦁","desc":"SEA-LION multilingual 27B","thinking":False,"type":"chat"},
    "cf_granite": {"provider":"cloudflare","model":"@cf/ibm-granite/granite-4.0-h-micro",              "label":"Granite 4.0 Micro",  "company":"IBM",     "emoji":"🪨","desc":"IBM Granite 4.0 micro","thinking":False,"type":"chat"},
    "cf_cybertron": {"provider":"cloudflare","model":"@cf/fblgit/una-cybertron-7b-v2-bf16",              "label":"Cybertron 7B",       "company":"FBL",     "emoji":"🤖","desc":"Una Cybertron 7B","thinking":False,"type":"chat"},
    "cf_phi2": {"provider":"cloudflare","model":"@cf/microsoft/phi-2",                               "label":"Phi-2 (CF)",         "company":"Microsoft","emoji":"🔷","desc":"Phi-2 — tiny & fast","thinking":False,"type":"chat"},
    "cf_hermes2": {"provider":"cloudflare","model":"@hf/nousresearch/hermes-2-pro-mistral-7b",         "label":"Hermes 2 Pro",       "company":"Nous",    "emoji":"🏛","desc":"Hermes 2 Pro Mistral 7B","thinking":False,"type":"chat"},
    "cf_llava": {"provider":"cloudflare","model":"@cf/llava-hf/llava-1.5-7b-hf",                    "label":"LLaVA 1.5 7B",       "company":"HF",      "emoji":"👁","desc":"LLaVA vision model","thinking":False,"type":"chat","vision":True},
    "cf_flux_klein4b":{"provider":"image_cf","model":"@cf/black-forest-labs/flux-2-klein-4b",             "label":"🖼 FLUX Klein 4B (CF)","company":"CF Image","emoji":"🎨","desc":"FLUX Klein 4B — fast on CF","thinking":False,"type":"image"},
    "cf_flux_klein9b":{"provider":"image_cf","model":"@cf/black-forest-labs/flux-2-klein-9b",             "label":"🖼 FLUX Klein 9B (CF)","company":"CF Image","emoji":"🎨","desc":"FLUX Klein 9B — quality on CF","thinking":False,"type":"image"},
    "cf_flux_dev": {"provider":"image_cf","model":"@cf/black-forest-labs/flux-2-dev",                  "label":"🖼 FLUX Dev (CF)",    "company":"CF Image","emoji":"🎨","desc":"FLUX Dev — highest quality","thinking":False,"type":"image"},
    "cf_phoenix": {"provider":"image_cf","model":"@cf/leonardo/phoenix-1.0",                          "label":"🖼 Phoenix 1.0",      "company":"CF Image","emoji":"🎨","desc":"Leonardo Phoenix on CF","thinking":False,"type":"image"},
    "cf_lucid": {"provider":"image_cf","model":"@cf/leonardo/lucid-origin",                         "label":"🖼 Lucid Origin",     "company":"CF Image","emoji":"🎨","desc":"Leonardo Lucid on CF","thinking":False,"type":"image"},
    "cf_dg_flux": {"provider":"image_cf","model":"@cf/deepgram/flux",                                  "label":"🖼 Deepgram FLUX",    "company":"CF Image","emoji":"🎨","desc":"Deepgram FLUX on CF","thinking":False,"type":"image"},
    "gh_img_gpt": {"provider":"image_gh","model":"openai/gpt-image-1","label":"🖼 GPT-image-1","company":"OpenAI Image","emoji":"🎨","desc":"GPT-image-1 — best quality","thinking":False,"type":"image"},
    "gh_img_dalle3": {"provider":"image_gh","model":"openai/dall-e-3","label":"🖼 DALL-E 3","company":"OpenAI Image","emoji":"🎨","desc":"DALL-E 3 — photorealistic","thinking":False,"type":"image"},
    # ── Direct provider models (use your own API key) ──────────────────────
    "oa_gpt4o": {"provider":"openai","model":"gpt-4o","label":"GPT-4o (Direct)","company":"OpenAI","emoji":"🤖","desc":"GPT-4o via your own OpenAI key","thinking":False,"type":"chat","vision":True},
    "oa_gpt4o_mini": {"provider":"openai","model":"gpt-4o-mini","label":"GPT-4o Mini (Direct)","company":"OpenAI","emoji":"🤖","desc":"GPT-4o Mini via your own OpenAI key","thinking":False,"type":"chat","vision":True},
    "oa_o3": {"provider":"openai","model":"o3","label":"o3 (Direct)","company":"OpenAI","emoji":"🧠","desc":"OpenAI o3 reasoning via your own key","thinking":True,"type":"chat"},
    "an_opus": {"provider":"anthropic","model":"claude-opus-4-1-20250805","label":"Claude Opus 4.1","company":"Anthropic","emoji":"🎭","desc":"Claude Opus via your own Anthropic key","thinking":True,"type":"chat","vision":True},
    "an_sonnet": {"provider":"anthropic","model":"claude-sonnet-4-5-20250929","label":"Claude Sonnet 4.5","company":"Anthropic","emoji":"🎭","desc":"Claude Sonnet via your own Anthropic key","thinking":False,"type":"chat","vision":True},
    "an_haiku": {"provider":"anthropic","model":"claude-3-5-haiku-20241022","label":"Claude Haiku 3.5","company":"Anthropic","emoji":"🎭","desc":"Claude Haiku — fast & cheap","thinking":False,"type":"chat","vision":True},
    "gm_flash": {"provider":"gemini","model":"gemini-2.0-flash","label":"Gemini 2.0 Flash","company":"Google","emoji":"💎","desc":"Gemini 2.0 Flash via your own key","thinking":False,"type":"chat","vision":True},
    "gm_pro": {"provider":"gemini","model":"gemini-2.5-pro","label":"Gemini 2.5 Pro","company":"Google","emoji":"💎","desc":"Gemini 2.5 Pro via your own key","thinking":True,"type":"chat","vision":True},
    "ds_chat": {"provider":"deepseek","model":"deepseek-chat","label":"DeepSeek Chat (Direct)","company":"DeepSeek","emoji":"🔵","desc":"DeepSeek V3 via your own key","thinking":False,"type":"chat"},
    "ds_reasoner": {"provider":"deepseek","model":"deepseek-reasoner","label":"DeepSeek R1 (Direct)","company":"DeepSeek","emoji":"🔵","desc":"DeepSeek R1 via your own key","thinking":True,"type":"chat"},
    "kimi_k2": {"provider":"moonshot","model":"kimi-k2-0711-preview","label":"Kimi K2","company":"Moonshot AI","emoji":"🌙","desc":"Kimi K2 via your own Moonshot key","thinking":False,"type":"chat"},
    "kimi_32k": {"provider":"moonshot","model":"moonshot-v1-32k","label":"Kimi Moonshot v1 32k","company":"Moonshot AI","emoji":"🌙","desc":"Moonshot v1 32k context","thinking":False,"type":"chat"},
    "mistral_large": {"provider":"mistral","model":"mistral-large-latest","label":"Mistral Large (Direct)","company":"Mistral","emoji":"🌀","desc":"Mistral Large via your own key","thinking":False,"type":"chat"},
    "mistral_small": {"provider":"mistral","model":"mistral-small-latest","label":"Mistral Small (Direct)","company":"Mistral","emoji":"🌀","desc":"Mistral Small via your own key","thinking":False,"type":"chat"},
    "xai_grok": {"provider":"xai","model":"grok-4","label":"Grok 4","company":"xAI","emoji":"✖️","desc":"Grok 4 via your own xAI key","thinking":True,"type":"chat"},
    "xai_grok_mini": {"provider":"xai","model":"grok-4-mini","label":"Grok 4 Mini","company":"xAI","emoji":"✖️","desc":"Grok 4 Mini — faster & cheaper","thinking":False,"type":"chat"},
    "cohere_cmdr": {"provider":"cohere","model":"command-r-plus","label":"Command R+","company":"Cohere","emoji":"🔮","desc":"Cohere Command R+ via your own key","thinking":False,"type":"chat"},
    "together_llama": {"provider":"together","model":"meta-llama/Llama-3.3-70B-Instruct-Turbo","label":"Llama 3.3 70B (Together)","company":"Together AI","emoji":"🦙","desc":"Llama 3.3 70B via Together AI","thinking":False,"type":"chat"},
    "perplexity_sonar": {"provider":"perplexity","model":"sonar-pro","label":"Perplexity Sonar Pro","company":"Perplexity","emoji":"🟣","desc":"Perplexity Sonar Pro — web-grounded","thinking":False,"type":"chat"},
}

FALLBACK_CHAINS = {
    "gh_gpt5":        ["gh_gpt4o_mini","gh_gpt4o","groq_llama33_70b","or_deepseek_v3"],
    "gh_gpt5_mini":   ["gh_gpt4o","gh_gpt4o_mini","groq_llama33_70b","or_lfm_instruct"],
    "gh_gpt4o":       ["gh_gpt4o","groq_llama4_scout","or_nemotron_12b_vl"],
    "gh_gpt4o_mini":  ["gh_gpt4o_mini","groq_llama31_8b","or_lfm_instruct"],
    "gh_o3":          ["gh_o4_mini","gh_o3_mini","gh_o1","groq_qwen3_32b","or_qwen3_235b"],
    "gh_o3_mini":     ["gh_o4_mini","gh_o3","gh_o1_mini","groq_qwen3_32b"],
    "gh_o4_mini":     ["gh_o3_mini","gh_o1_mini","groq_qwen3_32b","or_deepseek_r1"],
    "gh_o1":          ["gh_o3","gh_o4_mini","groq_gpt_oss_120b","or_deepseek_r1"],
    "gh_o1_mini":     ["gh_o4_mini","gh_o3_mini","groq_qwen3_32b"],
    "gh_deepseek_r1": ["or_deepseek_r1","groq_qwen3_32b","gh_o4_mini"],
    "gh_deepseek_v3": ["or_deepseek_v3","groq_llama33_70b","gh_gpt5_mini"],
    "gh_llama4_scout":["groq_llama4_scout","gh_gpt4o","or_nemotron_12b_vl"],
    "groq_compound":      ["groq_compound_mini","groq_llama33_70b","gh_gpt4o_mini","or_deepseek_v3"],
    "groq_compound_mini": ["groq_compound","groq_llama33_70b","groq_llama31_8b","or_lfm_instruct"],
    "groq_llama33_70b":   ["groq_compound","gh_gpt4o_mini","groq_gpt_oss_120b","or_deepseek_v3"],
    "groq_llama31_8b":    ["groq_compound_mini","groq_llama33_70b","or_lfm_instruct"],
    "groq_llama4_scout":  ["gh_llama4_scout","groq_llama33_70b","or_nemotron_12b_vl"],
    "groq_qwen3_32b":     ["gh_o4_mini","groq_llama33_70b","or_qwen3_235b"],
    "groq_gpt_oss_120b":  ["gh_o3","groq_llama33_70b","or_qwen3_235b"],
    "or_nemotron_super":  ["or_nemotron_70b","gh_gpt4o_mini","or_qwen3_235b"],
    "or_nemotron_70b":    ["or_nemotron_super","groq_llama33_70b","gh_gpt5_mini"],
    "or_qwen3_235b":      ["groq_qwen3_32b","or_qwen3_30b","gh_gpt5_mini"],
    "or_qwen3_30b":       ["or_qwen3_235b","groq_qwen3_32b","groq_llama33_70b"],
    "or_deepseek_r1":     ["gh_deepseek_r1","groq_qwen3_32b","gh_o4_mini"],
    "or_deepseek_v3":     ["gh_deepseek_v3","groq_llama33_70b","gh_gpt5_mini"],
    "or_mistral_small":   ["gh_gpt4o","groq_llama4_scout","or_nemotron_12b_vl"],
    "or_glm_45_air":      ["groq_llama31_8b","groq_compound_mini","or_lfm_instruct"],
    "or_lfm_instruct":    ["groq_llama31_8b","groq_compound_mini","or_glm_45_air"],
    "gh_phi4":            ["gh_phi4_mini","gh_gpt4o_mini","groq_llama33_70b"],
    "gh_phi4_mini":       ["gh_phi4","gh_gpt4o_mini","groq_llama31_8b"],
    "gh_phi4_mm":         ["gh_gpt4o","groq_llama4_scout","or_nemotron_12b_vl"],
    "gh_phi35_moe":       ["gh_phi4","groq_llama33_70b","or_lfm_instruct"],
}

MODEL_HEALTH = {}
def mark_fail(k): h=MODEL_HEALTH.setdefault(k,{"fails":0,"last_fail":0}); h["fails"]+=1; h["last_fail"]=_time.time()
def mark_ok(k): MODEL_HEALTH.get(k,{}).update({"fails":0}) if k in MODEL_HEALTH else None
def is_healthy(k):
    h=MODEL_HEALTH.get(k,{})
    if not h.get("fails",0): return True
    if _time.time()-h.get("last_fail",0)>300: MODEL_HEALTH[k]["fails"]=0; return True
    return h["fails"]<3

def healthy_chat_models(avail):
    PRI=["gh_gpt4o","gh_o4_mini","gh_o3_mini","gh_o3","gh_o1",
         "gh_deepseek_r1","gh_deepseek_v3","gh_phi4","gh_phi4_mini","gh_phi4_mm","gh_phi35_moe",
         "gh_llama4_scout","gh_gpt4o_mini","gh_o1_mini","groq_gpt_oss_120b","groq_llama33_70b",
         "groq_qwen3_32b","groq_llama4_scout","groq_compound","groq_compound_mini","groq_llama31_8b",
         "or_qwen3_235b","or_deepseek_r1","or_nemotron_super","or_nemotron_70b","or_deepseek_v3",
         "or_qwen3_30b","or_mistral_small","or_glm_45_air","or_lfm_instruct"]
    seen=set()
    for k in PRI:
        m=MODELS.get(k)
        if m and m.get("type","chat")=="chat" and m["provider"] in avail and is_healthy(k): seen.add(k); yield k
    for k,m in MODELS.items():
        if k not in seen and m.get("type","chat")=="chat" and m["provider"] in avail and is_healthy(k): yield k

def _call_overseer(sys_p, user_p, max_tokens=400, cheap=False, timeout=30):
    key=GROQ_KEY.strip() or OPENROUTER_KEY.strip()
    if not key: return None
    ep="https://api.groq.com/openai/v1/chat/completions" if GROQ_KEY.strip() else "https://openrouter.ai/api/v1/chat/completions"
    if GROQ_KEY.strip():
        model = (OVERSEER_ENHANCE_MODEL if cheap else OVERSEER_MODEL)
    else:
        model = "meta-llama/llama-3.1-8b-instruct:free" if cheap else "meta-llama/llama-3.3-70b-instruct:free"
    hdrs={"Authorization":f"Bearer {key}","Content-Type":"application/json"}
    if "openrouter" in ep: hdrs["HTTP-Referer"]="https://huggingface.co/spaces"
    try:
        r=req.post(ep,headers=hdrs,json={"model":model,"messages":[{"role":"system","content":sys_p},{"role":"user","content":user_p}],"max_tokens":max_tokens},timeout=timeout)
        if r.ok: return r.json()["choices"][0]["message"]["content"].strip()
        print(f"[overseer] {r.status_code}: {r.text[:200]}", flush=True)
    except Exception as _oe: print(f"[overseer] {type(_oe).__name__}: {_oe}", flush=True)
    return None

MODEL_PROFILES = {
    "groq_compound":      "Groq Compound Beta — has LIVE WEB SEARCH + agentic tools. Best for: news, current events, real-time data, anything needing internet. Blazing fast.",
    "groq_compound_mini": "Groq Compound Mini — fast web search + agentic. Best for: quick factual lookups, current info, light reasoning with web grounding.",
    "groq_llama33_70b":   "Llama 3.3 70B on Groq — ultra-fast 70B. Best for: general chat, creative writing, summarisation, quick Q&A. Very reliable.",
    "groq_llama31_8b":    "Llama 3.1 8B on Groq — fastest model. Best for: simple questions, high-speed responses, lightweight tasks.",
    "groq_llama4_scout":  "Llama 4 Scout on Groq — vision + multimodal. Best for: image analysis, visual questions, multilingual.",
    "groq_qwen3_32b":     "Qwen3 32B on Groq — extended thinking. Best for: complex math, step-by-step reasoning, science problems.",
    "groq_gpt_oss_120b":  "GPT-OSS 120B on Groq — huge model. Best for: complex analysis, nuanced writing, deep knowledge.",
    "gh_gpt4o":           "GPT-4o (GitHub) — OpenAI flagship. Best for: coding, vision, complex reasoning, professional writing.",
    "gh_gpt4o_mini":      "GPT-4o Mini (GitHub) — fast & cheap OpenAI. Best for: simple tasks, formatting, light coding.",
    "gh_o3":              "OpenAI o3 (GitHub) — strongest reasoning. Best for: advanced math proofs, hard science, complex logic.",
    "gh_o3_mini":         "OpenAI o3 Mini (GitHub) — fast reasoning. Best for: math, coding, logic with speed.",
    "gh_o4_mini":         "OpenAI o4 Mini (GitHub) — latest fast reasoning. Best for: coding, math, analysis.",
    "gh_o1":              "OpenAI o1 (GitHub) — deep thinker. Best for: research-level math, hard problems.",
    "gh_deepseek_r1":     "DeepSeek R1 (GitHub) — open reasoning. Best for: math, coding, chain-of-thought problems.",
    "gh_deepseek_v3":     "DeepSeek V3 (GitHub) — fast & capable. Best for: coding, analysis, writing.",
    "gh_llama4_scout":    "Llama 4 Scout (GitHub) — vision multimodal. Best for: image+text tasks.",
    "gh_phi4":            "Phi-4 (GitHub) — compact powerhouse. Best for: reasoning, coding, science.",
    "or_nemotron_super":  "NVIDIA Nemotron 253B (OpenRouter) — massive model. Best for: research, long documents, hard reasoning.",
    "or_qwen3_235b":      "Qwen3 235B MoE (OpenRouter) — huge MoE. Best for: deep knowledge, multilingual, complex analysis.",
    "or_deepseek_r1":     "DeepSeek R1 (OpenRouter) — reasoning + thinking. Best for: math, science, logic.",
    "or_deepseek_v3":     "DeepSeek V3 (OpenRouter) — fast chat. Best for: coding, writing, Q&A.",
    "or_mistral_small":   "Mistral Small 3.2 (OpenRouter) — vision + tools. Best for: structured tasks, image questions, JSON.",
    "cf_kimi_k2":         "Kimi K2.5 (Cloudflare) — powerful. Best for: reasoning, agentic tasks.",
    "cf_qwq_32b":         "QwQ 32B (Cloudflare) — deep thinking. Best for: math, puzzles, complex logic.",
    "cf_qwen25_coder":    "Qwen 2.5 Coder 32B (Cloudflare) — code specialist. Best for: programming, debugging, code review.",
    "cf_deepseek_r1":     "DeepSeek R1 distill (Cloudflare) — reasoning. Best for: math, step-by-step logic.",
}

def overseer_pick(msg, avail_provs, has_image=False):
    """Use Groq Compound Mini to intelligently pick the best model based on message content and model strengths."""
    def has(k):
        m=MODELS.get(k,{})
        return m.get("provider") in avail_provs and m.get("type","chat")=="chat" and is_healthy(k)

    # ── Simple query fast-path: cheapest/fastest model ──────────────────────
    _complex_kws=['code','python','javascript','typescript','debug','function','class','sql',
                  'math','equation','integral','derivative','proof','calculus','algebra',
                  'explain in detail','analyze','analyse','research','essay','article','story',
                  'poem','write a','create a document','generate a','build a','design a',
                  'compare','pros and cons','step by step','implement','algorithm','architecture',
                  '/doc','/ppt','/csv','/dataset','/report','fusionfile']
    _msg_l=msg.lower().strip()
    _is_simple=(len(msg)<130 and not any(kw in _msg_l for kw in _complex_kws) and msg.count('\n')<3)
    if _is_simple and not has_image:
        for k in ["groq_llama31_8b","or_lfm_instruct","gh_phi4_mini","or_glm_45_air","groq_compound_mini"]:
            if has(k): return {"key":k,"reason":f"{MODELS[k]['label']} · quick"}

    # Vision fast-path
    if has_image:
        for k in ["groq_llama4_scout","or_nemotron_12b_vl","or_mistral_small","gh_gpt4o","gh_llama4_scout","cf_llava"]:
            if has(k): return {"key":k,"reason":f"{MODELS[k]['label']} — vision"}

    cands=[k for k in MODEL_PROFILES if has(k)]
    if not cands:
        # Fall back to any healthy chat model
        for k in healthy_chat_models(avail_provs): return {"key":k,"reason":f"{MODELS[k]['label']} — only available"}
        return {"key":None,"reason":"No models available"}

    if len(cands)==1: return {"key":cands[0],"reason":f"{MODELS[cands[0]]['label']} — only option"}

    # Use Groq Compound Mini to pick — it has web search so it understands current context
    groq_key=(avail_provs.get("groq","") if isinstance(avail_provs,dict) else "")
    if not groq_key: groq_key=GROQ_KEY.strip()

    if groq_key and "groq" in avail_provs:
        profiles="\n".join(f'"{k}": {MODEL_PROFILES[k]}' for k in cands if k in MODEL_PROFILES)
        system="""You are a model router for Fusion.AI. Your ONLY job is to pick the single best model key for a given user message.

ROUTING RULES (follow strictly):
- If the message asks about news, current events, weather, prices, scores, real-time data → pick groq_compound or groq_compound_mini
- If the message is a math problem, equation, proof, calculus → pick groq_qwen3_32b or gh_o3_mini or gh_o4_mini  
- If the message is about coding/programming/debugging → pick cf_qwen25_coder or gh_gpt4o or gh_o4_mini
- If the message needs deep reasoning, logic puzzles, hard science → pick gh_o3 or gh_o3_mini or or_deepseek_r1
- If the message is creative writing, storytelling, poetry → pick groq_llama33_70b or gh_gpt4o
- If the message is a simple/short question or casual chat → pick groq_compound_mini or groq_llama33_70b
- If the message needs a very large model for complex analysis → pick or_nemotron_super or or_qwen3_235b
- Default for general questions: groq_compound (it has web search — always better grounded)

Reply with ONLY the model key. No explanation. No punctuation. Just the key."""
        try:
            r=req.post("https://api.groq.com/openai/v1/chat/completions",
                headers={"Authorization":f"Bearer {groq_key}","Content-Type":"application/json"},
                json={"model":"compound-beta-mini","messages":[
                    {"role":"system","content":system},
                    {"role":"user","content":f"User message: {msg[:800]}\n\nAvailable models:\n{profiles}\n\nPick the best key:"}
                ],"max_tokens":30,"stream":False,"temperature":0},
                timeout=8)
            if r.ok:
                picked=r.json()["choices"][0]["message"]["content"].strip().strip('"').split()[0]
                if picked in MODELS and has(picked):
                    return {"key":picked,"reason":f"{MODELS[picked]['label']} · AI-selected"}
        except Exception as e:
            print(f"[overseer] compound-mini failed: {e}",flush=True)

    # Fallback: smart rule-based (only reached if compound-mini unavailable)
    tl=msg.lower()
    priority=[
        (["news","latest","today","current","price","weather","score","stock","2025","2026"],"groq_compound"),
        (["code","python","javascript","debug","function","class","api","sql","script","program"],"cf_qwen25_coder"),
        (["math","equation","integral","derivative","proof","calculus","algebra","geometry","statistics"],"groq_qwen3_32b"),
        (["reason","logic","step by step","think","explain why","analyse","analyze"],"gh_o3_mini"),
        (["write","story","essay","poem","creative","blog","article"],"groq_llama33_70b"),
    ]
    for kws,model in priority:
        if model in cands and any(k in tl for k in kws):
            return {"key":model,"reason":f"{MODELS[model]['label']} — keyword match"}

    # Default: groq_compound if available (always better with web search), else best available
    for k in ["groq_compound","groq_compound_mini","groq_llama33_70b","gh_gpt4o"]:
        if k in cands: return {"key":k,"reason":f"{MODELS[k]['label']} — default"}
    return {"key":cands[0],"reason":f"{MODELS[cands[0]]['label']} — fallback"}

ENDPOINTS={"groq":"https://api.groq.com/openai/v1/chat/completions",
            "openrouter":"https://openrouter.ai/api/v1/chat/completions",
            "github":GITHUB_ENDPOINT,
            # ── Direct provider APIs (all OpenAI-compatible /chat/completions shape) ──
            "openai":"https://api.openai.com/v1/chat/completions",
            "deepseek":"https://api.deepseek.com/v1/chat/completions",
            "moonshot":"https://api.moonshot.ai/v1/chat/completions",
            "mistral":"https://api.mistral.ai/v1/chat/completions",
            "xai":"https://api.x.ai/v1/chat/completions",
            "together":"https://api.together.xyz/v1/chat/completions",
            "perplexity":"https://api.perplexity.ai/chat/completions",
            "cohere":"https://api.cohere.ai/compatibility/v1/chat/completions"}
            # NOTE: "anthropic" and "gemini" are NOT OpenAI-compatible — they have
            # dedicated call functions (_call_anthropic_api / _call_gemini_api) below.

# ── Anthropic (Claude) + Google Gemini — non-OpenAI-shaped direct APIs ────────
def _strip_system(msgs):
    """Pull system message(s) out of an OpenAI-style messages list."""
    sys_text=""; out=[]
    for m in msgs:
        if m.get("role")=="system": sys_text+=(m.get("content") or "")+"\n"
        else: out.append(m)
    return sys_text.strip(), out

def _call_anthropic_api(model,msgs,key,max_tokens=4096,image_b64=None,image_mime=None,timeout=90):
    """Call Anthropic /v1/messages. Returns (ok, text_or_error)."""
    sys_text,conv=_strip_system(msgs)
    amsgs=[]
    for m in conv:
        role=m.get("role")
        if role not in ("user","assistant"): continue
        content=m.get("content")
        if isinstance(content,str): amsgs.append({"role":role,"content":content})
        elif isinstance(content,list):
            blocks=[{"type":"text","text":c.get("text","")} for c in content if isinstance(c,dict) and c.get("type")=="text"]
            amsgs.append({"role":role,"content":blocks or [{"type":"text","text":""}]})
        else: amsgs.append({"role":role,"content":str(content)})
    if image_b64 and amsgs:
        for m in reversed(amsgs):
            if m["role"]=="user":
                txt=m["content"] if isinstance(m["content"],str) else "".join(b.get("text","") for b in m["content"] if isinstance(b,dict))
                m["content"]=[{"type":"image","source":{"type":"base64","media_type":image_mime or "image/jpeg","data":image_b64}},{"type":"text","text":txt}]
                break
    if not amsgs: amsgs=[{"role":"user","content":"Hello"}]
    body={"model":model,"max_tokens":max_tokens,"messages":amsgs}
    if sys_text: body["system"]=sys_text
    try:
        r=req.post("https://api.anthropic.com/v1/messages",
            headers={"x-api-key":key,"anthropic-version":"2023-06-01","content-type":"application/json"},
            json=body,timeout=timeout)
        if not r.ok:
            try: em=r.json().get("error",{}).get("message",r.text[:200])
            except: em=r.text[:200]
            return False,f"[{r.status_code}] {em}"
        data=r.json()
        txt="".join(b.get("text","") for b in data.get("content",[]) if isinstance(b,dict) and b.get("type")=="text")
        return True,txt
    except Exception as ex: return False,str(ex)

def _call_gemini_api(model,msgs,key,max_tokens=4096,image_b64=None,image_mime=None,timeout=90):
    """Call Google Gemini generateContent. Returns (ok, text_or_error)."""
    sys_text,conv=_strip_system(msgs)
    contents=[]
    for m in conv:
        role=m.get("role")
        if role not in ("user","assistant"): continue
        grole="model" if role=="assistant" else "user"
        content=m.get("content"); parts=[]
        if isinstance(content,str): parts.append({"text":content})
        elif isinstance(content,list):
            for c in content:
                if isinstance(c,dict) and c.get("type")=="text": parts.append({"text":c.get("text","")})
        else: parts.append({"text":str(content)})
        if parts: contents.append({"role":grole,"parts":parts})
    if image_b64 and contents:
        for c in reversed(contents):
            if c["role"]=="user":
                c["parts"].append({"inline_data":{"mime_type":image_mime or "image/jpeg","data":image_b64}})
                break
    if not contents: contents=[{"role":"user","content":[{"text":"Hello"}]}]
    body={"contents":contents,"generationConfig":{"maxOutputTokens":max_tokens}}
    if sys_text: body["systemInstruction"]={"parts":[{"text":sys_text}]}
    url=f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent?key={key}"
    try:
        r=req.post(url,headers={"Content-Type":"application/json"},json=body,timeout=timeout)
        if not r.ok:
            try: em=r.json().get("error",{}).get("message",r.text[:200])
            except: em=r.text[:200]
            return False,f"[{r.status_code}] {em}"
        data=r.json()
        try: txt="".join(p.get("text","") for p in data["candidates"][0]["content"]["parts"])
        except Exception: txt=""
        return True,txt
    except Exception as ex: return False,str(ex)

@app.post("/api/register")
async def register(request:Request):
    d=await request.json()
    u,p=d.get("username","").strip(),d.get("password","").strip()
    hcaptcha_tok=d.get("h-captcha-response","").strip()
    if not u or not p: return err("Username and password required")
    if len(u)<2: return err("Username must be 2+ chars")
    if len(p)<6: return err("Password must be 6+ chars")
    if HCAPTCHA_SECRET.strip():
        try:
            hc=req.post("https://api.hcaptcha.com/siteverify",data={"secret":HCAPTCHA_SECRET,"response":hcaptcha_tok},timeout=8)
            if not hc.json().get("success"): return err("Human verification failed. Please complete the CAPTCHA.")
        except: pass
    salt=uuid.uuid4().hex
    try:
        with db() as c: c.execute("INSERT INTO users(username,password,salt,created)VALUES(?,?,?,?)",(u,hash_pw(p,salt),salt,_now()))
        log_event(None,u,"register",request); return J({"ok":True})
    except sqlite3.IntegrityError: return err("Username taken",409)
@app.post("/api/login")
async def login(request:Request):
    d=await request.json()
    with db() as c: u=c.execute("SELECT * FROM users WHERE username=?",(d.get("username",""),)).fetchone()
    if not u or hash_pw(d.get("password",""),u["salt"])!=u["password"]: return err("Invalid credentials",401)
    tok=make_token(u["id"]); log_event(u["id"],u["username"],"login",request)
    with db() as c: pref=c.execute("SELECT * FROM preferences WHERE user_id=?",(u["id"],)).fetchone()
    return J({"ok":True,"username":u["username"],"token":tok,
              "theme":pref["theme"] if pref else "dark","model_key":pref["model_key"] if pref else "auto","is_dev":is_dev(u)})
@app.post("/api/guest")
async def guest_login(request:Request):
    gid=uuid.uuid4().hex[:12]; uname=f"guest_{gid}"; pw=uuid.uuid4().hex; salt=uuid.uuid4().hex
    with db() as c:
        c.execute("INSERT INTO users(username,password,salt,created)VALUES(?,?,?,?)",(uname,hash_pw(pw,salt),salt,_now()))
        uid=c.execute("SELECT id FROM users WHERE username=?",(uname,)).fetchone()["id"]
    tok=make_token(uid); log_event(uid,"Guest","guest_login",request)
    return J({"ok":True,"username":"Guest","token":tok,"theme":"system","model_key":"auto","is_dev":False})
@app.post("/api/logout")
async def logout(request:Request):
    tok=request.headers.get("x-auth-token","").strip()
    if tok:
        with db() as c:
            row=c.execute("SELECT user_id FROM tokens WHERE token=?",(tok,)).fetchone()
            if row:
                u2=c.execute("SELECT username FROM users WHERE id=?",(row["user_id"],)).fetchone()
                if u2 and u2["username"].startswith("guest_"):
                    for t in ["messages","tokens","api_keys","preferences","saved_items","memory","conv_messages","conversations"]:
                        c.execute(f"DELETE FROM {t} WHERE user_id=?",(row["user_id"],))
                    c.execute("DELETE FROM users WHERE id=?",(row["user_id"],))
                    return J({"ok":True})
            c.execute("DELETE FROM tokens WHERE token=?",(tok,))
    return J({"ok":True})
@app.get("/api/me")
async def whoami(request:Request):
    tok=request.headers.get("x-auth-token","").strip()
    u=get_user_by_token(tok)
    if not u: return J({"logged_in":False})
    # Auto-assign guest_id if missing
    import secrets as _sec
    with db() as c:
        row=c.execute("SELECT guest_id,avatar_emoji FROM users WHERE id=?",(u["id"],)).fetchone()
        gid=row["guest_id"] if row and row["guest_id"] else None
        aemoji=row["avatar_emoji"] if row and row["avatar_emoji"] else None
        if not gid:
            gid="guest_"+_sec.token_hex(4).upper()
            c.execute("UPDATE users SET guest_id=? WHERE id=?",(gid,u["id"]))
        if not aemoji:
            EMOJIS=["🦊","🐺","🦁","🐯","🦝","🐻","🦄","🐸","🦋","🌊","⚡","🌙","🔥","💎","🌸","🎯","🚀","🎸","🌈","🏔"]
            aemoji=EMOJIS[u["id"]%len(EMOJIS)]
            c.execute("UPDATE users SET avatar_emoji=? WHERE id=?",(aemoji,u["id"]))
    with db() as c: pref=c.execute("SELECT * FROM preferences WHERE user_id=?",(u["id"],)).fetchone()
    display="Guest" if u["username"].startswith("guest_") else (u["display_name"] if u["display_name"] else u["username"])
    avatar=u["avatar_url"] if u["avatar_url"] else ""
    return J({"logged_in":True,"username":display,"theme":pref["theme"] if pref else "system",
              "model_key":pref["model_key"] if pref else "auto","is_dev":is_dev(u),"avatar":avatar,"has_google":bool(u["google_id"])})
@app.post("/api/prefs")
async def save_prefs(request:Request):
    u=auth_user(request); d=await request.json()
    with db() as c: c.execute("INSERT OR REPLACE INTO preferences(user_id,theme,model_key)VALUES(?,?,?)",(u["id"],d.get("theme","system"),d.get("model_key","auto")))
    return J({"ok":True})
@app.get("/api/keys")
async def get_keys(request:Request):
    u=auth_user(request)
    with db() as c: rows=c.execute("SELECT provider FROM api_keys WHERE user_id=?",(u["id"],)).fetchall()
    return J({"providers":[r["provider"] for r in rows]})
@app.post("/api/keys")
async def save_key(request:Request):
    u=auth_user(request); d=await request.json()
    prov,key=d.get("provider"),d.get("key","").strip()
    ALLOWED=("groq","openrouter","huggingface","github","stability","cloudflare","extra","custom_endpoint")
    if prov not in ALLOWED: return err("Unknown provider")
    with db() as c:
        if not key: c.execute("DELETE FROM api_keys WHERE user_id=? AND provider=?",(u["id"],prov))
        else: c.execute("INSERT OR REPLACE INTO api_keys(user_id,provider,key_enc)VALUES(?,?,?)",(u["id"],prov,encrypt(key,u["salt"])))
    return J({"ok":True})
@app.post("/api/forgot-password")
async def forgot_password(request:Request):
    d=await request.json()
    uname,new_pw=d.get("username","").strip(),d.get("new_password","").strip()
    if not uname or not new_pw: return err("Username and new password required")
    if len(new_pw)<6: return err("Password must be 6+ chars")
    with db() as c: u=c.execute("SELECT * FROM users WHERE username=?",(uname,)).fetchone()
    if not u: return err("Username not found",404)
    new_salt=uuid.uuid4().hex
    with db() as c:
        c.execute("UPDATE users SET password=?,salt=? WHERE id=?",(hash_pw(new_pw,new_salt),new_salt,u["id"]))
        c.execute("DELETE FROM tokens WHERE user_id=?",(u["id"],))
    return J({"ok":True,"message":"Password updated."})
@app.get("/api/models")
async def get_models():
    path=os.path.join(_data_dir,"disabled_models.json"); disabled=[]
    if os.path.exists(path):
        try:
            with open(path) as f2: disabled=json.load(f2)
        except: pass
    models=[{"key":k,"label":v["label"],"company":v["company"],"emoji":v["emoji"],"desc":v["desc"],
             "provider":v["provider"],"vision":v.get("vision",False),"thinking":v.get("thinking",False),
             "type":v.get("type","chat"),"disabled":k in disabled} for k,v in MODELS.items()]
    return J({"models":models,"count":len(models)})
@app.post("/api/chat")
async def chat(request:Request):
    u=auth_user(request); d=await request.json()
    user_msg=d.get("message","").strip(); history=d.get("history",[])
    mkey=d.get("model_key"); image_b64=d.get("image_b64",""); image_mime=d.get("image_mime","image/jpeg")
    file_text=d.get("file_text","")
    if not user_msg and not image_b64 and not file_text: return err("Empty message")
    full_msg=user_msg
    if file_text: full_msg=(user_msg+"\n\n[File]:\n"+file_text[:80000]).strip() or "[File]:\n"+file_text[:80000]
    if not full_msg: full_msg="Analyse this image."
    avail=get_available(u["id"],u["salt"])
    has_img=bool(image_b64)
    if mkey and mkey in MODELS and MODELS[mkey]["provider"] in avail and MODELS[mkey].get("type","chat")=="chat": pick={"key":mkey,"reason":"Manual"}
    else: pick=overseer_pick(full_msg, avail, has_img)
    if not pick["key"]: return err("No API keys configured.")

    def build_chain(pk):
        chain=[pk]
        for fb in FALLBACK_CHAINS.get(pk,[]):
            if fb not in chain and fb in MODELS and MODELS[fb]["provider"] in avail and MODELS[fb].get("type","chat")=="chat": chain.append(fb)
        for k in ["groq_llama33_70b","groq_compound","groq_compound_mini","gh_phi4","or_deepseek_v3","groq_llama31_8b"]:
            if k not in chain and k in MODELS and MODELS[k]["provider"] in avail and MODELS[k].get("type","chat")=="chat": chain.append(k)
        return chain

    try_list=build_chain(pick["key"])
    with db() as c:
        mem=c.execute("SELECT key,value FROM memory WHERE user_id=? ORDER BY ts DESC LIMIT 40",(u["id"],)).fetchall()
    sys_parts=[f"You are Fusion.AI, an intelligent AI assistant built into an all-in-one AI platform. Today is {datetime.now().strftime('%A, %B %d, %Y')} and you're chatting with {u['username']}. PERSONALITY: Talk like a knowledgeable, thoughtful friend — warm, direct, never stiff. Use natural language; contractions are fine. Be concise on simple questions, go deep when warranted. Show genuine curiosity. Never start with 'Certainly!', 'Of course!', 'Great question!' or 'Absolutely!' — these are filler. Don't repeat the question back, just answer it. When you don't know something, say so plainly. FORMATTING: Use markdown (bold, code, headers, bullets) only when it genuinely helps. For math: $...$ inline, $$...$$ display block. For code: always complete runnable code — never placeholder comments. Keep responses tight — no padding, no 'In conclusion', no 'I hope this helps'. WEB SEARCH: When you get <web_context> tags, use that data and reference it naturally without saying 'According to my web search'. FILE GENERATION: When asked to make any file, wrap COMPLETE content EXACTLY as: <<<FUSIONFILE:filename.ext>>>\n[full content]\n<<<END_FUSIONFILE>>>. Types: .md .txt .csv .json .html .py .js .ts .css .sql .xml .yaml .sh .go .java .cpp .rb .php .toml .env .pptx.md .docx .xlsx. pptx.md: # Title per slide, --- between slides, 8+ slides. xlsx/docx: write CSV/markdown respectively (server converts to binary). ALWAYS complete — zero placeholders, zero truncation."]

    if mem: sys_parts.append("Remember about user:\n"+"\n".join(f"- {r['key']}: {r['value']}" for r in mem))
    sys_msg={"role":"system","content":" ".join(sys_parts)}
    base_msgs=[sys_msg]+history+[{"role":"user","content":full_msg}]
    with db() as c: c.execute("INSERT INTO messages(user_id,role,content,model,ts)VALUES(?,?,?,?,?)",(u["id"],"user",user_msg,MODELS[try_list[0]]["label"],_now()))

    def _call(key):
        m=MODELS[key]; prov=m["provider"]
        if prov in ("image_pol","image_gh","image_cf","image_or","image","video","video3d"): return False,"media",m["label"]
        if prov not in avail: return False,f"No key for {prov}",m["label"]
        # ── Anthropic / Gemini: non-OpenAI-shaped direct APIs, handled separately ──
        if prov=="anthropic":
            ok,res=_call_anthropic_api(m["model"],base_msgs,avail[prov],max_tokens=4096,image_b64=image_b64 or None,image_mime=image_mime)
            if not ok: mark_fail(key); return False,res,m["label"]
            mark_ok(key); return True,("__text__",res or "(no response)"),m["label"]
        if prov=="gemini":
            ok,res=_call_gemini_api(m["model"],base_msgs,avail[prov],max_tokens=4096,image_b64=image_b64 or None,image_mime=image_mime)
            if not ok: mark_fail(key); return False,res,m["label"]
            mark_ok(key); return True,("__text__",res or "(no response)"),m["label"]
        hdrs={"Authorization":f"Bearer {avail[prov]}","Content-Type":"application/json"}
        if prov=="openrouter": hdrs.update({"X-Title":"Fusion.AI","HTTP-Referer":"https://fusionai.space"})
        if image_b64 and prov in ("openrouter","groq","github"):
            use_msgs=[sys_msg]+history+[{"role":"user","content":[{"type":"image_url","image_url":{"url":f"data:{image_mime};base64,{image_b64}"}},{"type":"text","text":full_msg}]}]
        else: use_msgs=base_msgs
        body={"model":m["model"],"messages":use_msgs}
        try:
            if prov=="cloudflare":
                body.pop("model",None); body["max_tokens"]=4096
                # Retry cloudflare up to 2 times
                r=None
                for _cf_try in range(2):
                    try:
                        r=req.post(cf_ep(m["model"]),headers={"Authorization":f"Bearer {avail['cloudflare']}","Content-Type":"application/json"},json=body,timeout=60)
                        if r.ok: break
                    except req.exceptions.Timeout:
                        if _cf_try==1: raise
                        continue
                if not r.ok: mark_fail(key); return False,f"[{r.status_code}] {r.text[:80]}",m["label"]
                mark_ok(key)
                try: data=r.json(); txt=(data.get("result",{})or{}).get("response","") or (data.get("choices",[{}])[0].get("message",{})or{}).get("content","")
                except: txt=""
                return True,("__text__",txt or "(no response)"),m["label"]
            is_o=m.get("thinking",False) and prov=="github"
            if is_o:
                body["max_completion_tokens"]=8000
                r=req.post(ENDPOINTS[prov],headers=hdrs,json=body,timeout=180)
            else:
                body.update({"max_tokens":4096,"stream":True})
                # Compound Beta needs longer timeout due to tool calls
                _timeout=120 if "compound" in m.get("model","") else 90
                r=req.post(ENDPOINTS[prov],headers=hdrs,json=body,stream=True,timeout=_timeout)
            if not r.ok:
                try: e2=r.json().get("error",{}); em=e2.get("message",str(e2)) if isinstance(e2,dict) else str(e2)
                except: em=r.text[:200]
                if r.status_code==429: em="Rate limit — will try next model"
                elif r.status_code==413: em="Message too long for this model"
                elif r.status_code==503: em="Model temporarily unavailable"
                print(f"[fallback] {key} [{r.status_code}]: {em[:100]}",flush=True); mark_fail(key); return False,f"[{r.status_code}] {em}",m["label"]
            mark_ok(key)
            if is_o:
                try: txt=r.json()["choices"][0]["message"]["content"]
                except: txt=""
                return True,("__text__",txt),m["label"]
            return True,r,m["label"]
        except Exception as ex: print(f"[fallback] {key} exc: {ex}",flush=True); mark_fail(key); return False,str(ex),m["label"]

    import re as _re
    def gen_stream():
        resp=None; used_key=try_list[0]; used_label=MODELS[used_key]["label"]; tried=[]
        for i,k in enumerate(try_list):
            tried.append(MODELS[k]["label"]); ok,res,lbl=_call(k)
            if ok: resp=res; used_key=k; used_label=lbl; break
            if i+1<len(try_list):
                yield f"data: {json.dumps({'type':'retry','failed':lbl,'next':MODELS[try_list[i+1]]['label']})}\n\n"
        if resp is None:
            yield f"data: {json.dumps({'type':'error','message':'All models failed to respond. Try again in a moment, or check Settings → API Keys.'})}\n\n"
            return

        retry_note=f" (fallback: {', '.join(tried[:-1])})" if len(tried)>1 else ""

        # ── Non-streaming text (Cloudflare / o-series) ────────────────────────
        if isinstance(resp,tuple) and resp[0]=="__text__":
            txt=resp[1]
            yield f"data: {json.dumps({'type':'meta','model':used_label,'reason':MODELS[used_key]['provider']+retry_note})}\n\n"
            words=txt.split(' '); chunk=""
            for w in words:
                chunk+=w+" "
                if len(chunk)>40: yield f"data: {json.dumps({'type':'delta','text':chunk})}\n\n"; chunk=""
            if chunk: yield f"data: {json.dumps({'type':'delta','text':chunk.rstrip()})}\n\n"
            with db() as c2: c2.execute("INSERT INTO messages(user_id,role,content,model,ts)VALUES(?,?,?,?,?)",(u["id"],"assistant",txt,used_label,_now()))
            yield f"data: {json.dumps({'type':'done'})}\n\n"; return

        yield f"data: {json.dumps({'type':'meta','model':used_label,'reason':pick['reason']+retry_note})}\n\n"
        raw_buf=""; streamed=""; think_buf=""; think_on=False; think_done=False
        model_thinks=MODELS.get(used_key,{}).get("thinking",False)
        # Track Groq Compound tool call accumulation
        tool_args_buf=""; tool_name_buf=""; in_tool=False; tool_output_buf=""
        try:
            for line in resp.iter_lines():
                if not line: continue
                if isinstance(line,bytes): line=line.decode("utf-8","replace")
                if not line.startswith("data: ") or line=="data: [DONE]": continue
                try:
                    parsed=json.loads(line[6:])
                    ch=parsed.get("choices",[{}])[0]; delta=ch.get("delta",{})

                    # ── Groq Compound: tool_calls (web search results embedded) ──
                    tcs=delta.get("tool_calls") or []
                    for tc in tcs:
                        fn=tc.get("function",{})
                        tn=fn.get("name","") or ""
                        ta=fn.get("arguments","") or ""
                        if tn and not in_tool:
                            tool_name_buf=tn; in_tool=True
                            _stxt="🔍 *Searching ("+tn+")…*\n\n"
                            yield f"data: {json.dumps({'type':'delta','text':_stxt})}\n\n"
                            streamed+=_stxt
                        tool_args_buf+=ta

                    # ── tool result message (role=tool in subsequent message) ──
                    if delta.get("role")=="tool":
                        tool_output_buf+=delta.get("content","") or ""

                    rc=delta.get("reasoning_content") or ""
                    dc=delta.get("content") or ""
                    if rc:
                        if not think_on: yield f"data: {json.dumps({'type':'think_start'})}\n\n"; think_on=True
                        yield f"data: {json.dumps({'type':'think_delta','text':rc})}\n\n"
                    if dc:
                        raw_buf+=dc
                        if "<think>" in raw_buf and not think_on: yield f"data: {json.dumps({'type':'think_start'})}\n\n"; think_on=True
                        if think_on and not think_done:
                            if "</think>" in raw_buf:
                                m2=_re.search(r"<think>([\s\S]*?)</think>",raw_buf)
                                if m2:
                                    rem=m2.group(1).strip()[len(think_buf):]
                                    if rem: yield f"data: {json.dumps({'type':'think_delta','text':rem})}\n\n"
                                    think_done=True; yield f"data: {json.dumps({'type':'think_end'})}\n\n"
                            else:
                                ins=_re.search(r"<think>([\s\S]*?)$",raw_buf)
                                if ins:
                                    nw=ins.group(1)[len(think_buf):]
                                    if nw: think_buf+=nw; yield f"data: {json.dumps({'type':'think_delta','text':nw})}\n\n"
                                continue
                        if not think_on or think_done:
                            clean=_re.sub(r"<think>[\s\S]*?</think>","",raw_buf).strip()
                            if not model_thinks: clean=_re.sub(r"^(#+\s*)+","",clean).strip()
                            new=clean[len(streamed):]
                            if new: streamed+=new; yield f"data: {json.dumps({'type':'delta','text':new})}\n\n"

                    # ── finish_reason: tool_calls — Groq Compound done with search, do final non-stream call ──
                    fr=ch.get("finish_reason","")
                    if fr=="tool_calls" and in_tool:
                        # Proper Compound Beta follow-up: send assistant tool_call + tool result
                        try:
                            # Build the tool_call assistant message
                            asst_tool_msg={"role":"assistant","content":None,"tool_calls":[{
                                "id":"call_compound_"+tool_name_buf[:20],
                                "type":"function",
                                "function":{"name":tool_name_buf,"arguments":tool_args_buf or "{}"}
                            }]}
                            # Build tool result message
                            tool_result_msg={"role":"tool",
                                "tool_call_id":"call_compound_"+tool_name_buf[:20],
                                "content":tool_output_buf or "Search complete."}
                            msgs2=base_msgs+[asst_tool_msg,tool_result_msg]
                            body2={"model":MODELS[used_key]["model"],"messages":msgs2,
                                   "max_tokens":4096,"stream":True}
                            r2=req.post(ENDPOINTS["groq"],
                                headers={"Authorization":f"Bearer {avail['groq']}","Content-Type":"application/json"},
                                json=body2,stream=True,timeout=90)
                            if r2.ok:
                                follow_buf=""
                                for line2 in r2.iter_lines():
                                    if not line2: continue
                                    if isinstance(line2,bytes): line2=line2.decode("utf-8","replace")
                                    if not line2.startswith("data: ") or line2=="data: [DONE]": continue
                                    try:
                                        p2=json.loads(line2[6:])
                                        dc2=(p2.get("choices",[{}])[0].get("delta",{}) or {}).get("content") or ""
                                        if dc2:
                                            follow_buf+=dc2
                                            clean2=_re.sub(r"<think>[\s\S]*?</think>","",follow_buf).strip()
                                            new2=clean2[len(streamed):]
                                            if new2:
                                                streamed+=new2; raw_buf=clean2
                                                yield f"data: {json.dumps({'type':'delta','text':new2})}\n\n"
                                    except: pass
                            else:
                                # Fallback: non-streaming with just base_msgs + context note
                                fallback_msgs=base_msgs+[{"role":"user","content":"[Web search completed. Now provide a complete answer.]"}]
                                rb={"model":MODELS[used_key]["model"],"messages":fallback_msgs,"max_tokens":4096,"stream":False}
                                rf=req.post(ENDPOINTS["groq"],headers={"Authorization":f"Bearer {avail['groq']}","Content-Type":"application/json"},json=rb,timeout=60)
                                if rf.ok:
                                    final=(rf.json().get("choices",[{}])[0].get("message",{}) or {}).get("content") or ""
                                    if final:
                                        clean2=_re.sub(r"<think>[\s\S]*?</think>","",final).strip()
                                        new2=clean2[len(streamed):]
                                        if new2: streamed+=new2; raw_buf=clean2; yield f"data: {json.dumps({'type':'delta','text':new2})}\n\n"
                        except Exception as te:
                            print(f"[compound follow-up] {te}",flush=True)
                            # Last resort: emit error note
                            note="\n\n*(Web search completed but response retrieval failed — please retry)*"
                            yield f"data: {json.dumps({'type':'delta','text':note})}\n\n"
                        in_tool=False; tool_args_buf=""; tool_name_buf=""; tool_output_buf=""
                except Exception: pass

            if think_on and not think_done: yield f"data: {json.dumps({'type':'think_end'})}\n\n"
            full=_re.sub(r"<think>[\s\S]*?</think>","",raw_buf).strip()
            if not model_thinks: full=_re.sub(r"^(?:[#\-=]{2,}\s*\n)+","",full).strip()
            with db() as c2: c2.execute("INSERT INTO messages(user_id,role,content,model,ts)VALUES(?,?,?,?,?)",(u["id"],"assistant",full,used_label,_now()))
            try:
                um=full_msg.lower()
                if any(w in um for w in["my name","i am","i'm","i like","i prefer","i hate","i work","i live","remember"]):
                    mk=avail.get("groq","") or GROQ_KEY.strip()
                    if mk:
                        me2="https://api.groq.com/openai/v1/chat/completions"; mm="llama-3.1-8b-instant"
                        mh={"Authorization":f"Bearer {mk}","Content-Type":"application/json"}
                        mr=req.post(me2,headers=mh,json={"model":mm,"messages":[{"role":"user","content":f'User said: "{full_msg[:400]}"\nExtract up to 3 facts. Reply ONLY JSON:[{{"key":"...","value":"..."}}] or []'}],"max_tokens":200},timeout=8)
                        if mr.ok:
                            raw2=mr.json()["choices"][0]["message"]["content"].strip().strip("```json").strip("```").strip()
                            facts=json.loads(raw2)
                            if isinstance(facts,list):
                                with db() as cm:
                                    for f2 in facts[:3]:
                                        if isinstance(f2,dict) and f2.get("key") and f2.get("value"):
                                            cm.execute("INSERT OR REPLACE INTO memory(user_id,key,value,ts)VALUES(?,?,?,?)",(u["id"],str(f2["key"])[:100],str(f2["value"])[:500],_now()))
            except: pass
            yield f"data: {json.dumps({'type':'done'})}\n\n"
        except Exception as e: yield f"data: {json.dumps({'type':'error','message':str(e)})}\n\n"

    import asyncio
    loop = asyncio.get_event_loop()
    return StreamingResponse(
        _run_in_thread(gen_stream, loop),
        media_type="text/event-stream",
        headers={"Cache-Control":"no-cache","X-Accel-Buffering":"no"}
    )


def _run_in_thread(gen_fn, loop):
    """Wrap a blocking sync generator to run in a threadpool, yielding to async."""
    import concurrent.futures, queue, threading
    q = queue.Queue()
    sentinel = object()

    def run():
        try:
            for chunk in gen_fn():
                q.put(chunk)
        finally:
            q.put(sentinel)

    threading.Thread(target=run, daemon=True).start()

    async def aiter():
        while True:
            chunk = await loop.run_in_executor(None, q.get)
            if chunk is sentinel:
                break
            yield chunk
    return aiter()
@app.get("/api/history")
async def get_history(request:Request):
    u=auth_user(request)
    with db() as c:
        rows=c.execute("SELECT role,content,model,ts FROM messages WHERE user_id=? ORDER BY id DESC LIMIT 120",(u["id"],)).fetchall()
    return J({"messages":[dict(r) for r in reversed(rows)]})
@app.delete("/api/history")
async def clear_history(request:Request):
    u=auth_user(request)
    with db() as c: c.execute("DELETE FROM messages WHERE user_id=?",(u["id"],))
    return J({"ok":True})
@app.get("/api/conversations")
async def list_convs(request:Request):
    u=auth_user(request)
    with db() as c:
        rows=c.execute("SELECT id,title,created,updated FROM conversations WHERE user_id=? ORDER BY updated DESC LIMIT 50",(u["id"],)).fetchall()
    return J({"conversations":[dict(r) for r in rows]})
@app.post("/api/conversations")
async def new_conv(request:Request):
    u=auth_user(request); d=await request.json(); title=d.get("title","New Chat").strip()[:200]; now=_now()
    with db() as c:
        c.execute("INSERT INTO conversations(user_id,title,created,updated)VALUES(?,?,?,?)",(u["id"],title,now,now))
        cid=c.execute("SELECT last_insert_rowid()").fetchone()[0]
    return J({"ok":True,"id":cid,"title":title})
@app.get("/api/conversations/{cid}")
async def get_conv(cid:int,request:Request):
    u=auth_user(request)
    with db() as c:
        conv=c.execute("SELECT * FROM conversations WHERE id=? AND user_id=?",(cid,u["id"])).fetchone()
        if not conv: return err("Not found",404)
        msgs=c.execute("SELECT role,content,model,ts FROM conv_messages WHERE conv_id=? ORDER BY id",(cid,)).fetchall()
    return J({"id":cid,"title":conv["title"],"messages":[dict(m) for m in msgs]})
@app.delete("/api/conversations/{cid}")
async def del_conv(cid:int,request:Request):
    u=auth_user(request)
    with db() as c:
        c.execute("DELETE FROM conv_messages WHERE conv_id=? AND user_id=?",(cid,u["id"]))
        c.execute("DELETE FROM conversations WHERE id=? AND user_id=?",(cid,u["id"]))
    return J({"ok":True})
@app.post("/api/conversations/{cid}/rename")
async def rename_conv(cid:int,request:Request):
    u=auth_user(request); d=await request.json(); title=d.get("title","").strip()[:200]
    if not title: return err("title required")
    with db() as c: c.execute("UPDATE conversations SET title=? WHERE id=? AND user_id=?",(title,cid,u["id"]))
    return J({"ok":True})
@app.post("/api/conversations/{cid}/message")
async def save_conv_msg(cid:int,request:Request):
    u=auth_user(request); d=await request.json()
    role=d.get("role","user"); content=d.get("content","").strip(); model=d.get("model","")
    if not content: return err("content required"); now=_now()
    with db() as c:
        conv=c.execute("SELECT id FROM conversations WHERE id=? AND user_id=?",(cid,u["id"])).fetchone()
        if not conv: return err("Not found",404)
        now=datetime.now().isoformat()
        c.execute("INSERT INTO conv_messages(conv_id,user_id,role,content,model,ts)VALUES(?,?,?,?,?,?)",(cid,u["id"],role,content,model,now))
        cnt=c.execute("SELECT COUNT(*) as n FROM conv_messages WHERE conv_id=?",(cid,)).fetchone()["n"]
        if cnt<=2 and role=="user":
            t2=content[:60].strip().replace("\n"," ")
            c.execute("UPDATE conversations SET title=?,updated=? WHERE id=?",(t2,now,cid))
        else: c.execute("UPDATE conversations SET updated=? WHERE id=?",(now,cid))
    return J({"ok":True})
@app.get("/api/saved")
async def get_saved(request:Request):
    u=auth_user(request)
    with db() as c: rows=c.execute("SELECT id,title,content,ts FROM saved_items WHERE user_id=? ORDER BY id DESC",(u["id"],)).fetchall()
    return J({"items":[dict(r) for r in rows]})
@app.post("/api/saved")
async def add_saved(request:Request):
    u=auth_user(request); d=await request.json()
    title=d.get("title","").strip()[:200]; content=d.get("content","").strip()
    if not title or not content: return err("title and content required")
    with db() as c:
        c.execute("INSERT INTO saved_items(user_id,title,content,ts)VALUES(?,?,?,?)",(u["id"],title,content,_now()))
        iid=c.execute("SELECT last_insert_rowid()").fetchone()[0]
    return J({"ok":True,"id":iid})
@app.delete("/api/saved/{iid}")
async def del_saved(iid:int,request:Request):
    u=auth_user(request)
    with db() as c: c.execute("DELETE FROM saved_items WHERE id=? AND user_id=?",(iid,u["id"]))
    return J({"ok":True})
@app.get("/api/memory")
async def get_memory(request:Request):
    u=auth_user(request)
    with db() as c: rows=c.execute("SELECT key,value,ts FROM memory WHERE user_id=? ORDER BY ts DESC",(u["id"],)).fetchall()
    return J({"memory":[dict(r) for r in rows]})
@app.post("/api/memory")
async def set_memory(request:Request):
    u=auth_user(request); d=await request.json()
    key=d.get("key","").strip()[:100]; value=d.get("value","").strip()
    if not key: return err("key required")
    with db() as c:
        if not value: c.execute("DELETE FROM memory WHERE user_id=? AND key=?",(u["id"],key))
        else: c.execute("INSERT OR REPLACE INTO memory(user_id,key,value,ts)VALUES(?,?,?,?)",(u["id"],key,value,_now()))
    return J({"ok":True})
@app.delete("/api/memory/{key}")
async def del_memory(key:str,request:Request):
    u=auth_user(request)
    with db() as c: c.execute("DELETE FROM memory WHERE user_id=? AND key=?",(u["id"],key))
    return J({"ok":True})
@app.post("/api/detect-intent")
async def detect_intent(request:Request):
    auth_user(request); d=await request.json(); text=d.get("text","").strip(); tl=text.lower()
    if tl.startswith("/imagine "): return J({"intent":"image","prompt":text[9:].strip()})
    if tl.startswith("/img "): return J({"intent":"image","prompt":text[5:].strip()})
    if tl.startswith("/video "): return J({"intent":"video","prompt":text[7:].strip()})
    if tl.startswith("/3d "): return J({"intent":"video3d","prompt":text[4:].strip()})
    intent="chat"
    for kw in ["generate a 3d","create a 3d","make a 3d","3d model"]:
        if kw in tl: intent="video3d"; break
    for kw in ["generate a video","create a video","make a video"]:
        if kw in tl and intent=="chat": intent="video"; break
    # Only unambiguous image-creation phrases — NOT "draw a conclusion", "picture this" etc
    IMAGE_EXACT=[
        "generate an image","generate a photo","generate a picture",
        "create an image","create a photo","create a picture","create artwork",
        "make an image","make a photo","make a picture",
        "draw me a","draw me an","paint me a","paint me an",
        "sketch me a","sketch me an","illustrate a ","illustrate an ",
        "a picture of ","a photo of ","a painting of ","a drawing of ",
        "a sketch of ","an illustration of ",
        "show me a photo of","show me an image of","show me a picture of",
        "design a logo","design an icon","logo for my","icon for my",
        "photorealistic image of","digital art of","anime art of","pixel art of",
        "8k image of","wallpaper of ","generate art of ",
    ]
    for kw in IMAGE_EXACT:
        if kw in tl and intent=="chat": intent="image"; break
    # Override back to chat if explanation words present
    if intent=="image":
        if any(w in tl for w in ["explain","describe","what is","how does","tell me","why","when","who","code","script","function","bug","fix","error","calculate","solve"]): intent="chat"
    return J({"intent":intent,"prompt":text})
# ── Mega Free API Router ──────────────────────────────────────────────────────
_FREE_APIS = {
  "wiki":     lambda q: f"https://en.wikipedia.org/api/rest_v1/page/summary/{q.replace(' ','_')}",
  "ddg":      lambda q: f"http://openserp.alwaysdata.net/bing/search?text={q}",
  "dict":     lambda q: f"https://api.dictionaryapi.dev/api/v2/entries/en/{q}",
  "country":  lambda q: f"https://restcountries.com/v3.1/name/{q}?fullText=true",
  "uni":      lambda q: f"http://universities.hipolabs.com/search?name={q}",
  "books":    lambda q: f"https://openlibrary.org/search.json?q={q}&limit=3",
  "nasa":     lambda q: f"https://api.nasa.gov/planetary/apod?api_key=DEMO_KEY",
  "pokemon":  lambda q: f"https://pokeapi.co/api/v2/pokemon/{q.lower()}",
  "weather":  lambda q: f"https://api.open-meteo.com/v1/forecast?latitude={q.split(',')[0]}&longitude={q.split(',')[1]}&current_weather=true" if ',' in q else f"https://wttr.in/{q}?format=j1",
  "quake":    lambda q: "https://earthquake.usgs.gov/earthquakes/feed/v1.0/summary/significant_week.geojson",
  "sun":      lambda q: f"https://api.sunrise-sunset.org/json?lat={q.split(',')[0]}&lng={q.split(',')[1]}" if ',' in q else None,
  "airq":     lambda q: f"https://api.waqi.info/feed/{q}/?token=demo",
  "btc":      lambda q: "https://api.coindesk.com/v1/bpi/currentprice.json",
  "crypto":   lambda q: f"https://api.coingecko.com/api/v3/simple/price?ids={q}&vs_currencies=usd",
  "fx":       lambda q: "https://open.er-api.com/v6/latest/USD",
  "gold":     lambda q: "https://www.gold-api.com/price/XAU",
  "ip":       lambda q: f"http://ip-api.com/json/{q}" if q and q!="self" else "http://ip-api.com/json/",
  "zip":      lambda q: f"http://api.zippopotam.us/{q.split('/')[0]}/{q.split('/')[1]}" if '/' in q else None,
  "time":     lambda q: f"http://worldtimeapi.org/api/timezone/{q}" if q else "http://worldtimeapi.org/api/ip",
  "tv":       lambda q: f"https://api.tvmaze.com/search/shows?q={q}",
  "starwars": lambda q: f"https://swapi.dev/api/{q}/1/",
  "ghibli":   lambda q: "https://ghibliapi.vercel.app/films",
  "harrypotter": lambda q: "https://hp-api.onrender.com/api/characters",
  "space":    lambda q: "https://api.spaceflightnewsapi.net/v4/articles/?limit=5",
  "joke":     lambda q: "https://official-joke-api.appspot.com/jokes/random",
  "advice":   lambda q: "https://api.adviceslip.com/advice",
  "bored":    lambda q: "https://www.boredapi.com/api/activity",
  "gender":   lambda q: f"https://api.genderize.io?name={q}",
  "age":      lambda q: f"https://api.agify.io?name={q}",
  "nation":   lambda q: f"https://api.nationalize.io?name={q}",
  "number":   lambda q: f"http://numbersapi.com/{q}",
  "chuck":    lambda q: "https://api.chucknorris.io/jokes/random",
  "cat":      lambda q: "https://catfact.ninja/fact",
  "dog":      lambda q: "https://dog.ceo/api/breeds/image/random",
  "github":   lambda q: f"https://api.github.com/users/{q}",
  "gitrepos": lambda q: f"https://api.github.com/users/{q}/repos?per_page=5",
  "food":     lambda q: f"https://world.openfoodfacts.org/api/v0/product/{q}.json",
  "user":     lambda q: "https://randomuser.me/api/",
  "coffee":   lambda q: "https://fake-coffee-api.vercel.app/api?limit=1",
  "qr":       lambda q: f"https://api.qrserver.com/v1/create-qr-code/?data={q}&size=200x200",
  "population": lambda q: f"https://datausa.io/api/data?drilldowns=Nation&measures=Population&year=latest",
}

def _parse_free_api(text):
    """Detect if a message can be answered by a free API. Returns (api_key, query) or (None,None)."""
    tl=text.lower().strip()
    import re as _re, urllib.parse as _up
    # Wikipedia / definition / country / etc
    if _re.search(r"\bwhat is\b|\bwho is\b|\btell me about\b|\bdefine\b|\bexplain\b",tl):
        topic=_re.sub(r"(what is|who is|tell me about|define|explain)\s+","",tl).strip().rstrip("?.")
        if topic and len(topic.split())<6: return "wiki",_up.quote(topic)
    if "dictionary" in tl or "meaning of" in tl or "definition of" in tl:
        w=tl.replace("dictionary","").replace("meaning of","").replace("definition of","").strip().split()[0] if tl.replace("dictionary","").strip() else ""
        if w: return "dict",w
    if "country" in tl or "nation" in tl:
        m=_re.search(r"(country|nation|capital|currency|language).*?(?:of|called)?\s+([\w ]+)",tl)
        if m: return "country",_up.quote(m.group(2).strip())
    if "weather" in tl or "forecast" in tl or "temperature" in tl:
        city=_re.sub(r"(weather|forecast|temperature|in|at|for|today|tomorrow|°|degrees?)","",tl).strip()
        if city: return "weather",_up.quote(city)
    if "earthquake" in tl or "seismic" in tl: return "quake",""
    if "bitcoin" in tl or "btc price" in tl: return "btc",""
    if tl.startswith("price of ") or ("crypto" in tl and "price" in tl):
        coin=tl.replace("price of","").replace("crypto","").replace("price","").strip().split()[0]
        if coin: return "crypto",_up.quote(coin)
    if "exchange rate" in tl or "forex" in tl or "currency rate" in tl: return "fx",""
    if "gold price" in tl or "gold spot" in tl: return "gold",""
    if _re.search(r"\bpokemon\b",tl):
        m=_re.search(r"pokemon\s+([\w]+)",tl)
        if m: return "pokemon",m.group(1)
    if "star wars" in tl: return "starwars","people"
    if "ghibli" in tl or "studio ghibli" in tl: return "ghibli",""
    if "harry potter" in tl: return "harrypotter",""
    if "space news" in tl or "spaceflight" in tl or "rocket launch" in tl: return "space",""
    if "joke" in tl or "make me laugh" in tl or "funny" in tl: return "joke",""
    if "advice" in tl or "give me advice" in tl: return "advice",""
    if "bored" in tl or "activity" in tl or "something to do" in tl: return "bored",""
    if "chuck norris" in tl: return "chuck",""
    if "cat fact" in tl or "tell me about cats" in tl: return "cat",""
    if "dog" in tl and ("image" in tl or "picture" in tl or "photo" in tl): return "dog",""
    if tl.startswith("github ") or "github profile" in tl:
        u=tl.replace("github","").replace("profile","").strip().split()[0]
        if u: return "github",u
    if tl.startswith("what time is it") or "current time in" in tl:
        tz=_re.sub(r"(what time is it|current time in|time in)","",tl).strip().replace(" ","/")
        return "time",tz or ""
    if "my ip" in tl or "what is my ip" in tl: return "ip","self"
    if "number fact" in tl or "trivia about" in tl:
        m=_re.search(r"\d+",tl)
        if m: return "number",m.group(0)
    return None,None

def _call_free_api(api_key, query):
    """Call a free API and return a formatted response string."""
    import re as _re, json as _json, urllib.parse as _up, traceback as _tb
    fn=_FREE_APIS.get(api_key)
    if not fn: return None
    url=fn(query)
    if not url: return None
    try:
        r=req.get(url,timeout=8,headers={"User-Agent":"FusionAI/1.0"})
        if not r.ok: return None
        ct=r.headers.get("Content-Type","")
        if "image" in ct or url.endswith(".jpg") or url.endswith(".png"):
            return f"![Image]({url})"
        try: d=r.json()
        except: return r.text[:300] if r.text else None
        # Format response per API
        if api_key=="wiki":
            ex=d.get("extract",""); title=d.get("title","")
            thumb=d.get("thumbnail",{}).get("source","")
            out=f"**{title}**\n{ex[:600]}"
            if thumb: out+=f"\n![{title}]({thumb})"
            return out if ex else None
        if api_key=="dict":
            if isinstance(d,list) and d:
                entry=d[0]; word=entry.get("word",""); meanings=entry.get("meanings",[])
                out=f"**{word}**"
                for m in meanings[:2]:
                    pos=m.get("partOfSpeech",""); defs=m.get("definitions",[])
                    if defs: out+=f"\n*{pos}*: {defs[0].get('definition','')}"; ex=defs[0].get("example")
                    if ex: out+=f"\n> {ex}"
                return out
        if api_key=="country":
            if isinstance(d,list) and d:
                c=d[0]; n=c.get("name",{}).get("common",""); cap=(c.get("capital") or ["?"])[0]
                pop=c.get("population",0); cur=list((c.get("currencies") or {}).keys())
                langs=list((c.get("languages") or {}).values())
                flag=c.get("flags",{}).get("png","")
                out=f"**{n}** {'🏁'}\nCapital: {cap} | Population: {pop:,}\nCurrencies: {', '.join(cur)} | Languages: {', '.join(langs[:3])}"
                if flag: out+=f"\n![]({flag})"
                return out
        if api_key=="weather":
            if d.get("current_weather"):
                cw=d["current_weather"]; t=cw.get("temperature","?"); ws=cw.get("windspeed","?"); codes={0:"☀️ Clear",1:"🌤 Mostly clear",2:"⛅ Partly cloudy",3:"☁️ Overcast",61:"🌧 Rain",80:"🌦 Showers",95:"⛈ Thunderstorm"}
                code=cw.get("weathercode",0); desc=codes.get(code,f"Code {code}")
                return f"🌡️ **{t}°C** · {desc}\n💨 Wind: {ws} km/h"
            elif d.get("current_condition"):
                cc=d["current_condition"][0]; area=(d.get("nearest_area") or [{}])[0]
                city=(area.get("areaName") or [{}])[0].get("value","")
                return f"📍 **{city}** 🌡️ {cc.get('temp_C')}°C (feels {cc.get('FeelsLikeC')}°C)\n{cc.get('weatherDesc',[{}])[0].get('value','')} · 💧{cc.get('humidity')}%"
        if api_key=="quake":
            feats=d.get("features",[])[:3]
            if feats:
                out="🌍 **Recent Significant Earthquakes:**"
                for f in feats:
                    p=f.get("properties",{}); out+=f"\n• M{p.get('mag','?')} — {p.get('place','?')}"
                return out
        if api_key=="btc":
            bpi=d.get("bpi",{}).get("USD",{}); rate=bpi.get("rate","?")
            return f"₿ **Bitcoin**: ${rate} USD"
        if api_key=="crypto":
            for coin,prices in d.items():
                return f"💰 **{coin.title()}**: ${prices.get('usd','?')} USD"
        if api_key=="fx":
            rates=d.get("rates",{}); out="💱 **Exchange Rates (vs USD):**\n"
            for k,v in list(rates.items())[:8]: out+=f"{k}: {v}  "
            return out
        if api_key=="gold":
            return f"🥇 **Gold**: ${d.get('price','?')} / oz ({d.get('curr','XAU')})"
        if api_key=="pokemon":
            name=d.get("name","").title(); types=[t["type"]["name"] for t in d.get("types",[])]
            stats={s["stat"]["name"]:s["base_stat"] for s in d.get("stats",[][:6])}
            sprite=d.get("sprites",{}).get("front_default","")
            out=f"**{name}** — {' / '.join(types).title()}\nHP:{stats.get('hp','?')} ATK:{stats.get('attack','?')} DEF:{stats.get('defense','?')}"
            if sprite: out+=f"\n![]({sprite})"
            return out
        if api_key in ("joke","chuck"):
            j=d.get("joke") or d.get("value","") or (d.get("setup","")+" "+d.get("punchline",""))
            return f"😄 {j.strip()}"
        if api_key=="advice":
            return f"💡 {d.get('slip',{}).get('advice','')}"
        if api_key=="bored":
            return f"🎯 Try this: **{d.get('activity','')}** ({d.get('type','')})"
        if api_key=="cat":
            return f"🐱 Cat fact: {d.get('fact','')}"
        if api_key=="dog":
            url2=d.get("message","")
            return f"![🐶]({url2})" if url2 else None
        if api_key=="github":
            name=d.get("name") or d.get("login",""); bio=d.get("bio",""); repos=d.get("public_repos",0)
            followers=d.get("followers",0); avatar=d.get("avatar_url","")
            out=f"**{name}** (@{d.get('login','')})\n{bio}\n⭐ {repos} repos · 👥 {followers} followers"
            if avatar: out+=f"\n![]({avatar})"
            return out
        if api_key=="ghibli":
            if isinstance(d,list) and d:
                out="🎬 **Studio Ghibli Films:**"
                for f in d[:5]: out+=f"\n• **{f.get('title','')}** ({f.get('release_date','')[:4]}) — {f.get('description','')[:80]}…"
                return out
        if api_key=="harrypotter":
            if isinstance(d,list) and d:
                out="🧙 **Harry Potter Characters:**"
                for c in d[:6]: out+=f"\n• {c.get('name','')} — {c.get('house','?')}"
                return out
        if api_key=="space":
            out="🚀 **Latest Space News:**"
            for a in (d.get("results") or [])[:4]: out+=f"\n• {a.get('title','')}"
            return out if d.get("results") else None
        if api_key=="time":
            dt=d.get("datetime",""); tz=d.get("timezone","")
            return f"🕐 **{tz}**: {dt[:16].replace('T',' ')}"
        if api_key=="ip":
            return f"🌐 IP: {d.get('query','?')}\n📍 {d.get('city','')}, {d.get('country','')} ({d.get('isp','')})"
        if api_key=="user":
            res=(d.get("results") or [{}])[0]; n=res.get("name",{}); loc=res.get("location",{})
            return f"👤 Random user: {n.get('first','')} {n.get('last','')} from {loc.get('city','?')}, {loc.get('country','?')}"
        if api_key=="number":
            return f"🔢 {r.text}" if isinstance(r.text,str) and len(r.text)<200 else None
        if api_key=="starwars":
            items=d.get("results",[])
            if items: return f"⚔️ {items[0].get('name','')} — {items[0].get('birth_year','?')}, {items[0].get('homeworld','?')}"
        # Generic JSON fallback
        return _json.dumps(d,indent=2)[:500]
    except Exception as e:
        return None

@app.get("/api/freeapi/joke")
async def freeapi_joke(request:Request):
    auth_user(request)
    try:
        r=req.get("https://api.freeapi.app/api/v1/public/randomjokes?limit=1",timeout=8)
        if r.ok:
            d=r.json(); items=d.get("data",{}).get("data",[]) or d.get("data",[])
            if items: joke=items[0]; return J({"ok":True,"text":(joke.get("setup","")+" "+joke.get("punchline","")).strip() or joke.get("content","")})
    except: pass
    result=_call_free_api("joke","")
    return J({"ok":bool(result),"text":result or ""})

@app.get("/api/freeapi/quote")
async def freeapi_quote(request:Request):
    auth_user(request)
    try:
        r=req.get("https://api.freeapi.app/api/v1/public/quotes/quote/random",timeout=8)
        if r.ok:
            d=r.json(); q=d.get("data",{})
            return J({"ok":True,"text":f'"{q.get("content","")}" — {q.get("author","Unknown")}'})
    except: pass
    return J({"ok":False,"text":""})

@app.get("/api/freeapi/fact")
async def freeapi_fact(request:Request):
    auth_user(request)
    try:
        r=req.get("https://uselessfacts.jsph.pl/api/v2/facts/random?language=en",timeout=8)
        if r.ok: d=r.json(); return J({"ok":True,"text":d.get("text","")})
    except: pass
    return J({"ok":False,"text":""})

@app.post("/api/freeapi/query")
async def freeapi_query(request:Request):
    """Smart free-API router — called by frontend to answer without burning LLM tokens."""
    auth_user(request); d=await request.json()
    text=d.get("text","").strip()
    api_key,query=_parse_free_api(text)
    if not api_key: return J({"ok":False,"result":""})
    result=_call_free_api(api_key,query)
    return J({"ok":bool(result),"result":result or "","api":api_key})

# ── OpenSERP Bing Web Search (JSON, no API key needed) ────────────────────────
@app.post("/api/search/web")
async def web_search_proxy(request:Request):
    """Bing search via OpenSERP — returns clean JSON results."""
    import urllib.parse
    auth_user(request); d=await request.json()
    query=d.get("query","").strip()
    if not query: return J({"ok":False,"error":"query required","results":[]})
    try:
        url=f"http://openserp.alwaysdata.net/bing/search?text={urllib.parse.quote(query)}"
        r=req.get(url,headers={"User-Agent":"FusionAI/2.0","Accept":"application/json"},timeout=10)
        if not r.ok:
            return J({"ok":False,"error":f"Search API returned {r.status_code}","results":[]})
        raw=r.json()  # list of {title, url, description}
        results=[{"title":item.get("title",""),"url":item.get("url",""),"snippet":item.get("description","")} for item in (raw if isinstance(raw,list) else []) if item.get("url")]
        return J({"ok":True,"query":query,"instant":"","results":results[:10]})
    except Exception as e:
        return J({"ok":False,"error":str(e),"results":[]})

# Legacy alias — keeps old callers working
@app.post("/api/search/ddg")
async def ddg_search_alias(request:Request):
    return await web_search_proxy(request)


# ── Google OAuth ──────────────────────────────────────────────────────────────
@app.get("/api/auth/google")
async def google_auth_start(request:Request):
    if not GOOGLE_CLIENT_ID.strip(): return J({"ok":False,"error":"Google OAuth not configured"})
    import urllib.parse as _up
    state=uuid.uuid4().hex
    with db() as c: c.execute("INSERT OR REPLACE INTO oauth_state(state,created)VALUES(?,?)",(state,_now()))
    params={"client_id":GOOGLE_CLIENT_ID,"redirect_uri":str(request.base_url).rstrip("/")+"/api/auth/google/callback",
            "response_type":"code","scope":"openid email profile","state":state,"access_type":"offline","prompt":"select_account"}
    url="https://accounts.google.com/o/oauth2/v2/auth?"+_up.urlencode(params)
    return J({"ok":True,"url":url})

@app.get("/api/auth/google/callback")
async def google_auth_callback(request:Request):
    from fastapi.responses import RedirectResponse
    code=request.query_params.get("code",""); state=request.query_params.get("state","")
    if not code: return RedirectResponse("/?auth_error=no_code")
    with db() as c:
        row=c.execute("SELECT state FROM oauth_state WHERE state=?",(state,)).fetchone()
        if row: c.execute("DELETE FROM oauth_state WHERE state=?",(state,))
    try:
        r=req.post("https://oauth2.googleapis.com/token",data={"code":code,"client_id":GOOGLE_CLIENT_ID,"client_secret":GOOGLE_CLIENT_SECRET,
            "redirect_uri":str(request.base_url).rstrip("/")+"/api/auth/google/callback","grant_type":"authorization_code"},timeout=10)
        if not r.ok: return RedirectResponse("/?auth_error=token_fail")
        tokens=r.json(); access_token=tokens.get("access_token","")
        prof=req.get("https://www.googleapis.com/oauth2/v3/userinfo",headers={"Authorization":f"Bearer {access_token}"},timeout=8)
        if not prof.ok: return RedirectResponse("/?auth_error=profile_fail")
        p=prof.json(); gid=p.get("sub",""); email=p.get("email",""); name=p.get("name",email.split("@")[0]); avatar=p.get("picture","")
        with db() as c:
            existing=c.execute("SELECT id FROM users WHERE google_id=?",(gid,)).fetchone()
            if existing: uid=existing["id"]; c.execute("UPDATE users SET avatar_url=?,display_name=? WHERE id=?",(avatar,name,uid))
            else:
                uname=f"g_{email.split('@')[0]}_{uuid.uuid4().hex[:6]}"; salt=uuid.uuid4().hex; pw=uuid.uuid4().hex
                c.execute("INSERT OR IGNORE INTO users(username,password,salt,created,google_id,avatar_url,display_name)VALUES(?,?,?,?,?,?,?)",(uname,hash_pw(pw,salt),salt,_now(),gid,avatar,name))
                uid=c.execute("SELECT id FROM users WHERE google_id=?",(gid,)).fetchone()["id"]
        tok=make_token(uid); log_event(uid,name,"google_login",request)
        return RedirectResponse(f"/?g_token={tok}&g_name={_up.quote(name)}&g_avatar={_up.quote(avatar)}")
    except Exception as e:
        import traceback; traceback.print_exc()
        return RedirectResponse(f"/?auth_error={str(e)[:60]}")





# ── Weather & GeoIP (free) ────────────────────────────────────────────────────
@app.get("/api/weather")
async def get_weather(request:Request):
    auth_user(request)
    lat=request.query_params.get("lat","") or request.headers.get("x-user-lat","")
    lon=request.query_params.get("lon","") or request.headers.get("x-user-lon","")
    city=request.query_params.get("city","") or request.headers.get("x-user-city","")
    # Use real client IP for auto-detection if no lat/lon/city provided
    if not lat and not lon and not city:
        client_ip=(request.headers.get("x-forwarded-for","") or "").split(",")[0].strip() or \
                   request.headers.get("x-real-ip","") or request.headers.get("cf-connecting-ip","") or ""
        try:
            import ipaddress as _ia; _pip=_ia.ip_address(client_ip)
            if _pip.is_private or _pip.is_loopback: client_ip=""
        except: pass
        if client_ip:
            try:
                gr=req.get(f"https://ip-api.com/json/{client_ip}",timeout=5)
                if gr.ok:
                    gd=gr.json(); lat=str(gd.get("lat","")); lon=str(gd.get("lon","")); city=gd.get("city","")
            except: pass
    try:
        if OPENWEATHER_KEY.strip():
            q=f"lat={lat}&lon={lon}" if lat and lon else f"q={city or 'London'}"
            r=req.get(f"https://api.openweathermap.org/data/2.5/weather?{q}&appid={OPENWEATHER_KEY}&units=metric",timeout=8)
            if r.ok:
                d=r.json(); return J({"ok":True,"city":d.get("name",""),"temp":round(d["main"]["temp"]),"feels_like":round(d["main"]["feels_like"]),"humidity":d["main"]["humidity"],"desc":d["weather"][0]["description"],"wind_speed":d.get("wind",{}).get("speed",0),"icon":f"https://openweathermap.org/img/wn/{d['weather'][0]['icon']}@2x.png"})
        q2=f"{lat},{lon}" if lat and lon else (city or "auto")
        r2=req.get(f"https://wttr.in/{q2}?format=j1",timeout=8)
        if r2.ok:
            d=r2.json(); cur=d["current_condition"][0]; area=(d.get("nearest_area") or [{}])[0]
            city_name=(area.get("areaName") or [{}])[0].get("value","")
            return J({"ok":True,"city":city_name,"temp":int(cur["temp_C"]),"feels_like":int(cur["FeelsLikeC"]),"humidity":int(cur["humidity"]),"desc":cur["weatherDesc"][0]["value"],"wind_speed":round(float(cur.get("windspeedKmph",0))/3.6,1),"icon":""})
    except: pass
    return J({"ok":False,"error":"Weather unavailable"})

@app.get("/api/geo/ip")
async def geo_ip(request:Request):
    auth_user(request)
    # Get real client IP from forwarded headers (not server IP)
    client_ip = (request.headers.get("x-forwarded-for","") or "").split(",")[0].strip() or \
                request.headers.get("x-real-ip","") or \
                request.headers.get("cf-connecting-ip","") or \
                (request.client.host if request.client else "")
    # Strip localhost/private IPs — fall back to blank so API auto-detects
    import ipaddress
    try:
        _ip=ipaddress.ip_address(client_ip)
        if _ip.is_private or _ip.is_loopback: client_ip=""
    except: pass
    apis=[
        (f"https://ipapi.co/{client_ip}/json/" if client_ip else "https://ipapi.co/json/",
         lambda d: {"city":d.get("city",""),"country":d.get("country_name",""),"lat":d.get("latitude",""),"lon":d.get("longitude",""),"timezone":d.get("timezone",""),"ip":d.get("ip","")}),
        (f"https://ip-api.com/json/{client_ip}" if client_ip else "https://ip-api.com/json/",
         lambda d: {"city":d.get("city",""),"country":d.get("country",""),"lat":d.get("lat",""),"lon":d.get("lon",""),"timezone":d.get("timezone",""),"ip":d.get("query","")}),
        (f"https://ipinfo.io/{client_ip}/json" if client_ip else "https://ipinfo.io/json",
         lambda d: {"city":d.get("city",""),"country":d.get("country",""),"lat":(d.get("loc","0,0").split(",")+[""])[0],"lon":(d.get("loc","0,0").split(",")+["",""])[1],"timezone":d.get("timezone",""),"ip":d.get("ip","")}),
    ]
    for url,parser in apis:
        try:
            r=req.get(url,timeout=6,headers={"Accept":"application/json"})
            if r.ok:
                d=r.json(); parsed=parser(d)
                if parsed.get("city"): return J({"ok":True,**parsed,"client_ip":client_ip})
        except: continue
    return J({"ok":False})

# ── Real 3D Model generation ──────────────────────────────────────────────────
@app.post("/api/generate/3d/model")
async def gen_3d_model(request:Request):
    import base64 as _b64, urllib.parse as _up
    auth_user(request); d=await request.json(); prompt=d.get("prompt","").strip()
    if not prompt: return err("Prompt required")
    if HF_TOKEN.strip():
        try:
            r=req.post("https://api-inference.huggingface.co/models/openai/shap-e",headers={"Authorization":f"Bearer {HF_TOKEN}","Content-Type":"application/json"},json={"inputs":prompt},timeout=120)
            if r.ok and len(r.content)>1000:
                ct=r.headers.get("Content-Type","application/octet-stream")
                return J({"ok":True,"b64":_b64.b64encode(r.content).decode(),"mime":ct.split(";")[0],"format":"glb","backend":"HuggingFace · Shap-E"})
        except: pass
    enhanced=f"{prompt}, 3D model render, turntable view, studio lighting, octane render, high detail mesh, 8k"
    url=f"https://image.pollinations.ai/prompt/{_up.quote(enhanced,safe='')}?model=flux&width=1024&height=1024&nologo=true&enhance=true"
    try:
        r3=req.get(url,timeout=90,allow_redirects=True)
        if r3.ok and r3.headers.get("Content-Type","").startswith("image"):
            return J({"ok":True,"b64":_b64.b64encode(r3.content).decode(),"mime":"image/png","format":"render","backend":"AI 3D Engine","enhanced_prompt":enhanced})
    except: pass
    return J({"ok":False,"error":"3D generation unavailable"},502)

@app.post("/api/overseer/enhance")
async def overseer_enhance(request:Request):
    u=auth_user(request); d=await request.json()
    prompt=d.get("prompt","").strip(); mode=d.get("mode","image")
    answers=d.get("answers",{}); wc=len(prompt.split())
    if answers or wc>=12:
        ans_text="\n".join(f"- {q}: {a}" for q,a in answers.items()) if answers else ""
        if mode=="chat": sys_p="Expand this short message into a clear specific request. Same intent. Max 80 words. Return ONLY the expanded prompt."; up=f"Message: {prompt}"
        elif mode=="image": sys_p="Write ONE highly detailed image generation prompt. No preamble. Max 120 words."; up=f"Base: {prompt}\nPrefs:\n{ans_text}" if ans_text else f"Base: {prompt}"
        elif mode=="video3d": sys_p="Write ONE detailed 3D object description. Include materials, style, colors. Max 100 words."; up=f"Base: {prompt}\nPrefs:\n{ans_text}" if ans_text else f"Base: {prompt}"
        else: sys_p="Write ONE cinematic video prompt. Include action, camera, lighting, mood. Max 100 words."; up=f"Base: {prompt}\nPrefs:\n{ans_text}" if ans_text else f"Base: {prompt}"
        enhanced=_call_overseer(sys_p,up,200,cheap=True) or prompt
        return J({"ok":True,"action":"generate","enhanced_prompt":enhanced})
    if mode=="chat":
        enhanced=_call_overseer("Expand this short message. Same intent. Max 80 words. Return ONLY the prompt.",f"Message: {prompt}",150) or prompt
        return J({"ok":True,"action":"generate","enhanced_prompt":enhanced})
    Q={
        "image":'Ask 2 quick questions. Reply ONLY valid JSON:[{"q":"Style?","options":["Photorealistic","Digital art","Anime","Cinematic","Sketch"]},{"q":"Aspect ratio?","options":["Square 1:1","Landscape 16:9","Portrait 9:16"]}]',
        "video3d":'Ask 2 quick questions. Reply ONLY valid JSON:[{"q":"Style?","options":["Realistic","Stylized","Low-poly","Cartoon"]},{"q":"Detail?","options":["Simple","Medium","Highly detailed"]}]',
        "video":'Ask 2 quick questions. Reply ONLY valid JSON:[{"q":"Camera style?","options":["Cinematic","Drone shot","Close-up","Wide angle"]},{"q":"Mood?","options":["Dramatic","Calm","Energetic","Mysterious"]}]',
    }
    raw=_call_overseer(Q.get(mode,Q["image"]),f'Prompt: "{prompt}"',350)
    if not raw: return J({"ok":True,"action":"generate","enhanced_prompt":prompt})
    try:
        qs=json.loads(raw.strip().lstrip("```json").lstrip("```").rstrip("```").strip())
        if isinstance(qs,list) and qs: return J({"ok":True,"action":"ask","questions":qs[:3]})
    except: pass
    return J({"ok":True,"action":"generate","enhanced_prompt":prompt})

_CF_IMG_MODELS=["@cf/black-forest-labs/flux-2-klein-4b","@cf/black-forest-labs/flux-2-klein-9b",
                "@cf/black-forest-labs/flux-2-dev","@cf/leonardo/phoenix-1.0","@cf/deepgram/flux"]
_POL_TO_CF={"flux":"@cf/black-forest-labs/flux-2-klein-4b","flux-realism":"@cf/black-forest-labs/flux-2-dev",
            "flux-anime":"@cf/black-forest-labs/flux-2-klein-9b","flux-3d":"@cf/black-forest-labs/flux-2-dev",
            "turbo":"@cf/black-forest-labs/flux-2-klein-4b","img_flux":"@cf/black-forest-labs/flux-2-klein-4b",
            "img_flux_realism":"@cf/black-forest-labs/flux-2-dev","img_flux_anime":"@cf/black-forest-labs/flux-2-klein-9b",
            "img_flux_3d":"@cf/black-forest-labs/flux-2-dev","img_turbo":"@cf/black-forest-labs/flux-2-klein-4b"}
@app.post("/api/generate/image")
async def gen_image(request:Request):
    import base64 as _b64,traceback as _tb
    auth_user(request); d=await request.json()
    prompt=d.get("prompt","").strip(); model=d.get("model","")
    if not prompt: return err("Prompt required")
    if not (CF_ACCOUNT_ID.strip() and CF_KEY.strip()) and not (CF_ACCOUNT_ID2.strip() and CF_KEY2.strip()): return err("CF_ACCOUNT_ID and CF_KEY required for image generation",400)
    if model and model.startswith("@cf/"): try_models=[model]+[m for m in _CF_IMG_MODELS if m!=model]
    elif model in _POL_TO_CF: cf=_POL_TO_CF[model]; try_models=[cf]+[m for m in _CF_IMG_MODELS if m!=cf]
    else: try_models=_CF_IMG_MODELS
    last_err="Unknown error"
    for cf_model in try_models:
        try:
            _acc=cf_acc_for(cf_model); _key=cf_key_for(cf_model)
            if not _acc.strip() or not _key.strip(): continue
            hdrs={"Authorization":f"Bearer {_key}"}
            r=req.post(cf_ep(cf_model,_acc),headers=hdrs,files={"prompt":(None,prompt)},timeout=90)
            ct=r.headers.get("Content-Type","")
            if r.ok:
                if "image" in ct: return J({"ok":True,"b64":_b64.b64encode(r.content).decode(),"mime":ct.split(";")[0],"url":"","backend":f"CF · {cf_model}"})
                try:
                    data=r.json(); img=data.get("result",{}).get("image","") or data.get("result",{}).get("b64_json","")
                    if img: return J({"ok":True,"b64":img,"mime":"image/png","url":"","backend":f"CF · {cf_model}"})
                except: pass
            last_err=f"HTTP {r.status_code}"
        except Exception as e: last_err=str(e); _tb.print_exc()
    return J({"ok":False,"error":"Image generation failed","dev_error":last_err,"b64":"","mime":"image/png","url":""},502)
@app.post("/api/generate/image/cloudflare")
async def gen_image_cf(request:Request):
    import base64 as _b64,traceback as _tb
    auth_user(request); d=await request.json()
    prompt=d.get("prompt","").strip(); model=d.get("model","@cf/black-forest-labs/flux-2-klein-4b")
    if not prompt: return err("Prompt required")
    if not CF_ACCOUNT_ID.strip() or not CF_KEY.strip(): return err("CF keys not set",400)
    hdrs={"Authorization":f"Bearer {CF_KEY}"}
    try:
        r=req.post(cf_ep(model),headers=hdrs,files={"prompt":(None,prompt)},timeout=90)
        ct=r.headers.get("Content-Type","")
        if r.ok:
            if "image" in ct: return J({"ok":True,"b64":_b64.b64encode(r.content).decode(),"mime":ct.split(";")[0],"url":"","backend":f"CF · {model}"})
            try:
                data=r.json(); img=data.get("result",{}).get("image","") or data.get("result",{}).get("b64_json","")
                if img: return J({"ok":True,"b64":img,"mime":"image/png","url":"","backend":f"CF · {model}"})
            except: pass
        return J({"ok":False,"error":f"CF image failed (HTTP {r.status_code})","b64":"","mime":"image/png","url":""},502)
    except Exception as e: _tb.print_exc(); return J({"ok":False,"error":str(e),"b64":"","mime":"image/png","url":""},502)
@app.post("/api/generate/image/github")
async def gen_image_github(request:Request):
    import base64 as _b64,traceback as _tb
    auth_user(request); d=await request.json()
    prompt=d.get("prompt","").strip(); model=d.get("model","gpt-image-1"); size=d.get("size","1024x1024")
    if not prompt: return err("Prompt required")
    if not GITHUB_TOKEN.strip(): return err("GITHUB_TOKEN not set",400)
    hdrs={"Authorization":f"Bearer {GITHUB_TOKEN}","Content-Type":"application/json"}
    try:
        r=req.post("https://models.github.ai/inference/images/generations",headers=hdrs,
                   json={"model":model,"prompt":prompt,"n":1,"size":size,"response_format":"b64_json"},timeout=120)
        if r.ok:
            items=r.json().get("data",[])
            if items and items[0].get("b64_json"):
                return J({"ok":True,"b64":items[0]["b64_json"],"mime":"image/png","url":"","backend":f"GitHub · {model}"})
        return J({"ok":False,"error":f"GitHub image failed (HTTP {r.status_code})","b64":"","mime":"image/png","url":""},502)
    except Exception as e: _tb.print_exc(); return J({"ok":False,"error":str(e),"b64":"","mime":"image/png","url":""},502)
@app.post("/api/generate/image/worker")
async def gen_image_worker(request:Request):
    import base64 as _b64,traceback as _tb
    auth_user(request); d=await request.json()
    prompt=d.get("prompt","").strip(); model_key=d.get("model_key","sdxl"); want_enrich=d.get("want_enrich",True)
    if not prompt: return err("Prompt required")
    wurl=WORKER_URL.strip() or "https://fusionai.pantathagat.workers.dev/"
    hdrs={"Content-Type":"application/json"}
    if WORKER_KEY.strip(): hdrs["Authorization"]=f"Bearer {WORKER_KEY}"
    try:
        r=req.post(wurl,headers=hdrs,json={"prompt":prompt,"model_key":model_key,"want_enrich":want_enrich},timeout=120)
        ct=r.headers.get("Content-Type","")
        if r.ok and "image" in ct:
            return J({"ok":True,"b64":_b64.b64encode(r.content).decode(),"mime":ct.split(";")[0],"url":"","backend":f"CF Worker · {model_key}"})
        return J({"ok":False,"error":f"Worker image failed (HTTP {r.status_code})","b64":"","mime":"image/jpeg","url":""},502)
    except Exception as e: _tb.print_exc(); return J({"ok":False,"error":str(e),"b64":"","mime":"image/jpeg","url":""},502)
@app.post("/api/generate/image/stability")
async def gen_image_stability(request:Request):
    import base64 as _b64,traceback as _tb
    auth_user(request); d=await request.json()
    prompt=d.get("prompt","").strip(); model=d.get("model","sd3-medium")
    if not prompt: return err("Prompt required")
    if not STABILITY_KEY.strip(): return err("STABILITY_KEY not set",400)
    EP={"sd3-medium":"https://api.stability.ai/v2beta/stable-image/generate/sd3",
        "sd3-large":"https://api.stability.ai/v2beta/stable-image/generate/sd3",
        "sd3-large-turbo":"https://api.stability.ai/v2beta/stable-image/generate/sd3",
        "stable-image-core":"https://api.stability.ai/v2beta/stable-image/generate/core",
        "stable-image-ultra":"https://api.stability.ai/v2beta/stable-image/generate/ultra"}
    ep=EP.get(model,"https://api.stability.ai/v2beta/stable-image/generate/core")
    hdrs={"Authorization":f"Bearer {STABILITY_KEY}","Accept":"image/*"}
    files={"prompt":(None,prompt),"output_format":(None,"png")}
    if model.startswith("sd3"): files["model"]=(None,model)
    try:
        r=req.post(ep,headers=hdrs,files=files,timeout=90)
        ct=r.headers.get("Content-Type","")
        if r.ok and r.content:
            mime=ct.split(";")[0] if ct.startswith("image") else "image/png"
            return J({"ok":True,"b64":_b64.b64encode(r.content).decode(),"mime":mime,"url":"","backend":f"Stability · {model}"})
        return J({"ok":False,"error":f"Stability failed (HTTP {r.status_code})","b64":"","mime":"image/png","url":""},502)
    except Exception as e: _tb.print_exc(); return J({"ok":False,"error":str(e),"b64":"","mime":"image/png","url":""},502)

def _pol_get(url,timeout=120): return req.get(url,timeout=timeout,allow_redirects=True)
@app.post("/api/generate/video")
async def gen_video(request:Request):
    import base64 as _b64,traceback as _tb
    auth_user(request); d=await request.json()
    prompt=d.get("prompt","").strip(); model=d.get("model","animatediff")
    w=min(int(d.get("width",512)),1024); h=min(int(d.get("height",512)),1024)
    if not prompt: return err("Prompt required")
    import urllib.parse
    pm={"animatediff":"animatediff","ltx":"ltx","wan":"wan"}.get(model,"animatediff")
    url=f"https://video.pollinations.ai/prompt/{urllib.parse.quote(prompt,safe='')}?model={pm}&width={w}&height={h}&nologo=true"
    try:
        r=_pol_get(url,180); ct=r.headers.get("Content-Type","")
        if r.ok and (ct.startswith("video/") or len(r.content)>10000):
            mime=ct.split(";")[0] if ct.startswith("video/") else "video/mp4"
            return J({"ok":True,"b64":_b64.b64encode(r.content).decode(),"mime":mime,"url":url,"backend":"AI Video Engine"})
        return J({"ok":False,"error":f"Video gen failed (HTTP {r.status_code})","b64":"","mime":"video/mp4","url":url},502)
    except Exception as e: _tb.print_exc(); return J({"ok":False,"error":str(e),"b64":"","mime":"video/mp4","url":url},502)
@app.post("/api/generate/3d")
async def gen_3d(request:Request):
    import base64 as _b64,traceback as _tb,time as _t,urllib.parse
    auth_user(request); d=await request.json()
    prompt=d.get("prompt","").strip(); style=d.get("style","3d render")
    w=min(int(d.get("width",1024)),2048); h=min(int(d.get("height",1024)),2048)
    if not prompt: return err("Prompt required")
    enhanced=f"{prompt}, {style}, studio lighting, high detail, 8k, professional CGI"
    if CF_ACCOUNT_ID.strip() and CF_KEY.strip():
        hdrs={"Authorization":f"Bearer {CF_KEY}"}
        for cf_model in _CF_IMG_MODELS:
            try:
                r=req.post(cf_ep(cf_model),headers=hdrs,files={"prompt":(None,enhanced)},timeout=90)
                ct=r.headers.get("Content-Type","")
                if r.ok and "image" in ct:
                    return J({"ok":True,"b64":_b64.b64encode(r.content).decode(),"mime":ct.split(";")[0],"url":"","backend":f"CF · {style}","enhanced_prompt":enhanced})
                try:
                    data=r.json(); img=data.get("result",{}).get("image","")
                    if img: return J({"ok":True,"b64":img,"mime":"image/png","url":"","backend":f"CF · {style}","enhanced_prompt":enhanced})
                except: pass
            except: pass
    url=f"https://image.pollinations.ai/prompt/{urllib.parse.quote(enhanced,safe='')}?model=flux&width={w}&height={h}&nologo=true&enhance=true"
    for _ in range(3):
        try:
            r=_pol_get(url,120); ct=r.headers.get("Content-Type","")
            if r.ok and ct.startswith("image"):
                return J({"ok":True,"b64":_b64.b64encode(r.content).decode(),"mime":ct.split(";")[0],"url":url,"backend":f"Pollinations · {style}","enhanced_prompt":enhanced})
            if r.status_code==429: _t.sleep(5); continue
            break
        except: break
    return J({"ok":False,"error":"3D gen failed","b64":"","mime":"image/jpeg","url":url},502)
@app.post("/api/generate/audio")
async def gen_audio(request:Request):
    u=auth_user(request); d=await request.json()
    text=d.get("text","").strip(); voice=d.get("voice","alloy")
    if not text: return err("Text required")
    avail=get_available(u["id"],u["salt"])
    api_key=avail.get("openrouter","") or avail.get("groq","")
    if not api_key: return err("API key required for TTS",400)
    hdrs={"Authorization":f"Bearer {api_key}","Content-Type":"application/json","X-Title":"Fusion.AI"}
    try:
        r=req.post("https://openrouter.ai/api/v1/audio/speech",headers=hdrs,json={"model":"openai/tts-1","input":text[:4096],"voice":voice},timeout=30)
        if r.ok:
            import base64 as _b64
            return J({"ok":True,"b64":_b64.b64encode(r.content).decode(),"mime":r.headers.get("Content-Type","audio/mpeg")})
        return err("TTS unavailable",502)
    except Exception as e: return err(str(e),500)
@app.get("/api/dev/stats")
async def dev_stats(request:Request):
    dev_check(request)
    with db() as c:
        users=c.execute("SELECT COUNT(*) as n FROM users WHERE username NOT LIKE 'guest_%'").fetchone()["n"]
        guests=c.execute("SELECT COUNT(*) as n FROM users WHERE username LIKE 'guest_%'").fetchone()["n"]
        msgs=c.execute("SELECT COUNT(*) as n FROM messages").fetchone()["n"]
        msgs_today=c.execute("SELECT COUNT(*) as n FROM messages WHERE ts>=date('now')").fetchone()["n"]
        visits_24h=c.execute("SELECT COUNT(DISTINCT username) as n FROM visitor_log WHERE event IN('login','guest_login') AND ts>=datetime('now','-1 day')").fetchone()["n"]
        total_visits=c.execute("SELECT COUNT(*) as n FROM visitor_log WHERE event IN('login','guest_login')").fetchone()["n"]
        online_now=c.execute("SELECT COUNT(*) as n FROM tokens WHERE created>=datetime('now','-30 minutes')").fetchone()["n"]
        recent_users=c.execute("SELECT username,created FROM users WHERE username NOT LIKE 'guest_%' ORDER BY id DESC LIMIT 20").fetchall()
        top_models=c.execute("SELECT model,COUNT(*) as cnt FROM messages WHERE model IS NOT NULL GROUP BY model ORDER BY cnt DESC LIMIT 15").fetchall()
        recent_msgs=c.execute("SELECT u.username,m.role,m.content,m.model,m.ts FROM messages m JOIN users u ON u.id=m.user_id ORDER BY m.id DESC LIMIT 30").fetchall()
        visit_log=c.execute("SELECT username,event,ip,ts FROM visitor_log ORDER BY id DESC LIMIT 50").fetchall()
    return J({"total_users":users,"total_guests":guests,"total_messages":msgs,"messages_today":msgs_today,
              "visits_24h":visits_24h,"total_visits":total_visits,"online_now":online_now,
              "recent_users":[dict(r) for r in recent_users],"top_models":[dict(r) for r in top_models],
              "recent_messages":[{"username":r["username"],"role":r["role"],"content":r["content"][:120],"model":r["model"],"ts":r["ts"]} for r in recent_msgs],
              "visitor_log":[dict(r) for r in visit_log],
              "env_keys":{"GROQ_KEY":bool(GROQ_KEY),"OPENROUTER_KEY":bool(OPENROUTER_KEY),"HF_TOKEN":bool(HF_TOKEN),
                          "GITHUB_TOKEN":bool(GITHUB_TOKEN),"STABILITY_KEY":bool(STABILITY_KEY),"CF_ACCOUNT_ID":bool(CF_ACCOUNT_ID),
                          "CF_KEY":bool(CF_KEY),"WORKER_KEY":bool(WORKER_KEY),"SECRET_KEY":bool(_raw_secret),
                          "DEV_PASSWORD":bool(DEV_PASSWORD),"TP_USERNAME":bool(TP_USERNAME),
                          "CF_ACCOUNT_ID2":bool(CF_ACCOUNT_ID2),"CF_KEY2":bool(CF_KEY2),"WORKER_URL2":bool(WORKER_URL2),
                          "OPENAI_KEY":bool(OPENAI_KEY),"ANTHROPIC_KEY":bool(ANTHROPIC_KEY),"GEMINI_KEY":bool(GEMINI_KEY),
                          "DEEPSEEK_KEY":bool(DEEPSEEK_KEY),"MOONSHOT_KEY":bool(MOONSHOT_KEY),"MISTRAL_KEY":bool(MISTRAL_KEY),
                          "XAI_KEY":bool(XAI_KEY),"COHERE_KEY":bool(COHERE_KEY),"TOGETHER_KEY":bool(TOGETHER_KEY),
                          "PERPLEXITY_KEY":bool(PERPLEXITY_KEY)}})
@app.get("/api/dev/users")
async def dev_all_users(request:Request):
    dev_check(request)
    with db() as c:
        users=c.execute("""SELECT u.id,u.username,u.created,COUNT(m.id) as msg_count,MAX(m.ts) as last_active,
            GROUP_CONCAT(DISTINCT ak.provider) as providers,
            (SELECT value FROM memory WHERE user_id=u.id AND key='__banned__' LIMIT 1) as banned
            FROM users u LEFT JOIN messages m ON m.user_id=u.id LEFT JOIN api_keys ak ON ak.user_id=u.id
            GROUP BY u.id ORDER BY u.id DESC""").fetchall()
    return J({"users":[dict(u) for u in users]})
@app.post("/api/dev/user/{uid}/ban")
async def dev_ban(uid:int,request:Request):
    dev_check(request); d=await request.json(); reason=d.get("reason","No reason")[:200]
    with db() as c:
        u=c.execute("SELECT username FROM users WHERE id=?",(uid,)).fetchone()
        if not u: return err("Not found",404)
        c.execute("INSERT OR REPLACE INTO memory(user_id,key,value,ts)VALUES(?,?,?,?)",(uid,"__banned__",reason,_now()))
        c.execute("DELETE FROM tokens WHERE user_id=?",(uid,))
    return J({"ok":True,"message":f"User {u['username']} banned"})
@app.post("/api/dev/user/{uid}/unban")
async def dev_unban(uid:int,request:Request):
    dev_check(request)
    with db() as c:
        u=c.execute("SELECT username FROM users WHERE id=?",(uid,)).fetchone()
        if not u: return err("Not found",404)
        c.execute("DELETE FROM memory WHERE user_id=? AND key='__banned__'",(uid,))
    return J({"ok":True,"message":f"User {u['username']} unbanned"})
@app.delete("/api/dev/user/{uid}/delete")
async def dev_del_user(uid:int,request:Request):
    dev_check(request)
    with db() as c:
        u=c.execute("SELECT username FROM users WHERE id=?",(uid,)).fetchone()
        if not u: return err("Not found",404)
        for t in ["messages","tokens","api_keys","preferences","saved_items","memory","conv_messages","conversations","visitor_log"]:
            c.execute(f"DELETE FROM {t} WHERE user_id=?",(uid,))
        c.execute("DELETE FROM users WHERE id=?",(uid,))
    return J({"ok":True,"message":f"User {u['username']} deleted"})
@app.post("/api/dev/user/{uid}/reset_keys")
async def dev_reset_keys(uid:int,request:Request):
    dev_check(request)
    with db() as c:
        u=c.execute("SELECT username FROM users WHERE id=?",(uid,)).fetchone()
        if not u: return err("Not found",404)
        c.execute("DELETE FROM api_keys WHERE user_id=?",(uid,))
    return J({"ok":True,"message":f"Keys for {u['username']} cleared"})
@app.get("/api/dev/user/{uid}/history")
async def dev_user_hist(uid:int,request:Request):
    dev_check(request)
    with db() as c:
        u=c.execute("SELECT username FROM users WHERE id=?",(uid,)).fetchone()
        if not u: return err("Not found",404)
        msgs=c.execute("SELECT role,content,model,ts FROM messages WHERE user_id=? ORDER BY id DESC LIMIT 50",(uid,)).fetchall()
    return J({"username":u["username"],"messages":[dict(m) for m in msgs]})
@app.get("/api/dev/request-logs")
async def dev_logs(request:Request):
    dev_check(request)
    with db() as c:
        logs=c.execute("SELECT m.id,u.username,m.role,m.content,m.model,m.ts FROM messages m JOIN users u ON u.id=m.user_id ORDER BY m.id DESC LIMIT 200").fetchall()
    return J({"logs":[dict(l) for l in logs]})
@app.get("/api/dev/search-logs")
async def dev_search_logs(request:Request):
    dev_check(request); q=request.query_params.get("q","").strip().lower()
    with db() as c:
        logs=c.execute("SELECT m.id,u.username,m.role,m.content,m.model,m.ts FROM messages m JOIN users u ON u.id=m.user_id WHERE LOWER(m.content) LIKE ? OR LOWER(u.username) LIKE ? ORDER BY m.id DESC LIMIT 100",(f"%{q}%",f"%{q}%")).fetchall()
    return J({"logs":[dict(l) for l in logs],"query":q})
@app.get("/api/dev/export-logs")
async def dev_export_logs(request:Request):
    dev_check(request)
    with db() as c:
        logs=c.execute("SELECT m.id,u.username,m.role,m.content,m.model,m.ts FROM messages m JOIN users u ON u.id=m.user_id ORDER BY m.id DESC LIMIT 5000").fetchall()
    return J({"logs":[dict(l) for l in logs],"exported":len(logs)})

def _load_json_file(path,default):
    if os.path.exists(path):
        try:
            with open(path) as f: d=json.load(f); default.update(d)
        except: pass
    return default
@app.get("/api/dev/feature-flags")
async def dev_get_flags(request:Request):
    dev_check(request)
    return J(_load_json_file(os.path.join(_data_dir,"feature_flags.json"),
             {"ai_requests_enabled":True,"image_gen_enabled":True,"video_gen_enabled":True,
              "chat_streaming_enabled":True,"worker_image_enabled":True,"slow_mode":False}))
@app.post("/api/dev/feature-flags")
async def dev_set_flags(request:Request):
    dev_check(request); d=await request.json()
    with open(os.path.join(_data_dir,"feature_flags.json"),"w") as f: json.dump(d,f)
    return J({"ok":True})
@app.get("/api/dev/rate-limits")
async def dev_get_limits(request:Request):
    dev_check(request)
    return J(_load_json_file(os.path.join(_data_dir,"rate_limits.json"),
             {"global_rpm":120,"per_user_rpm":20,"image_per_user_hour":30,"emergency_stop":False}))
@app.post("/api/dev/rate-limits")
async def dev_set_limits(request:Request):
    dev_check(request); d=await request.json()
    with open(os.path.join(_data_dir,"rate_limits.json"),"w") as f: json.dump(d,f)
    return J({"ok":True})
@app.post("/api/dev/test-prompt")
async def dev_test_prompt(request:Request):
    u=dev_check(request); d=await request.json()
    prompt=d.get("prompt","").strip(); mkey=d.get("model_key","groq_llama33_70b")
    if not prompt: return err("Prompt required")
    avail=get_available(u["id"],u["salt"]); m=MODELS.get(mkey)
    if not m: return err("Unknown model")
    prov=m["provider"]
    if prov not in avail: return err(f"No key for {prov}")
    hdrs={"Authorization":f"Bearer {avail[prov]}","Content-Type":"application/json"}
    if prov=="openrouter": hdrs.update({"X-Title":"Fusion.AI"})
    body={"messages":[{"role":"user","content":prompt}],"max_tokens":1024}
    try:
        if prov=="cloudflare":
            r=req.post(cf_ep(m["model"]),headers={"Authorization":f"Bearer {avail['cloudflare']}","Content-Type":"application/json"},json=body,timeout=30)
            data=r.json(); reply=(data.get("result",{})or{}).get("response","") or (data.get("choices",[{}])[0].get("message",{})or{}).get("content","")
        elif prov in ENDPOINTS:
            body["model"]=m["model"]; r=req.post(ENDPOINTS[prov],headers=hdrs,json=body,timeout=30)
            data=r.json(); reply=data["choices"][0]["message"]["content"]
        else: return err("Unsupported provider")
        return J({"ok":True,"reply":reply,"model":m["label"],"status":r.status_code})
    except Exception as e: return J({"ok":False,"error":str(e),"status":500},500)
@app.get("/api/dev/disabled-models")
async def dev_disabled_models(request:Request):
    dev_check(request)
    path=os.path.join(_data_dir,"disabled_models.json")
    disabled=[]
    if os.path.exists(path):
        try:
            with open(path) as f: disabled=json.load(f)
        except: pass
    return J({"disabled":disabled})
@app.post("/api/dev/model-enable")
async def dev_model_enable(request:Request):
    dev_check(request); d=await request.json()
    key=d.get("key",""); enabled=d.get("enabled",True)
    if key not in MODELS: return err("Unknown model")
    path=os.path.join(_data_dir,"disabled_models.json"); disabled=[]
    if os.path.exists(path):
        try:
            with open(path) as f: disabled=json.load(f)
        except: pass
    if not enabled and key not in disabled: disabled.append(key)
    elif enabled and key in disabled: disabled.remove(key)
    with open(path,"w") as f: json.dump(disabled,f)
    return J({"ok":True,"disabled_models":disabled})

from concurrent.futures import ThreadPoolExecutor, as_completed

def _test_one(key,m,api_key):
    import time as _t
    prov=m["provider"]
    if prov not in ("groq","openrouter"): return {"key":key,"label":m["label"],"provider":prov,"model":m["model"],"status":"no_key","ms":0}
    ep=ENDPOINTS.get(prov); hdrs={"Authorization":f"Bearer {api_key}","Content-Type":"application/json"}
    if prov=="openrouter": hdrs.update({"X-Title":"Fusion.AI"})
    t0=_t.time()
    try:
        r=req.post(ep,headers=hdrs,json={"model":m["model"],"messages":[{"role":"user","content":"Reply with exactly: OK"}],"max_tokens":10,"stream":False},timeout=20)
        ms=round((_t.time()-t0)*1000)
        if r.ok:
            try: reply=r.json()["choices"][0]["message"]["content"].strip()[:60]
            except: reply="OK"
            return {"key":key,"label":m["label"],"provider":prov,"model":m["model"],"status":"ok","reply":reply,"ms":ms}
        return {"key":key,"label":m["label"],"provider":prov,"model":m["model"],"status":"error","error":f"HTTP {r.status_code}","http":r.status_code,"ms":ms}
    except Exception as e:
        return {"key":key,"label":m["label"],"provider":prov,"model":m["model"],"status":"exception","error":str(e),"ms":round((_time.time()-t0)*1000)}
@app.get("/api/test")
async def test_models(request:Request):
    dev_check(request)
    u=auth_user(request); avail=get_available(u["id"],u["salt"])
    # Test ALL chat models across ALL providers with keys
    testable_provs={"groq","openrouter","github","cloudflare"}
    tasks=[(k,m,avail[m["provider"]]) for k,m in MODELS.items()
           if m.get("type","chat")=="chat" and m["provider"] in avail and m["provider"] in testable_provs]
    results=[]
    with ThreadPoolExecutor(max_workers=10) as ex:
        futs={ex.submit(_test_one,*t):t for t in tasks}
        for f in as_completed(futs): results.append(f.result())
    # Add no-key models
    for k,m in MODELS.items():
        if m.get("type","chat")=="chat" and (m["provider"] not in avail or m["provider"] not in testable_provs):
            results.append({"key":k,"label":m["label"],"provider":m["provider"],"model":m["model"],"status":"no_key","ms":0})
    results.sort(key=lambda x:({"ok":0,"error":1,"exception":1,"no_key":2}.get(x["status"],3),x.get("label","")))
    ok=sum(1 for r in results if r["status"]=="ok")
    bad=sum(1 for r in results if r["status"] in("error","exception"))
    failed=[r for r in results if r["status"] in("error","exception")]
    report="All tested models OK." if not failed else "Failed:\n"+"\n".join(f"- {r['label']}: {r.get('error','?')}" for r in failed)
    return J({"summary":{"ok":ok,"failed":bad,"no_key":len(results)-ok-bad,"total":len(results)},"overseer_report":report,"models":results,"health":MODEL_HEALTH})

# ── Deep Research: ALL available models in parallel with timing ──────────────
@app.post("/api/deep-research")
async def deep_research(request:Request):
    import time as _t2, urllib.parse as _up2, re as _re2
    t_start=_t2.time()
    u=auth_user(request); d=await request.json()
    prompt=d.get("prompt","").strip()
    mode=d.get("mode","general")   # general | math | code
    if not prompt: return err("Prompt required")
    avail=get_available(u["id"],u["salt"])

    # Mode-specific model sets
    MATH_KEYS={"gh_o3","gh_o3_mini","gh_o4_mini","gh_o1","gh_deepseek_r1","groq_qwen3_32b","cf_qwq_32b","cf_deepseek_r1","or_deepseek_r1","groq_compound"}
    CODE_KEYS={"cf_qwen25_coder","gh_gpt4o","gh_gpt4o_mini","gh_o4_mini","gh_deepseek_v3","or_deepseek_v3","gh_phi4","or_mistral_small","groq_compound","groq_llama33_70b"}

    def _mode_ok(k):
        if mode=="math": return k in MATH_KEYS
        if mode=="code": return k in CODE_KEYS
        return True

    # Mode-specific system prompts
    if mode=="math":
        sys_p=("You are a world-class mathematician. Solve the following problem step-by-step. Show all working. Use LaTeX notation: $...$ for inline math, $$...$$ for display equations. Give the exact final answer clearly.")
    elif mode=="code":
        sys_p=("You are an expert software engineer. Write clean, complete, production-ready code for the following request. Include brief explanations. Use proper formatting. No placeholders or truncation — write every line.")
    else:
        sys_p=("You are a specialist AI researcher. Answer the following question concisely and accurately in 3-5 sentences from YOUR unique model perspective and training. Be direct. Do not repeat the question.")

    # 1. Grab Bing web context via OpenSERP (skip for math/code — not needed)
    web_ctx=""
    if mode=="general":
        try:
            sr=req.get(f"http://openserp.alwaysdata.net/bing/search?text={_up2.quote(prompt)}",
                       headers={"User-Agent":"FusionAI/2.0","Accept":"application/json"},timeout=8)
            if sr.ok:
                items=sr.json() if isinstance(sr.json(),list) else []
                for item in items[:5]:
                    snip=item.get("description","") or item.get("snippet","")
                    if snip: web_ctx+=snip[:250]+"\n"
        except: pass
    if web_ctx: sys_p+=f"\n\nLive web context:\n{web_ctx[:800]}"

    base_msg=[{"role":"system","content":sys_p},{"role":"user","content":prompt}]

    # 2. All chat models filtered by mode
    all_chat=[k for k,m in MODELS.items()
              if m.get("type","chat")=="chat"
              and m["provider"] in avail
              and m["provider"] in ("groq","openrouter","github","cloudflare")
              and _mode_ok(k)]

    def _fetch_agent(key):
        m=MODELS[key]; prov=m["provider"]; mdl=m["model"]; label=m["label"]
        t0=_t2.time()
        try:
            key_val=avail[prov]
            hdrs={"Authorization":f"Bearer {key_val}","Content-Type":"application/json"}
            if prov=="openrouter": hdrs.update({"X-Title":"Fusion.AI","HTTP-Referer":"https://fusionai.space"})
            body={"messages":base_msg,"max_tokens":350,"stream":False}
            if prov=="cloudflare":
                r=req.post(cf_ep(mdl),headers={"Authorization":f"Bearer {key_val}","Content-Type":"application/json"},json=body,timeout=22)
                dj=r.json(); txt=(dj.get("result",{})or{}).get("response","") or (dj.get("choices",[{}])[0].get("message",{})or{}).get("content","")
            else:
                body["model"]=mdl; ep=ENDPOINTS.get(prov,"")
                r=req.post(ep,headers=hdrs,json=body,timeout=22)
                if not r.ok: return None
                txt=r.json()["choices"][0]["message"]["content"]
            import re as _re3; txt=_re3.sub(r"<think>[\s\S]*?</think>","",txt).strip()
            ms=round((_t2.time()-t0)*1000)
            return {"key":key,"model":label,"prov":prov,"text":txt,"ms":ms} if txt else None
        except: return None

    with ThreadPoolExecutor(max_workers=20) as ex:
        futs=[ex.submit(_fetch_agent,k) for k in all_chat]
        raw=[f.result() for f in as_completed(futs,timeout=28)]
    perspectives=[r for r in raw if r and r.get("text")]

    if not perspectives:
        return J({"ok":False,"error":"All agents failed","perspectives":[],"total_ms":round((_t2.time()-t_start)*1000)})

    # 3. Synthesise with Groq Compound (has web search) or best available
    synth_key=avail.get("groq","") or avail.get("openrouter","") or avail.get("github","")
    if avail.get("groq"): synth_ep=ENDPOINTS["groq"]; synth_model="compound-beta"
    elif avail.get("openrouter"): synth_ep=ENDPOINTS["openrouter"]; synth_model="deepseek/deepseek-chat-v3-0324:free"
    elif avail.get("github"): synth_ep=ENDPOINTS["github"]; synth_model="openai/gpt-4o"
    else: synth_ep=None
    combined="\n\n".join(f"[{p['model']} · {p['ms']}ms]: {p['text']}" for p in sorted(perspectives,key=lambda x:x['ms']))
    synthesis=""
    if synth_ep and synth_key:
        try:
            if mode=="math":
                synth_sys="You are the world's best mathematician. Review all model solutions below. Identify the correct approach. Write a definitive, step-by-step solution using LaTeX ($$...$$). Flag any errors in the individual solutions."
            elif mode=="code":
                synth_sys="You are a senior software engineer. Review all code solutions below. Write the single best, most complete, production-ready implementation. Clean up any bugs. Include comments."
            else:
                synth_sys="You are a master research synthesiser. Combine these AI perspectives into one definitive answer. Use headers for structure. Be direct. Do not say 'the models say' — write as a confident synthesis."
            synth_msg=[
                {"role":"system","content":synth_sys},
                {"role":"user","content":f"Question: {prompt}\n\nAll AI responses ({len(perspectives)} models):\n{combined}\n\nWrite the master synthesis:"}
            ]
            h2={"Authorization":f"Bearer {synth_key}","Content-Type":"application/json"}
            if "openrouter" in synth_ep: h2.update({"X-Title":"Fusion.AI","HTTP-Referer":"https://fusionai.space"})
            max_tok=1800 if mode in ("math","code") else 1200
            sr=req.post(synth_ep,headers=h2,json={"model":synth_model,"messages":synth_msg,"max_tokens":max_tok,"stream":False},timeout=40)
            if sr.ok: synthesis=sr.json()["choices"][0]["message"]["content"].strip()
        except: pass
    if not synthesis:
        synthesis="\n\n".join(p["text"] for p in perspectives[:5])

    total_ms=round((_t2.time()-t_start)*1000)
    return J({"ok":True,"synthesis":synthesis,"perspectives":perspectives,
              "web_ctx":web_ctx.strip(),"agents_used":len(perspectives),
              "total_ms":total_ms,"models_queried":len(all_chat),"mode":mode})
@app.get("/api/test/image")
async def test_image(request:Request):
    dev_check(request)
    import base64 as _b64,time as _t,traceback as _tb
    if not CF_ACCOUNT_ID.strip() or not CF_KEY.strip(): return J({"ok":False,"error":"CF_ACCOUNT_ID and CF_KEY not set"})
    model="@cf/black-forest-labs/flux-2-klein-4b"
    t0=_t.time()
    try:
        r=req.post(cf_ep(model),headers={"Authorization":f"Bearer {CF_KEY}"},files={"prompt":(None,"a red circle")},timeout=60)
        ms=round((_t.time()-t0)*1000); ct=r.headers.get("Content-Type","")
        if r.ok and ("image" in ct or len(r.content)>100):
            return J({"ok":True,"status":r.status_code,"content_type":ct,"size_bytes":len(r.content),"ms":ms,"model":model,"b64_preview":_b64.b64encode(r.content[:512]).decode()})
        try:
            data=r.json(); img=data.get("result",{}).get("image","")
            if img: return J({"ok":True,"status":r.status_code,"ms":ms,"model":model,"size_bytes":len(img)})
        except: pass
        return J({"ok":False,"status":r.status_code,"ms":ms,"model":model,"body":r.text[:300],"error":f"HTTP {r.status_code}"})
    except Exception as e: return J({"ok":False,"ms":round((_t.time()-t0)*1000),"model":model,"error":f"{type(e).__name__}: {e}","traceback":_tb.format_exc()[-400:]})
@app.get("/api/health")
async def model_health_route(request:Request):
    dev_check(request)
    status=[{"key":k,"label":m.get("label",k),"type":m.get("type","chat"),"provider":m["provider"],
             "healthy":is_healthy(k),"fails":MODEL_HEALTH.get(k,{}).get("fails",0)} for k,m in MODELS.items()]
    status.sort(key=lambda x:(x["healthy"],x["fails"]))
    return J({"models":status,"raw":MODEL_HEALTH})


# ── KaTeX static files (served locally to avoid CDN tracking-prevention warnings) ──
import threading as _threading
_katex_cache = {}
_katex_lock = _threading.Lock()

def _fetch_katex():
    """Download KaTeX files once at startup and cache in memory."""
    files = {
        "katex.css":       "https://cdn.jsdelivr.net/npm/katex@0.16.9/dist/katex.min.css",
        "katex.js":        "https://cdn.jsdelivr.net/npm/katex@0.16.9/dist/katex.min.js",
        "auto-render.js":  "https://cdn.jsdelivr.net/npm/katex@0.16.9/dist/contrib/auto-render.min.js",
    }
    for name, url in files.items():
        try:
            r = req.get(url, timeout=15)
            if r.ok:
                with _katex_lock:
                    _katex_cache[name] = (r.content, r.headers.get("Content-Type","text/plain"))
        except: pass

_threading.Thread(target=_fetch_katex, daemon=True).start()

@app.get("/static/katex.css")
async def serve_katex_css():
    from fastapi.responses import Response
    with _katex_lock:
        data = _katex_cache.get("katex.css")
    if data:
        return Response(content=data[0], media_type="text/css",
                       headers={"Cache-Control":"public,max-age=86400"})
    return Response(content="/* KaTeX not loaded */", media_type="text/css")

@app.get("/static/katex.js")
async def serve_katex_js():
    from fastapi.responses import Response
    with _katex_lock:
        data = _katex_cache.get("katex.js")
    if data:
        return Response(content=data[0], media_type="application/javascript",
                       headers={"Cache-Control":"public,max-age=86400"})
    return Response(content="", media_type="application/javascript")

@app.get("/static/auto-render.js")
async def serve_autorender_js():
    from fastapi.responses import Response
    with _katex_lock:
        data = _katex_cache.get("auto-render.js")
    if data:
        return Response(content=data[0], media_type="application/javascript",
                       headers={"Cache-Control":"public,max-age=86400"})
    return Response(content="", media_type="application/javascript")

# ── Serve HTML ────────────────────────────────────────────────────────────────
_HTML = r"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8"><meta name="theme-color" content="#000000" id="metaThemeColor"/>
<meta name="viewport" content="width=device-width,initial-scale=1,maximum-scale=1,interactive-widget=resizes-content">
<link rel="icon" href="data:image/svg+xml,<svg xmlns='http://www.w3.org/2000/svg' viewBox='0 0 100 100'><text y='.9em' font-size='90'>⚡</text></svg>">
<title>Fusion.AI</title>
<link rel="stylesheet" href="/static/katex.css">
<script defer src="/static/katex.js"></script>
<script defer src="/static/auto-render.js" onload="window._katexReady=true;setTimeout(function(){document.querySelectorAll('.bbl').forEach(function(e){_renderMath(e);});},100);"></script>
<link href="https://fonts.googleapis.com/css2?family=Bebas+Neue&family=DM+Sans:wght@300;400;500;600&family=DM+Mono:wght@400;500&display=swap" rel="stylesheet">
<style>
:root{
 /* ── New FusionAI dark purple/cyan theme ── */
 --red:#7c3aed;--red2:#6d28d9;--redg:rgba(124,58,237,.28);
 --blue:#06b6d4;--blueg:rgba(6,182,212,.22);
 --grad:linear-gradient(135deg,#7c3aed,#06b6d4);
 --gradt:linear-gradient(90deg,#7c3aed,#06b6d4);
 --green:#22c55e;--purple:#a78bfa;
 --accent:#7c3aed;--accent2:#06b6d4;
 --glass-bg:rgba(10,10,20,0.55);--glass-bg2:rgba(12,12,24,0.65);
 --glass-surf:rgba(16,16,30,0.60);--glass-surf2:rgba(20,20,38,0.65);
 --glass-bdr:rgba(124,58,237,0.15);--glass-bdr2:rgba(124,58,237,0.28);
 --blur:blur(24px) saturate(180%);--blur2:blur(16px) saturate(160%);
 --tx:#eeeeff;--tx2:#8888bb;--tx3:#44446a;--shad:rgba(0,0,0,.8);--inp:rgba(8,8,18,0.65);
 --bg-body:#08080f;
 --sidebar-w:242px;
 --hdr-h:52px;
}
*{margin:0;padding:0;box-sizing:border-box}
html,body{height:100%;overflow:hidden}
body{font-family:'DM Sans',sans-serif;color:var(--tx);background:var(--bg-body);transition:background .3s,color .3s}
.bg{position:fixed;inset:0;z-index:0;overflow:hidden;pointer-events:none}
.bg-orb{position:absolute;border-radius:50%;filter:blur(80px);animation:orb 18s ease-in-out infinite alternate}
.bg-orb1{width:700px;height:700px;top:-200px;left:-200px;background:radial-gradient(circle,rgba(124,58,237,.3),transparent 65%);animation-duration:16s}
.bg-orb2{width:800px;height:800px;bottom:-250px;right:-250px;background:radial-gradient(circle,rgba(6,182,212,.22),transparent 65%);animation-duration:22s;animation-delay:-8s}
.bg-orb3{width:500px;height:500px;top:40%;left:40%;background:radial-gradient(circle,rgba(100,20,200,.18),transparent 65%);animation-duration:28s;animation-delay:-4s}
.bg-grid{position:absolute;inset:0;background-image:linear-gradient(rgba(124,58,237,.04) 1px,transparent 1px),linear-gradient(90deg,rgba(124,58,237,.04) 1px,transparent 1px);background-size:48px 48px;mask-image:radial-gradient(ellipse at center,black 40%,transparent 75%)}
/* Matrix canvas background */
#matrixCanvas{position:fixed;inset:0;z-index:0;opacity:0;transition:opacity .5s;pointer-events:none}
#matrixCanvas.active{opacity:.18}
body.light #matrixCanvas.active{opacity:.07}
/* Right vertical quick-tools bar */
.rtoolbar{position:fixed;right:0;top:50%;transform:translateY(-50%);z-index:100;display:flex;flex-direction:column;gap:6px;padding:10px 6px;background:rgba(10,10,20,.7);backdrop-filter:blur(20px);border:1px solid var(--glass-bdr);border-right:none;border-radius:14px 0 0 14px;box-shadow:-4px 0 24px rgba(0,0,0,.4)}
.rtool-btn{width:38px;height:38px;background:rgba(255,255,255,.06);border:1px solid var(--glass-bdr);border-radius:9px;display:flex;align-items:center;justify-content:center;cursor:pointer;font-size:16px;transition:all .2s;color:var(--tx2);position:relative}
.rtool-btn:hover{border-color:var(--red);color:var(--tx);background:rgba(124,58,237,.12);transform:translateX(-3px)}
.rtool-btn.active{border-color:var(--blue);background:rgba(6,182,212,.15);color:var(--blue)}
.rtool-tip{position:absolute;right:48px;top:50%;transform:translateY(-50%);background:rgba(10,10,20,.95);border:1px solid var(--glass-bdr2);border-radius:7px;padding:4px 10px;font-size:11px;color:var(--tx);white-space:nowrap;pointer-events:none;opacity:0;transition:opacity .15s}
.rtool-btn:hover .rtool-tip{opacity:1}
/* Left settings sidebar */
.lsidebar{position:fixed;left:0;top:52px;bottom:0;width:240px;background:rgba(6,8,18,.92);backdrop-filter:blur(22px);border-right:1px solid var(--glass-bdr);z-index:90;transform:translateX(-100%);transition:transform .25s cubic-bezier(.16,1,.3,1);overflow-y:auto;display:flex;flex-direction:column}
.lsidebar.open{transform:translateX(0)}
.lsidebar-header{padding:16px 16px 10px;border-bottom:1px solid var(--glass-bdr);flex-shrink:0}
.lsidebar-header h3{font-size:13px;font-weight:700;color:var(--tx);margin-bottom:3px}
.lsidebar-header p{font-size:11px;color:var(--tx3)}
.lsidebar-close{position:absolute;top:12px;right:12px;background:none;border:none;color:var(--tx3);font-size:16px;cursor:pointer;width:28px;height:28px;display:flex;align-items:center;justify-content:center;border-radius:6px}
.lsidebar-close:hover{color:var(--tx);background:rgba(255,255,255,.07)}
.lsidebar-section{padding:12px 14px;border-bottom:1px solid var(--glass-bdr)}
.lsidebar-section h4{font-size:9px;font-weight:700;color:var(--tx3);text-transform:uppercase;letter-spacing:1px;margin-bottom:9px}
.lsidebar-item{display:flex;align-items:center;gap:9px;padding:8px 10px;border-radius:9px;cursor:pointer;transition:background .15s;color:var(--tx2);font-size:13px}
.lsidebar-item:hover{background:rgba(255,255,255,.07);color:var(--tx)}
.lsidebar-item.active{background:rgba(124,58,237,.1);color:var(--tx);border:1px solid rgba(124,58,237,.2)}
.lsidebar-item span{font-size:15px}
/* Google sign-in button */
.google-btn{width:100%;padding:10px;background:rgba(255,255,255,.07);border:1.5px solid rgba(255,255,255,.15);border-radius:10px;color:var(--tx);font-family:'DM Sans',sans-serif;font-size:13px;font-weight:500;cursor:pointer;display:flex;align-items:center;justify-content:center;gap:9px;transition:all .2s;margin-bottom:10px}
.google-btn:hover{background:rgba(255,255,255,.12);border-color:rgba(255,255,255,.3);transform:translateY(-1px)}
.google-btn svg{width:18px;height:18px;flex-shrink:0}
/* Weather widget */
.weather-widget{display:none;background:rgba(6,182,212,.08);border:1px solid rgba(6,182,212,.2);border-radius:11px;padding:10px 13px;font-size:12px;color:var(--tx2);margin-bottom:10px;animation:fi .3s ease}
.weather-widget.show{display:flex;align-items:center;gap:10px}
.weather-temp{font-size:22px;font-weight:700;color:var(--tx);font-family:'Bebas Neue',sans-serif}
.weather-info{flex:1}
/* Camera button */
#cameraBtn{display:none}
@media(max-width:768px){.rtoolbar{display:none}.lsidebar{width:100vw}}
@media(pointer:fine){#cameraBtn{display:flex}}
body.light{background:#f5f6fa!important;color:#111827}
body.light header{background:rgba(255,255,255,0.97)!important;border-bottom:1px solid rgba(180,190,220,.45)}
body.light .msg.ai .bbl{background:rgba(255,255,255,0.97);color:#111827;border:1px solid rgba(180,190,220,.4)}
body.light .msg.user .bbl{background:var(--grad);color:#fff}
body.light .iz{background:rgba(255,255,255,0.93);border-top:1px solid rgba(180,190,220,.4)}
body.light .ibox{background:#fff;border:1.5px solid rgba(180,190,220,.55)}
body.light .ibox textarea{color:#111827}
body.light .ibox textarea::placeholder{color:#9ca3af}
body.light .sp{background:rgba(248,249,252,0.99)}
body.light .krow{background:rgba(240,242,248,0.85)}
body.light .msel{background:#fff;color:#111827}
body.light .msel option{background:#fff;color:#111827}
body.light .msel optgroup{background:#fff}
body.light .drop{background:rgba(248,249,252,0.99)}
body.light .ditem{color:#374151}
body.light .saved-item{background:rgba(240,242,248,0.85)}
body.light .mem-item{background:rgba(240,242,248,0.85)}
body.light .thought{background:rgba(240,242,248,0.9)}
body.light .chat-opts-panel{background:rgba(248,249,252,0.99)}
body.light .ov-modal{background:rgba(248,249,252,0.99)}
body.light code,body.light pre{background:rgba(220,225,240,.7);color:#1e3a5f}
@keyframes orb{from{transform:translate(0,0) scale(1)}to{transform:translate(60px,80px) scale(1.1)}}
/* KaTeX / Math rendering */
.katex-display{overflow-x:auto;overflow-y:hidden;padding:6px 0;margin:6px 0}
.katex{font-size:1.05em;color:var(--tx)}
.katex-display>.katex{font-size:1.15em}
.bbl .katex-display{background:rgba(255,255,255,.04);border-radius:8px;padding:10px 14px;border:1px solid var(--glass-bdr);display:block;text-align:center}

.page{display:none;flex-direction:column;height:100%;position:fixed;inset:0;z-index:20;overflow:visible}
.page.active{display:flex}
#authPage{align-items:center;justify-content:center;overflow-y:auto}
#chatPage.active{display:flex;flex-direction:column}
.card{background:rgba(6,9,20,0.95);backdrop-filter:blur(28px) saturate(180%);-webkit-backdrop-filter:blur(28px) saturate(180%);border:1px solid rgba(80,120,255,0.28);border-radius:24px;padding:40px 36px;width:420px;max-width:95vw;box-shadow:0 40px 100px rgba(0,0,0,.85),0 0 80px rgba(6,182,212,.08),0 0 0 1px rgba(255,255,255,.05),inset 0 1px 0 rgba(255,255,255,.07);position:relative;overflow:hidden;animation:slideup .45s cubic-bezier(.16,1,.3,1)}
.card::before{content:'';position:absolute;top:0;left:0;right:0;height:3px;background:var(--grad)}
@keyframes slideup{from{opacity:0;transform:translateY(28px) scale(.97)}to{opacity:1;transform:none}}
.card-logo{display:flex;align-items:center;gap:13px;margin-bottom:26px}
.lmark{width:46px;height:46px;background:var(--grad);border-radius:14px;display:flex;align-items:center;justify-content:center;font-size:23px;box-shadow:0 8px 24px var(--redg)}
.ltxt{font-family:'Bebas Neue',sans-serif;font-size:32px;letter-spacing:2.5px;background:var(--gradt);-webkit-background-clip:text;-webkit-text-fill-color:transparent}
.tabs{display:flex;background:rgba(255,255,255,.04);border:1px solid var(--glass-bdr);border-radius:14px;padding:4px;margin-bottom:22px;gap:4px}
.tab{flex:1;padding:10px;border:none;background:none;border-radius:10px;color:var(--tx2);font-family:'DM Sans',sans-serif;font-size:13px;font-weight:500;cursor:pointer;transition:all .22s}
.tab.active{background:var(--grad);color:#fff;box-shadow:0 4px 14px var(--redg)}
.field{margin-bottom:13px}
.field label{display:block;font-size:10.5px;font-weight:600;color:var(--tx2);margin-bottom:6px;text-transform:uppercase;letter-spacing:1px}
.field input{width:100%;background:var(--inp);backdrop-filter:var(--blur2);border:1.5px solid var(--glass-bdr);border-radius:12px;padding:12px 15px;color:var(--tx);font-family:'DM Sans',sans-serif;font-size:14px;outline:none;transition:border-color .2s,box-shadow .2s}
.field input:focus{border-color:var(--blue);box-shadow:0 0 0 3.5px var(--blueg)}
.field input::placeholder{color:var(--tx3)}
.ferr{background:rgba(124,58,237,.12);backdrop-filter:var(--blur2);border:1px solid rgba(124,58,237,.3);border-radius:10px;padding:10px 14px;font-size:12.5px;color:var(--red);margin-bottom:13px;display:none}
.btn{width:100%;background:var(--grad);border:none;border-radius:12px;padding:13px;color:#fff;font-family:'DM Sans',sans-serif;font-size:15px;font-weight:600;cursor:pointer;transition:all .2s;letter-spacing:.3px}
.btn:hover:not(:disabled){opacity:.9;transform:translateY(-2px);box-shadow:0 10px 28px var(--redg)}
.btn:disabled{opacity:.5;cursor:not-allowed}
.divider{display:flex;align-items:center;gap:12px;margin:16px 0;color:var(--tx3);font-size:11px}
.divider::before,.divider::after{content:'';flex:1;height:1px;background:var(--glass-bdr)}
.ghost-btn{width:100%;background:rgba(255,255,255,.05);backdrop-filter:var(--blur2);border:1.5px solid var(--glass-bdr);border-radius:12px;padding:12px;color:var(--tx2);font-family:'DM Sans',sans-serif;font-size:14px;cursor:pointer;transition:all .2s}
.ghost-btn:hover{border-color:var(--blue);color:var(--tx)}
.tlink{text-align:center;margin-top:14px;font-size:12.5px;color:var(--tx3)}
.tlink a{color:var(--blue);cursor:pointer;text-decoration:none}
header{display:flex;align-items:center;justify-content:space-between;padding:10px 16px;background:rgba(6,8,18,0.7);backdrop-filter:var(--blur);-webkit-backdrop-filter:var(--blur);border-bottom:1px solid var(--glass-bdr);flex-shrink:0;position:relative;z-index:10;box-shadow:0 1px 0 rgba(255,255,255,.05),0 4px 24px rgba(0,0,0,.4)}
.hl{display:flex;align-items:center;gap:9px}
.hlm{width:30px;height:30px;background:var(--grad);border-radius:8px;display:flex;align-items:center;justify-content:center;font-size:14px}
.htxt{font-family:'Bebas Neue',sans-serif;font-size:20px;letter-spacing:2px;background:var(--gradt);-webkit-background-clip:text;-webkit-text-fill-color:transparent}
.hr{display:flex;align-items:center;gap:7px}
.mpill{display:flex;align-items:center;gap:6px;background:rgba(255,255,255,.06);backdrop-filter:var(--blur2);border:1px solid var(--glass-bdr);border-radius:20px;padding:5px 12px;font-size:11px;color:var(--tx2);font-family:'DM Mono',monospace;max-width:180px}
.mpill span{overflow:hidden;text-overflow:ellipsis;white-space:nowrap}
.ldot{width:7px;height:7px;background:var(--green);border-radius:50%;box-shadow:0 0 6px var(--green);animation:pulse 2s infinite;flex-shrink:0}
@keyframes pulse{0%,100%{opacity:1}50%{opacity:.35}}
.ibtn{width:30px;height:30px;background:rgba(255,255,255,.06);backdrop-filter:var(--blur2);border:1px solid var(--glass-bdr);border-radius:8px;display:flex;align-items:center;justify-content:center;cursor:pointer;font-size:13px;transition:all .2s;color:var(--tx2)}
.ibtn:hover{border-color:var(--blue);color:var(--blue);background:rgba(6,182,212,.1)}
.ibtn.dev-btn{border-color:rgba(124,58,237,.4);color:var(--red);background:rgba(124,58,237,.08)}
.ibtn.dev-btn:hover{border-color:var(--red);background:rgba(124,58,237,.18)}
.uchip{display:flex;align-items:center;gap:6px;background:rgba(255,255,255,.06);backdrop-filter:var(--blur2);border:1px solid var(--glass-bdr);border-radius:20px;padding:4px 10px 4px 4px;cursor:pointer;font-size:12px;color:var(--tx2);transition:all .2s;position:relative;user-select:none}
.uchip:hover{border-color:var(--red)}
.uav{width:22px;height:22px;background:var(--grad);border-radius:50%;display:flex;align-items:center;justify-content:center;font-size:10px;color:#fff;font-weight:700;flex-shrink:0}
.drop{position:absolute;top:calc(100% + 8px);right:0;background:rgba(10,10,20,0.95);backdrop-filter:var(--blur);-webkit-backdrop-filter:var(--blur);border:1px solid var(--glass-bdr2);border-radius:14px;box-shadow:0 18px 48px rgba(0,0,0,.7);min-width:170px;z-index:200;overflow:hidden;display:none}
.drop.open{display:block;animation:slideup .14s ease}
.ditem{padding:10px 14px;font-size:13px;color:var(--tx2);cursor:pointer;display:flex;align-items:center;gap:9px;transition:background .14s}
.ditem:hover{background:rgba(255,255,255,.07);color:var(--tx)}
.ditem.danger:hover{background:rgba(124,58,237,.12);color:var(--red)}
.ditem.dev-item{color:var(--red);opacity:.85}
.ditem.dev-item:hover{background:rgba(124,58,237,.12);opacity:1}
.dsep{height:1px;background:var(--glass-bdr);margin:3px 0}
.ovl{position:fixed;inset:0;background:rgba(0,0,0,.4);z-index:250;display:none;backdrop-filter:blur(4px)}
.ovl.open{display:block}
.sp{position:fixed;top:0;right:-520px;width:520px;max-width:100vw;height:100vh;height:100dvh;background:rgba(6,10,22,0.97);backdrop-filter:blur(28px);border-left:1px solid var(--glass-bdr);z-index:300;box-shadow:-12px 0 48px rgba(0,0,0,.5);transition:right .28s cubic-bezier(.16,1,.3,1);display:flex;flex-direction:column}
.sp.open{right:0}
.sp-head{padding:14px 16px 10px;border-bottom:1px solid var(--glass-bdr);flex-shrink:0;display:flex;align-items:center;justify-content:space-between}
.sp-head h3{font-size:13px;font-weight:700;color:var(--tx)}
.sp-close{background:none;border:none;color:var(--tx3);font-size:18px;cursor:pointer;padding:4px 7px;border-radius:6px;line-height:1}
.sp-close:hover{color:var(--tx);background:rgba(255,255,255,.07)}
.sp-inner{display:flex;flex:1;overflow:hidden;min-height:0}
.sp-tabs{display:flex;flex-direction:column;gap:2px;padding:10px 6px;width:108px;min-width:108px;border-right:1px solid var(--glass-bdr);overflow-y:auto;flex-shrink:0;scrollbar-width:none;background:rgba(255,255,255,.015)}
.sp-tabs::-webkit-scrollbar{display:none}
.sp-tab{flex:none;display:flex;align-items:center;gap:7px;padding:8px 10px;background:none;border:none;border-left:3px solid transparent;color:var(--tx3);font-family:'DM Sans',sans-serif;font-size:11px;cursor:pointer;transition:all .15s;text-align:left;white-space:nowrap;width:100%;border-radius:0 8px 8px 0}
.sp-tab:hover{background:rgba(255,255,255,.07);color:var(--tx2);border-left-color:rgba(255,255,255,.15)}
.sp-tab.active{background:rgba(124,58,237,.12);color:#fff;font-weight:600;border-left-color:var(--red)}
.sp-tab .st-ico{font-size:14px;flex-shrink:0;width:18px;text-align:center}
.sp-tab .st-lbl{font-size:11px;flex:1}
.sp-tab-sep{height:1px;background:var(--glass-bdr);margin:4px 8px;flex-shrink:0}
.sp-body{flex:1;overflow-y:auto;padding:4px 14px 20px;scrollbar-width:thin}
.sp-body::-webkit-scrollbar{width:3px}
.sp-body::-webkit-scrollbar-thumb{background:rgba(100,140,220,.2);border-radius:3px}
.ksave{background:var(--grad);border:none;border-radius:9px;padding:8px 13px;color:#fff;font-size:11px;font-weight:600;cursor:pointer;white-space:nowrap}
.ksave:hover{opacity:.85}
.kdel{background:rgba(124,58,237,.12);border:1px solid rgba(124,58,237,.25);border-radius:9px;padding:8px 10px;color:var(--red);font-size:11px;cursor:pointer}
.msel{width:100%;background:var(--inp);border:1.5px solid var(--glass-bdr);border-radius:11px;padding:9px 12px;color:var(--tx);font-family:'DM Mono',monospace;font-size:11px;outline:none;cursor:pointer;transition:border-color .2s}
.msel:focus{border-color:var(--blue)}
.msel option,.msel optgroup{background:#0a0e1c}
.saved-item{background:rgba(255,255,255,.04);border:1px solid var(--glass-bdr);border-radius:12px;padding:12px 14px;margin-bottom:9px;cursor:pointer;transition:all .2s;position:relative}
.saved-item:hover{border-color:var(--blue);background:rgba(6,182,212,.06)}
.si-title{font-size:13px;font-weight:600;margin-bottom:4px;padding-right:28px;overflow:hidden;text-overflow:ellipsis;white-space:nowrap}
.si-preview{font-size:11px;color:var(--tx2);overflow:hidden;text-overflow:ellipsis;white-space:nowrap}
.si-ts{font-size:10px;color:var(--tx3);margin-top:4px}
.si-del{position:absolute;top:10px;right:10px;background:none;border:none;color:var(--tx3);cursor:pointer;font-size:15px;padding:2px;opacity:0;transition:opacity .2s}
.saved-item:hover .si-del{opacity:1}
.si-del:hover{color:var(--red)}
.empty-state{text-align:center;padding:32px 16px;color:var(--tx3);font-size:13px}
.empty-state .ei{font-size:32px;display:block;margin-bottom:8px}
.mem-item{display:flex;align-items:flex-start;gap:10px;padding:9px 12px;background:rgba(255,255,255,.04);border:1px solid var(--glass-bdr);border-radius:10px;margin-bottom:7px}
.mem-key{font-size:11px;font-weight:600;color:var(--blue);min-width:80px;flex-shrink:0;font-family:'DM Mono',monospace;padding-top:1px}
.mem-val{font-size:12px;color:var(--tx2);flex:1;line-height:1.5}
.mem-del{background:none;border:none;color:var(--tx3);cursor:pointer;font-size:13px;flex-shrink:0;padding:2px;transition:color .2s}
.mem-del:hover{color:var(--red)}
.mem-add-row{display:flex;gap:7px;margin-top:12px}
.mem-add-row input{flex:1;background:var(--inp);border:1.5px solid var(--glass-bdr);border-radius:9px;padding:8px 10px;color:var(--tx);font-family:'DM Mono',monospace;font-size:11px;outline:none}
.mem-add-row input:focus{border-color:var(--blue)}
/* Image Gen Panel */
.img-gen-panel{background:rgba(255,255,255,.04);border:1px solid var(--glass-bdr);border-radius:12px;padding:14px;margin-bottom:12px}
.img-gen-panel h4{font-size:11px;font-weight:700;color:var(--tx2);margin-bottom:10px;text-transform:uppercase;letter-spacing:.8px}
.ig-prompt{width:100%;background:var(--inp);border:1.5px solid var(--glass-bdr);border-radius:9px;padding:9px 12px;color:var(--tx);font-family:'DM Sans',sans-serif;font-size:12px;outline:none;resize:vertical;min-height:60px;margin-bottom:8px}
.ig-prompt:focus{border-color:var(--blue)}
.ig-row{display:flex;gap:7px;margin-bottom:8px;flex-wrap:wrap}
.ig-sel{flex:1;min-width:80px;background:var(--inp);border:1.5px solid var(--glass-bdr);border-radius:9px;padding:7px 10px;color:var(--tx);font-family:'DM Mono',monospace;font-size:11px;outline:none;cursor:pointer}
.ig-btn{background:var(--grad);border:none;border-radius:9px;padding:9px 16px;color:#fff;font-size:12px;font-weight:600;cursor:pointer;white-space:nowrap;width:100%}
.ig-btn:hover{opacity:.88}
.ig-btn:disabled{opacity:.4;cursor:not-allowed}
.ig-result{margin-top:10px;display:none!important}
.gen-img-card{border-radius:14px;overflow:hidden;border:1px solid var(--glass-bdr2);background:var(--glass-bg2);max-width:480px;margin:4px 0}
.gen-img-preview{width:100%;display:block;cursor:zoom-in;transition:opacity .2s}
.gen-img-preview:hover{opacity:.92}
.gen-img-actions{display:flex;align-items:center;gap:8px;padding:10px 12px;background:rgba(0,0,0,.22);flex-wrap:wrap}
.gen-img-btn{background:var(--grad);border:none;border-radius:8px;padding:7px 14px;color:#fff;font-size:12px;font-weight:600;cursor:pointer;transition:opacity .2s}
.gen-img-btn:hover{opacity:.85}
.gen-img-btn-sec{background:rgba(255,255,255,.12);border:1px solid var(--glass-bdr2)}
.gen-img-meta{font-size:10px;color:var(--tx3);margin-left:auto;font-family:'DM Mono',monospace}
.chat-opts-panel{position:fixed;top:52px;right:170px;background:rgba(10,10,20,0.97);backdrop-filter:var(--blur);-webkit-backdrop-filter:var(--blur);border:1px solid var(--glass-bdr2);border-radius:14px;box-shadow:0 18px 48px rgba(0,0,0,.7);min-width:170px;z-index:200;overflow:hidden;animation:slideup .14s ease}
.co-item{padding:10px 16px;font-size:13px;color:var(--tx2);cursor:pointer;transition:background .14s}
.co-item:hover{background:rgba(255,255,255,.07);color:var(--tx)}
.co-sep{height:1px;background:var(--glass-bdr);margin:3px 0}
.gen-stop-btn{background:rgba(124,58,237,.15);border:1px solid rgba(124,58,237,.4);color:var(--red);border-radius:9px;padding:8px 14px;font-size:12px;font-weight:600;cursor:pointer;width:100%;margin-top:6px;display:none}
.gen-stop-btn:hover{background:rgba(124,58,237,.28)}
.audio-result{margin-top:10px;display:none}
.audio-result.show{display:block}
.audio-result audio{width:100%;border-radius:8px;margin-bottom:6px}
.ig-result img{max-width:100%;border-radius:10px;border:1px solid var(--glass-bdr2);display:block;margin-bottom:8px}
.ig-result.show{display:block!important}
.ig-img{width:100%;border-radius:10px;border:1px solid var(--glass-bdr2);display:block}
.ig-dl{display:flex;gap:6px;margin-top:7px}
.ig-dl a,.ig-dl button{flex:1;text-align:center;padding:7px;background:rgba(255,255,255,.06);border:1px solid var(--glass-bdr);border-radius:8px;color:var(--tx2);font-size:11px;cursor:pointer;text-decoration:none;transition:all .2s}
.ig-dl a:hover,.ig-dl button:hover{border-color:var(--blue);color:var(--tx)}
/* Audio */
.audio-panel{background:rgba(255,255,255,.04);border:1px solid var(--glass-bdr);border-radius:12px;padding:14px;margin-bottom:12px}
.audio-panel h4{font-size:11px;font-weight:700;color:var(--tx2);margin-bottom:10px;text-transform:uppercase;letter-spacing:.8px}
.audio-result{margin-top:10px;display:none}
.audio-result.show{display:block}
audio{width:100%;border-radius:8px;margin-top:6px}
/* 3D panel */
.threed-viewer-wrap{background:rgba(34,197,94,.05);border:1px solid rgba(34,197,94,.15);border-radius:10px;padding:14px;text-align:center;margin-bottom:8px}
.dur-row{display:flex;align-items:center;gap:10px;margin-bottom:8px}
.dur-row label{font-size:10px;color:var(--tx3);white-space:nowrap}
.dur-row input[type=range]{flex:1;accent-color:var(--red)}
.dur-lbl{font-size:11px;color:var(--tx2);min-width:28px;font-family:'DM Mono',monospace}
/* Voice chat */
.voice-chat-panel{background:rgba(124,58,237,.06);border:1px solid rgba(124,58,237,.2);border-radius:12px;padding:14px;margin-bottom:12px}
.vc-status{font-size:11px;color:var(--tx3);margin-bottom:10px;line-height:1.6}
.vc-btn{width:100%;padding:11px;background:rgba(124,58,237,.12);border:1.5px solid rgba(124,58,237,.3);border-radius:10px;color:var(--red);font-family:'DM Sans',sans-serif;font-size:13px;font-weight:600;cursor:pointer;transition:all .2s;display:flex;align-items:center;justify-content:center;gap:8px}
.vc-btn:hover{background:rgba(124,58,237,.2);border-color:var(--red)}
.vc-btn.active{background:var(--red);color:#fff;animation:vcpulse 1.5s infinite}
@keyframes vcpulse{0%,100%{box-shadow:0 0 0 0 rgba(124,58,237,.4)}50%{box-shadow:0 0 0 8px rgba(124,58,237,0)}}
.vc-transcript{min-height:36px;background:rgba(255,255,255,.04);border:1px solid var(--glass-bdr);border-radius:8px;padding:8px 11px;font-size:12px;color:var(--tx2);margin-top:8px;font-style:italic}
/* Chat */
.chat-body{flex:1;overflow-y:auto;padding:22px 16px;scroll-behavior:smooth}
.chat-body::-webkit-scrollbar{width:3px}
.chat-body::-webkit-scrollbar-thumb{background:rgba(100,140,220,.2);border-radius:3px}
.msgs{max-width:780px;margin:0 auto}
.welcome{display:flex;flex-direction:column;align-items:center;text-align:center;padding:30px 0 16px;animation:slideup .5s ease}
.wico{width:72px;height:72px;background:var(--grad);border-radius:20px;display:flex;align-items:center;justify-content:center;font-size:34px;margin-bottom:16px;box-shadow:0 16px 48px var(--redg);animation:float 4s ease-in-out infinite}
@keyframes float{0%,100%{transform:translateY(0)}50%{transform:translateY(-8px)}}
.welcome h1{font-family:'Bebas Neue',sans-serif;font-size:42px;letter-spacing:4px;background:var(--gradt);-webkit-background-clip:text;-webkit-text-fill-color:transparent;margin-bottom:8px}
.welcome>p{color:var(--tx2);font-size:13px;max-width:400px;line-height:1.75;margin-bottom:6px}
.wsubt{font-family:'Bebas Neue',sans-serif;font-size:17px;letter-spacing:2px;background:var(--gradt);-webkit-background-clip:text;-webkit-text-fill-color:transparent;margin-bottom:24px}
.chips{display:flex;flex-wrap:wrap;gap:8px;justify-content:center;max-width:600px}
.chip{background:rgba(255,255,255,.05);backdrop-filter:var(--blur2);border:1px solid var(--glass-bdr);border-radius:22px;padding:8px 15px;font-size:12px;color:var(--tx2);cursor:pointer;transition:all .2s;white-space:nowrap}
.chip:hover{border-color:var(--red);color:var(--tx);background:rgba(124,58,237,.1);transform:translateY(-2px)}
.msg{display:flex;gap:10px;margin-bottom:20px;animation:fi .28s ease}
@keyframes fi{from{opacity:0;transform:translateY(8px)}to{opacity:1;transform:none}}
.msg.user{flex-direction:row-reverse}
.mav{width:30px;height:30px;border-radius:8px;display:flex;align-items:center;justify-content:center;font-size:12px;flex-shrink:0;align-self:flex-start;margin-top:2px}
.msg.ai .mav{background:var(--grad);box-shadow:0 4px 12px var(--redg)}
.msg.user .mav{background:rgba(255,255,255,.08);border:1px solid var(--glass-bdr);font-size:10px;font-weight:700;color:var(--tx2)}
.mcont{flex:1;min-width:0}
.mname{font-size:10px;font-weight:600;color:var(--tx3);text-transform:uppercase;letter-spacing:.9px;margin-bottom:4px;display:flex;align-items:center;gap:7px}
.msg.user .mname{justify-content:flex-end}
.mtag{background:rgba(255,255,255,.07);border:1px solid var(--glass-bdr2);border-radius:20px;padding:2px 8px;font-size:9px;font-weight:500;color:var(--tx2);font-family:'DM Mono',monospace}
.bbl{display:inline-block;padding:11px 15px;border-radius:15px;font-size:14px;line-height:1.78;max-width:100%;word-break:break-word}
.msg.ai .bbl{background:rgba(10,16,36,0.92);backdrop-filter:blur(20px) saturate(170%);-webkit-backdrop-filter:blur(20px) saturate(170%);border:1px solid rgba(60,100,200,0.22);border-top-left-radius:6px;box-shadow:inset 0 1px 0 rgba(255,255,255,.06),0 2px 16px rgba(0,0,0,.28)}
.msg.user .bbl{background:linear-gradient(135deg,var(--red2),var(--red));color:#fff;border-top-right-radius:4px;box-shadow:0 4px 16px var(--redg);float:right}
.msg.user .mcont{display:flex;flex-direction:column;align-items:flex-end}
.bbl pre{background:rgba(0,0,0,.5);border:1px solid rgba(124,58,237,.2);border-radius:9px;padding:12px;margin:9px 0;overflow-x:auto;font-size:12px;font-family:'DM Mono',monospace;line-height:1.65}
.bbl code{background:rgba(6,182,212,.15);color:#70b5fa;padding:2px 5px;border-radius:5px;font-family:'DM Mono',monospace;font-size:12px}
.bbl pre code{background:none;color:var(--tx);padding:0}
.rtag{font-size:10px;color:var(--tx3);margin-top:4px;font-style:italic}
.msg-img{max-width:240px;max-height:180px;border-radius:9px;border:1px solid var(--glass-bdr);margin-bottom:7px;display:block}
.msg-actions{display:flex;gap:6px;margin-top:5px;opacity:0;transition:opacity .2s;flex-wrap:wrap}
.msg:hover .msg-actions{opacity:1}
.mact{background:rgba(255,255,255,.05);border:1px solid var(--glass-bdr);border-radius:7px;padding:4px 9px;font-size:11px;color:var(--tx3);cursor:pointer;transition:all .2s}
.mact:hover{color:var(--tx);border-color:var(--blue)}
.gen-img-msg{border-radius:12px;border:1px solid var(--glass-bdr2);max-width:320px;display:block;margin:8px 0}
/* Thinking */
.thinking-wrap{display:flex;align-items:center;gap:10px;padding:10px 0}
.thinking-dots{display:flex;gap:5px;align-items:center}
.thinking-dots span{width:8px;height:8px;border-radius:50%;animation:td 1.4s infinite ease-in-out}
.thinking-dots span:nth-child(1){background:var(--red);animation-delay:0s}
.thinking-dots span:nth-child(2){background:#9b45be;animation-delay:.2s}
.thinking-dots span:nth-child(3){background:var(--blue);animation-delay:.4s}
@keyframes td{0%,60%,100%{transform:translateY(0) scale(1);opacity:.4}30%{transform:translateY(-9px) scale(1.1);opacity:1}}
.thinking-txt{font-size:12px;color:var(--tx3);font-style:italic;animation:blink 2s infinite}
@keyframes blink{0%,100%{opacity:.5}50%{opacity:1}}
.streaming .bbl::after{content:'\25ae';display:inline-block;animation:cur .7s infinite;color:var(--blue);margin-left:1px;font-size:13px}
@keyframes cur{0%,100%{opacity:1}50%{opacity:0}}
.thought{margin-bottom:10px;border-radius:12px;overflow:hidden;border:1px solid rgba(120,80,220,.25);background:rgba(80,40,160,.08)}
.thought-hdr{display:flex;align-items:center;gap:8px;padding:8px 13px;cursor:pointer;user-select:none;transition:background .2s}
.thought-hdr:hover{background:rgba(120,80,220,.1)}
.thought-icon{font-size:14px}
.thought-label{font-size:11px;font-weight:600;color:rgba(160,120,255,.85);letter-spacing:.5px;flex:1}
.thought-toggle{font-size:11px;color:rgba(160,120,255,.5);transition:transform .25s}
.thought-toggle.open{transform:rotate(180deg)}
.thought-body{max-height:0;overflow:hidden;transition:max-height .35s cubic-bezier(.16,1,.3,1)}
.thought-body.open{max-height:1200px}
.thought-body-open{max-height:1200px!important}
.thought-txt{padding:10px 13px 13px;font-size:12px;color:rgba(180,150,255,.7);line-height:1.75;font-family:'DM Mono',monospace;white-space:pre-wrap;border-top:1px solid rgba(120,80,220,.15)}
.thought-streaming .thought-hdr{background:rgba(100,60,200,.12)}
.thought-dots span{animation:dotpulse 1.2s infinite;opacity:0}
.thought-dots span:nth-child(2){animation-delay:.2s}
.thought-dots span:nth-child(3){animation-delay:.4s}
@keyframes dotpulse{0%,80%,100%{opacity:0}40%{opacity:1}}
/* Input zone */
.iz{flex-shrink:0;padding:10px 16px 14px;background:rgba(5,7,16,0.8);backdrop-filter:var(--blur);-webkit-backdrop-filter:var(--blur);border-top:1px solid var(--glass-bdr);box-shadow:0 -1px 0 rgba(255,255,255,.04),0 -8px 32px rgba(0,0,0,.2)}
.ibox{max-width:800px;margin:0 auto;background:rgba(8,12,26,0.75);backdrop-filter:var(--blur2);-webkit-backdrop-filter:var(--blur2);border:1.5px solid var(--glass-bdr);border-radius:20px;display:flex;align-items:flex-end;gap:8px;padding:8px 10px;transition:border-color .25s,box-shadow .25s}
.ibox:focus-within{border-color:rgba(6,182,212,.5);box-shadow:0 0 0 3px rgba(6,182,212,.1)}
textarea{flex:1;background:none;border:none;color:var(--tx);font-family:'DM Sans',sans-serif;font-size:14px;resize:none;outline:none;min-height:24px;max-height:140px;line-height:1.65;padding:3px 0}
textarea::placeholder{color:var(--tx3)}
.ib-row{display:flex;align-items:flex-end;gap:5px}
.iibtn{width:32px;height:32px;background:rgba(255,255,255,.07);border:1px solid var(--glass-bdr);border-radius:8px;display:flex;align-items:center;justify-content:center;cursor:pointer;font-size:14px;transition:all .2s;color:var(--tx2);flex-shrink:0}
.iibtn:hover{border-color:var(--blue);color:var(--blue);background:rgba(6,182,212,.1)}
.iibtn.active{border-color:var(--red);color:var(--red);background:rgba(124,58,237,.12);animation:micpulse 1s infinite}
@keyframes micpulse{0%,100%{box-shadow:0 0 0 0 rgba(124,58,237,.4)}50%{box-shadow:0 0 0 5px rgba(124,58,237,0)}}
.sbtn{width:36px;height:36px;background:var(--grad);border:none;border-radius:10px;color:#fff;font-size:15px;cursor:pointer;display:flex;align-items:center;justify-content:center;flex-shrink:0;transition:all .2s;box-shadow:0 4px 16px var(--redg)}
.sbtn:hover{transform:scale(1.08)}
.sbtn:disabled{opacity:.28;cursor:not-allowed;transform:none;box-shadow:none}
.attach-preview{max-width:780px;margin:0 auto 6px;display:none;align-items:center;gap:10px;padding:7px 12px;background:rgba(255,255,255,.04);border:1px solid var(--glass-bdr);border-radius:12px}
.attach-preview.show{display:flex}
.ap-img{width:48px;height:48px;object-fit:cover;border-radius:7px;border:1px solid var(--glass-bdr2)}
.ap-name{font-size:12px;color:var(--tx2);flex:1;overflow:hidden;text-overflow:ellipsis;white-space:nowrap}
.ap-del{background:none;border:none;color:var(--tx3);font-size:16px;cursor:pointer;transition:color .2s;padding:0 3px}
.ap-del:hover{color:var(--red)}
.ifooter{max-width:780px;margin:4px auto 0;display:flex;justify-content:space-between;align-items:center}
.hint{font-size:10px;color:var(--tx3)}
.tbtn{background:none;border:none;color:var(--tx3);font-family:'DM Sans',sans-serif;font-size:11px;cursor:pointer;padding:2px 6px;border-radius:5px;transition:all .2s}
.tbtn:hover{color:var(--red)}
.voice-bar{display:none;max-width:780px;margin:0 auto 6px;align-items:center;gap:10px;padding:7px 13px;background:rgba(124,58,237,.08);border:1px solid rgba(124,58,237,.25);border-radius:12px}
.voice-bar.show{display:flex}
.vbars{display:flex;gap:3px;align-items:center;height:18px}
.vbar{width:3px;background:var(--red);border-radius:2px;animation:vwave .8s ease-in-out infinite}
.vbar:nth-child(1){height:5px}.vbar:nth-child(2){height:12px}.vbar:nth-child(3){height:18px}.vbar:nth-child(4){height:12px}.vbar:nth-child(5){height:5px}
@keyframes vwave{0%,100%{transform:scaleY(.4)}50%{transform:scaleY(1)}}
.vtxt{font-size:12px;color:var(--red);flex:1}
.vstop{background:rgba(124,58,237,.12);border:none;border-radius:6px;color:var(--red);font-size:12px;cursor:pointer;padding:3px 9px}
.err-overlay{position:fixed;top:68px;left:50%;transform:translateX(-50%);max-width:520px;width:calc(100% - 28px);background:rgba(18,4,4,0.9);backdrop-filter:blur(20px);-webkit-backdrop-filter:blur(20px);border:1px solid rgba(124,58,237,.35);border-radius:13px;padding:12px 16px;display:flex;align-items:center;gap:11px;z-index:400;pointer-events:none;opacity:0;transition:opacity .3s,transform .3s;transform:translateX(-50%) translateY(-8px);box-shadow:0 8px 32px rgba(124,58,237,.2)}
.err-overlay.show{opacity:1;transform:translateX(-50%) translateY(0);pointer-events:auto}
.err-overlay .et{font-size:13px;color:#ffb3b3;flex:1;line-height:1.5}
.err-overlay .ec{background:none;border:none;color:rgba(255,180,180,.6);font-size:16px;cursor:pointer;flex-shrink:0}
.toggle-switch input{opacity:0;width:0;height:0}
input:checked+
input:checked+
.toast{position:fixed;bottom:24px;left:50%;transform:translateX(-50%) translateY(14px);background:rgba(10,15,30,0.95);backdrop-filter:var(--blur);border:1px solid var(--glass-bdr2);border-radius:12px;padding:10px 20px;font-size:13px;color:var(--tx);box-shadow:0 12px 36px rgba(0,0,0,.6);z-index:600;opacity:0;transition:all .28s;pointer-events:none;white-space:nowrap}
.toast.show{opacity:1;transform:translateX(-50%) translateY(0)}
/* Dev modal */
.modal-bg{position:fixed;inset:0;background:rgba(0,0,0,.65);z-index:500;display:none;align-items:center;justify-content:center;backdrop-filter:blur(8px)}
.modal-bg.open{display:flex;animation:slideup .18s ease}
.modal{background:rgba(6,10,22,0.98);backdrop-filter:var(--blur);border:1px solid var(--glass-bdr2);border-radius:20px;padding:0;width:780px;max-width:96vw;max-height:92vh;overflow:hidden;position:relative;box-shadow:0 40px 80px rgba(0,0,0,.7);display:flex;flex-direction:column}
.modal-header{padding:20px 24px 0;border-bottom:1px solid var(--glass-bdr);flex-shrink:0}
.modal-header h2{font-size:18px;font-weight:700;margin-bottom:3px;background:var(--gradt);-webkit-background-clip:text;-webkit-text-fill-color:transparent}
.modal-header p{color:var(--tx2);font-size:12px;margin-bottom:14px}
.modal-close{position:absolute;top:16px;right:18px;background:none;border:none;color:var(--tx2);font-size:20px;cursor:pointer}
.modal-tabs{display:flex;gap:0;border-bottom:1px solid var(--glass-bdr);flex-shrink:0;padding:0 24px;overflow-x:auto}
.modal-tab{padding:10px 14px;background:none;border:none;border-bottom:2px solid transparent;color:var(--tx3);font-family:'DM Sans',sans-serif;font-size:12px;font-weight:600;cursor:pointer;transition:all .2s;white-space:nowrap}
.modal-tab.active{border-bottom-color:var(--red);color:var(--tx)}
.modal-body{flex:1;overflow-y:auto;padding:20px 24px;scrollbar-width:thin}
.modal-body::-webkit-scrollbar{width:4px}
.modal-body::-webkit-scrollbar-thumb{background:rgba(100,140,220,.2)}
.stat-grid{display:grid;grid-template-columns:repeat(auto-fit,minmax(130px,1fr));gap:10px;margin-bottom:18px}
.stat-card{background:rgba(255,255,255,.04);border:1px solid var(--glass-bdr);border-radius:12px;padding:14px 16px}
.stat-card label{font-size:9px;color:var(--tx3);text-transform:uppercase;letter-spacing:1px;font-weight:600;display:block;margin-bottom:5px}
.stat-card .val{font-size:26px;font-weight:700;font-family:'Bebas Neue',sans-serif;letter-spacing:1px;background:var(--gradt);-webkit-background-clip:text;-webkit-text-fill-color:transparent}
.stat-card .sub{font-size:10px;color:var(--tx3);margin-top:2px}
.dev-list{background:rgba(255,255,255,.04);border:1px solid var(--glass-bdr);border-radius:11px;padding:12px 14px;margin-bottom:12px}
.dev-list h5{font-size:9px;color:var(--tx3);text-transform:uppercase;letter-spacing:.8px;margin-bottom:10px;font-weight:600}
.dev-row{display:flex;justify-content:space-between;align-items:center;font-size:12px;padding:6px 0;border-bottom:1px solid rgba(124,58,237,.08);color:var(--tx2);gap:8px}
.dev-row:last-child{border-bottom:none}
.dev-row .k{font-family:'DM Mono',monospace;color:var(--tx);font-size:11px;flex:1;overflow:hidden;text-overflow:ellipsis;white-space:nowrap}
.dev-row .v{font-size:11px;color:var(--tx3);flex-shrink:0}
.dev-row .ok{color:var(--green);font-weight:600}
.dev-row .no{color:var(--red);font-weight:600}
.model-test-row{display:flex;align-items:center;gap:10px;padding:8px 10px;background:rgba(255,255,255,.03);border:1px solid var(--glass-bdr);border-radius:9px;margin-bottom:6px}
.model-test-row.ok{border-color:rgba(34,197,94,.2)}
.model-test-row.error,.model-test-row.exception{border-color:rgba(124,58,237,.2)}
.model-test-row.no_key{opacity:.4}
.mtbadge{font-size:9px;font-weight:700;padding:2px 8px;border-radius:20px;white-space:nowrap;flex-shrink:0}
.mtbadge.ok{background:rgba(34,197,94,.15);color:#22c55e;border:1px solid rgba(34,197,94,.3)}
.mtbadge.error,.mtbadge.exception{background:rgba(124,58,237,.15);color:var(--red);border:1px solid rgba(124,58,237,.3)}
.mtbadge.no_key{background:rgba(100,120,160,.12);color:var(--tx3)}
.mtlabel{font-size:12px;font-weight:600;flex:1}
.mtdetail{font-size:10px;color:var(--tx3);font-family:'DM Mono',monospace;flex:2;overflow:hidden;text-overflow:ellipsis;white-space:nowrap}
.mtms{font-size:10px;color:var(--tx3);width:50px;text-align:right;flex-shrink:0}
.overseer-box{background:rgba(80,40,160,.1);border:1px solid rgba(120,80,220,.25);border-radius:12px;padding:14px 16px;margin-bottom:16px}
.overseer-box h4{font-size:10px;font-weight:700;color:rgba(160,120,255,.9);text-transform:uppercase;letter-spacing:1px;margin-bottom:8px}
.overseer-box p{font-size:12px;color:rgba(200,180,255,.8);line-height:1.75;white-space:pre-wrap;font-family:'DM Mono',monospace}
.run-test-btn{background:var(--grad);border:none;border-radius:10px;padding:10px 20px;color:#fff;font-size:13px;font-weight:600;cursor:pointer;margin-bottom:16px;transition:all .2s}
.run-test-btn:hover{opacity:.88}
.run-test-btn:disabled{opacity:.4;cursor:not-allowed}
/* Theme Panel */
.theme-mode-btn{flex:1;padding:9px;border-radius:9px;border:1.5px solid var(--glass-bdr);background:rgba(255,255,255,.06);color:var(--tx2);cursor:pointer;font-size:12px;font-weight:600;transition:all .2s}
.theme-mode-btn.active{border-color:var(--red)!important;color:var(--tx)!important;background:rgba(255,255,255,.1)!important}
.accent-swatch{width:30px;height:30px;border-radius:8px;cursor:pointer;border:2px solid transparent;transition:all .2s}
.accent-swatch.active{border-color:white!important;transform:scale(1.15)!important}
.bg-swatch{height:50px;border-radius:9px;cursor:pointer;border:2px solid transparent;font-size:10px;font-weight:600;letter-spacing:.5px;display:flex;align-items:center;justify-content:center;transition:all .2s}
.bg-swatch.active{border-color:var(--red)!important;transform:scale(1.04)!important}
/* Overseer Q&A */
.overseer-qa{background:rgba(120,80,220,.1);border:1px solid rgba(120,80,220,.25);border-radius:14px;padding:14px 16px;margin-bottom:8px;animation:fi .3s ease}
.overseer-qa-title{font-size:11px;font-weight:700;color:rgba(180,150,255,.9);text-transform:uppercase;letter-spacing:.8px;margin-bottom:10px;display:flex;align-items:center;gap:7px}
.overseer-q{margin-bottom:12px}
.overseer-q-text{font-size:13px;color:var(--tx);margin-bottom:7px;font-weight:500}
.overseer-chips{display:flex;flex-wrap:wrap;gap:6px}
.overseer-chip{padding:5px 13px;background:rgba(255,255,255,.07);border:1.5px solid rgba(120,80,220,.3);border-radius:20px;font-size:12px;color:var(--tx2);cursor:pointer;transition:all .18s;user-select:none}
.overseer-chip:hover{border-color:rgba(160,120,255,.7);color:var(--tx);background:rgba(120,80,220,.15)}
.overseer-chip.selected{border-color:rgba(160,120,255,.9);color:#fff;background:rgba(120,80,220,.35)}
.overseer-generate-btn{width:100%;margin-top:10px;padding:10px;background:linear-gradient(135deg,rgba(120,80,220,.85),rgba(60,20,160,.85));border:none;border-radius:10px;color:#fff;font-family:'DM Sans',sans-serif;font-size:13px;font-weight:600;cursor:pointer;transition:all .2s;letter-spacing:.3px}
.overseer-generate-btn:hover{opacity:.9;transform:translateY(-1px)}
.overseer-skip{background:none;border:none;color:var(--tx3);font-size:11px;cursor:pointer;margin-top:6px;text-decoration:underline;width:100%;text-align:center;font-family:'DM Sans',sans-serif}
.overseer-skip:hover{color:var(--tx2)}
/* Overseer Modal */
.ov-modal-bg{position:fixed;inset:0;background:rgba(0,0,0,.72);z-index:800;display:flex;align-items:center;justify-content:center;backdrop-filter:blur(12px);animation:fadeIn .2s ease}
@keyframes fadeIn{from{opacity:0}to{opacity:1}}
.ov-modal{background:rgba(8,10,22,0.97);backdrop-filter:var(--blur);border:1px solid rgba(120,80,220,.4);border-radius:22px;width:520px;max-width:94vw;max-height:88vh;overflow:hidden;box-shadow:0 40px 100px rgba(0,0,0,.8),0 0 0 1px rgba(160,120,255,.08);display:flex;flex-direction:column;animation:slideup .25s cubic-bezier(.16,1,.3,1)}
.ov-modal-hdr{padding:22px 24px 0;flex-shrink:0}
.ov-modal-badge{display:inline-flex;align-items:center;gap:6px;background:rgba(120,80,220,.15);border:1px solid rgba(120,80,220,.35);border-radius:20px;padding:4px 12px;font-size:10px;font-weight:700;color:rgba(180,150,255,.9);text-transform:uppercase;letter-spacing:.8px;margin-bottom:12px}
.ov-modal-title{font-size:18px;font-weight:700;color:var(--tx);margin-bottom:5px}
.ov-modal-sub{font-size:12px;color:var(--tx2);margin-bottom:14px;line-height:1.6}
.ov-modal-prompt{background:rgba(120,80,220,.08);border:1px solid rgba(120,80,220,.2);border-radius:10px;padding:10px 13px;font-size:12px;color:rgba(180,150,255,.85);font-family:'DM Mono',monospace;margin-bottom:18px;word-break:break-word;line-height:1.5}
.ov-modal-body{flex:1;overflow-y:auto;padding:4px 24px 16px;scrollbar-width:thin}
.ov-modal-body::-webkit-scrollbar{width:3px}
.ov-modal-body::-webkit-scrollbar-thumb{background:rgba(120,80,220,.3);border-radius:3px}
.ov-q{margin-bottom:18px}
.ov-q-num{font-size:9px;font-weight:700;color:rgba(160,120,255,.55);text-transform:uppercase;letter-spacing:.8px;margin-bottom:5px}
.ov-q-text{font-size:13px;color:var(--tx);font-weight:600;margin-bottom:9px}
.ov-chips{display:flex;flex-wrap:wrap;gap:7px}
.ov-chip{padding:6px 14px;background:rgba(255,255,255,.06);border:1.5px solid rgba(120,80,220,.25);border-radius:22px;font-size:12px;color:var(--tx2);cursor:pointer;transition:all .18s;user-select:none}
.ov-chip:hover{border-color:rgba(160,120,255,.65);color:var(--tx);background:rgba(120,80,220,.12)}
.ov-chip.sel{border-color:rgba(160,120,255,.9);color:#fff;background:rgba(120,80,220,.42);box-shadow:0 0 0 3px rgba(120,80,220,.18)}
.ov-modal-footer{padding:14px 24px 22px;flex-shrink:0;display:flex;flex-direction:column;gap:9px;border-top:1px solid rgba(120,80,220,.15)}
.ov-gen-btn{padding:13px;background:linear-gradient(135deg,#7b50dc,#4a20a0);border:none;border-radius:12px;color:#fff;font-family:'DM Sans',sans-serif;font-size:14px;font-weight:600;cursor:pointer;transition:all .2s;letter-spacing:.3px;box-shadow:0 4px 20px rgba(120,80,220,.45)}
.ov-gen-btn:hover{opacity:.92;transform:translateY(-2px);box-shadow:0 8px 28px rgba(120,80,220,.55)}
.ov-skip-btn{padding:10px;background:none;border:1px solid var(--glass-bdr);border-radius:10px;color:var(--tx3);font-family:'DM Sans',sans-serif;font-size:12px;cursor:pointer;transition:all .2s}
.ov-skip-btn:hover{color:var(--tx2);border-color:var(--glass-bdr2)}
@media(max-width:600px){
 .card{padding:28px 20px;border-radius:20px}
 header{padding:7px 10px}
 .mpill{display:none}
 .ibtn{width:28px;height:28px;font-size:12px}
 #uName{max-width:55px;overflow:hidden;text-overflow:ellipsis;white-space:nowrap}
 .sp{width:100vw;right:-100vw}
 .sp.open{right:0}
 .chat-body{padding:12px 8px}
 .msg-actions{opacity:1}
 .iz{padding:6px 8px 8px}
 .ibox{padding:5px 7px;border-radius:14px}
 .toast{bottom:16px;font-size:12px;max-width:calc(100vw - 32px);white-space:normal;text-align:center}
 .modal{border-radius:16px}
 .stat-grid{grid-template-columns:1fr 1fr}
}
@media(hover:none){.msg-actions{opacity:1}.si-del{opacity:1}}
/* Code Blocks */
.code-block-wrap{border-radius:12px;border:1px solid var(--glass-bdr2);overflow:hidden;margin:8px 0;background:rgba(6,10,20,0.7)}
.code-block-header{display:flex;align-items:center;gap:8px;padding:7px 12px;background:rgba(255,255,255,.04);border-bottom:1px solid var(--glass-bdr);flex-wrap:wrap}
.code-lang-badge{font-family:'DM Mono',monospace;font-size:9px;font-weight:700;letter-spacing:1px;padding:2px 8px;border-radius:20px;background:var(--redg);border:1px solid var(--glass-bdr2);color:var(--tx2)}
.code-lines{font-size:10px;color:var(--tx3);font-family:'DM Mono',monospace;flex:1}
.code-copy-btn,.code-expand-btn{font-size:10px;background:rgba(255,255,255,.06);border:1px solid var(--glass-bdr);border-radius:6px;color:var(--tx2);cursor:pointer;padding:3px 9px;transition:all .15s;font-family:'DM Sans',sans-serif}
.code-copy-btn:hover,.code-expand-btn:hover{background:rgba(255,255,255,.12);color:var(--tx)}
.code-block-body{overflow:hidden;transition:max-height .3s ease}
.code-block-body.code-collapsed{max-height:0}
.code-block-wrap pre{margin:0;overflow-x:auto;padding:14px 16px}
.code-block-wrap code.hlcode{font-family:'DM Mono',monospace;font-size:13px;line-height:1.7;display:block}
code.inline-code{font-family:'DM Mono',monospace;font-size:12.5px;background:rgba(255,255,255,.09);border:1px solid var(--glass-bdr);border-radius:5px;padding:1px 6px}
/* ── Compact card for huge code blocks (>60 lines) ─────────────────────── */
.code-card{border-radius:12px;border:1px solid var(--glass-bdr2);overflow:hidden;margin:8px 0;background:rgba(6,10,20,0.7);cursor:pointer;transition:all .15s}
.code-card:hover{border-color:var(--accent);box-shadow:0 4px 18px rgba(124,58,237,.18);transform:translateY(-1px)}
.code-card-top{display:flex;align-items:center;gap:10px;padding:12px 14px;border-bottom:1px solid var(--glass-bdr)}
.code-card-icon{font-size:22px;flex-shrink:0}
.code-card-info{flex:1;min-width:0}
.code-card-title{font-size:13px;font-weight:700;color:var(--tx)}
.code-card-sub{font-size:10.5px;color:var(--tx3);font-family:'DM Mono',monospace;margin-top:2px;overflow:hidden;text-overflow:ellipsis;white-space:nowrap}
.code-card-open{font-size:11px;font-weight:700;color:var(--accent2);flex-shrink:0;padding:5px 10px;border:1px solid rgba(6,182,212,.35);border-radius:8px;background:rgba(6,182,212,.08)}
.code-card-preview{padding:10px 14px;background:rgba(0,0,0,.2);max-height:70px;overflow:hidden;position:relative}
.code-card-preview pre{margin:0;font-family:'DM Mono',monospace;font-size:11px;color:var(--tx3);line-height:1.6;white-space:pre-wrap;word-break:break-all}
.code-card-preview::after{content:'';position:absolute;bottom:0;left:0;right:0;height:30px;background:linear-gradient(transparent,rgba(6,10,20,0.9))}
/* ── Floating annotated code window ────────────────────────────────────── */
.code-win-overlay{position:fixed;inset:0;z-index:9995;background:rgba(2,4,12,.72);backdrop-filter:blur(8px);display:flex;align-items:center;justify-content:center;animation:fadeIn .15s ease}
.code-win{position:relative;width:min(880px,92vw);height:min(640px,84vh);background:rgba(8,8,16,.98);border:1px solid var(--glass-bdr2);border-radius:18px;box-shadow:0 40px 100px rgba(0,0,0,.85);display:flex;flex-direction:column;overflow:hidden;animation:slideup .18s cubic-bezier(.16,1,.3,1)}
.code-win-head{display:flex;align-items:center;justify-content:space-between;padding:12px 16px;background:rgba(124,58,237,.08);border-bottom:1px solid var(--glass-bdr);cursor:move;flex-shrink:0;user-select:none}
.code-win-title{font-size:13px;font-weight:700;color:var(--tx);display:flex;align-items:center;gap:8px}
.code-win-lines{font-size:10.5px;font-weight:600;color:var(--tx3);font-family:'DM Mono',monospace;background:rgba(255,255,255,.06);padding:2px 8px;border-radius:10px}
.code-win-actions{display:flex;gap:6px}
.cw-btn{font-size:11px;font-weight:600;background:rgba(255,255,255,.06);border:1px solid var(--glass-bdr);border-radius:7px;color:var(--tx2);cursor:pointer;padding:5px 11px;transition:all .15s;font-family:'DM Sans',sans-serif}
.cw-btn:hover{background:rgba(124,58,237,.18);color:var(--tx);border-color:var(--glass-bdr2)}
.cw-btn.cw-close:hover{background:rgba(239,68,68,.18);color:#f87171;border-color:rgba(239,68,68,.35)}
.code-win-outline{display:flex;flex-wrap:wrap;gap:6px;padding:8px 16px;border-bottom:1px solid var(--glass-bdr);background:rgba(0,0,0,.15);flex-shrink:0}
.cw-outline-chip{font-size:10px;font-family:'DM Mono',monospace;color:var(--accent2);background:rgba(6,182,212,.08);border:1px solid rgba(6,182,212,.2);padding:2px 8px;border-radius:10px}
.code-win-body{flex:1;overflow:auto;padding:14px 0}
.cw-pre{margin:0;padding:0 16px;font-family:'DM Mono',monospace;font-size:12.5px;line-height:1.75;color:var(--tx)}
.cw-lnum{display:inline-block;width:38px;color:var(--tx3);user-select:none;text-align:right;margin-right:14px;font-size:11px}
/* Syntax highlighting colors */
.hl-kw{color:#c792ea}.hl-str{color:#c3e88d}.hl-num{color:#f78c6c}.hl-cmt{color:#546e7a;font-style:italic}.hl-attr{color:#82aaff}
/* ══ Compact code card — replaces huge inline blocks, opens floating window ══ */
.code-card{display:flex;align-items:center;gap:12px;margin:10px 0;padding:13px 15px;border-radius:14px;border:1px solid var(--glass-bdr2);background:linear-gradient(135deg,rgba(124,58,237,.08),rgba(6,182,212,.05));cursor:pointer;transition:all .18s}
.code-card:hover{border-color:rgba(124,58,237,.5);background:linear-gradient(135deg,rgba(124,58,237,.14),rgba(6,182,212,.08));transform:translateY(-1px);box-shadow:0 6px 20px rgba(124,58,237,.15)}
.code-card-icon{width:38px;height:38px;border-radius:10px;background:var(--grad);display:flex;align-items:center;justify-content:center;font-size:17px;flex-shrink:0;font-family:'DM Mono',monospace;font-weight:800;color:#fff}
.code-card-body{flex:1;min-width:0}
.code-card-title{font-size:13px;font-weight:700;color:var(--tx);display:flex;align-items:center;gap:8px}
.code-card-title .clines{font-size:10px;font-weight:600;color:var(--tx3);font-family:'DM Mono',monospace;background:rgba(255,255,255,.06);padding:1px 7px;border-radius:20px}
.code-card-preview{font-size:11px;color:var(--tx3);font-family:'DM Mono',monospace;overflow:hidden;text-overflow:ellipsis;white-space:nowrap;margin-top:3px}
.code-card-actions{display:flex;gap:6px;flex-shrink:0}
.code-card-actions button{font-size:10.5px;background:rgba(255,255,255,.06);border:1px solid var(--glass-bdr);border-radius:8px;color:var(--tx2);cursor:pointer;padding:6px 11px;transition:all .15s;font-family:'DM Sans',sans-serif;font-weight:600}
.code-card-actions button:hover{background:rgba(124,58,237,.2);color:var(--tx);border-color:rgba(124,58,237,.4)}
.code-card-actions .cc-open{background:var(--grad);color:#fff;border:none}
.code-card-actions .cc-open:hover{opacity:.88}
/* Floating annotated code window */
.code-float-overlay{position:fixed;inset:0;z-index:9995;background:rgba(2,4,12,.75);backdrop-filter:blur(10px);display:none;align-items:center;justify-content:center;animation:fadeIn .18s}
.code-float-overlay.show{display:flex}
.code-float-win{width:min(920px,94vw);height:min(680px,86vh);background:rgba(8,8,16,0.98);border:1px solid rgba(124,58,237,.3);border-radius:18px;overflow:hidden;display:flex;flex-direction:column;box-shadow:0 50px 120px rgba(0,0,0,.7);animation:slideup .2s cubic-bezier(.16,1,.3,1)}
.code-float-head{display:flex;align-items:center;gap:10px;padding:12px 16px;background:rgba(124,58,237,.07);border-bottom:1px solid var(--glass-bdr2);cursor:move;flex-shrink:0;user-select:none}
.code-float-lang{font-family:'DM Mono',monospace;font-size:10px;font-weight:800;letter-spacing:1px;padding:3px 10px;border-radius:20px;background:var(--grad);color:#fff}
.code-float-lines{font-size:11px;color:var(--tx3);font-family:'DM Mono',monospace;flex:1}
.code-float-actions{display:flex;gap:6px}
.code-float-actions button{font-size:11px;background:rgba(255,255,255,.06);border:1px solid var(--glass-bdr);border-radius:8px;color:var(--tx2);cursor:pointer;padding:6px 12px;transition:all .15s;font-weight:600}
.code-float-actions button:hover{background:rgba(124,58,237,.2);color:var(--tx)}
.code-float-actions .cf-close:hover{background:rgba(239,68,68,.18);color:#f87171;border-color:rgba(239,68,68,.4)}
.code-float-body{flex:1;display:flex;overflow:hidden;min-height:0}
.code-float-outline{width:200px;flex-shrink:0;border-right:1px solid var(--glass-bdr);overflow-y:auto;padding:10px;background:rgba(0,0,0,.15)}
.code-float-outline-title{font-size:9px;font-weight:700;letter-spacing:1px;color:var(--tx3);text-transform:uppercase;margin-bottom:8px;padding:0 4px}
.cf-outline-item{display:flex;align-items:center;gap:6px;padding:6px 8px;border-radius:8px;cursor:pointer;font-size:11px;color:var(--tx2);transition:all .12s;margin-bottom:1px}
.cf-outline-item:hover{background:rgba(124,58,237,.14);color:var(--tx)}
.cf-outline-item .cft{font-size:9px;color:var(--purple);font-family:'DM Mono',monospace;flex-shrink:0}
.cf-outline-empty{font-size:10.5px;color:var(--tx3);padding:8px;line-height:1.6}
.code-float-code-wrap{flex:1;overflow:auto;padding:16px 18px}
.code-float-code-wrap pre{margin:0}
.code-float-code-wrap code.hlcode{font-family:'DM Mono',monospace;font-size:12.5px;line-height:1.75;display:block}
.cf-line-hl{background:rgba(124,58,237,.18);display:block;margin:0 -18px;padding:0 18px;border-left:2px solid var(--purple)}
@media(max-width:768px){.code-float-outline{display:none}.code-float-win{width:96vw;height:88vh}}
/* Thinking block improvements */
.thought{border-radius:12px;border:1px solid rgba(120,80,220,.2);background:rgba(80,40,180,.06);margin-bottom:10px;overflow:hidden}
.thought-hdr{display:flex;align-items:center;gap:8px;padding:9px 13px;cursor:pointer;user-select:none;background:rgba(120,80,220,.06)}
.thought-icon{font-size:14px}
.thought-label{font-size:11px;font-weight:700;color:rgba(180,150,255,.9);flex:1;text-transform:uppercase;letter-spacing:.7px}
.thought-dots span{animation:tdot 1.2s infinite;opacity:0;display:inline-block;margin-left:1px}
.thought-dots span:nth-child(2){animation-delay:.2s}.thought-dots span:nth-child(3){animation-delay:.4s}
@keyframes tdot{0%,80%,100%{opacity:0}40%{opacity:1}}
.thought-streaming .thought-dots{display:inline}
.thought:not(.thought-streaming) .thought-dots{display:none}
.thought-toggle{font-size:10px;color:var(--tx3);transition:transform .2s}
.thought-toggle.open{transform:rotate(180deg)}
.thought-body{max-height:0;overflow:hidden;transition:max-height .4s ease}
.thought-body.open,.thought-body.thought-body-open{max-height:2000px}
.thought-txt{padding:10px 14px;font-size:12px;color:var(--tx2);line-height:1.75;font-family:'DM Mono',monospace;white-space:pre-wrap;word-break:break-word}
/* ── Matrix BG Canvas ── */
#matrixCanvas{position:fixed;inset:0;z-index:0;opacity:0;transition:opacity .6s;pointer-events:none}
#matrixCanvas.vis{opacity:.16}
body.light #matrixCanvas.vis{opacity:.06;filter:hue-rotate(180deg) invert(1)}
/* ── Compact sizes ── */
.msg{margin-bottom:14px}
.bbl{padding:11px 15px;font-size:13.5px}
.chat-body{padding:16px 12px}
/* ── Right vertical quick-tool bar ── */
.rtb{position:fixed;right:0;top:50%;transform:translateY(-50%);z-index:110;display:flex;flex-direction:column;gap:5px;padding:8px 5px;background:rgba(6,8,18,.85);backdrop-filter:blur(18px);border:1px solid var(--glass-bdr);border-right:none;border-radius:12px 0 0 12px}
.rtb-btn{width:36px;height:36px;background:rgba(255,255,255,.05);border:1px solid var(--glass-bdr);border-radius:8px;display:flex;align-items:center;justify-content:center;cursor:pointer;font-size:15px;transition:all .18s;color:var(--tx2);position:relative}
.rtb-btn:hover{border-color:var(--red);background:rgba(124,58,237,.1);color:var(--tx);transform:translateX(-2px)}
.rtb-btn.on{border-color:var(--blue);background:rgba(6,182,212,.12);color:var(--blue)}
.rtb-tip{position:absolute;right:44px;top:50%;transform:translateY(-50%);background:rgba(6,8,18,.97);border:1px solid var(--glass-bdr2);border-radius:7px;padding:3px 9px;font-size:11px;color:var(--tx);white-space:nowrap;pointer-events:none;opacity:0;transition:opacity .15s;z-index:200}
.rtb-btn:hover .rtb-tip{opacity:1}
.rtb-sep{height:1px;background:var(--glass-bdr);margin:2px 0}
/* ── Left settings sidebar ── */
.lsb{position:fixed;left:0;top:52px;bottom:0;width:240px;background:rgba(5,7,16,.95);backdrop-filter:blur(22px);border-right:1px solid var(--glass-bdr);z-index:95;transform:translateX(-100%);transition:transform .25s cubic-bezier(.16,1,.3,1);overflow-y:auto;display:flex;flex-direction:column}
.lsb.open{transform:translateX(0)}
.lsb-hdr{padding:14px 14px 10px;border-bottom:1px solid var(--glass-bdr);flex-shrink:0;display:flex;align-items:center;justify-content:space-between}
.lsb-hdr h3{font-size:13px;font-weight:700;color:var(--tx)}
.lsb-x{background:none;border:none;color:var(--tx3);font-size:16px;cursor:pointer;padding:3px 6px;border-radius:5px}
.lsb-x:hover{color:var(--tx);background:rgba(255,255,255,.07)}
.lsb-sec{padding:10px 12px;border-bottom:1px solid var(--glass-bdr)}
.lsb-sec h4{font-size:9px;font-weight:700;color:var(--tx3);text-transform:uppercase;letter-spacing:1px;margin-bottom:8px}
.lsb-item{display:flex;align-items:center;gap:8px;padding:7px 9px;border-radius:8px;cursor:pointer;transition:background .15s;color:var(--tx2);font-size:12.5px;user-select:none}
.lsb-item:hover{background:rgba(255,255,255,.06);color:var(--tx)}
.lsb-item.active{background:rgba(124,58,237,.09);color:var(--tx);border-left:2px solid var(--red)}
.lsb-item .ei2{font-size:14px}
.lsb-overlay{position:fixed;inset:0;z-index:94;display:none}
.lsb-overlay.open{display:block}
/* ── Google button ── */
.google-btn{width:100%;padding:10px;background:rgba(255,255,255,.06);border:1.5px solid rgba(255,255,255,.13);border-radius:10px;color:var(--tx);font-family:'DM Sans',sans-serif;font-size:13px;font-weight:500;cursor:pointer;display:flex;align-items:center;justify-content:center;gap:9px;transition:all .2s;margin-bottom:10px}
.google-btn:hover{background:rgba(255,255,255,.11);border-color:rgba(255,255,255,.28);transform:translateY(-1px)}
/* ── Weather badge ── */
.wx-badge{display:none;align-items:center;gap:8px;background:rgba(6,182,212,.07);border:1px solid rgba(6,182,212,.18);border-radius:10px;padding:7px 12px;font-size:12px;color:var(--tx2);margin-bottom:10px;cursor:pointer}
.wx-badge.show{display:flex}
.wx-temp{font-size:20px;font-weight:700;color:var(--tx);font-family:'Bebas Neue',sans-serif;letter-spacing:1px}
/* ── Langsearch result card ── */
.search-card{background:rgba(255,255,255,.04);border:1px solid var(--glass-bdr);border-radius:10px;padding:10px 13px;margin-bottom:7px}
.search-card-title{font-size:12px;font-weight:600;color:var(--blue);margin-bottom:3px}
.search-card-snip{font-size:11px;color:var(--tx2);line-height:1.6}
/* ── Camera btn ── */
.cam-btn{display:none}
@media(any-pointer:coarse),(max-width:700px){.cam-btn{display:flex}.rtb{display:none}}
/* ── Compact mobile ── */
/* Mobile keyboard & layout fixes */
.page{flex-direction:column;height:100%;height:100dvh}
#chatPage{flex-direction:column;height:100dvh}
#chatPage header{flex-shrink:0;position:sticky;top:0;z-index:50}
.chat-body{flex:1;overflow-y:auto;-webkit-overflow-scrolling:touch}
.iz{flex-shrink:0;position:sticky;bottom:0;z-index:50;padding-bottom:max(12px,env(safe-area-inset-bottom));overflow:visible}
@media(max-width:600px){
  .bbl{font-size:13px;padding:8px 11px}
  .msg{margin-bottom:10px}
  .ibox{padding:5px 7px;border-radius:13px;gap:5px}
  header{padding:6px 10px}
  .sp{width:100vw;right:-100vw;top:0;height:100dvh;border-radius:0}.sp-tabs{width:82px;min-width:82px;padding:6px 3px}.sp-tab{padding:7px 6px;font-size:10px}.sp-tab .st-lbl{font-size:10px}.sp-tab .st-ico{font-size:12px}
  .chat-body{padding:8px 6px}
  .toast{font-size:12px;max-width:calc(100vw - 24px);white-space:normal;text-align:center;bottom:calc(60px + env(safe-area-inset-bottom))}
  .welcome h1{font-size:28px;letter-spacing:2px}
  .welcome>p{font-size:12px}
  .chips{gap:5px;justify-content:center}
  .chip{padding:6px 10px;font-size:11.5px}
  .rtb{display:none!important}
  .lsb{width:90vw}
  .mpill{display:none}
  .msg-actions{opacity:1}
  .wico{width:56px;height:56px;font-size:26px}
  .ibtn{width:30px;height:30px;font-size:13px}
  .iibtn{width:30px;height:30px;font-size:13px}
  .sbtn{width:32px;height:32px}
  textarea{font-size:16px}
  .sp-hdr{padding:12px 14px 8px}
  .sp-tab{padding:6px 8px;font-size:10px}
}
@media(min-width:601px){
  .page,.chat-body{height:100vh}
}
/* ── API key new row ── */
.krow-extra{margin-top:8px}
/* ── Sleek additions ── */
.card{box-shadow:0 40px 100px rgba(0,0,0,.7),0 0 0 1px rgba(255,255,255,.05);backdrop-filter:blur(36px)}
.btn{font-size:14px;letter-spacing:.3px;box-shadow:0 4px 16px var(--redg)}
.btn:hover{transform:translateY(-1px);box-shadow:0 6px 24px var(--redg)}
.google-btn:active,.btn:active,.ghost-btn:active{transform:scale(.97)}
.ibtn{transition:background .15s,border-color .15s,transform .15s}
.ibtn:hover{transform:scale(1.08)}
.ibtn:active{transform:scale(.93)}
.iibtn:active{transform:scale(.9)}
.rtb-btn:active{transform:translateX(-4px) scale(.93)}
.lsb-item:active{transform:scale(.97)}
.chip:active{transform:translateY(-1px) scale(.97)}
header{backdrop-filter:blur(24px) saturate(200%)}
.sp{box-shadow:-12px 0 48px rgba(0,0,0,.5)}
/* ── Liquid Glass & Premium Transparency ── */
.glass-card{background:rgba(8,14,30,0.42);backdrop-filter:blur(32px) saturate(200%) brightness(1.08);-webkit-backdrop-filter:blur(32px) saturate(200%) brightness(1.08);border:1px solid rgba(120,160,255,0.13);box-shadow:0 8px 48px rgba(0,0,0,0.5),inset 0 1px 0 rgba(255,255,255,0.07)}
/* Apple-style smooth transitions on everything */
*{transition:background-color 0.25s ease,border-color 0.25s ease,box-shadow 0.25s ease,opacity 0.2s ease}
/* Override for transforms — don't slow them */
.chip,.btn,.ksave,.mact,.sp-tab,.lsidebar-item,.rtool-btn{transition:all 0.18s cubic-bezier(0.4,0,0.2,1)!important}
header{background:rgba(6,8,18,0.72)!important;backdrop-filter:blur(32px) saturate(220%)!important;-webkit-backdrop-filter:blur(32px) saturate(220%)!important;border-bottom:1px solid rgba(80,120,200,0.13)!important}
.sp{background:rgba(5,8,18,0.88)!important;backdrop-filter:blur(44px) saturate(210%)!important;-webkit-backdrop-filter:blur(44px) saturate(210%)!important;border-left:1px solid rgba(80,120,200,0.14)!important}
.sp-tabs{background:rgba(4,6,14,0.55)!important;backdrop-filter:blur(20px)!important}
.lsidebar{background:rgba(4,6,14,0.91)!important;backdrop-filter:blur(36px) saturate(200%)!important;-webkit-backdrop-filter:blur(36px) saturate(200%)!important}
.msg.ai .bbl{background:rgba(10,16,34,0.62)!important;backdrop-filter:blur(22px) saturate(180%)!important;-webkit-backdrop-filter:blur(22px) saturate(180%)!important;border:1px solid rgba(80,120,200,0.16)!important;box-shadow:0 4px 28px rgba(0,0,0,0.3),inset 0 1px 0 rgba(255,255,255,0.05)!important}
.iz{background:rgba(4,6,16,0.85)!important;backdrop-filter:blur(32px) saturate(200%)!important;-webkit-backdrop-filter:blur(32px) saturate(200%)!important;border-top:1px solid rgba(80,120,200,0.12)!important}
.ibox{background:rgba(8,12,28,0.72)!important;backdrop-filter:blur(24px) saturate(180%)!important;-webkit-backdrop-filter:blur(24px) saturate(180%)!important;border:1.5px solid rgba(80,120,200,0.2)!important}
.ibox:focus-within{border-color:rgba(124,58,237,.5)!important;box-shadow:0 0 0 3px rgba(124,58,237,.1),0 4px 24px rgba(0,0,0,.4)!important}
.card{background:rgba(8,14,30,0.78)!important;backdrop-filter:blur(44px) saturate(220%)!important;-webkit-backdrop-filter:blur(44px) saturate(220%)!important;border:1px solid rgba(120,160,255,0.12)!important;box-shadow:0 40px 100px rgba(0,0,0,.75),inset 0 1px 0 rgba(255,255,255,0.07)!important}
.rtoolbar{background:rgba(6,8,18,0.82)!important;backdrop-filter:blur(28px)!important;border:1px solid rgba(80,120,200,0.14)!important}
.drop,.chat-opts-panel,.modal{background:rgba(5,8,18,0.96)!important;backdrop-filter:blur(48px) saturate(220%)!important;-webkit-backdrop-filter:blur(48px) saturate(220%)!important;border:1px solid rgba(80,120,200,0.14)!important}
.thought{background:rgba(80,40,160,0.1)!important;backdrop-filter:blur(12px)!important}
.krow{background:rgba(8,14,30,0.5);backdrop-filter:blur(16px);border:1px solid rgba(80,120,200,0.12);border-radius:12px;padding:12px}
/* Light mode full overhaul */
body.light *{--glass-bg:rgba(255,255,255,0.75);--glass-bdr:rgba(0,0,0,0.09);--inp:rgba(240,242,250,0.9);--tx:#0f172a;--tx2:#475569;--tx3:#94a3b8;--bg-body:#f1f5f9}
body.light .msg.ai .bbl{background:rgba(255,255,255,0.92)!important;border:1px solid rgba(200,210,230,.5)!important;color:#0f172a!important}
body.light .ibox{background:rgba(255,255,255,0.88)!important;border-color:rgba(180,190,220,.5)!important}
body.light .iz{background:rgba(248,250,252,0.92)!important}
body.light header{background:rgba(255,255,255,0.88)!important}
body.light .sp{background:rgba(248,250,252,0.97)!important}
body.light .drop,.drop.open{background:rgba(248,250,252,0.98)!important}
body.light code{background:rgba(220,228,250,.7)!important;color:#1e3a5f!important}
body.light pre{background:rgba(220,228,250,.5)!important;color:#1e3a5f!important}
/* ══ Tool bar ══════════════════════════════════════════════════════════════ */
.tbar{max-width:800px;margin:6px auto 0;display:flex;align-items:center;gap:6px;overflow:visible;scrollbar-width:none;padding:2px 4px 4px;position:relative}
.tbar::-webkit-scrollbar{display:none}
.tbar-btn{display:inline-flex;align-items:center;gap:5px;padding:6px 13px;border-radius:20px;font-size:11.5px;font-weight:600;color:rgba(150,150,210,.75);background:rgba(16,16,30,0.7);border:1px solid rgba(50,80,160,.2);cursor:pointer;white-space:nowrap;flex-shrink:0;transition:all .18s cubic-bezier(.4,0,.2,1);letter-spacing:.2px;line-height:1;backdrop-filter:blur(8px)}
.tbar-btn span{font-size:11px}
.tbar-btn:hover{border-color:rgba(6,182,212,.5);color:#d8d8f5;background:rgba(6,182,212,.12);transform:translateY(-1px);box-shadow:0 4px 14px rgba(6,182,212,.15)}
.tbar-btn.active,.tbar-btn.tbar-open{background:rgba(6,182,212,.18);border-color:rgba(6,182,212,.55);color:#d0e8ff;box-shadow:0 2px 10px rgba(6,182,212,.2)}
.tbar-btn.web-active{border-color:rgba(34,197,94,.4);color:#4ade80;background:rgba(34,197,94,.09);box-shadow:0 2px 8px rgba(34,197,94,.1)}
.tbar-sep{width:1px;height:16px;background:rgba(50,80,160,.2);flex-shrink:0;margin:0 3px}
/* Token counter */
.tok-ctr{display:inline-flex;align-items:center;gap:4px;margin-left:auto;flex-shrink:0;padding:0 2px}
.tok-num{font-size:11px;font-weight:600;color:var(--blue);font-family:'DM Mono',monospace;min-width:20px;text-align:right;transition:color .2s}
.tok-bar{width:40px;height:2.5px;background:rgba(80,120,200,.13);border-radius:3px;overflow:hidden}
.tok-fill{height:100%;background:linear-gradient(90deg,var(--blue),var(--red));border-radius:3px;transition:width .12s ease,background .2s}
.tok-ctr.warn .tok-num{color:var(--red)}
/* ══ Tool panels ═══════════════════════════════════════════════════════════ */
.tpanel{max-width:800px;margin:4px auto 0;background:rgba(8,13,28,.72);backdrop-filter:blur(28px) saturate(180%);-webkit-backdrop-filter:blur(28px) saturate(180%);border:1px solid rgba(80,120,200,.18);border-radius:14px;padding:14px 16px 13px;animation:tpIn .17s cubic-bezier(.16,1,.3,1)}
@keyframes tpIn{from{opacity:0;transform:translateY(5px) scale(.99)}to{opacity:1;transform:none}}
.tpanel-hd{display:flex;align-items:center;gap:8px;margin-bottom:7px}
.tpanel-ico{font-size:15px;line-height:1}
.tpanel-title{font-size:12px;font-weight:700;color:var(--tx);letter-spacing:.2px}
.tpanel-badge{background:rgba(6,182,212,.18);border:1px solid rgba(6,182,212,.3);border-radius:7px;padding:1px 7px;font-size:9px;font-weight:700;color:var(--blue);letter-spacing:.5px;text-transform:uppercase}
.tpanel-badge.green{background:rgba(34,197,94,.12);border-color:rgba(34,197,94,.3);color:#4ade80}
.tpanel-x{margin-left:auto;background:none;border:none;color:var(--tx3);font-size:13px;cursor:pointer;padding:0 2px;line-height:1;transition:color .15s}
.tpanel-x:hover{color:var(--tx)}
.tpanel-desc{font-size:11px;color:var(--tx3);line-height:1.65;margin-bottom:9px}
.tpanel-desc strong{color:var(--tx2)}
.tpanel-row{display:flex;gap:7px;align-items:stretch}
.tpanel-input{flex:1;background:rgba(6,10,22,.7);backdrop-filter:blur(12px);border:1.5px solid rgba(80,120,200,.18);border-radius:10px;padding:9px 13px;color:var(--tx);font-size:13px;outline:none;font-family:'DM Sans',sans-serif;transition:border-color .18s,box-shadow .18s;min-width:0}
.tpanel-input:focus{border-color:rgba(6,182,212,.5);box-shadow:0 0 0 3px rgba(6,182,212,.09)}
.tpanel-input::placeholder{color:var(--tx3)}
.tpanel-go{background:var(--grad);border:none;border-radius:10px;padding:9px 16px;color:#fff;font-size:12px;font-weight:700;cursor:pointer;white-space:nowrap;flex-shrink:0;letter-spacing:.2px;transition:opacity .15s,transform .15s}
.tpanel-go:hover{opacity:.88;transform:translateY(-1px)}
.tpanel-go:active{transform:translateY(0)}
.tpanel-sel{background:rgba(6,10,22,.8);border:1.5px solid rgba(80,120,200,.18);border-radius:8px;padding:6px 10px;color:var(--tx);font-size:12px;outline:none;cursor:pointer;min-width:100px;font-family:'DM Sans',sans-serif}
.tpanel-lvl{padding:4px 11px;border-radius:14px;font-size:11px;font-weight:500;color:var(--tx3);background:rgba(255,255,255,.04);border:1px solid rgba(255,255,255,.07);cursor:pointer;transition:all .15s}
.tpanel-lvl.active{background:rgba(6,182,212,.18);border-color:rgba(6,182,212,.4);color:#c8d8ff}
.tpanel-lvl:hover:not(.active){background:rgba(255,255,255,.08);color:var(--tx)}
/* Deep research live card grid */
.dr-header{display:flex;align-items:center;gap:8px;margin-bottom:8px;flex-wrap:wrap}
.dr-title-ico{font-size:16px;line-height:1}
.dr-title-txt{font-size:13px;font-weight:700;color:var(--tx)}
.dr-live-badge{background:rgba(124,58,237,.15);border:1px solid rgba(124,58,237,.3);border-radius:6px;padding:1px 7px;font-size:9px;font-weight:700;color:var(--red);letter-spacing:.5px;animation:pulse 1.2s infinite}
.dr-stat{font-size:11px;color:var(--tx3);font-family:'DM Mono',monospace}
.dr-prompt{font-size:12px;color:var(--tx2);font-style:italic;margin-bottom:4px;line-height:1.5;opacity:.8}
.dr-web-ctx{font-size:10px;color:#4ade80;background:rgba(34,197,94,.06);border:1px solid rgba(34,197,94,.15);border-radius:7px;padding:4px 10px;margin-bottom:10px;display:inline-block}
.dr-wrap{display:grid;grid-template-columns:repeat(auto-fill,minmax(240px,1fr));gap:8px;margin-top:4px}
.dr-card{background:rgba(6,10,22,.65);backdrop-filter:blur(16px);border:1px solid rgba(80,120,200,.15);border-radius:11px;padding:11px 13px;cursor:pointer;transition:border-color .2s}
.dr-card.done{border-color:rgba(34,197,94,.22)}
.dr-card.fail{border-color:rgba(124,58,237,.18);opacity:.55}
.dr-card.pend{border-color:rgba(80,120,200,.12)}
.dr-card.expanded .dr-card-txt{max-height:none!important;mask-image:none!important}
.dr-card-hd{display:flex;align-items:center;gap:6px;margin-bottom:6px}
.dr-card-dot{width:6px;height:6px;border-radius:50%;flex-shrink:0;transition:background .3s}
.dr-card.done .dr-card-dot{background:#4ade80}
.dr-card.fail .dr-card-dot{background:var(--red)}
.dr-card.pend .dr-card-dot{background:rgba(255,200,0,.7);animation:pulse 1s infinite}
.dr-card-name{font-size:9.5px;font-weight:700;color:var(--tx3);font-family:'DM Mono',monospace;text-transform:uppercase;letter-spacing:.5px;flex:1;overflow:hidden;text-overflow:ellipsis;white-space:nowrap}
.dr-card-ms{font-size:9px;font-family:'DM Mono',monospace;flex-shrink:0}
.dr-card-txt{font-size:12px;color:var(--tx);line-height:1.7;max-height:108px;overflow:hidden;mask-image:linear-gradient(to bottom,black 55%,transparent)}
.dr-synth{background:rgba(6,182,212,.07);border:1px solid rgba(6,182,212,.18);border-radius:12px;padding:15px 17px;margin-top:12px}
.dr-synth-hd{font-size:10px;font-weight:700;color:var(--blue);text-transform:uppercase;letter-spacing:.8px;margin-bottom:9px}
.dr-synth-txt{font-size:14px;color:var(--tx);line-height:1.85}
/* Fb-chip in browser */
.fb-chip{padding:4px 10px;border-radius:14px;font-size:11px;font-weight:500;color:var(--tx3);background:rgba(255,255,255,.04);border:1px solid rgba(255,255,255,.07);cursor:pointer;transition:all .15s;white-space:nowrap}
.fb-chip:hover{border-color:var(--blue);color:var(--tx);background:rgba(6,182,212,.1)}
/* Fusion Browser results */
.fb-result{background:rgba(8,13,28,.62);backdrop-filter:blur(16px);border:1px solid rgba(80,120,200,.13);border-radius:11px;padding:11px 14px;margin-bottom:7px;cursor:pointer;transition:all .16s}
.fb-result:hover{border-color:rgba(6,182,212,.38);background:rgba(6,182,212,.07);transform:translateX(2px)}
.fb-result-title{font-size:13px;font-weight:600;color:var(--blue);margin-bottom:3px}
.fb-result-url{font-size:10px;color:var(--tx3);margin-bottom:4px;font-family:'DM Mono',monospace;overflow:hidden;text-overflow:ellipsis;white-space:nowrap}
.fb-result-snip{font-size:12px;color:var(--tx2);line-height:1.65}
.fb-instant{background:rgba(6,182,212,.07);border:1px solid rgba(6,182,212,.18);backdrop-filter:blur(12px);border-radius:11px;padding:12px 15px;margin-bottom:11px}


/* Sleek input */
.ibox{box-shadow:0 2px 16px rgba(0,0,0,.3)}
.ibox:focus-within{box-shadow:0 4px 24px rgba(6,182,212,.15),0 0 0 2px rgba(6,182,212,.15)}
/* reCAPTCHA badge positioning */
.grecaptcha-badge{visibility:hidden}
/* Compact header pill */
.mpill{max-width:120px;overflow:hidden;text-overflow:ellipsis}



/* ══ Conversations Sidebar ═══════════════════════════════════════════════ */
.conv-sidebar{position:fixed;left:0;top:0;bottom:0;width:242px;background:rgba(9,9,17,0.97);border-right:1px solid rgba(124,58,237,.16);z-index:200;display:flex;flex-direction:column;transform:translateX(0);transition:transform .25s cubic-bezier(.4,0,.2,1);backdrop-filter:blur(28px)}
.conv-sidebar.closed{transform:translateX(-242px)}
.conv-sidebar-overlay{position:fixed;inset:0;z-index:199;background:rgba(0,0,0,.5);backdrop-filter:blur(4px);opacity:0;pointer-events:none;transition:opacity .25s}
.conv-sidebar-overlay.open{opacity:1;pointer-events:all}
#chatPage{transition:padding-left .25s cubic-bezier(.4,0,.2,1)}
@media(min-width:769px){
  #chatPage{padding-left:242px}
  #chatPage.sidebar-closed{padding-left:0}
  .conv-sidebar-overlay{display:none}
}
@media(max-width:768px){
  .conv-sidebar{transform:translateX(-242px)}
  .conv-sidebar.open{transform:translateX(0)}
  .conv-sidebar.closed{transform:translateX(-242px)}
}
.conv-sb-head{display:flex;align-items:center;padding:16px 14px 12px;border-bottom:1px solid rgba(124,58,237,.2);gap:8px;flex-shrink:0}
.conv-sb-logo{font-size:18px}
.conv-sb-title{font-size:14px;font-weight:800;color:#eeeeff;flex:1;letter-spacing:.2px}
.conv-sb-close{width:26px;height:26px;border-radius:50%;border:1px solid rgba(124,58,237,.3);background:rgba(20,20,38,.6);color:rgba(150,150,210,.7);font-size:12px;cursor:pointer;display:flex;align-items:center;justify-content:center;transition:all .15s}
.conv-sb-close:hover{background:rgba(124,58,237,.15);color:#ff8090;border-color:rgba(124,58,237,.4)}
.conv-sb-new{display:flex;align-items:center;gap:6px;margin:10px 12px;padding:9px 14px;background:linear-gradient(135deg,rgba(6,182,212,.2),rgba(124,58,237,.15));border:1px solid rgba(6,182,212,.35);border-radius:12px;cursor:pointer;font-size:13px;font-weight:700;color:#d8d8f5;transition:all .18s}
.conv-sb-new:hover{background:linear-gradient(135deg,rgba(6,182,212,.3),rgba(124,58,237,.22));border-color:rgba(6,182,212,.6);transform:translateY(-1px)}
.conv-sb-search{margin:0 12px 8px;position:relative}
.conv-sb-search input{width:100%;background:rgba(16,16,30,.7);border:1px solid rgba(124,58,237,.25);border-radius:10px;padding:7px 10px 7px 30px;font-size:12px;color:#d8d8f5;outline:none;font-family:'DM Sans',sans-serif;transition:border-color .15s}
.conv-sb-search input:focus{border-color:rgba(6,182,212,.5)}
.conv-sb-search input::placeholder{color:rgba(60,100,160,.5)}
.conv-sb-search-icon{position:absolute;left:10px;top:50%;transform:translateY(-50%);font-size:11px;color:rgba(60,100,160,.5)}
.conv-sb-list{flex:1;overflow-y:auto;padding:4px 8px 20px}
.conv-sb-list::-webkit-scrollbar{width:3px}
.conv-sb-list::-webkit-scrollbar-thumb{background:rgba(6,182,212,.2);border-radius:2px}
.conv-sb-day{font-size:9px;font-weight:700;letter-spacing:1px;color:rgba(60,100,160,.6);text-transform:uppercase;padding:10px 8px 4px;margin-top:4px}
.conv-sb-item{display:flex;align-items:center;gap:8px;padding:9px 10px;border-radius:11px;cursor:pointer;border:1px solid transparent;transition:all .14s;margin-bottom:2px}
.conv-sb-item:hover{background:rgba(6,182,212,.08);border-color:rgba(6,182,212,.15)}
.conv-sb-item.active{background:rgba(6,182,212,.14);border-color:rgba(6,182,212,.3)}
.conv-sb-item-icon{font-size:14px;flex-shrink:0;width:22px;text-align:center}
.conv-sb-item-body{flex:1;min-width:0}
.conv-sb-item-title{font-size:12px;font-weight:600;color:#d8d8f5;overflow:hidden;text-overflow:ellipsis;white-space:nowrap}
.conv-sb-item-preview{font-size:10px;color:rgba(130,130,180,.65);overflow:hidden;text-overflow:ellipsis;white-space:nowrap;margin-top:1px}
.conv-sb-item-del{width:20px;height:20px;border-radius:5px;border:none;background:transparent;color:rgba(130,130,180,.4);font-size:11px;cursor:pointer;display:flex;align-items:center;justify-content:center;opacity:0;transition:all .12s;flex-shrink:0}
.conv-sb-item:hover .conv-sb-item-del{opacity:1}
.conv-sb-item-del:hover{background:rgba(124,58,237,.15);color:#f87171}
/* ══ Profile Panel ═══════════════════════════════════════════════════════ */
.profile-panel{position:fixed;inset:0;z-index:9998;display:flex;align-items:center;justify-content:center;background:rgba(2,4,12,.8);backdrop-filter:blur(12px);animation:fadeIn .18s ease}
.profile-panel.hidden{display:none}
.profile-box{background:rgba(5,9,22,0.98);border:1px solid rgba(124,58,237,.3);border-radius:24px;width:420px;max-width:95vw;overflow:hidden;box-shadow:0 40px 100px rgba(0,0,0,.8);animation:slideup .2s cubic-bezier(.16,1,.3,1)}
.profile-cover{height:100px;background:linear-gradient(135deg,rgba(6,182,212,.3),rgba(124,58,237,.25));position:relative;display:flex;align-items:flex-end;padding:0 20px 0}
.profile-avatar-wrap{position:absolute;bottom:-30px;left:20px;cursor:pointer}
.profile-avatar-big{width:64px;height:64px;border-radius:50%;background:linear-gradient(135deg,#06b6d4,#7c3aed);border:3px solid rgba(5,9,22,.98);display:flex;align-items:center;justify-content:center;font-size:28px;transition:all .2s;box-shadow:0 4px 20px rgba(6,182,212,.3)}
.profile-avatar-big:hover{transform:scale(1.05);box-shadow:0 6px 28px rgba(6,182,212,.4)}
.profile-avatar-edit{position:absolute;bottom:-2px;right:-2px;width:20px;height:20px;background:#06b6d4;border-radius:50%;border:2px solid rgba(5,9,22,.98);display:flex;align-items:center;justify-content:center;font-size:9px;color:#fff}
.profile-body{padding:44px 20px 20px}
.profile-name{font-size:20px;font-weight:800;color:#eeeeff;margin-bottom:2px}
.profile-gid{font-size:11px;color:rgba(130,130,180,.6);font-family:'DM Mono',monospace;margin-bottom:12px}
.profile-stats{display:flex;gap:20px;padding:12px 0;border-top:1px solid rgba(124,58,237,.2);border-bottom:1px solid rgba(124,58,237,.2);margin-bottom:14px}
.profile-stat{text-align:center;flex:1}
.profile-stat-val{font-size:20px;font-weight:800;color:#a78bfa;font-family:'DM Mono',monospace}
.profile-stat-lbl{font-size:9px;color:rgba(130,130,180,.6);text-transform:uppercase;letter-spacing:.8px;margin-top:2px}
.profile-actions{display:flex;gap:8px}
.profile-btn{flex:1;padding:9px;border-radius:11px;border:1px solid rgba(124,58,237,.3);background:rgba(16,16,30,.6);color:rgba(120,180,255,.8);font-size:12px;font-weight:600;cursor:pointer;transition:all .15s;font-family:'DM Sans',sans-serif}
.profile-btn:hover{background:rgba(6,182,212,.15);border-color:rgba(6,182,212,.4);color:#d8d8f5}
.profile-close-btn{position:absolute;top:12px;right:14px;width:28px;height:28px;border-radius:50%;border:1px solid rgba(255,255,255,.1);background:rgba(255,255,255,.07);color:rgba(200,220,255,.6);font-size:13px;cursor:pointer;display:flex;align-items:center;justify-content:center}
.profile-close-btn:hover{background:rgba(124,58,237,.15);color:#f87171}
/* Emoji picker */
.emoji-picker-wrap{padding:14px 20px;border-top:1px solid rgba(124,58,237,.2)}
.emoji-picker-title{font-size:10px;font-weight:700;letter-spacing:1px;color:rgba(130,130,180,.6);text-transform:uppercase;margin-bottom:8px}
.emoji-picker-grid{display:grid;grid-template-columns:repeat(10,1fr);gap:4px}
.emoji-opt{width:32px;height:32px;border-radius:8px;border:1px solid rgba(124,58,237,.2);background:rgba(16,16,30,.5);font-size:18px;cursor:pointer;display:flex;align-items:center;justify-content:center;transition:all .12s}
.emoji-opt:hover{background:rgba(6,182,212,.15);border-color:rgba(6,182,212,.4);transform:scale(1.15)}
.emoji-opt.selected{background:rgba(6,182,212,.25);border-color:rgba(6,182,212,.7);box-shadow:0 0 0 2px rgba(6,182,212,.3)}
/* ══ AI Tools Hub ════════════════════════════════════════════════════════ */
.tools-hub-overlay{position:fixed;inset:0;z-index:9990;background:rgba(2,4,12,.85);backdrop-filter:blur(14px);display:flex;align-items:center;justify-content:center;animation:fadeIn .18s}
.tools-hub-overlay.hidden{display:none}
.tools-hub{background:rgba(4,7,18,.98);border:1px solid rgba(124,58,237,.25);border-radius:24px;width:680px;max-width:96vw;max-height:88vh;overflow:hidden;display:flex;flex-direction:column;box-shadow:0 40px 100px rgba(0,0,0,.85);animation:slideup .2s cubic-bezier(.16,1,.3,1)}
.tools-hub-head{display:flex;align-items:center;padding:18px 22px 14px;border-bottom:1px solid rgba(124,58,237,.2)}
.tools-hub-title{font-size:16px;font-weight:800;color:#eeeeff;flex:1}
.tools-hub-close{width:28px;height:28px;border-radius:50%;border:1px solid rgba(124,58,237,.3);background:rgba(20,20,38,.6);color:rgba(150,150,210,.7);font-size:14px;cursor:pointer;display:flex;align-items:center;justify-content:center;transition:all .15s}
.tools-hub-close:hover{background:rgba(124,58,237,.15);color:#f87171}
.tools-hub-body{padding:16px 22px;overflow-y:auto;flex:1}
.tools-hub-grid{display:grid;grid-template-columns:repeat(auto-fill,minmax(150px,1fr));gap:10px;margin-bottom:16px}
.tool-card{background:rgba(16,16,30,.7);border:1px solid rgba(124,58,237,.2);border-radius:14px;padding:16px 14px;cursor:pointer;transition:all .18s;text-align:center}
.tool-card:hover{background:rgba(6,182,212,.1);border-color:rgba(6,182,212,.4);transform:translateY(-2px);box-shadow:0 6px 20px rgba(6,182,212,.12)}
.tool-card-icon{font-size:28px;margin-bottom:8px}
.tool-card-name{font-size:12px;font-weight:700;color:#d8d8f5}
.tool-card-desc{font-size:10px;color:rgba(130,130,180,.7);margin-top:3px;line-height:1.4}
/* Tool panels */
.tool-panel{display:none;flex-direction:column;gap:12px}
.tool-panel.active{display:flex}
.tool-panel-back{display:flex;align-items:center;gap:6px;font-size:12px;color:rgba(130,130,180,.7);cursor:pointer;padding:4px 0;transition:color .15s;width:fit-content}
.tool-panel-back:hover{color:#d8d8f5}
.tool-panel-title{font-size:15px;font-weight:800;color:#eeeeff;margin-bottom:2px}
.tool-textarea{width:100%;background:rgba(6,10,26,.8);border:1px solid rgba(124,58,237,.25);border-radius:12px;padding:10px 13px;font-size:13px;color:#eeeeff;outline:none;font-family:'DM Sans',sans-serif;resize:vertical;transition:border-color .15s;min-height:100px}
.tool-textarea:focus{border-color:rgba(6,182,212,.5);box-shadow:0 0 0 3px rgba(6,182,212,.07)}
.tool-textarea::placeholder{color:rgba(60,100,160,.5)}
.tool-select{background:rgba(6,10,26,.8);border:1px solid rgba(124,58,237,.25);border-radius:10px;padding:8px 12px;font-size:12px;color:#d8d8f5;outline:none;font-family:'DM Sans',sans-serif;cursor:pointer;transition:border-color .15s}
.tool-select:focus{border-color:rgba(6,182,212,.5)}
.tool-run-btn{background:linear-gradient(135deg,#06b6d4,#7c3aed);border:none;border-radius:11px;padding:10px 24px;font-size:13px;font-weight:700;color:#fff;cursor:pointer;transition:all .18s;font-family:'DM Sans',sans-serif;align-self:flex-end}
.tool-run-btn:hover{opacity:.9;transform:translateY(-1px);box-shadow:0 6px 20px rgba(6,182,212,.35)}
.tool-run-btn:disabled{opacity:.4;cursor:not-allowed;transform:none}
.tool-result-box{background:rgba(6,10,26,.8);border:1px solid rgba(40,80,160,.2);border-radius:12px;padding:14px;font-size:13px;color:rgba(180,220,255,.9);line-height:1.7;white-space:pre-wrap;max-height:300px;overflow-y:auto;position:relative}
.tool-result-copy{position:absolute;top:8px;right:10px;background:rgba(6,182,212,.12);border:1px solid rgba(6,182,212,.3);border-radius:7px;padding:4px 9px;font-size:10px;color:rgba(100,180,255,.8);cursor:pointer;font-family:'DM Sans',sans-serif}
.tool-result-copy:hover{background:rgba(6,182,212,.25);color:#fff}
/* Quiz */
.quiz-q{font-size:14px;font-weight:700;color:#d0e8ff;margin-bottom:10px;line-height:1.5}
.quiz-opts{display:flex;flex-direction:column;gap:6px;margin-bottom:10px}
.quiz-opt-btn{background:rgba(16,16,30,.7);border:1px solid rgba(124,58,237,.25);border-radius:10px;padding:10px 14px;font-size:13px;color:rgba(160,210,255,.85);cursor:pointer;text-align:left;transition:all .15s;font-family:'DM Sans',sans-serif}
.quiz-opt-btn:hover{background:rgba(6,182,212,.12);border-color:rgba(6,182,212,.4);color:#d8d8f5}
.quiz-opt-btn.correct{background:rgba(34,197,94,.12);border-color:rgba(34,197,94,.5);color:#4ade80}
.quiz-opt-btn.wrong{background:rgba(239,68,68,.1);border-color:rgba(239,68,68,.4);color:#f87171}
.quiz-progress{font-size:11px;color:rgba(80,130,200,.6);font-family:'DM Mono',monospace;margin-bottom:10px}
.quiz-score{text-align:center;padding:20px;font-size:22px;font-weight:800;color:#a78bfa}
/* ══ Quick action bar (declutter) ════════════════════════════════════════ */
.qa-bar{display:flex;gap:6px;overflow-x:auto;padding:8px 0 4px;scrollbar-width:none;max-width:800px;margin:0 auto}
.qa-bar::-webkit-scrollbar{display:none}
.qa-chip{display:inline-flex;align-items:center;gap:4px;padding:5px 12px;border-radius:20px;font-size:11px;font-weight:600;color:rgba(100,160,220,.75);background:rgba(16,16,30,.6);border:1px solid rgba(124,58,237,.2);cursor:pointer;white-space:nowrap;flex-shrink:0;transition:all .15s;backdrop-filter:blur(8px)}
.qa-chip:hover{background:rgba(6,182,212,.14);border-color:rgba(6,182,212,.45);color:#d8d8f5;transform:translateY(-1px)}
/* ══ Hamburger / sidebar toggle ══════════════════════════════════════════ */
.sidebar-toggle{width:34px;height:34px;border-radius:10px;border:1px solid rgba(124,58,237,.25);background:rgba(16,16,30,.6);color:rgba(100,160,220,.7);font-size:15px;cursor:pointer;display:flex;align-items:center;justify-content:center;transition:all .15s;flex-shrink:0}
.sidebar-toggle:hover{background:rgba(6,182,212,.14);border-color:rgba(6,182,212,.4);color:#d8d8f5}

/* ══ Toolbox dropdown ══════════════════════════════════════════════════ */
.toolbox-wrap{position:relative;flex-shrink:0}
.toolbox-btn{display:inline-flex;align-items:center;gap:5px;padding:6px 14px;border-radius:20px;font-size:11.5px;font-weight:700;color:#d0e8ff;background:linear-gradient(135deg,rgba(6,182,212,.22),rgba(120,60,220,.18));border:1px solid rgba(80,140,255,.4);cursor:pointer;white-space:nowrap;transition:all .18s cubic-bezier(.4,0,.2,1);letter-spacing:.3px;backdrop-filter:blur(10px)}
.toolbox-btn:hover{background:linear-gradient(135deg,rgba(6,182,212,.32),rgba(120,60,220,.26));border-color:rgba(100,170,255,.6);box-shadow:0 4px 18px rgba(6,182,212,.25);transform:translateY(-1px)}
.toolbox-btn.open{background:linear-gradient(135deg,rgba(6,182,212,.35),rgba(120,60,220,.28));border-color:rgba(120,180,255,.65);box-shadow:0 4px 20px rgba(6,182,212,.3)}
.toolbox-menu{position:absolute;bottom:calc(100% + 8px);top:auto;left:0;z-index:9999;background:rgba(6,10,26,0.97);border:1px solid rgba(60,100,200,.35);border-radius:16px;padding:8px;width:230px;box-shadow:0 20px 60px rgba(0,0,0,.7),0 0 0 1px rgba(80,140,255,.08);backdrop-filter:blur(28px);animation:dropIn .16s cubic-bezier(.16,1,.3,1)}
@keyframes dropIn{from{opacity:0;transform:translateY(8px) scale(.97)}to{opacity:1;transform:translateY(0) scale(1)}}
.toolbox-menu.hidden{display:none}
.toolbox-section{font-size:9px;font-weight:700;letter-spacing:1.2px;color:rgba(80,120,200,.6);text-transform:uppercase;padding:6px 8px 3px;margin-top:2px}
.toolbox-item{display:flex;align-items:center;gap:8px;padding:8px 10px;border-radius:10px;cursor:pointer;transition:all .14s;border:1px solid transparent}
.toolbox-item:hover{background:rgba(6,182,212,.12);border-color:rgba(6,182,212,.2);color:#d0e8ff}
.toolbox-item-icon{font-size:15px;width:22px;text-align:center;flex-shrink:0}
.toolbox-item-label{font-size:12px;font-weight:600;color:rgba(160,200,255,.85)}
.toolbox-item-desc{font-size:10px;color:rgba(130,130,180,.7);margin-top:1px}
.toolbox-divider{height:1px;background:rgba(124,58,237,.3);margin:5px 4px}
/* ══ File Creator Modal ════════════════════════════════════════════════ */
.file-modal-overlay{position:fixed;inset:0;z-index:99999;background:rgba(2,4,12,0.8);backdrop-filter:blur(12px);display:flex;align-items:center;justify-content:center;animation:fadeIn .18s ease}
.file-modal-overlay.hidden{display:none}
.file-modal{background:rgba(6,10,26,0.98);border:1px solid rgba(60,100,200,.35);border-radius:22px;width:520px;max-width:95vw;max-height:88vh;overflow:hidden;display:flex;flex-direction:column;box-shadow:0 40px 100px rgba(0,0,0,.8),0 0 80px rgba(6,182,212,.08);animation:slideup .2s cubic-bezier(.16,1,.3,1)}
.file-modal-head{display:flex;align-items:center;justify-content:space-between;padding:20px 22px 14px;border-bottom:1px solid rgba(124,58,237,.25)}
.file-modal-title{font-size:16px;font-weight:700;color:#eeeeff;letter-spacing:.2px}
.file-modal-close{width:28px;height:28px;border-radius:50%;border:1px solid rgba(124,58,237,.3);background:rgba(20,20,38,.6);color:rgba(150,150,210,.7);font-size:14px;cursor:pointer;display:flex;align-items:center;justify-content:center;transition:all .15s}
.file-modal-close:hover{background:rgba(124,58,237,.15);border-color:rgba(124,58,237,.4);color:#ff8090}
.file-modal-body{padding:16px 22px;overflow-y:auto;flex:1}
.file-type-grid{display:grid;grid-template-columns:repeat(3,1fr);gap:6px;margin-bottom:16px;max-height:320px;overflow-y:auto;padding-right:4px}
.ftype-card{border:1px solid rgba(124,58,237,.25);border-radius:12px;padding:10px 8px;cursor:pointer;transition:all .16s;text-align:center;background:rgba(16,16,30,.5)}
.ftype-card:hover{border-color:rgba(6,182,212,.5);background:rgba(6,182,212,.1);transform:translateY(-2px);box-shadow:0 6px 18px rgba(6,182,212,.15)}
.ftype-card.selected{border-color:rgba(6,182,212,.7);background:rgba(6,182,212,.16);box-shadow:0 0 0 1px rgba(6,182,212,.3)}
.ftype-card-icon{font-size:22px;margin-bottom:4px}
.ftype-card-name{font-size:11px;font-weight:700;color:#d8d8f5}
.ftype-card-ext{font-size:9px;color:rgba(130,130,180,.7);margin-top:1px;font-family:'DM Mono',monospace}
.file-topic-wrap{margin-bottom:12px}
.file-topic-label{font-size:11px;font-weight:600;color:rgba(150,150,210,.8);margin-bottom:6px;letter-spacing:.3px;text-transform:uppercase}
.file-topic-input{width:100%;background:rgba(16,16,30,.7);border:1px solid rgba(124,58,237,.3);border-radius:10px;padding:10px 13px;font-size:13px;color:#eeeeff;outline:none;font-family:'DM Sans',sans-serif;transition:border-color .15s;resize:none}
.file-topic-input:focus{border-color:rgba(6,182,212,.55);box-shadow:0 0 0 3px rgba(6,182,212,.08)}
.file-topic-input::placeholder{color:rgba(60,90,140,.6)}
.file-modal-foot{padding:12px 22px 18px;display:flex;gap:8px;justify-content:flex-end;border-top:1px solid rgba(124,58,237,.2)}
.file-cancel-btn{padding:9px 18px;border-radius:10px;border:1px solid rgba(124,58,237,.3);background:rgba(16,16,30,.5);color:rgba(150,150,210,.7);font-size:13px;font-weight:600;cursor:pointer;transition:all .15s;font-family:'DM Sans',sans-serif}
.file-cancel-btn:hover{border-color:rgba(60,100,200,.5);color:#d8d8f5}
.file-gen-btn{padding:9px 22px;border-radius:10px;border:none;background:linear-gradient(135deg,#06b6d4,#7c3aed);color:#fff;font-size:13px;font-weight:700;cursor:pointer;transition:all .18s;font-family:'DM Sans',sans-serif;letter-spacing:.2px}
.file-gen-btn:hover{opacity:.9;transform:translateY(-1px);box-shadow:0 6px 20px rgba(6,182,212,.4)}
.file-gen-btn:disabled{opacity:.4;cursor:not-allowed;transform:none}
/* ── Fusion File Card ──────────────────────────────────── */
.fusion-file-card{background:rgba(16,16,30,0.9);border:1px solid rgba(80,180,255,0.3);border-radius:14px;overflow:hidden;margin:10px 0;box-shadow:0 4px 24px rgba(6,182,212,.12)}
.ffc-header{display:flex;align-items:center;gap:10px;padding:12px 14px;background:rgba(6,182,212,.08);border-bottom:1px solid rgba(124,58,237,.2)}
.ffc-icon{font-size:22px;flex-shrink:0}
.ffc-info{flex:1;min-width:0}
.ffc-name{font-size:13px;font-weight:700;color:#eeeeff;font-family:'DM Mono',monospace;overflow:hidden;text-overflow:ellipsis;white-space:nowrap}
.ffc-meta{font-size:10px;color:rgba(100,160,255,.7);margin-top:2px}
.ffc-dl{display:inline-flex;align-items:center;gap:5px;background:var(--grad);border:none;border-radius:8px;padding:6px 12px;font-size:11px;font-weight:600;color:#fff;cursor:pointer;text-decoration:none;transition:all .2s;white-space:nowrap;flex-shrink:0}
.ffc-dl:hover{opacity:.88;transform:translateY(-1px);box-shadow:0 4px 12px rgba(6,182,212,.4)}
.ffc-copy{background:rgba(255,255,255,.07);border:1px solid rgba(80,120,200,.3);border-radius:8px;padding:6px 10px;font-size:11px;color:rgba(150,190,255,.8);cursor:pointer;transition:all .2s;white-space:nowrap;flex-shrink:0;font-family:'DM Sans',sans-serif}
.ffc-copy:hover{background:rgba(6,182,212,.15);color:#fff}
.ffc-preview{padding:10px 14px 12px;overflow-x:auto}
.ffc-table{width:100%;border-collapse:collapse;font-size:11px;font-family:'DM Mono',monospace}
.ffc-table th{background:rgba(6,182,212,.12);color:#a0c4ff;padding:5px 8px;text-align:left;border:1px solid rgba(124,58,237,.2);font-weight:700;font-size:10px;text-transform:uppercase;letter-spacing:.5px}
.ffc-table td{color:rgba(200,220,255,.85);padding:4px 8px;border:1px solid rgba(124,58,237,.2)}
.ffc-table tr:hover td{background:rgba(6,182,212,.05)}
.ffc-code-pre{font-size:11px;color:rgba(180,220,255,.8);margin:0;background:rgba(0,0,0,.35);border-radius:8px;padding:10px 13px;max-height:140px;overflow:auto;line-height:1.6;font-family:'DM Mono',monospace}
.ffc-more{font-size:10px;color:rgba(100,150,255,.6);margin-top:5px;font-style:italic}
.ffc-md-prev{font-size:12px;color:rgba(180,220,255,.85);line-height:1.7;max-height:140px;overflow:hidden;mask-image:linear-gradient(to bottom,#fff 70%,transparent)}
.ffc-html-prev{border-radius:8px;overflow:hidden;border:1px solid rgba(124,58,237,.3)}
.ffc-pptx-prev{padding:10px;background:rgba(6,182,212,.07);border-radius:8px;color:rgba(180,220,255,.85);font-size:13px}
.ffc-json{color:#7dd3fc}
.ffc-code-pre .kw{color:#c084fc}.ffc-code-pre .str{color:#86efac}.ffc-code-pre .num{color:#fb923c}.ffc-code-pre .cm{color:#6b7280;font-style:italic}

/* ══ Multi-Model Arena ════════════════════════════════════════════════════ */
.arena-wrap{display:flex;flex-direction:column;gap:10px;width:100%}
.arena-header{display:flex;align-items:center;gap:8px;padding:10px 14px;background:rgba(6,182,212,.07);border:1px solid rgba(6,182,212,.2);border-radius:14px}
.arena-badge{font-size:9px;font-weight:800;letter-spacing:1.5px;text-transform:uppercase;background:linear-gradient(90deg,#06b6d4,#7c3aed);color:#fff;padding:3px 8px;border-radius:20px}
.arena-prompt{font-size:12px;color:rgba(150,190,255,.7);flex:1;overflow:hidden;text-overflow:ellipsis;white-space:nowrap}
.arena-grid{display:grid;grid-template-columns:repeat(auto-fit,minmax(280px,1fr));gap:10px}
.arena-card{background:rgba(16,16,30,0.88);border:1px solid rgba(124,58,237,.2);border-radius:14px;overflow:hidden;cursor:pointer;transition:all .2s;position:relative}
.arena-card:hover{border-color:rgba(6,182,212,.5);transform:translateY(-2px);box-shadow:0 8px 28px rgba(6,182,212,.15)}
.arena-card.selected{border-color:rgba(6,182,212,.8);box-shadow:0 0 0 2px rgba(6,182,212,.3),0 8px 28px rgba(6,182,212,.2)}
.arena-card.selected::before{content:"✓ Selected";position:absolute;top:8px;right:10px;font-size:9px;font-weight:800;color:#06b6d4;background:rgba(6,182,212,.15);border:1px solid rgba(6,182,212,.4);border-radius:20px;padding:2px 8px;letter-spacing:.5px}
.arena-card-head{display:flex;align-items:center;gap:8px;padding:10px 12px;border-bottom:1px solid rgba(124,58,237,.15)}
.arena-card-emoji{font-size:18px}
.arena-card-name{font-size:12px;font-weight:700;color:#d8d8f5;flex:1}
.arena-card-ms{font-size:10px;color:rgba(130,130,180,.6);font-family:'DM Mono',monospace}
.arena-card-body{padding:10px 12px;font-size:12.5px;color:rgba(180,220,255,.85);line-height:1.65;max-height:200px;overflow-y:auto}
.arena-card-body::-webkit-scrollbar{width:3px}
.arena-card-body::-webkit-scrollbar-thumb{background:rgba(6,182,212,.3);border-radius:2px}
.arena-actions{display:flex;gap:8px;justify-content:flex-end;padding:6px 0 2px}
.arena-use-btn{background:linear-gradient(135deg,#06b6d4,#7c3aed);border:none;border-radius:10px;padding:8px 20px;font-size:12px;font-weight:700;color:#fff;cursor:pointer;transition:all .18s;font-family:'DM Sans',sans-serif}
.arena-use-btn:hover{opacity:.9;transform:translateY(-1px);box-shadow:0 4px 16px rgba(6,182,212,.4)}
.arena-use-btn:disabled{opacity:.4;cursor:not-allowed}
.arena-model-sel{display:flex;flex-wrap:wrap;gap:6px;margin:8px 0}
.arena-model-chip{display:inline-flex;align-items:center;gap:4px;padding:4px 10px;border-radius:20px;border:1px solid rgba(124,58,237,.3);background:rgba(16,16,30,.5);font-size:11px;color:rgba(120,180,255,.75);cursor:pointer;transition:all .15s;font-family:'DM Sans',sans-serif}
.arena-model-chip:hover{border-color:rgba(6,182,212,.5);color:#d0e8ff}
.arena-model-chip.sel{background:rgba(6,182,212,.18);border-color:rgba(6,182,212,.6);color:#d0e8ff}
.arena-loading{display:flex;align-items:center;gap:10px;padding:18px;color:rgba(100,160,255,.7);font-size:13px}
.arena-spinner{width:18px;height:18px;border:2px solid rgba(6,182,212,.2);border-top-color:#06b6d4;border-radius:50%;animation:spin .8s linear infinite}
@keyframes spin{to{transform:rotate(360deg)}}
/* ══ Extreme Deep Think ══════════════════════════════════════════════════ */
.edt-wrap{display:flex;flex-direction:column;gap:12px;width:100%}
.edt-header{display:flex;align-items:center;gap:10px;padding:12px 16px;background:linear-gradient(135deg,rgba(124,58,237,.1),rgba(6,182,212,.08));border:1px solid rgba(124,58,237,.3);border-radius:14px}
.edt-badge{font-size:9px;font-weight:800;letter-spacing:1.5px;text-transform:uppercase;background:linear-gradient(90deg,#7c3aed,#06b6d4);color:#fff;padding:3px 9px;border-radius:20px}
.edt-meta{font-size:11px;color:rgba(150,180,255,.6)}
.edt-phase-bar{display:flex;gap:4px;flex-wrap:wrap}
.edt-phase{display:flex;align-items:center;gap:4px;padding:4px 10px;border-radius:20px;font-size:10px;font-weight:600;border:1px solid rgba(124,58,237,.3);background:rgba(16,16,30,.5);color:rgba(80,130,200,.6);transition:all .3s}
.edt-phase.active{background:rgba(6,182,212,.15);border-color:rgba(6,182,212,.5);color:#a78bfa}
.edt-phase.done{background:rgba(34,197,94,.08);border-color:rgba(34,197,94,.3);color:#4ade80}
.edt-phase-dot{width:6px;height:6px;border-radius:50%;background:currentColor;opacity:.7}
.edt-synthesis{background:rgba(16,16,30,.9);border:1px solid rgba(80,140,255,.25);border-radius:14px;padding:16px 18px;line-height:1.75;font-size:13px;color:rgba(200,225,255,.9)}
.edt-perspectives{display:grid;grid-template-columns:repeat(auto-fit,minmax(240px,1fr));gap:8px}
.edt-persp-card{background:rgba(6,10,22,.85);border:1px solid rgba(30,60,130,.2);border-radius:11px;overflow:hidden}
.edt-persp-head{display:flex;align-items:center;gap:6px;padding:8px 11px;background:rgba(6,182,212,.05);border-bottom:1px solid rgba(30,60,130,.15)}
.edt-persp-name{font-size:11px;font-weight:700;color:#a8c8f0;flex:1}
.edt-persp-rounds{font-size:9px;color:rgba(130,130,180,.6);font-family:'DM Mono',monospace}
.edt-persp-body{padding:9px 11px;font-size:11.5px;color:rgba(160,200,255,.8);line-height:1.6;max-height:150px;overflow:auto}
.edt-stats{display:flex;gap:12px;flex-wrap:wrap;padding:10px 14px;background:rgba(6,10,22,.7);border:1px solid rgba(30,60,130,.2);border-radius:11px;font-size:11px;color:rgba(130,130,180,.7)}
.edt-stat{display:flex;flex-direction:column;gap:2px}
.edt-stat-val{font-size:18px;font-weight:800;color:#a78bfa;font-family:'DM Mono',monospace}
.edt-stat-lbl{font-size:9px;text-transform:uppercase;letter-spacing:.8px}
.edt-loading-phases{display:flex;flex-direction:column;gap:6px;padding:8px 0}
.edt-loading-row{display:flex;align-items:center;gap:8px;font-size:11.5px;color:rgba(120,170,255,.7)}
.edt-loading-row .arena-spinner{width:12px;height:12px}
.edt-loading-row.done{color:rgba(74,222,128,.7)}
.edt-loading-row.done .arena-spinner{display:none}
.edt-loading-row.done::before{content:"✓";color:#4ade80;font-weight:800;width:12px;text-align:center}

/* ── Upload file preview improvements ── */
.upload-preview-box{background:rgba(16,16,30,0.88);border:1px solid rgba(124,58,237,.25);border-radius:11px;padding:10px 13px;margin:6px 0;display:flex;align-items:flex-start;gap:10px}
.upload-preview-img{max-width:120px;max-height:90px;border-radius:7px;object-fit:cover;border:1px solid rgba(124,58,237,.2)}
.upload-preview-meta{flex:1;min-width:0}
.upload-preview-name{font-size:12px;font-weight:600;color:var(--tx);overflow:hidden;text-overflow:ellipsis;white-space:nowrap}
.upload-preview-size{font-size:10px;color:var(--tx3);margin-top:2px}
.upload-preview-text{font-size:11px;color:var(--tx2);margin-top:5px;font-family:'DM Mono',monospace;max-height:80px;overflow:hidden;line-height:1.5;white-space:pre-wrap;word-break:break-all}
.upload-dl-btn{background:rgba(6,182,212,.12);border:1px solid rgba(6,182,212,.3);border-radius:7px;padding:5px 10px;font-size:10px;color:rgba(100,180,255,.9);cursor:pointer;text-decoration:none;transition:all .2s;font-family:'DM Sans',sans-serif}
.upload-dl-btn:hover{background:rgba(6,182,212,.25);color:#fff}
/* ── Improved code copy button ── */
.code-copy-btn{font-size:10px;background:rgba(6,182,212,.12);border:1px solid rgba(6,182,212,.3);border-radius:6px;color:rgba(100,180,255,.9);cursor:pointer;padding:3px 10px;transition:all .15s;font-family:'DM Sans',sans-serif}
.code-copy-btn:hover{background:rgba(6,182,212,.28);color:#fff;border-color:rgba(6,182,212,.6)}
/* ── Human-like AI avatar ── */
.msg.ai .mav{background:linear-gradient(135deg,#06b6d4,#7c3aed);box-shadow:0 4px 14px rgba(6,182,212,.4)}
/* ── Toolbar doc button ── */
.tbar-doc-btn{display:flex;align-items:center;gap:5px;padding:5px 11px;background:rgba(6,182,212,.1);border:1px solid rgba(6,182,212,.3);border-radius:20px;font-size:11px;color:rgba(100,180,255,.9);cursor:pointer;transition:all .2s;white-space:nowrap;font-family:'DM Sans',sans-serif}
.tbar-doc-btn:hover{background:rgba(6,182,212,.2);color:#fff;border-color:rgba(6,182,212,.6)}
/* ── Welcome screen doc chips ── */
.chip.doc-chip{border-color:rgba(6,182,212,.3);background:rgba(6,182,212,.07)}
.chip.doc-chip:hover{border-color:rgba(6,182,212,.6);background:rgba(6,182,212,.15)}
/* ── Thinking animation improvement ── */
.thinking-txt{font-size:12px;color:rgba(120,180,255,.75);font-style:italic;animation:blink 2s infinite}
/* ── Model pill improvement ── */
.mpill{background:rgba(10,20,50,.7);border-color:rgba(60,100,200,.25)}
/* ── Better header glow ── */
header{box-shadow:0 1px 0 rgba(60,100,200,.2),0 8px 40px rgba(0,0,0,.6),0 0 80px rgba(6,182,212,.04)}

</style>
<script src="https://js.hcaptcha.com/1/api.js" async defer></script>
</head>
<body>
<canvas id="matrixCanvas"></canvas>
<canvas id="particleCanvas"></canvas>
<div class="bg" id="bgDiv"><div class="bg-orb bg-orb1"></div><div class="bg-orb bg-orb2"></div><div class="bg-orb bg-orb3"></div><div class="bg-grid"></div></div>
<div class="err-overlay" id="errOverlay"><span>&#x26A0;&#xFE0F;</span><span class="et" id="errOverlayMsg"></span><button class="ec" onclick="hideErrOverlay()">&#x2715;</button></div>


<!-- AUTH PAGE -->
<div class="page active" id="authPage">
  <div class="card" id="authCard">
    <div class="card-logo"><div class="lmark">&#x26A1;</div><span class="ltxt">Fusion.AI</span></div>
    <div class="tabs">
      <button class="tab active" id="tabLogin" onclick="switchTab('login')">Sign In</button>
      <button class="tab" id="tabReg" onclick="switchTab('register')">Register</button>
      <button class="tab" id="tabFP" onclick="switchTab('forgot')">Reset</button>
    </div>
    <div id="authErr" class="ferr"></div>
    <div id="loginForm">
      <button class="google-btn" id="googleSignInBtn" onclick="doGoogleSignIn()">
        <svg viewBox="0 0 48 48"><path fill="#EA4335" d="M24 9.5c3.54 0 6.71 1.22 9.21 3.6l6.85-6.85C35.9 2.38 30.47 0 24 0 14.62 0 6.51 5.38 2.56 13.22l7.98 6.19C12.43 13.72 17.74 9.5 24 9.5z"/><path fill="#4285F4" d="M46.98 24.55c0-1.57-.15-3.09-.38-4.55H24v9.02h12.94c-.58 2.96-2.26 5.48-4.78 7.18l7.73 6c4.51-4.18 7.09-10.36 7.09-17.65z"/><path fill="#FBBC05" d="M10.53 28.59c-.48-1.45-.76-2.99-.76-4.59s.27-3.14.76-4.59l-7.98-6.19C.92 16.46 0 20.12 0 24c0 3.88.92 7.54 2.56 10.78l7.97-6.19z"/><path fill="#34A853" d="M24 48c6.48 0 11.93-2.13 15.89-5.81l-7.73-6c-2.18 1.48-4.97 2.36-8.16 2.36-6.26 0-11.57-4.22-13.47-9.91l-7.98 6.19C6.51 42.62 14.62 48 24 48z"/><path fill="none" d="M0 0h48v48H0z"/></svg>
        Continue with Google
      </button>
      <div class="divider">or</div>
      <div class="field"><label>Username</label><input type="text" id="lUser" placeholder="your username" autocomplete="username"/></div>
      <div class="field"><label>Password</label><input type="password" id="lPass" placeholder="••••••••" autocomplete="current-password" onkeydown="if(event.key==='Enter')doLogin()"/></div>
      <button class="btn" id="loginBtn" onclick="doLogin()">Sign In &#x2192;</button>
      <div class="divider">or</div>
      <button class="ghost-btn" id="guestBtn" onclick="guestLogin()">Continue as Guest</button>
    </div>
    <div id="registerForm" style="display:none">
      <div class="field"><label>Username</label><input type="text" id="rUser" placeholder="choose a username" autocomplete="username"/></div>
      <div class="field"><label>Password</label><input type="password" id="rPass" placeholder="min 6 characters" autocomplete="new-password"/></div>
      <div class="field"><label>Confirm</label><input type="password" id="rPass2" placeholder="repeat password" autocomplete="new-password" onkeydown="if(event.key==='Enter')doRegister()"/></div>
      <div id="hcaptchaWrap" style="margin-bottom:12px;transform-origin:left top">
        <div class="h-captcha" id="hcaptchaWidget" data-sitekey="__HSITE__" data-theme="dark" data-size="normal" data-callback="onHcaptchaSuccess" data-error-callback="onHcaptchaError"></div>
        <div id="captchaErr" style="font-size:11px;color:var(--red);margin-top:5px;display:none">&#x26A0; Please complete the human verification.</div>
      </div>
      <button class="btn" id="regBtn" onclick="doRegister()">Create Account &#x2192;</button>
    </div>
    <div id="forgotForm" style="display:none">
      <p style="font-size:12.5px;color:var(--tx2);margin-bottom:14px;line-height:1.65">Enter your username and a new password.</p>
      <div class="field"><label>Username</label><input type="text" id="fpUser" placeholder="your username"/></div>
      <div class="field"><label>New Password</label><input type="password" id="fpPass" placeholder="new password (min 6)"/></div>
      <div class="field"><label>Confirm</label><input type="password" id="fpPass2" placeholder="repeat new password" onkeydown="if(event.key==='Enter')doForgotPw()"/></div>
      <button class="btn" id="fpBtn" onclick="doForgotPw()">Reset Password &#x2192;</button>
    </div>
  </div>
</div>

<!-- CHAT PAGE -->
<div class="page" id="chatPage">

<!-- ══ Conversations Sidebar ══════════════════════════════════════════════ -->
<div class="conv-sidebar-overlay" id="convSbOverlay" onclick="closeConvSidebar()"></div>
<div class="conv-sidebar" id="convSidebar">
  <div class="conv-sb-head">
    <span class="conv-sb-logo">⚡</span>
    <span class="conv-sb-title">Fusion.AI</span>
    <button class="conv-sb-close" onclick="closeConvSidebar()">&#x2715;</button>
  </div>
  <div class="conv-sb-new" onclick="if(window.innerWidth<=768)closeConvSidebar();newConv()">
    <span style="font-size:16px">✏️</span> New Chat
  </div>
  <div class="conv-sb-search">
    <span class="conv-sb-search-icon">🔍</span>
    <input type="text" placeholder="Search conversations…" id="convSbSearch" oninput="filterConvs(this.value)">
  </div>
  <div class="conv-sb-list" id="convSbList">
    <div style="font-size:11px;color:rgba(60,100,160,.5);text-align:center;padding:20px">Loading…</div>
  </div>
</div>
<!-- ══ Profile Panel ═══════════════════════════════════════════════════════ -->
<div class="profile-panel hidden" id="profilePanel">
  <div class="profile-box" style="position:relative">
    <button class="profile-close-btn" onclick="closeProfile()">&#x2715;</button>
    <div class="profile-cover" id="profileCover"></div>
    <div class="profile-avatar-wrap" onclick="toggleEmojiPicker()">
      <div class="profile-avatar-big" id="profileAvatarBig">😊</div>
      <div class="profile-avatar-edit">✏️</div>
    </div>
    <div class="profile-body">
      <div class="profile-name" id="profileName">—</div>
      <div class="profile-gid" id="profileGid">guest_????</div>
      <div class="profile-stats">
        <div class="profile-stat"><div class="profile-stat-val" id="pStatConvs">—</div><div class="profile-stat-lbl">Chats</div></div>
        <div class="profile-stat"><div class="profile-stat-val" id="pStatMsgs">—</div><div class="profile-stat-lbl">Messages</div></div>
        <div class="profile-stat"><div class="profile-stat-val" id="pStatMem">—</div><div class="profile-stat-lbl">Memories</div></div>
      </div>
      <div class="profile-actions">
        <button class="profile-btn" onclick="exportChat()">📥 Export Chat</button>
        <button class="profile-btn" onclick="closeProfile();openSP('memory')">🧠 Memory</button>
        <button class="profile-btn" onclick="closeProfile();logout()">🚪 Sign Out</button>
      </div>
    </div>
    <div class="emoji-picker-wrap" id="emojiPickerWrap" style="display:none">
      <div class="emoji-picker-title">Choose your avatar</div>
      <div class="emoji-picker-grid" id="emojiGrid"></div>
    </div>
  </div>
</div>
<!-- ══ AI Tools Hub ════════════════════════════════════════════════════════ -->
<div class="tools-hub-overlay hidden" id="toolsHubOverlay" onclick="toolsHubClickOut(event)">
  <div class="tools-hub">
    <div class="tools-hub-head">
      <div class="tools-hub-title" id="toolsHubTitle">🛠 AI Tools</div>
      <button class="tools-hub-close" onclick="closeToolsHub()">&#x2715;</button>
    </div>
    <div class="tools-hub-body" id="toolsHubBody">
      <!-- populated by JS -->
    </div>
  </div>
</div>

<!-- LEFT SETTINGS SIDEBAR -->
<div class="lsb-overlay" id="lsbOverlay" onclick="closeLSB()"></div>
<div class="lsb" id="lsb">
  <div class="lsb-hdr"><h3>⚙️ Settings</h3><button class="lsb-x" onclick="closeLSB()">✕</button></div>
  <div class="lsb-sec">
    <h4>Navigation</h4>
    <div class="lsb-item" onclick="closeLSB();openSP()"><span class="ei2">🔑</span> API Keys</div>
    <div class="lsb-item" onclick="closeLSB();openSP('convs')"><span class="ei2">💬</span> Conversations</div>
    <div class="lsb-item" onclick="closeLSB();openSP('memory')"><span class="ei2">🧠</span> Memory</div>
    <div class="lsb-item" onclick="closeLSB();openSP('saved')"><span class="ei2">🔖</span> Saved Items</div>
    <div class="lsb-item" onclick="closeLSB();openSP('theme')"><span class="ei2">🎨</span> Theme</div>
    <div class="lsb-item" onclick="closeLSB();openSP('extra')"><span class="ei2">🔗</span> Custom Endpoint</div>
    <div class="lsb-item" onclick="closeLSB();openSP('scraper')"><span class="ei2">🦆</span> Web Search</div>
    <div class="lsb-item" onclick="closeLSB();openSP('tts')"><span class="ei2">🔊</span> Voice / TTS</div>
    <div class="lsb-item" onclick="closeLSB();openSP('voicechat')"><span class="ei2">🎤</span> Voice Chat</div>
    <div class="lsb-item" onclick="closeLSB();openSP('audio')"><span class="ei2">🎵</span> Audio Gen</div>
  </div>
  <div class="lsb-sec">
    <h4>Quick Actions</h4>
    <div class="lsb-item" onclick="closeLSB();clearChat()"><span class="ei2">🗑️</span> Clear Chat</div>
    <div class="lsb-item" onclick="closeLSB();exportChat()"><span class="ei2">📥</span> Export Chat</div>
    <div class="lsb-item" id="lsbDevBtn" style="display:none" onclick="closeLSB();openDevModal()"><span class="ei2">🛠️</span> Dev Dashboard</div>
  </div>
  <div class="lsb-sec" id="lsbWxWrap">
    <h4>🌤 Weather</h4>
    <div id="wxCard" style="display:none;background:rgba(6,182,212,.07);backdrop-filter:blur(16px);border:1px solid rgba(6,182,212,.18);border-radius:14px;padding:14px 16px;animation:fi .3s ease">
      <div style="display:flex;align-items:flex-start;gap:12px">
        <div id="wxIcon" style="font-size:36px;line-height:1;filter:drop-shadow(0 2px 8px rgba(6,182,212,.4))">🌤</div>
        <div style="flex:1;min-width:0">
          <div style="display:flex;align-items:baseline;gap:6px">
            <span id="wxTemp" style="font-size:28px;font-weight:800;color:var(--tx);font-family:'Bebas Neue',sans-serif;letter-spacing:1px">--°</span>
            <span id="wxFeels" style="font-size:11px;color:var(--tx3)">Feels --°</span>
          </div>
          <div id="wxCity" style="font-size:12px;font-weight:700;color:var(--tx);margin-top:1px;overflow:hidden;text-overflow:ellipsis;white-space:nowrap">Loading…</div>
          <div id="wxDesc" style="font-size:11px;color:var(--tx2);margin-top:2px;text-transform:capitalize"></div>
        </div>
      </div>
      <div style="display:flex;gap:10px;margin-top:10px;padding-top:10px;border-top:1px solid rgba(80,120,200,.1)">
        <div style="flex:1;text-align:center"><div style="font-size:10px;color:var(--tx3);margin-bottom:2px">HUMIDITY</div><div id="wxHumid" style="font-size:13px;font-weight:600;color:var(--tx)">--%</div></div>
        <div style="flex:1;text-align:center"><div style="font-size:10px;color:var(--tx3);margin-bottom:2px">WIND</div><div id="wxWind" style="font-size:13px;font-weight:600;color:var(--tx)">-- km/h</div></div>
        <div style="flex:1;text-align:center"><div style="font-size:10px;color:var(--tx3);margin-bottom:2px">UPDATED</div><div id="wxTime" style="font-size:11px;font-weight:600;color:var(--tx)">--:--</div></div>
      </div>
    </div>
    <div id="wxLoading" style="font-size:11px;color:var(--tx3);padding:8px 0">⏳ Fetching weather…</div>
  </div>
  <div style="padding:12px 12px 0">
    <button class="ghost-btn" style="width:100%;font-size:12px;padding:8px;color:var(--red)" onclick="doLogout()">🚪 Sign Out</button>
  </div>
</div>

<!-- RIGHT QUICK-TOOL BAR -->
<div class="rtb" id="rtb">
  <div class="rtb-btn" onclick="openSP('imagine')" title=""><span>🖼</span><span class="rtb-tip">Image Gen</span></div>
  <div class="rtb-btn" onclick="openSP('video')" title=""><span>🎬</span><span class="rtb-tip">Video</span></div>
  <div class="rtb-btn" onclick="openSP('threed')" title=""><span>🧊</span><span class="rtb-tip">3D Model</span></div>
  <div class="rtb-sep"></div>
  <div class="rtb-btn" onclick="openSP('audio')" title=""><span>🎵</span><span class="rtb-tip">Audio</span></div>
  <div class="rtb-btn" onclick="openSP('voicechat')" title=""><span>🎤</span><span class="rtb-tip">Voice</span></div>
  <div class="rtb-sep"></div>
  <div class="rtb-btn" onclick="openSP('model')" title=""><span>🤖</span><span class="rtb-tip">Model</span></div>
  <div class="rtb-btn" onclick="openSP('convs')" title=""><span>💬</span><span class="rtb-tip">History</span></div>
  <div class="rtb-btn" onclick="openSP('memory')" title=""><span>🧠</span><span class="rtb-tip">Memory</span></div>
  <div class="rtb-sep"></div>
  <div class="rtb-btn" onclick="openSP('theme')" title=""><span>🎨</span><span class="rtb-tip">Theme</span></div>
  <div class="rtb-btn" onclick="openSP('scraper')" title=""><span>🌐</span><span class="rtb-tip">Fusion.Browser</span></div>
  <div class="rtb-btn" onclick="doDeepResearch(document.getElementById('msgIn').value.trim()||'latest AI breakthroughs')" title=""><span>🔬</span><span class="rtb-tip">Deep Research</span></div>
  <div class="rtb-sep"></div>
  <div class="rtb-btn" onclick="toggleLSB()" title=""><span>⚙️</span><span class="rtb-tip">Settings</span></div>
</div>

  <header>
    <div class="hl" style="gap:8px;display:flex;align-items:center">
      <button class="sidebar-toggle" onclick="toggleConvSidebar()" title="Conversations">&#x2630;</button>
      <div class="hlm">&#x26A1;</div><span class="htxt">Fusion.AI</span>
    </div>
    <div class="hr">
      <div class="mpill"><div class="ldot"></div><span id="mlabel">Auto</span></div>
      <div class="ibtn" id="themeBtn" onclick="toggleTheme()" title="Toggle theme">&#x1F319;</div>
      <div class="ibtn" id="chatOptBtn" onclick="toggleChatOpts()" title="Chat options">&#x1F4AC;</div>
      <div class="ibtn" id="voiceChatHeaderBtn" onclick="openSP('voicechat')" title="Voice Chat">&#x1F3A4;</div>
      <div class="ibtn" onclick="openSP()" title="Settings">&#x2699;&#xFE0F;</div>
      <div class="ibtn" id="devHeaderBtn" onclick="openDevModal()" title="Dev Dashboard" style="display:none">&#x1F6E0;&#xFE0F;</div>
      <div style="position:relative">
        <div class="uchip" id="uChip" onclick="toggleDrop()">
          <div class="uav" id="uAv" onclick="openProfile()" style="cursor:pointer;font-size:16px" title="View profile">?</div>
          <span id="uName">Guest</span>
          <span style="color:var(--tx3);font-size:9px;margin-left:2px">&#x25BE;</span>
        </div>
        <div class="drop" id="uDrop">
          <div class="ditem" onclick="openSP()">&#x2699;&#xFE0F; Settings</div>
          <div class="ditem" onclick="openSP('audio')">&#x1F3B5; Audio / Music</div>
          <div class="ditem" onclick="openSP('voicechat')">&#x1F3A4; Voice Chat</div>
          <div class="ditem" onclick="openSP('saved')">&#x1F516; Saved Items</div>
          <div class="ditem" onclick="openSP('memory')">&#x1F9E0; Memory</div>
          <div class="ditem" onclick="openSP('theme')">&#x1F3A8; Theme</div>
          <div class="ditem" onclick="clearChat()">&#x1F5D1;&#xFE0F; Clear Chat</div>
          <div class="dsep" id="devDropSep" style="display:none"></div>
          <div class="ditem dev-item" id="devDropItem" onclick="openDevModal()" style="display:none">&#x1F6E0;&#xFE0F; Dev Dashboard</div>
          <div class="dsep"></div>
          <div class="ditem danger" onclick="doLogout()">&#x1F6AA; Sign Out</div>
        </div>
      </div>
    </div>

<!-- CHAT OPTIONS PANEL -->
<div class="chat-opts-panel" id="chatOptsPanel" style="display:none">
  <div class="co-item" onclick="setChatFontSize('small')">&#x1F524; Small text</div>
  <div class="co-item" onclick="setChatFontSize('normal')">&#x1F4AC; Normal text</div>
  <div class="co-item" onclick="setChatFontSize('large')">&#x1F535; Large text</div>
  <div class="co-sep"></div>
  <div class="co-item" onclick="exportChat()">&#x1F4E5; Export chat</div>
  <div class="co-item" onclick="clearChat();closeChatOpts()">&#x1F5D1;&#xFE0F; Clear chat</div>
</div>
</header>
  <div class="chat-body" id="chatBody"><div class="msgs" id="msgs"></div></div>
  <div class="voice-bar" id="voiceBar">
    <div class="vbars"><div class="vbar"></div><div class="vbar"></div><div class="vbar"></div><div class="vbar"></div><div class="vbar"></div></div>
    <span class="vtxt" id="voiceTxt">Listening&#x2026;</span>
    <button class="vstop" onclick="stopVoice()">Stop</button>
  </div>
  <div class="attach-preview" id="imgPrevWrap">
    <img class="ap-img" id="imgPrev" src="" alt=""/>
    <span class="ap-name" id="imgPrevName"></span>
    <button class="ap-del" onclick="clearImage()">&#x2715;</button>
  </div>
  <div class="attach-preview" id="filePrevWrap">
    <span id="fileIcon" style="font-size:20px">&#x1F4C4;</span>
    <div style="flex:1;min-width:0"><div class="ap-name" id="fileName"></div><div style="font-size:10px;color:var(--tx3);margin-top:2px" id="fileSize"></div></div>
    <button class="ap-del" onclick="clearFile()">&#x2715;</button>
  </div>
  <div class="iz">
    <div class="ibox">
      <textarea id="msgIn" placeholder="Ask anything — web search always on · /imagine · /video · /search [query]" rows="1" onkeydown="handleKey(event)" oninput="ar(this);updateTokenCount(this.value)"></textarea>
      <div class="ib-row">
        <input type="file" id="imgInput" accept="image/*" style="display:none" onchange="handleImageFile(event)"/>
        <input type="file" id="fileInput" accept="*/*" style="display:none" onchange="handleFileUpload(event)"/>
        <input type="file" id="cameraInput" accept="image/*" capture="environment" style="display:none" onchange="handleImageFile(event)"/>
        <div class="iibtn" onclick="document.getElementById('imgInput').click()" title="Attach Image">🖼️</div>
        <div class="iibtn cam-btn" onclick="document.getElementById('cameraInput').click()" title="Take Photo">📷</div>
        <div class="iibtn" onclick="document.getElementById('fileInput').click()" title="Attach File">📎</div>
        <div class="iibtn" id="micBtn" onclick="toggleVoice()" title="Voice Input">🎤</div>
        <button class="sbtn stop-chat-btn" id="stopChatBtn" onclick="stopChat()" title="Stop" style="display:none;background:rgba(124,58,237,.85)">⏹</button>
        <button class="sbtn" id="sbtn" onclick="sendMsg()">➤</button>
      </div>
    </div>
    <!-- Quick Action Chips -->
    <div class="qa-bar" id="qaBar">
      <div class="qa-chip" onclick="openToolsHub()">🛠 Tools Hub</div>
      <div class="qa-chip" onclick="sendQuickMsg('Summarise this for me')">📋 Summarise</div>
      <div class="qa-chip" onclick="sendQuickMsg('Translate this to English')">🌍 Translate</div>
      <div class="qa-chip" onclick="sendQuickMsg('Rewrite this more professionally')">✍️ Rewrite</div>
      <div class="qa-chip" onclick="sendQuickMsg('Debug this code and explain the issues')">🐛 Debug</div>
      <div class="qa-chip" onclick="sendQuickMsg('Explain this step by step')">🧪 Explain</div>
      <div class="qa-chip" onclick="openFileCreator()">📁 Create File</div>
      <div class="qa-chip" onclick="openArena()">🏆 Arena</div>
    </div>
    <!-- Tool bar — below chatbox, icons expand inline panels -->
    <div class="tbar" id="tbar">
      <!-- Decluttered toolbar: Hub + Arena + Deep Think + File -->
      <div class="toolbox-wrap" id="toolboxWrap">
        <button class="toolbox-btn" id="toolboxBtn" onclick="toggleToolbox(event)">⚡ Tools ▾</button>
        <div class="toolbox-menu hidden" id="toolboxMenu">
          <div class="toolbox-section">Search & Research</div>
          <div class="toolbox-item" onclick="tbarToggle('web');closeToolbox()"><span class="toolbox-item-icon">🌐</span><div><div class="toolbox-item-label">Web Search</div><div class="toolbox-item-desc">Live internet results</div></div></div>
          <div class="toolbox-item" onclick="tbarToggle('deep');closeToolbox()"><span class="toolbox-item-icon">🔬</span><div><div class="toolbox-item-label">Deep Research</div><div class="toolbox-item-desc">Multi-source analysis</div></div></div>
          <div class="toolbox-item" onclick="tbarToggle('browser');closeToolbox()"><span class="toolbox-item-icon">🗺</span><div><div class="toolbox-item-label">Browser</div><div class="toolbox-item-desc">Fetch any webpage</div></div></div>
          <div class="toolbox-divider"></div>
          <div class="toolbox-section">AI Modes</div>
          <div class="toolbox-item" onclick="tbarToggle('math');closeToolbox()"><span class="toolbox-item-icon">📐</span><div><div class="toolbox-item-label">Math Mode</div><div class="toolbox-item-desc">Equations & proofs</div></div></div>
          <div class="toolbox-item" onclick="tbarToggle('code');closeToolbox()"><span class="toolbox-item-icon">💻</span><div><div class="toolbox-item-label">Code Mode</div><div class="toolbox-item-desc">Programming & debug</div></div></div>
          <div class="toolbox-item" onclick="tbarToggle('write');closeToolbox()"><span class="toolbox-item-icon">✍️</span><div><div class="toolbox-item-label">Write Mode</div><div class="toolbox-item-desc">Essays, blogs, stories</div></div></div>
          <div class="toolbox-item" onclick="tbarToggle('sum');closeToolbox()"><span class="toolbox-item-icon">⚡</span><div><div class="toolbox-item-label">Summarise</div><div class="toolbox-item-desc">Condense any content</div></div></div>
          <div class="toolbox-item" onclick="tbarToggle('trans');closeToolbox()"><span class="toolbox-item-icon">🌍</span><div><div class="toolbox-item-label">Translate</div><div class="toolbox-item-desc">Any language</div></div></div>
          <div class="toolbox-item" onclick="tbarToggle('explain');closeToolbox()"><span class="toolbox-item-icon">🧪</span><div><div class="toolbox-item-label">Explain</div><div class="toolbox-item-desc">Break down anything</div></div></div>
          <div class="toolbox-item" onclick="tbarToggle('compare');closeToolbox()"><span class="toolbox-item-icon">⚖️</span><div><div class="toolbox-item-label">Compare</div><div class="toolbox-item-desc">Side by side analysis</div></div></div>
          <div class="toolbox-item" onclick="tbarToggle('debug');closeToolbox()"><span class="toolbox-item-icon">🐛</span><div><div class="toolbox-item-label">Debug</div><div class="toolbox-item-desc">Find & fix errors</div></div></div>
          <div class="toolbox-divider"></div>
          <div class="toolbox-section">Multi-Model</div>
          <div class="toolbox-item" onclick="openArena();closeToolbox()"><span class="toolbox-item-icon">🏆</span><div><div class="toolbox-item-label">Model Arena</div><div class="toolbox-item-desc">5 AIs answer, you pick</div></div></div>
          <div class="toolbox-item" onclick="openExtremeThink();closeToolbox()"><span class="toolbox-item-icon">🧠</span><div><div class="toolbox-item-label">Extreme Deep Think</div><div class="toolbox-item-desc">All models · 5 rounds · 12 searches</div></div></div>
          <div class="toolbox-divider"></div>
          <div class="toolbox-section">Generate</div>
          <div class="toolbox-item" onclick="openSP('imagine');closeToolbox()"><span class="toolbox-item-icon">🎨</span><div><div class="toolbox-item-label">Generate Image</div><div class="toolbox-item-desc">AI image creation</div></div></div>
          <div class="toolbox-item" onclick="openSvgGen();closeToolbox()"><span class="toolbox-item-icon">🖼</span><div><div class="toolbox-item-label">AI Art (Free)</div><div class="toolbox-item-desc">LLM → SVG → PNG, no API cost</div></div></div>
          <div class="toolbox-divider"></div>
          <div class="toolbox-section">Power Tools</div>
          <div class="toolbox-item" onclick="openFusionOS();closeToolbox()"><span class="toolbox-item-icon">⬛</span><div><div class="toolbox-item-label">FusionOS</div><div class="toolbox-item-desc">AI desktop — terminal, files, code, agent</div></div></div>
          <div class="toolbox-item" onclick="openComputer();closeToolbox()"><span class="toolbox-item-icon">🖥</span><div><div class="toolbox-item-label">AI Computer</div><div class="toolbox-item-desc">Web search → cited AI answer</div></div></div>
          <div class="toolbox-item" onclick="openFileCreator();closeToolbox()"><span class="toolbox-item-icon">📁</span><div><div class="toolbox-item-label">Create File</div><div class="toolbox-item-desc">Doc, CSV, PPT, code…</div></div></div>
        </div>
      </div>
      <div class="tbar-sep"></div>
      <button class="tbar-btn" id="tbtn-os" onclick="openFusionOS()" style="background:linear-gradient(135deg,rgba(74,158,255,.14),rgba(168,85,247,.14));border-color:rgba(168,85,247,.28)">⬛ <span>FusionOS</span></button>
      <button class="tbar-btn" id="tbtn-file" onclick="openFileCreator()">📁 <span>Files</span></button>
      <div class="tbar-sep"></div>
      <div class="tok-ctr" id="tokCtr" title="~tokens in message">
        <span class="tok-num" id="tokNum">0</span><span style="font-size:9px;color:var(--tx3)">tok</span>
        <div class="tok-bar"><div class="tok-fill" id="tokFill" style="width:0%"></div></div>
      </div>
    </div>

    <!-- ── Tool panels ─────────────────────────────────────── -->
    <div class="tpanel" id="tpanel-web" style="display:none">
      <div class="tpanel-hd"><span class="tpanel-ico">🌐</span><span class="tpanel-title">Web Search</span><span class="tpanel-badge green">Always On</span><button class="tpanel-x" onclick="tbarClose()">✕</button></div>
      <p class="tpanel-desc">Every message is automatically web-grounded via Bing search. Use this for a quick direct search sent straight to chat.</p>
      <div class="tpanel-row">
        <input class="tpanel-input" id="tinput-web" placeholder="Search anything…" onkeydown="if(event.key==='Enter'){_tpSend('web');}"/>
        <button class="tpanel-go" onclick="_tpSend('web')">Search ↵</button>
      </div>
    </div>

    <div class="tpanel" id="tpanel-deep" style="display:none">
      <div class="tpanel-hd"><span class="tpanel-ico">🔬</span><span class="tpanel-title">Deep Research</span><span class="tpanel-badge">Every AI</span><button class="tpanel-x" onclick="tbarClose()">✕</button></div>
      <p class="tpanel-desc">Sends your question to <strong>every available model simultaneously</strong> — Groq, OpenRouter, GitHub, Cloudflare — then synthesises a master answer with per-model timing.</p>
      <div class="tpanel-row">
        <input class="tpanel-input" id="tinput-deep" placeholder="Research question — all models answer in parallel…" onkeydown="if(event.key==='Enter'){_tpDeep();}"/>
        <button class="tpanel-go" onclick="_tpDeep()">Research ↵</button>
      </div>
    </div>

    <div class="tpanel" id="tpanel-math" style="display:none">
      <div class="tpanel-hd"><span class="tpanel-ico">📐</span><span class="tpanel-title">Math Solver</span><span class="tpanel-badge">All Math AIs</span><button class="tpanel-x" onclick="tbarClose()">✕</button></div>
      <p class="tpanel-desc">Routes to all best math models (o3, o4-mini, QwQ, Qwen3, DeepSeek R1) simultaneously. Renders with KaTeX — use $…$ inline or $$…$$ for display equations.</p>
      <div class="tpanel-row">
        <input class="tpanel-input" id="tinput-math" placeholder="Solve, prove, integrate, derive… e.g. ∫₀^∞ e^{-x²} dx" onkeydown="if(event.key==='Enter'){_tpMath();}"/>
        <button class="tpanel-go" onclick="_tpMath()">Solve ↵</button>
      </div>
    </div>

    <div class="tpanel" id="tpanel-code" style="display:none">
      <div class="tpanel-hd"><span class="tpanel-ico">💻</span><span class="tpanel-title">Code Engine</span><span class="tpanel-badge">All Code AIs</span><button class="tpanel-x" onclick="tbarClose()">✕</button></div>
      <p class="tpanel-desc">Hits every available coding model (Qwen Coder, GPT-4o, DeepSeek V3, Phi-4, Llama) in parallel. Returns best-voted answer with syntax highlighting.</p>
      <div class="tpanel-row">
        <input class="tpanel-input" id="tinput-code" placeholder="What to build, fix, or explain… any language" onkeydown="if(event.key==='Enter'){_tpCode();}"/>
        <button class="tpanel-go" onclick="_tpCode()">Code ↵</button>
      </div>
    </div>

    <div class="tpanel" id="tpanel-write" style="display:none">
      <div class="tpanel-hd"><span class="tpanel-ico">✍️</span><span class="tpanel-title">Writing Assistant</span><span class="tpanel-badge">Creative</span><button class="tpanel-x" onclick="tbarClose()">✕</button></div>
      <p class="tpanel-desc">Essays, emails, stories, posts, scripts, cover letters — specify tone and format in your prompt.</p>
      <div class="tpanel-row">
        <input class="tpanel-input" id="tinput-write" placeholder="Write a… [format] about [topic] in [tone] style" onkeydown="if(event.key==='Enter'){_tpSendPre('write','[Write] ');}"/>
        <button class="tpanel-go" onclick="_tpSendPre('write','[Write] ')">Write ↵</button>
      </div>
    </div>

    <div class="tpanel" id="tpanel-sum" style="display:none">
      <div class="tpanel-hd"><span class="tpanel-ico">⚡</span><span class="tpanel-title">Summarise</span><span class="tpanel-badge">Fast</span><button class="tpanel-x" onclick="tbarClose()">✕</button></div>
      <p class="tpanel-desc">Paste text, a URL, or describe content. Get a concise, structured summary with key points.</p>
      <div class="tpanel-row">
        <input class="tpanel-input" id="tinput-sum" placeholder="Paste text or describe what to summarise…" onkeydown="if(event.key==='Enter'){_tpSendPre('sum','Summarise this concisely with key bullet points:\n');}"/>
        <button class="tpanel-go" onclick="_tpSendPre('sum','Summarise this concisely with key bullet points:\n')">Summarise ↵</button>
      </div>
    </div>

    <div class="tpanel" id="tpanel-trans" style="display:none">
      <div class="tpanel-hd"><span class="tpanel-ico">🌍</span><span class="tpanel-title">Translate</span><span class="tpanel-badge">100+ langs</span><button class="tpanel-x" onclick="tbarClose()">✕</button></div>
      <div style="display:flex;gap:7px;margin-bottom:8px;align-items:center">
        <select id="transFrom" class="tpanel-sel"><option value="auto">Auto-detect</option><option>English</option><option>Arabic</option><option>French</option><option>Spanish</option><option>German</option><option>Chinese</option><option>Japanese</option><option>Korean</option><option>Portuguese</option><option>Russian</option><option>Italian</option><option>Hindi</option><option>Turkish</option></select>
        <span style="color:var(--tx3);font-size:14px">→</span>
        <select id="transTo" class="tpanel-sel"><option>English</option><option>Arabic</option><option>French</option><option>Spanish</option><option>German</option><option>Chinese</option><option>Japanese</option><option>Korean</option><option>Portuguese</option><option>Russian</option><option>Italian</option><option>Hindi</option><option>Turkish</option></select>
      </div>
      <div class="tpanel-row">
        <input class="tpanel-input" id="tinput-trans" placeholder="Text to translate…" onkeydown="if(event.key==='Enter'){_tpTrans();}"/>
        <button class="tpanel-go" onclick="_tpTrans()">Translate ↵</button>
      </div>
    </div>

    <div class="tpanel" id="tpanel-debug" style="display:none">
      <div class="tpanel-hd"><span class="tpanel-ico">🐛</span><span class="tpanel-title">Debug</span><span class="tpanel-badge">Fix It</span><button class="tpanel-x" onclick="tbarClose()">✕</button></div>
      <p class="tpanel-desc">Paste your broken code + error message. Gets routed to all coding models and returns a fixed version with explanation.</p>
      <div class="tpanel-row">
        <input class="tpanel-input" id="tinput-debug" placeholder="Paste code or describe the bug + error message…" onkeydown="if(event.key==='Enter'){_tpSendPre('debug','Debug and fix this code. Explain the root cause and provide the corrected version:\n');}"/>
        <button class="tpanel-go" onclick="_tpSendPre('debug','Debug and fix this code. Explain the root cause and provide the corrected version:\n')">Debug ↵</button>
      </div>
    </div>

    <div class="tpanel" id="tpanel-explain" style="display:none">
      <div class="tpanel-hd"><span class="tpanel-ico">🧪</span><span class="tpanel-title">Explain</span><span class="tpanel-badge">Deep Dive</span><button class="tpanel-x" onclick="tbarClose()">✕</button></div>
      <div style="display:flex;gap:6px;margin-bottom:8px;flex-wrap:wrap">
        <button class="tpanel-lvl active" onclick="_setLevel(this,'eli5')">👶 Simple</button>
        <button class="tpanel-lvl" onclick="_setLevel(this,'normal')">🎓 Normal</button>
        <button class="tpanel-lvl" onclick="_setLevel(this,'expert')">🔬 Expert</button>
        <button class="tpanel-lvl" onclick="_setLevel(this,'analogy')">💡 Analogy</button>
      </div>
      <div class="tpanel-row">
        <input class="tpanel-input" id="tinput-explain" placeholder="Concept, topic, or paste text to explain…" onkeydown="if(event.key==='Enter'){_tpExplain();}"/>
        <button class="tpanel-go" onclick="_tpExplain()">Explain ↵</button>
      </div>
    </div>

    <div class="tpanel" id="tpanel-compare" style="display:none">
      <div class="tpanel-hd"><span class="tpanel-ico">⚖️</span><span class="tpanel-title">Compare</span><span class="tpanel-badge">Pro/Con</span><button class="tpanel-x" onclick="tbarClose()">✕</button></div>
      <p class="tpanel-desc">Compare two things side by side — technologies, products, ideas, frameworks — with pros, cons, and a recommendation.</p>
      <div style="display:flex;gap:7px;margin-bottom:8px">
        <input class="tpanel-input" id="tinput-cmp1" placeholder="Option A…" style="flex:1"/>
        <span style="color:var(--tx3);align-self:center;font-size:12px;font-weight:700">vs</span>
        <input class="tpanel-input" id="tinput-cmp2" placeholder="Option B…" style="flex:1"/>
      </div>
      <div class="tpanel-row">
        <input class="tpanel-input" id="tinput-cmpctx" placeholder="Context (optional) — for what use case?" style="flex:1"/>
        <button class="tpanel-go" onclick="_tpCompare()">Compare ↵</button>
      </div>
    </div>

    <div class="tpanel" id="tpanel-browser" style="display:none">
      <div class="tpanel-hd"><span class="tpanel-ico">🌐</span><span class="tpanel-title">Fusion.Browser</span><span class="tpanel-badge">Private</span><button class="tpanel-x" onclick="tbarClose()">✕</button></div>
      <p class="tpanel-desc">Search and browse the web. Results open in the full Fusion.Browser panel with AI summarisation.</p>
      <div class="tpanel-row">
        <input class="tpanel-input" id="tinput-browser" placeholder="Search or enter URL…" onkeydown="if(event.key==='Enter'){_tpBrowser();}"/>
        <button class="tpanel-go" onclick="_tpBrowser()">Browse ↵</button>
      </div>
    </div>
    <div class="ifooter">
      <span class="hint">Enter ↵ · Shift+Enter newline · 🌐 Web always on</span>
      <button class="tbtn" onclick="clearChat()">✕ Clear</button>
    </div>
  </div>
</div>

<!-- SIDE PANEL -->
<div class="ovl" id="ovl" onclick="closeSP()"></div>
<div class="sp" id="sp">
  <div class="sp-head"><h3>&#x2699;&#xFE0F; Settings</h3><button class="sp-close" onclick="closeSP()">&#x2715;</button></div>
  <div class="sp-inner">
  <div class="sp-tabs">
    <button class="sp-tab active" onclick="switchSPTab('model')" id="spt-model" title="Model"><span class="st-ico">🤖</span><span class="st-lbl">Model</span></button>
    <button class="sp-tab" onclick="switchSPTab('imagine')" id="spt-imagine"><span class="st-ico">🖼</span><span class="st-lbl">Image</span></button>
    <button class="sp-tab" onclick="switchSPTab('video')" id="spt-video"><span class="st-ico">🎬</span><span class="st-lbl">Video</span></button>
    <button class="sp-tab" onclick="switchSPTab('threed')" id="spt-threed"><span class="st-ico">🧊</span><span class="st-lbl">3D</span></button>
    <button class="sp-tab" onclick="switchSPTab('audio')" id="spt-audio"><span class="st-ico">🎵</span><span class="st-lbl">Audio</span></button>
    <button class="sp-tab" onclick="switchSPTab('voicechat')" id="spt-voicechat"><span class="st-ico">🎤</span><span class="st-lbl">Voice</span></button>
    <button class="sp-tab" onclick="switchSPTab('tts')" id="spt-tts"><span class="st-ico">🔊</span><span class="st-lbl">TTS</span></button>
    <div class="sp-tab-sep"></div>
    <button class="sp-tab" onclick="switchSPTab('convs')" id="spt-convs"><span class="st-ico">💬</span><span class="st-lbl">Chats</span></button>
    <button class="sp-tab" onclick="switchSPTab('saved')" id="spt-saved"><span class="st-ico">🔖</span><span class="st-lbl">Saved</span></button>
    <button class="sp-tab" onclick="switchSPTab('memory')" id="spt-memory"><span class="st-ico">🧠</span><span class="st-lbl">Memory</span></button>
    <div class="sp-tab-sep"></div>
    <button class="sp-tab" onclick="switchSPTab('theme')" id="spt-theme"><span class="st-ico">🎨</span><span class="st-lbl">Theme</span></button>
    <button class="sp-tab" onclick="switchSPTab('keys')" id="spt-keys"><span class="st-ico">⚙️</span><span class="st-lbl">Info</span></button>
    <button class="sp-tab" onclick="switchSPTab('scraper')" id="spt-scraper" title="Fusion.Browser"><span class="st-ico">🌐</span><span class="st-lbl">Browser</span></button>
    <button class="sp-tab" onclick="switchSPTab('extra')" id="spt-extra"><span class="st-ico">🔗</span><span class="st-lbl">Custom EP</span></button>
  </div>
  <div class="sp-body">
    <!-- INFO — clean user-facing panel, no API key details -->
    <div id="spTab-keys">
      <div class="sp-sec" style="margin-top:16px">
        <div style="background:var(--glass-surf);border:1px solid var(--glass-bdr2);border-radius:14px;padding:22px 18px;text-align:center">
          <div style="font-size:36px;margin-bottom:12px">⚡</div>
          <div style="font-size:16px;font-weight:700;color:var(--tx);margin-bottom:8px">Fusion.AI</div>
          <div style="font-size:12px;color:var(--tx2);line-height:1.8;margin-bottom:18px">
            Powered by our proprietary AI infrastructure.<br>
            All models and generation services run on our hardware.
          </div>
          <div style="display:flex;flex-direction:column;gap:8px;text-align:left">
            <div style="background:rgba(34,197,94,.08);border:1px solid rgba(34,197,94,.2);border-radius:10px;padding:11px 14px;font-size:12px;color:var(--tx2);display:flex;align-items:center;gap:8px">
              <span style="font-size:16px">🖼️</span><div><strong style="color:var(--tx)">Image Generation</strong><div style="margin-top:2px;color:var(--tx3)">High-quality multi-model AI generation</div></div>
            </div>
            <div style="background:rgba(34,197,94,.08);border:1px solid rgba(34,197,94,.2);border-radius:10px;padding:11px 14px;font-size:12px;color:var(--tx2);display:flex;align-items:center;gap:8px">
              <span style="font-size:16px">🎬</span><div><strong style="color:var(--tx)">Video Generation</strong><div style="margin-top:2px;color:var(--tx3)">Fusion Video Engine — fast diffusion</div></div>
            </div>
            <div style="background:rgba(34,197,94,.08);border:1px solid rgba(34,197,94,.2);border-radius:10px;padding:11px 14px;font-size:12px;color:var(--tx2);display:flex;align-items:center;gap:8px">
              <span style="font-size:16px">💬</span><div><strong style="color:var(--tx)">AI Chat</strong><div style="margin-top:2px;color:var(--tx3)">50+ models — GPT-5, DeepSeek, Llama 4 &amp; more</div></div>
            </div>
            <div style="background:rgba(34,197,94,.08);border:1px solid rgba(34,197,94,.2);border-radius:10px;padding:11px 14px;font-size:12px;color:var(--tx2);display:flex;align-items:center;gap:8px">
              <span style="font-size:16px">🧊</span><div><strong style="color:var(--tx)">3D Generation</strong><div style="margin-top:2px;color:var(--tx3)">Fusion 3D Engine — photorealistic renders</div></div>
            </div>
          </div>
        </div>
      </div>
      <div class="sp-sec">
        <button class="ghost-btn" style="font-size:13px;padding:10px" onclick="doLogout()">🚪 Sign Out</button>
      </div>
    </div>
    <!-- MODEL -->
    <div id="spTab-model" style="display:none">
      <div class="sp-sec" style="margin-top:16px">
        <h4>Model / Mode Selection</h4>
        <p style="font-size:11px;color:var(--tx2);margin-bottom:10px;line-height:1.6">Select a chat model. <strong>Auto</strong> mode picks the best model for each request. Use the 🖼 Image tab for image generation.</p>
        <select class="msel" id="mOverride" onchange="savePref()">
          <option value="auto">&#x26A1; Auto — AI Overseer picks best model</option>
        </select>
        <div id="imgModeNote" style="display:none;margin-top:10px;padding:10px 13px;background:rgba(6,182,212,.1);border:1px solid rgba(6,182,212,.2);border-radius:10px;font-size:11px;color:var(--tx2);line-height:1.6">
          &#x1F3A8; <strong>Image Mode active.</strong> Type your image prompt and press Send.
        </div>
        <p style="font-size:11px;color:var(--tx3);margin-top:10px;line-height:1.6">In Auto mode, Fusion.AI selects the best available model for your query automatically.</p>
      </div>
    </div>
    <!-- VOICE CHAT -->
    <div id="spTab-voicechat" style="display:none">
      <div class="sp-sec" style="margin-top:16px">
        <h4>&#x1F3A4; Voice Chat Mode</h4>
        <p style="font-size:11px;color:var(--tx2);margin-bottom:12px;line-height:1.6">Speak to Fusion.AI and hear responses. Uses browser speech recognition + TTS.</p>
        <div class="voice-chat-panel">
          <div class="vc-status" id="vcStatus">Click to start voice conversation. Speak, then AI responds aloud.</div>
          <button class="vc-btn" id="vcBtn" onclick="toggleVoiceChat()">&#x1F3A4; Start Voice Chat</button>
          <div class="vc-transcript" id="vcTranscript">Your speech appears here...</div>
        </div>
        <div class="krow" style="margin-top:12px">
          <div class="krow-top"><span class="kprov">&#x1F50A; Voice</span></div>
          <select class="msel" id="vcVoiceSelect" onchange="saveVoiceChatPref()" style="margin-bottom:8px">
            <option value="">Default Voice</option>
          </select>
          <div class="dur-row">
            <label>Speed</label>
            <input type="range" id="vcRate" min="0.5" max="2" step="0.1" value="1" oninput="document.getElementById('vcRateLbl').textContent=parseFloat(this.value).toFixed(1)+'x';saveVoiceChatPref()"/>
            <span class="dur-lbl" id="vcRateLbl">1.0x</span>
          </div>
        </div>
      </div>
    </div>
    <!-- KEYS/INFO Tab — API + endpoint settings -->
    <div id="spTab-extra" style="display:none">
      <div class="sp-sec" style="margin-top:16px">
        <h4>&#x1F517; Connections</h4>
        <div class="krow">
          <div class="krow-top"><span class="kprov">&#x1F517; Custom API Endpoint</span><span class="kstat" id="customEndpointStat">Not set</span></div>
          <p style="font-size:10px;color:var(--tx3);margin:0 0 8px">Add your own AI endpoint (OpenAI-compatible API URL)</p>
          <input type="url" id="customEndpointUrl" placeholder="https://your-endpoint.com/v1/chat/completions" autocomplete="off" style="width:100%;background:var(--inp);border:1.5px solid var(--glass-bdr);border-radius:9px;padding:8px 10px;color:var(--tx);font-family:'DM Mono',monospace;font-size:11px;outline:none;margin-bottom:6px;box-sizing:border-box"/>
          <input type="text" id="customEndpointModel" placeholder="Model name (e.g. gpt-4o, llama-3)" autocomplete="off" style="width:100%;background:var(--inp);border:1.5px solid var(--glass-bdr);border-radius:9px;padding:8px 10px;color:var(--tx);font-family:'DM Mono',monospace;font-size:11px;outline:none;margin-bottom:6px;box-sizing:border-box"/>
          <div style="display:flex;gap:6px;margin-bottom:6px">
            <input type="password" id="customEndpointKey" placeholder="API Key (optional)" autocomplete="off" style="flex:1;background:var(--inp);border:1.5px solid var(--glass-bdr);border-radius:9px;padding:8px 10px;color:var(--tx);font-family:'DM Mono',monospace;font-size:11px;outline:none"/>
            <button class="ksave" onclick="saveCustomEndpoint()">Save</button>
          </div>
          <div style="font-size:10px;color:var(--tx3);line-height:1.6">Saved endpoint appears as <strong>Custom</strong> in model selector.</div>
        </div>
        <div class="krow" style="margin-top:8px">
          <div class="krow-top"><span class="kprov">&#x1F517; Backup API</span><span class="kstat" id="extraStat">Not set</span></div>
          <p style="font-size:10px;color:var(--tx3);margin:0 0 8px">Custom endpoint — your own model, proxy, or service</p>
          <input type="text" id="extraUrlIn" placeholder="https://your-api.com/v1/chat/completions" style="width:100%;background:var(--inp);border:1.5px solid var(--glass-bdr);border-radius:9px;padding:8px 10px;color:var(--tx);font-family:'DM Mono',monospace;font-size:11px;outline:none;margin-bottom:6px"/>
          <div style="display:flex;gap:6px">
            <input type="password" id="extraIn" placeholder="API key" autocomplete="off" style="flex:1;background:var(--inp);border:1.5px solid var(--glass-bdr);border-radius:9px;padding:8px 10px;color:var(--tx);font-family:'DM Mono',monospace;font-size:11px;outline:none"/>
            <button class="ksave" onclick="saveKey('extra')">Save</button>
          </div>
        </div>
      </div>
    </div>
<!-- IMAGE GEN -->
    <div id="spTab-imagine" style="display:none">
      <div class="sp-sec" style="margin-top:16px">
        <h4>&#x1F3A8; AI Image Generation</h4>
        <div style="background:rgba(34,197,94,.08);border:1px solid rgba(34,197,94,.2);border-radius:10px;padding:10px 13px;margin-bottom:12px;font-size:11px;color:var(--tx2);line-height:1.7">
          &#x2705; High-quality AI image generation — multiple models available.
        </div>
        <div class="img-gen-panel">
          <h4>&#x2728; Generate Image</h4>
          <textarea class="ig-prompt" id="igPrompt" placeholder="Describe the image you want to create..."></textarea>
          <div class="ig-row">
            <select class="ig-sel" id="igModel">
              <optgroup label="☁️ Enhanced (server configured)">
                <option value="@cf/black-forest-labs/flux-2-klein-4b">FLUX Klein 4B — fast</option>
                <option value="@cf/black-forest-labs/flux-2-klein-9b">FLUX Klein 9B — quality</option>
                <option value="@cf/black-forest-labs/flux-2-dev">FLUX Dev — highest quality</option>
                <option value="@cf/leonardo/phoenix-1.0">Leonardo Phoenix 1.0</option>
                <option value="@cf/leonardo/lucid-origin">Leonardo Lucid Origin</option>
                <option value="@cf/deepgram/flux">Deepgram FLUX</option>
              </optgroup>
              <optgroup label="🏭 CF Worker (custom endpoint)">
                <option value="worker:sdxl">Worker · SDXL</option>
                <option value="worker:dreamshaper">Worker · DreamShaper</option>
                <option value="worker:lightning">Worker · Lightning (fast)</option>
                <option value="worker:fast">Worker · Fast (turbo)</option>
              </optgroup>
              <optgroup label="✨ Premium (server configured)">
                <option value="openai/gpt-image-1">GPT-image-1 — best quality</option>
                <option value="openai/dall-e-3">DALL-E 3 — photorealistic</option>
              </optgroup>
              <optgroup label="🎨 Standard (always available)">
                <option value="flux" selected>FLUX — great quality</option>
                <option value="flux-realism">FLUX Realism — photorealistic</option>
                <option value="flux-anime">FLUX Anime — anime style</option>
                <option value="flux-3d">FLUX 3D — 3D renders</option>
                <option value="turbo">Turbo — fastest</option>
              </optgroup>
              <optgroup label="🔮 Stability AI (needs STABILITY_KEY)" id="stabilityImgGroup">
                <option value="stable-image-core">Stable Image Core</option>
                <option value="stable-image-ultra">Stable Image Ultra — highest quality</option>
                <option value="sd3-medium">SD3 Medium</option>
                <option value="sd3-large">SD3 Large</option>
                <option value="sd3-large-turbo">SD3 Large Turbo — fast</option>
              </optgroup>
            </select>
            <select class="ig-sel" id="igSize">
              <option value="1024x1024">1024×1024 Square</option>
              <option value="1280x720">1280×720 Wide</option>
              <option value="720x1280">720×1280 Tall</option>
              <option value="1920x1080">1920×1080 HD</option>
            </select>
          </div>
          <button class="ig-btn" id="igGenBtn" onclick="generateImage()" style="margin-bottom:6px">&#x2728; Generate Image</button>
          <button class="gen-stop-btn" id="imgStopBtn" onclick="stopGeneration('igGenBtn','imgStopBtn')">&#x23F9;&#xFE0F; Stop</button>
          <div class="ig-result" id="igResult">
            <img class="ig-img" id="igImg" src="" alt="Generated Image" style="cursor:zoom-in" onclick="openImgFull(this.src)"/>
            <div id="igBackendTag" style="font-size:10px;color:var(--tx3);margin:4px 0 6px;text-align:right"></div>
            <div class="ig-dl">
              <button id="igDlBtn" onclick="downloadGenImg()">&#x2B07;&#xFE0F; Download</button>
              <button onclick="sendGenImgToChat()">&#x1F4AC; Send to Chat</button>
            </div>
          </div>
        </div>
      </div>
    </div>
    <!-- VIDEO GEN -->
    <div id="spTab-video" style="display:none">
      <div class="sp-sec" style="margin-top:16px">
        <h4>🎬 AI Video Generation</h4>
        <div style="background:rgba(34,197,94,.08);border:1px solid rgba(34,197,94,.2);border-radius:10px;padding:10px 13px;margin-bottom:12px;font-size:11px;color:var(--tx2);line-height:1.7">
          ✅ Free video generation — no API key needed.
        </div>
        <div class="img-gen-panel">
          <h4>🎬 Generate Video</h4>
          <textarea class="ig-prompt" id="videoPrompt" placeholder="Describe your video scene... e.g. 'A majestic eagle soaring over snow-capped mountains, cinematic, golden hour'"></textarea>
          <div class="ig-row">
            <select class="ig-sel" id="videoModel">
              <option value="animatediff">AnimateDiff — smooth motion</option>
              <option value="ltx">LTX — high quality</option>
              <option value="wan">Wan — fast</option>
            </select>
            <select class="ig-sel" id="videoSize">
              <option value="512x512">512×512 Square</option>
              <option value="768x432">768×432 Wide</option>
              <option value="432x768">432×768 Tall</option>
            </select>
          </div>
          <button class="ig-btn" id="videoGenBtn" onclick="generateVideo()" style="margin-bottom:6px">🎬 Generate Video</button>
          <button class="gen-stop-btn" id="videoStopBtn" onclick="stopVideoGen()">⏹️ Stop</button>
          <div id="videoResult" style="display:none;margin-top:10px">
            <video id="videoEl" controls style="width:100%;border-radius:10px;border:1px solid var(--glass-bdr2);display:block;margin-bottom:8px"></video>
            <div id="videoBackendTag" style="font-size:10px;color:var(--tx3);margin-bottom:8px"></div>
            <div class="ig-dl">
              <button onclick="downloadVideo()">⬇️ Download</button>
              <button onclick="sendVideoToChat()">💬 Send to Chat</button>
            </div>
          </div>
        </div>
      </div>
    </div>
    <!-- 3D GEN -->
    <div id="spTab-threed" style="display:none">
      <div class="sp-sec" style="margin-top:16px">
        <h4>🧊 AI 3D Generation</h4>
        <div style="background:rgba(34,197,94,.08);border:1px solid rgba(34,197,94,.2);border-radius:10px;padding:10px 13px;margin-bottom:12px;font-size:11px;color:var(--tx2);line-height:1.7">
          ✅ Free 3D generation — renders a photorealistic preview image.
        </div>
        <div class="img-gen-panel">
          <h4>🧊 Generate 3D Model</h4>
          <textarea class="ig-prompt" id="threedPrompt" placeholder="Describe the 3D object... e.g. 'A futuristic sports car, metallic red, studio lighting, 3D render'" style="min-height:60px"></textarea>
          <div class="ig-row">
            <select class="ig-sel" id="threedStyle">
              <option value="3d render">3D Render — realistic</option>
              <option value="low poly 3d">Low Poly — stylized</option>
              <option value="clay render 3d">Clay Render — matte</option>
              <option value="isometric 3d">Isometric — game-ready</option>
              <option value="blender 3d">Blender Style — detailed</option>
            </select>
            <select class="ig-sel" id="threedSize">
              <option value="1024x1024">1024×1024</option>
              <option value="1280x720">1280×720 Wide</option>
            </select>
          </div>
          <button class="ig-btn" id="threedGenBtn" onclick="generate3D()" style="margin-bottom:6px">🧊 Generate 3D</button>
          <button class="gen-stop-btn" id="threedStopBtn" onclick="stop3DGen()">⏹️ Stop</button>
          <div id="threedResult" style="display:none;margin-top:10px">
            <div id="threedModelViewer" style="display:none;margin-bottom:8px"></div>
            <img id="threedImg" src="" style="width:100%;border-radius:10px;border:1px solid var(--glass-bdr2);display:block;margin-bottom:8px;cursor:zoom-in" onclick="openImgFull(this.src)"/>
            <div id="threedModelTag" style="font-size:10px;color:var(--tx3);margin-bottom:8px"></div>
            <div class="ig-dl">
              <button onclick="download3D()">⬇️ Download</button>
              <button onclick="send3DToChat()">💬 Send to Chat</button>
            </div>
          </div>
        </div>
      </div>
    </div>
    <div id="spTab-audio" style="display:none">
      <div class="sp-sec" style="margin-top:16px">
        <h4>&#x1F3B5; Audio &amp; Music Generation</h4>
        <div style="background:rgba(167,139,250,.08);border:1px solid rgba(167,139,250,.25);border-radius:10px;padding:10px 13px;margin-bottom:12px;font-size:11px;color:var(--tx2);line-height:1.7">
          Music and speech generation coming soon.
        </div>
        <div class="audio-panel">
          <h4>&#x1F3B6; Music Generation</h4>
          <textarea class="ig-prompt" id="musicPrompt" placeholder="e.g. 'upbeat jazz piano, 120bpm' or 'calm ambient nature sounds'"></textarea>
          <button class="ig-btn" id="musicGenBtn" onclick="generateMusic()">&#x1F3B5; Generate Music</button>
          <button class="gen-stop-btn" id="audioStopBtn" onclick="stopGeneration('musicGenBtn','audioStopBtn')">&#x23F9;&#xFE0F; Stop</button>
          <div class="audio-result" id="musicResult"><audio id="musicAudio" controls></audio><button id="musicDlBtn" style="display:none;margin-top:6px;width:100%" class="ig-btn" onclick="downloadAudio('musicAudio','music.mp3')">⬇️ Download Music</button></div>
        </div>
        <div class="audio-panel">
          <h4>&#x1F4E3; AI Speech Synthesis</h4>
          <textarea class="ig-prompt" id="speechText" placeholder="Enter text to convert to speech..."></textarea>
          <div style="display:flex;gap:7px">
            <button class="ig-btn" id="speechGenBtn" onclick="generateSpeech()" style="flex:2">&#x1F4E3; AI Speech</button>
            <button class="ig-btn" onclick="browserTTS()" style="flex:1;background:rgba(6,182,212,.8)">&#x1F50A; Browser</button>
          </div>
          <div style="display:flex;gap:7px;margin-top:6px"><select class="ig-sel" id="speechVoice"><option value="alloy">Alloy</option><option value="echo">Echo</option><option value="fable">Fable</option><option value="onyx">Onyx</option><option value="nova">Nova</option><option value="shimmer">Shimmer</option></select></div>
          <button class="gen-stop-btn" id="speechStopBtn" onclick="stopGeneration('speechGenBtn','speechStopBtn')">&#x23F9;&#xFE0F; Stop</button>
          <div class="audio-result" id="speechResult"><audio id="speechAudio" controls></audio><button id="speechDlBtn" style="display:none;margin-top:6px;width:100%" class="ig-btn" onclick="downloadAudio('speechAudio','speech.mp3')">⬇️ Download Audio</button></div>
        </div>
      </div>
    </div>
    <!-- TTS -->
    <div id="spTab-tts" style="display:none">
      <div class="sp-sec" style="margin-top:16px">
        <h4>&#x1F50A; Auto-Read Responses</h4>
        <p style="font-size:11px;color:var(--tx2);margin-bottom:14px;line-height:1.6">Automatically read AI responses aloud using browser TTS.</p>
        <div class="krow">
          <div class="krow-top"><span class="kprov">AI Voice</span><span class="kstat" id="voiceStatus" style="display:none">Enabled</span></div>
          <select class="msel" id="voiceSelect" onchange="saveVoicePref()" style="margin-bottom:10px">
            <option value="">&#x1F507; Off (no voice)</option>
          </select>
          <div style="display:flex;gap:8px;margin-top:6px">
            <button class="ghost-btn" style="font-size:12px;padding:8px;flex:1" onclick="previewVoice()">&#x25B6; Preview</button>
            <button class="ghost-btn" style="font-size:12px;padding:8px;flex:1" onclick="stopSpeech()">&#x23F9; Stop</button>
          </div>
          <div style="margin-top:12px">
            <label style="font-size:10px;font-weight:600;color:var(--tx2);text-transform:uppercase;letter-spacing:1px;display:block;margin-bottom:6px">Speed</label>
            <div style="display:flex;align-items:center;gap:10px">
              <input type="range" id="voiceRate" min="0.5" max="2" step="0.1" value="1" oninput="updateRateLabel(this);saveVoicePref()" style="flex:1;accent-color:var(--red)"/>
              <span id="voiceRateLabel" style="font-size:11px;color:var(--tx2);min-width:28px;font-family:'DM Mono',monospace">1.0x</span>
            </div>
          </div>
        </div>
      </div>
    </div>
    <!-- SAVED -->
    <div id="spTab-saved" style="display:none">
      <div class="sp-sec" style="margin-top:16px">
        <h4>Saved Items</h4>
        <div id="savedList"><div class="empty-state"><span class="ei">&#x1F516;</span>Nothing saved yet.</div></div>
      </div>
    </div>
    <!-- CONVERSATIONS -->
    <div id="spTab-convs" style="display:none">
      <div class="sp-sec" style="margin-top:16px">
        <div style="display:flex;align-items:center;justify-content:space-between;margin-bottom:12px">
          <h4 style="margin:0">💬 Saved Conversations</h4>
          <button onclick="newConv()" style="background:var(--grad);border:none;border-radius:8px;padding:5px 12px;color:#fff;font-size:11px;font-weight:600;cursor:pointer">+ New</button>
        </div>
        <p style="font-size:11px;color:var(--tx2);margin-bottom:12px;line-height:1.6">Your chat history is saved automatically. Pick up any conversation where you left off.</p>
        <div id="convList"><div class="empty-state"><span class="ei">💬</span>No saved conversations yet.</div></div>
      </div>
    </div>
    <!-- MEMORY -->
    <div id="spTab-memory" style="display:none">
      <div class="sp-sec" style="margin-top:16px">
        <h4>AI Memory</h4>
        <p style="font-size:11px;color:var(--tx2);margin-bottom:12px;line-height:1.6">The AI remembers these facts about you in every conversation.</p>
        <div id="memoryList"><div class="empty-state"><span class="ei">&#x1F9E0;</span>No memories yet.</div></div>
        <div class="mem-add-row">
          <input type="text" id="memKey" placeholder="key (e.g. name)" style="max-width:110px"/>
          <input type="text" id="memVal" placeholder="value" onkeydown="if(event.key==='Enter')addMemory()"/>
          <button class="ksave" onclick="addMemory()">Add</button>
        </div>
      </div>
    </div>
    <!-- THEME -->
    <div id="spTab-theme" style="display:none">
      <div class="sp-sec" style="margin-top:16px">
        <h4>&#x1F3A8; Theme &amp; Background</h4>
        <div class="krow" style="margin-bottom:10px">
          <div class="krow-top"><span class="kprov">Color Mode</span></div>
          <div style="display:flex;gap:8px;margin-top:6px;flex-wrap:wrap">
            <button class="theme-mode-btn" id="tmb-system" onclick="setThemeMode('system')">&#x1F4BB; System</button>
            <button class="theme-mode-btn" id="tmb-dark" onclick="setThemeMode('dark')">&#x1F319; Dark</button>
            <button class="theme-mode-btn" id="tmb-light" onclick="setThemeMode('light')">&#x2600;&#xFE0F; Light</button>
          </div>
        </div>
        <div class="krow" style="margin-bottom:10px">
          <div class="krow-top"><span class="kprov">Accent Color</span></div>
          <div style="display:flex;flex-wrap:wrap;gap:8px;margin-top:6px" id="accentPicker">
            <div class="accent-swatch" data-accent="red" onclick="setAccent('red')" style="background:linear-gradient(135deg,#7c3aed,#06b6d4)" title="Red/Blue"></div>
            <div class="accent-swatch" data-accent="purple" onclick="setAccent('purple')" style="background:linear-gradient(135deg,#9333ea,#06b6d4)" title="Purple/Cyan"></div>
            <div class="accent-swatch" data-accent="green" onclick="setAccent('green')" style="background:linear-gradient(135deg,#16a34a,#0ea5e9)" title="Green/Blue"></div>
            <div class="accent-swatch" data-accent="orange" onclick="setAccent('orange')" style="background:linear-gradient(135deg,#ea580c,#eab308)" title="Orange/Yellow"></div>
            <div class="accent-swatch" data-accent="pink" onclick="setAccent('pink')" style="background:linear-gradient(135deg,#db2777,#9333ea)" title="Pink/Purple"></div>
            <div class="accent-swatch" data-accent="teal" onclick="setAccent('teal')" style="background:linear-gradient(135deg,#0f766e,#1d4ed8)" title="Teal/Indigo"></div>
          </div>
        </div>
        <div class="krow" style="margin-bottom:10px">
          <div class="krow-top"><span class="kprov">Background Style</span></div>
          <div style="display:grid;grid-template-columns:1fr 1fr;gap:8px;margin-top:8px" id="bgPicker">
            <div class="bg-swatch" data-bg="orbs" onclick="setBg('orbs')" style="background:radial-gradient(circle at 30% 40%,rgba(124,58,237,.5),transparent 50%),radial-gradient(circle at 70% 60%,rgba(6,182,212,.4),transparent 50%),#060a12;color:#fff">&#x1F525; Orbs</div>
            <div class="bg-swatch" data-bg="mesh" onclick="setBg('mesh')" style="background:linear-gradient(135deg,#0a0a1a,#0d1a2e,#0a1520);color:#6a90c0">&#x1F4A0; Mesh</div>
            <div class="bg-swatch" data-bg="aurora" onclick="setBg('aurora')" style="background:linear-gradient(135deg,#001a0a,#0a001a,#001a1a);color:#4ade80">&#x1F30C; Aurora</div>
            <div class="bg-swatch" data-bg="plain" onclick="setBg('plain')" style="background:#060810;color:#3a5070">&#x25A0; Plain</div>
            <div class="bg-swatch" data-bg="sunset" onclick="setBg('sunset')" style="background:linear-gradient(135deg,#1a0010,#2a0820,#0a1030);color:#f472b6">&#x1F307; Sunset</div>
            <div class="bg-swatch" data-bg="forest" onclick="setBg('forest')" style="background:linear-gradient(135deg,#001a08,#0a1a10);color:#4ade80">&#x1F332; Forest</div>
            <div class="bg-swatch" data-bg="matrix" onclick="setBg('matrix')" style="background:#000900;color:#00ff41;font-family:'DM Mono',monospace">&#x1F4BB; Matrix</div>
          </div>
        </div>
        <div class="krow">
          <div class="krow-top"><span class="kprov">Custom Gradient Colors</span></div>
          <div style="display:flex;gap:8px;align-items:center;margin-top:6px">
            <input type="color" id="customBgColor" value="#060a12" oninput="setCustomBg(this.value)" style="width:44px;height:36px;border:none;border-radius:8px;cursor:pointer;padding:0;background:none"/>
            <input type="color" id="customBgColor2" value="#0a0e1c" oninput="setCustomBg2(this.value)" style="width:44px;height:36px;border:none;border-radius:8px;cursor:pointer;padding:0;background:none"/>
            <span style="font-size:11px;color:var(--tx3);flex:1;line-height:1.4">Pick 2 colors for gradient background</span>
          </div>
          <button class="ghost-btn" style="font-size:12px;padding:8px;margin-top:10px" onclick="resetTheme()">&#x21BA; Reset to Default</button>
        </div>
      </div>
    </div>

    <!-- FUSION BROWSER -->
    <div id="spTab-scraper" style="display:none;flex-direction:column;padding:0;height:calc(100vh - 52px)">
      <!-- Browser chrome bar -->
      <div style="padding:10px 12px 8px;border-bottom:1px solid var(--glass-bdr);background:rgba(6,10,22,0.6);backdrop-filter:blur(20px);flex-shrink:0">
        <div style="display:flex;align-items:center;gap:6px;margin-bottom:8px">
          <div style="width:10px;height:10px;border-radius:50%;background:#ff5f57"></div>
          <div style="width:10px;height:10px;border-radius:50%;background:#febc2e"></div>
          <div style="width:10px;height:10px;border-radius:50%;background:#28c840"></div>
          <span style="font-size:11px;font-weight:700;color:var(--tx);margin-left:6px;letter-spacing:.5px">Fusion.Browser</span>
          <span style="margin-left:auto;font-size:10px;color:var(--tx3)">🔒 Private · AI-Powered</span>
        </div>
        <div style="display:flex;gap:6px;align-items:center">
          <button onclick="fbBack()" style="background:rgba(255,255,255,.06);border:1px solid var(--glass-bdr);border-radius:7px;width:28px;height:28px;color:var(--tx2);cursor:pointer;font-size:13px;display:flex;align-items:center;justify-content:center;flex-shrink:0" title="Back">‹</button>
          <button onclick="fbForward()" style="background:rgba(255,255,255,.06);border:1px solid var(--glass-bdr);border-radius:7px;width:28px;height:28px;color:var(--tx2);cursor:pointer;font-size:13px;display:flex;align-items:center;justify-content:center;flex-shrink:0" title="Forward">›</button>
          <button onclick="fbRefresh()" style="background:rgba(255,255,255,.06);border:1px solid var(--glass-bdr);border-radius:7px;width:28px;height:28px;color:var(--tx2);cursor:pointer;font-size:13px;display:flex;align-items:center;justify-content:center;flex-shrink:0" title="Refresh">↻</button>
          <div style="flex:1;display:flex;align-items:center;background:rgba(8,14,30,0.7);backdrop-filter:blur(16px);border:1.5px solid var(--glass-bdr);border-radius:9px;padding:0 10px;gap:7px;height:32px">
            <span id="fbSecure" style="font-size:11px;color:#22c55e">🔒</span>
            <input id="fbUrl" type="text" value="https://www.bing.com/" placeholder="Search or enter URL…"
              style="flex:1;background:none;border:none;outline:none;color:var(--tx);font-size:12px;font-family:'DM Mono',monospace"
              onkeydown="if(event.key==='Enter')fbNavigate()"/>
            <button onclick="fbNavigate()" style="background:var(--grad);border:none;border-radius:6px;padding:3px 10px;color:#fff;font-size:11px;font-weight:700;cursor:pointer">Go</button>
          </div>
          <button onclick="fbSendToAI()" style="background:rgba(6,182,212,.15);border:1px solid rgba(6,182,212,.3);border-radius:7px;padding:0 10px;height:28px;color:var(--blue);cursor:pointer;font-size:11px;font-weight:600;flex-shrink:0;white-space:nowrap">💬 Ask AI</button>
        </div>
        <!-- Quick search bar -->
        <div style="display:flex;gap:5px;margin-top:8px;flex-wrap:wrap">
          <button onclick="fbSearch('news today')" class="fb-chip">📰 News</button>
          <button onclick="fbSearch('weather today')" class="fb-chip">🌤 Weather</button>
          <button onclick="fbLoadUrl('https://www.bing.com/')" class="fb-chip">🔎 Bing</button>
          <button onclick="fbLoadUrl('https://en.m.wikipedia.org/')" class="fb-chip">📖 Wikipedia</button>
          <button onclick="fbLoadUrl('https://news.ycombinator.com/')" class="fb-chip">🔶 HN</button>
        </div>
      </div>
      <!-- Loading bar -->
      <div id="fbLoadBar" style="height:2px;background:var(--grad);width:0%;transition:width .3s;flex-shrink:0"></div>
      <!-- Status bar -->
      <div id="fbStatus" style="padding:3px 12px;font-size:10px;color:var(--tx3);background:rgba(6,10,22,0.4);flex-shrink:0;display:none"></div>
      <!-- Browser results area (DDG search rendered natively) -->
      <div id="fbContent" style="flex:1;overflow-y:auto;padding:12px">
        <!-- Welcome screen -->
        <div id="fbWelcome" style="text-align:center;padding:32px 16px">
          <div style="font-size:48px;margin-bottom:12px">🌐</div>
          <div style="font-size:18px;font-weight:700;color:var(--tx);margin-bottom:6px;font-family:'Bebas Neue',sans-serif;letter-spacing:2px">FUSION.BROWSER</div>
          <div style="font-size:12px;color:var(--tx2);margin-bottom:24px;line-height:1.7">AI-powered browser · Private · No tracking<br>Search, browse, and instantly ask AI about any page</div>
          <div style="display:flex;gap:8px;justify-content:center;flex-wrap:wrap">
            <button onclick="fbSearch('latest AI news')" style="background:rgba(6,182,212,.12);border:1px solid rgba(6,182,212,.25);border-radius:10px;padding:8px 14px;color:var(--tx2);font-size:12px;cursor:pointer">🤖 AI News</button>
            <button onclick="fbSearch('cryptocurrency prices')" style="background:rgba(124,58,237,.1);border:1px solid rgba(124,58,237,.2);border-radius:10px;padding:8px 14px;color:var(--tx2);font-size:12px;cursor:pointer">📈 Crypto</button>
            <button onclick="fbSearch('science breakthroughs 2025')" style="background:rgba(34,197,94,.08);border:1px solid rgba(34,197,94,.18);border-radius:10px;padding:8px 14px;color:var(--tx2);font-size:12px;cursor:pointer">🔬 Science</button>
          </div>
        </div>
        <!-- Search results -->
        <div id="fbResults" style="display:none"></div>
      </div>
    </div>
  </div>
  </div><!-- /sp-inner -->
</div>

<!-- DEV DASHBOARD MODAL -->
<div class="modal-bg" id="devModal">
  <div class="modal">
    <div class="modal-header">
      <h2>&#x1F6E0;&#xFE0F; Dev Dashboard</h2>
      <p>Real-time stats for Fusion.AI</p>
      <button class="modal-close" onclick="closeDevModal()">&#x2715;</button>
    </div>
    <div class="modal-tabs">
      <button class="modal-tab active" id="mdt-overview" onclick="switchModalTab('overview')">&#x1F4CA; Overview</button>
      <button class="modal-tab" id="mdt-users" onclick="switchModalTab('users')">&#x1F464; Users</button>
      <button class="modal-tab" id="mdt-logs" onclick="switchModalTab('logs')">&#x1F4DC; Logs</button>
      <button class="modal-tab" id="mdt-models" onclick="switchModalTab('models')">&#x1F916; Models</button>
      <button class="modal-tab" id="mdt-limits" onclick="switchModalTab('limits')">&#x1F6E1; Limits</button>
      <button class="modal-tab" id="mdt-flags" onclick="switchModalTab('flags')">&#x1F6A9; Flags</button>
      <button class="modal-tab" id="mdt-console" onclick="switchModalTab('console')">&#x1F9EA; Console</button>
      <button class="modal-tab" id="mdt-tokens" onclick="switchModalTab('tokens')">&#x1F511; Keys</button>
      <button class="modal-tab" id="mdt-activity" onclick="switchModalTab('activity')">&#x1F4AC; Activity</button>
      <button class="modal-tab" id="mdt-visitors" onclick="switchModalTab('visitors')">&#x1F4CD; Visitors</button>
      <button class="modal-tab" id="mdt-allmodels" onclick="switchModalTab('allmodels')">&#x1F916; All Models</button>
      <button class="modal-tab" id="mdt-menutoggle" onclick="switchModalTab('menutoggle')">&#x1F6A7; Menu</button>
    </div>
    <div class="modal-body" id="devModalBody">
      <div style="text-align:center;padding:40px;color:var(--tx3)">&#x23F3; Loading&#x2026;</div>
    </div>
  </div>
</div>

<div class="toast" id="toast"></div>


<!-- ══ Arena Modal ═══════════════════════════════════════════════════════ -->
<div class="file-modal-overlay hidden" id="arenaModalOverlay" onclick="arenaClickOut(event)">
  <div class="file-modal" id="arenaModal" style="width:680px;max-width:96vw">
    <div class="file-modal-head">
      <div class="file-modal-title">🏆 Model Arena</div>
      <button class="file-modal-close" onclick="closeArena()">&#x2715;</button>
    </div>
    <div class="file-modal-body">
      <div style="font-size:12px;color:rgba(120,180,255,.7);margin-bottom:10px">Ask a question — up to 5 AI models answer simultaneously. Pick the best response.</div>
      <div class="file-topic-label">Your Question</div>
      <textarea class="file-topic-input" id="arenaPromptInput" rows="2" placeholder="e.g. Explain quantum entanglement in simple terms…"></textarea>
      <div class="file-topic-label" style="margin-top:10px">Models (pick up to 5, or leave empty for auto)</div>
      <div class="arena-model-sel" id="arenaModelSel"></div>
      <div style="font-size:10px;color:rgba(60,100,160,.6);margin-top:4px" id="arenaSelCount">0 selected · auto-pick 5</div>
    </div>
    <div class="file-modal-foot">
      <button class="file-cancel-btn" onclick="closeArena()">Cancel</button>
      <button class="file-gen-btn" id="arenaRunBtn" onclick="runArena()">🏆 Run Arena</button>
    </div>
  </div>
</div>
<!-- ══ Extreme Deep Think Modal ════════════════════════════════════════════ -->
<div class="file-modal-overlay hidden" id="edtModalOverlay" onclick="edtClickOut(event)">
  <div class="file-modal" id="edtModal" style="width:560px;max-width:96vw">
    <div class="file-modal-head">
      <div class="file-modal-title">🧠 Extreme Deep Think</div>
      <button class="file-modal-close" onclick="closeEDT()">&#x2715;</button>
    </div>
    <div class="file-modal-body">
      <div style="font-size:12px;color:rgba(150,180,255,.7);margin-bottom:12px;line-height:1.6">
        12 web searches · up to 12 AI models · 5 rounds of thinking each · master synthesis.<br>
        <span style="color:rgba(100,140,200,.5);font-size:10px">Takes 30–90 seconds. Best for complex questions requiring deep analysis.</span>
      </div>
      <div class="file-topic-label">Your Question</div>
      <textarea class="file-topic-input" id="edtPromptInput" rows="3" placeholder="e.g. What are the long-term economic effects of universal basic income? Analyse all perspectives."></textarea>
    </div>
    <div class="file-modal-foot">
      <button class="file-cancel-btn" onclick="closeEDT()">Cancel</button>
      <button class="file-gen-btn" id="edtRunBtn" onclick="runEDT()">🧠 Start Deep Think</button>
    </div>
  </div>
</div>

<!-- ══ File Creator Modal ════════════════════════════════════════════════ -->
<div class="file-modal-overlay hidden" id="fileModalOverlay" onclick="fileModalClickOut(event)">
  <div class="file-modal" id="fileModal">
    <div class="file-modal-head">
      <div class="file-modal-title">🗂 Create a File</div>
      <button class="file-modal-close" onclick="closeFileCreator()">&#x2715;</button>
    </div>
    <div class="file-modal-body">
      <div class="file-type-grid" id="fileTypeGrid"></div>
      <div class="file-topic-wrap">
        <div class="file-topic-label">Topic / Description</div>
        <textarea class="file-topic-input" id="fileTopicInput" rows="3" placeholder="e.g. Sales data for Q1 2026 with product, region and revenue columns..."></textarea>
      </div>
    </div>
    <div class="file-modal-foot">
      <button class="file-cancel-btn" onclick="closeFileCreator()">Cancel</button>
      <button class="file-gen-btn" id="fileGenBtn" onclick="submitFileCreator()">&#x2728; Generate File</button>
    </div>
  </div>
</div>

<script>
var uName='', hist=[], loading=false, isDark=true, currentMKey='auto', isDevUser=false, userAvatar='';
// ── Persistent login via cookies (365-day) ───────────────────────────────────
function _setCookie(n,v,days){var e=new Date();e.setDate(e.getDate()+(days||365));document.cookie=n+'='+encodeURIComponent(v)+';expires='+e.toUTCString()+';path=/;SameSite=Lax';}
function _getCookie(n){var ca=document.cookie.split(';');for(var i=0;i<ca.length;i++){var c=ca[i].trim();if(c.startsWith(n+'='))return decodeURIComponent(c.slice(n.length+1));}return '';}
function _delCookie(n){document.cookie=n+'=;expires=Thu, 01 Jan 1970 00:00:00 UTC;path=/';}
var authToken=_getCookie('fusion_token')||localStorage.getItem('fusion_token')||'';
var pendingImageB64='', pendingImageMime='', pendingImageName='';
var pendingFileText='', pendingFileName='', pendingFileSize=0;
var recognition=null, isRecording=false;
var devStats=null, devTestResults=null, currentModalTab='overview';
var ttsEnabled=false;
var voicePref = JSON.parse(localStorage.getItem('fusion_voice')||'{}');
var vcPref = JSON.parse(localStorage.getItem('fusion_vc')||'{}');
var vcActive=false, vcRecog=null;
var lastGenImgB64='', lastGenImgMime='image/jpeg';
var allModels=[];

// ── Menu toggles — declared early so _applyMenuToggles() can be called any time ──
var _devMenuToggles = JSON.parse(localStorage.getItem('dev_menu_toggles')||'{}');
function _applyMenuToggles(){
  var toggleMap={
    'rtb_imagegen':document.querySelector('.rtb-btn[onclick*="imagine"]'),
    'rtb_video':document.querySelector('.rtb-btn[onclick*="video"]'),
    'rtb_threed':document.querySelector('.rtb-btn[onclick*="threed"]'),
    'rtb_audio':document.querySelector('.rtb-btn[onclick*="audio"]'),
    'rtb_voice':document.querySelector('.rtb-btn[onclick*="voicechat"]'),
    'sp_imagetab':document.getElementById('spt-imagine'),
    'sp_videotab':document.getElementById('spt-video'),
    'sp_threedtab':document.getElementById('spt-threed'),
    'sp_audiotab':document.getElementById('spt-audio'),
    'sp_voicetab':document.getElementById('spt-voicechat'),
    'sp_ttstab':document.getElementById('spt-tts'),
    'sp_scrapetab':document.getElementById('spt-scraper'),
    'header_voice':document.getElementById('voiceChatHeaderBtn'),
    'lsb_weather':document.getElementById('lsbWxWrap'),
  };
  Object.keys(toggleMap).forEach(function(k){
    var el=toggleMap[k];
    if(el){
      var hidden=_devMenuToggles[k]===false;
      el.style.display=hidden?'none':'';
    }
  });
}

// Theme state
var themePref = JSON.parse(_getCookie('fusion_theme')||localStorage.getItem('fusion_theme_pref')||'{"mode":"dark","accent":"red","bg":"plain"}');

function apiFetch(url, opts) {
  opts = opts||{}; opts.headers = opts.headers||{};
  if(authToken) opts.headers['X-Auth-Token'] = authToken;
  if(_userGeo&&_userGeo.lat){
    opts.headers['X-User-Lat']=String(_userGeo.lat||'');
    opts.headers['X-User-Lon']=String(_userGeo.lon||'');
    // City can contain non-Latin chars — strip to ASCII-safe for HTTP headers
    var safeCity=(_userGeo.city||'').replace(/[^\x20-\x7E]/g,'').trim().substring(0,80);
    opts.headers['X-User-City']=safeCity;
  }
  return fetch(url, opts);
}

var ACCENTS = {
  red:    {r:'#7c3aed',b:'#06b6d4',rg:'rgba(124,58,237,.25)',bg:'rgba(6,182,212,.22)',r2:'#6d28d9'},
  purple: {r:'#9333ea',b:'#06b6d4',rg:'rgba(147,51,234,.25)',bg:'rgba(6,182,212,.22)',r2:'#7e22ce'},
  green:  {r:'#16a34a',b:'#0ea5e9',rg:'rgba(22,163,74,.25)',bg:'rgba(14,165,233,.22)',r2:'#15803d'},
  orange: {r:'#ea580c',b:'#eab308',rg:'rgba(234,88,12,.25)',bg:'rgba(234,179,8,.22)',r2:'#c2410c'},
  pink:   {r:'#db2777',b:'#9333ea',rg:'rgba(219,39,119,.25)',bg:'rgba(147,51,234,.22)',r2:'#be185d'},
  teal:   {r:'#0f766e',b:'#1d4ed8',rg:'rgba(15,118,110,.25)',bg:'rgba(29,78,216,.22)',r2:'#0d6960'},
};
var BG_STYLES = {
  orbs:   {body:'#07090f', o1:'rgba(124,58,237,.35)', o2:'rgba(6,182,212,.3)',  o3:'rgba(120,40,200,.2)'},
  mesh:   {body:'#040810', o1:'rgba(26,60,120,.2)',   o2:'rgba(40,80,180,.18)', o3:'rgba(60,40,160,.12)'},
  aurora: {body:'#020c08', o1:'rgba(16,163,74,.3)',   o2:'rgba(56,189,248,.22)',o3:'rgba(147,51,234,.2)'},
  plain:  {body:'#000000', o1:'rgba(0,0,0,0)',         o2:'rgba(0,0,0,0)',       o3:'rgba(0,0,0,0)'},
  sunset: {body:'#0a0010', o1:'rgba(219,39,119,.35)', o2:'rgba(147,51,234,.25)',o3:'rgba(234,88,12,.2)'},
  forest: {body:'#020c06', o1:'rgba(22,163,74,.3)',   o2:'rgba(15,118,110,.22)',o3:'rgba(5,150,105,.18)'},
  matrix: {body:'#000900', o1:'rgba(0,255,0,.06)',    o2:'rgba(0,200,50,.04)',  o3:'rgba(0,150,30,.04)'},
};
var DARK_MODES = {
  dark:{
    'glass-bg':'rgba(10,10,20,0.45)','glass-bg2':'rgba(10,15,28,0.58)',
    'glass-surf':'rgba(16,16,30,0.52)','glass-surf2':'rgba(20,20,38,0.58)',
    'glass-bdr':'rgba(124,58,237,0.18)','glass-bdr2':'rgba(80,120,200,0.28)',
    'tx':'#e8f0ff','tx2':'#7a95c0','tx3':'#3a5070',
    'inp':'rgba(10,10,20,0.6)','shad':'rgba(0,0,0,.7)'
  },
  light:{
    'glass-bg':'rgba(255,255,255,0.98)','glass-bg2':'rgba(248,249,252,1.0)',
    'glass-surf':'rgba(240,242,248,0.96)','glass-surf2':'rgba(235,238,245,1.0)',
    'glass-bdr':'rgba(180,190,220,0.45)','glass-bdr2':'rgba(160,175,215,0.6)',
    'tx':'#0f172a','tx2':'#374151','tx3':'#9ca3af',
    'inp':'rgba(255,255,255,1.0)','shad':'rgba(0,0,0,.08)'
  },
  system:{
    'glass-bg':'rgba(10,10,20,0.45)','glass-bg2':'rgba(10,15,28,0.58)',
    'glass-surf':'rgba(16,16,30,0.52)','glass-surf2':'rgba(20,20,38,0.58)',
    'glass-bdr':'rgba(124,58,237,0.18)','glass-bdr2':'rgba(80,120,200,0.28)',
    'tx':'#e8f0ff','tx2':'#7a95c0','tx3':'#3a5070',
    'inp':'rgba(10,10,20,0.6)','shad':'rgba(0,0,0,.7)'
  },
};

var _matrixRunning=false,_matrixRaf=null;
function startMatrix(isDark){
  var c=document.getElementById('matrixCanvas');
  if(!c)return; c.classList.add('vis');
  if(_matrixRunning)return; _matrixRunning=true;
  c.width=window.innerWidth; c.height=window.innerHeight;
  var ctx=c.getContext('2d');
  var cols=Math.floor(c.width/14); var drops=Array(cols).fill(1);
  var chars='ABCDEFGHIJKLMNOPQRSTUVWXYZ0123456789@#$%^&*()アイウエオカキクケコサシスセソタチツテトナニヌネノ';
  var color=isDark?'#00ff41':'#003300';
  function draw(){
    ctx.fillStyle=isDark?'rgba(0,9,0,.05)':'rgba(240,255,240,.05)';
    ctx.fillRect(0,0,c.width,c.height);
    ctx.fillStyle=color; ctx.font='13px DM Mono,monospace';
    for(var i=0;i<drops.length;i++){
      var ch=chars[Math.floor(Math.random()*chars.length)];
      ctx.fillText(ch,i*14,drops[i]*14);
      if(drops[i]*14>c.height&&Math.random()>.975)drops[i]=0;
      drops[i]++;
    }
    _matrixRaf=requestAnimationFrame(draw);
  }
  draw();
}
function stopMatrix(){
  var c=document.getElementById('matrixCanvas');
  if(c)c.classList.remove('vis');
  _matrixRunning=false;
  if(_matrixRaf){cancelAnimationFrame(_matrixRaf);_matrixRaf=null;}
  var ctx=c&&c.getContext('2d');if(ctx)ctx.clearRect(0,0,c.width,c.height);
}
window.addEventListener('resize',function(){
  var c=document.getElementById('matrixCanvas');
  if(c&&_matrixRunning){c.width=window.innerWidth;c.height=window.innerHeight;}
});

function applyTheme(p) {
  if(!p) p = themePref;
  var root = document.documentElement.style;
  var rawMode = p.mode || 'system';
  var mode;
  if(rawMode === 'system') {
    mode = (window.matchMedia && window.matchMedia('(prefers-color-scheme: light)').matches) ? 'light' : 'dark';
  } else {
    mode = rawMode;
  }
  var acc = ACCENTS[p.accent||'red'];
  var bgKey = p.bg || 'plain';
  var bg = BG_STYLES[bgKey] || BG_STYLES['plain'];
  if(bgKey==='matrix'){startMatrix(mode!=='light');}else{stopMatrix();}
  var dm = DARK_MODES[mode] || DARK_MODES['dark'];
  // Accent
  root.setProperty('--red',   acc.r);root.setProperty('--blue',  acc.b);root.setProperty('--red2',  acc.r2);root.setProperty('--redg',  acc.rg);
  root.setProperty('--blueg', acc.bg);root.setProperty('--grad',  'linear-gradient(135deg,'+acc.r+','+acc.b+')');
  root.setProperty('--gradt', 'linear-gradient(90deg,'+acc.r+','+acc.b+')');
  // Dark/light vars
  Object.keys(dm).forEach(function(k){ root.setProperty('--'+k, dm[k]); });
  // Body background
  var bodyBg;
  if(p.customBg && p.customBg2) {
    bodyBg = 'linear-gradient(135deg,'+p.customBg+','+p.customBg2+')';
  } else if(p.customBg) {
    bodyBg = p.customBg;
  } else {
    bodyBg = bg.body;
  }
  root.setProperty('--bg-body', mode==='light'?'#ffffff':'#000000');
  var orb1=document.querySelector('.bg-orb1'),orb2=document.querySelector('.bg-orb2'),orb3=document.querySelector('.bg-orb3');
  var bgDiv=document.getElementById('bgDiv');
  if(mode==='light'){
    document.body.style.background='#ffffff';
    if(orb1)orb1.style.opacity='0'; if(orb2)orb2.style.opacity='0'; if(orb3)orb3.style.opacity='0';
    if(bgDiv)bgDiv.style.opacity='0';
  } else {
    document.body.style.background='#000000';
    // Hide orbs if bg is plain or matrix (transparent orb colors)
    var showOrbs = (bgKey !== 'plain' && bgKey !== 'matrix');
    if(orb1)orb1.style.opacity=showOrbs?'1':'0';
    if(orb2)orb2.style.opacity=showOrbs?'1':'0';
    if(orb3)orb3.style.opacity=showOrbs?'1':'0';
    if(bgDiv)bgDiv.style.opacity='1';
  }
  // Orbs
  var o1=document.querySelector('.bg-orb1');
  var o2=document.querySelector('.bg-orb2');
  var o3=document.querySelector('.bg-orb3');
  if(o1) o1.style.background='radial-gradient(circle,'+bg.o1+',transparent 65%)';if(o2) o2.style.background='radial-gradient(circle,'+bg.o2+',transparent 65%)';
  if(o3) o3.style.background='radial-gradient(circle,'+bg.o3+',transparent 65%)';
  // Light class
  document.body.classList.toggle('light', mode==='light');isDark = mode !== 'light';
  var tb=document.getElementById('themeBtn');
  if(tb) tb.textContent = isDark ? '\u{1F319}' : '\u2600\uFE0F';
}

function saveTheme() {
  localStorage.setItem('fusion_theme_pref', JSON.stringify(themePref));_setCookie('fusion_theme',JSON.stringify(themePref),365);applyTheme(themePref);updateThemeUI();
  var mtc=document.getElementById('metaThemeColor');if(mtc)mtc.content=(themePref.mode==='light'?'#ffffff':'#000000');
}

function setThemeMode(m){ themePref.mode=m; saveTheme(); try{var w=document.getElementById('hcaptchaWidget');if(w&&typeof hcaptcha!=='undefined'){w.setAttribute('data-theme',m==='light'?'light':'dark');hcaptcha.reset();}}catch(e){}  ['dark','light','system'].forEach(function(x){var b=document.getElementById('tmb-'+x);if(b)b.classList.toggle('active',x===m);}); }
function setAccent(a)    { themePref.accent=a; saveTheme(); }
function setBg(b)        { themePref.bg=b; delete themePref.customBg; delete themePref.customBg2; saveTheme(); }
function setCustomBg(v)  { themePref.customBg=v; themePref.bg='custom'; saveTheme(); }
function setCustomBg2(v) { themePref.customBg2=v; themePref.bg='custom'; saveTheme(); }
function resetTheme()    { themePref={mode:'dark',accent:'red',bg:'plain'}; saveTheme(); showToast('\u21BA Theme reset'); }

// ── Left sidebar toggle ────────────────────────────────────────────────────
function toggleLSB(){var el=document.getElementById('lsb'),ov=document.getElementById('lsbOverlay');el.classList.toggle('open');ov.classList.toggle('open');}
function closeLSB(){var el=document.getElementById('lsb'),ov=document.getElementById('lsbOverlay');el.classList.remove('open');ov.classList.remove('open');}
function openLSB(){var el=document.getElementById('lsb'),ov=document.getElementById('lsbOverlay');el.classList.add('open');ov.classList.add('open');}

// ── Google Sign-In ────────────────────────────────────────────────────────
async function doGoogleSignIn(){
  try{
    var r=await fetch('/api/auth/google');
    var d=await r.json();
    if(d.ok&&d.url){window.location.href=d.url;}
    else{showToast('Google sign-in not configured. Use username/password.',5000);}
  }catch(e){showToast('Google sign-in unavailable',4000);}
}

// ── Weather loader ────────────────────────────────────────────────────────
var _userGeo = JSON.parse(_getCookie('fusion_geo')||localStorage.getItem('fusion_geo')||'{}');
var _WX_ICONS={'clear sky':'☀️','few clouds':'🌤','scattered clouds':'⛅','broken clouds':'☁️','shower rain':'🌧','rain':'🌧','thunderstorm':'⛈','snow':'❄️','mist':'🌫','light rain':'🌦','overcast clouds':'☁️','moderate rain':'🌧'};
function _wxIcon(desc){desc=(desc||'').toLowerCase();for(var k in _WX_ICONS){if(desc.includes(k))return _WX_ICONS[k];}return '🌤';}
async function _loadWeather(){
  try{
    var geo=await apiFetch('/api/geo/ip');
    var gd=await geo.json();
    var lat=gd.lat||'',lon=gd.lon||'',city=gd.city||'';
    if(lat&&lon){_userGeo={lat:lat,lon:lon,city:city,ts:Date.now()};localStorage.setItem('fusion_geo',JSON.stringify(_userGeo));_setCookie('fusion_geo',JSON.stringify(_userGeo),1);}
    var qs=lat&&lon?'?lat='+lat+'&lon='+lon+'&city='+encodeURIComponent(city):'?city='+encodeURIComponent(city);
    var wr=await apiFetch('/api/weather'+qs);
    var wd=await wr.json();
    var lod=document.getElementById('wxLoading'); if(lod)lod.style.display='none';
    if(wd.ok){
      var now=new Date(); var hm=now.getHours().toString().padStart(2,'0')+':'+now.getMinutes().toString().padStart(2,'0');
      var desc=wd.desc||''; var icon=_wxIcon(desc);
      var el=function(id){return document.getElementById(id);};
      if(el('wxIcon'))el('wxIcon').textContent=icon;
      if(el('wxTemp'))el('wxTemp').textContent=wd.temp+'°C';
      if(el('wxFeels'))el('wxFeels').textContent='Feels '+(wd.feels_like||wd.temp)+'°';
      if(el('wxCity'))el('wxCity').textContent=wd.city||city||'Local';
      if(el('wxDesc'))el('wxDesc').textContent=desc;
      if(el('wxHumid'))el('wxHumid').textContent=(wd.humidity||'--')+'%';
      if(el('wxWind'))el('wxWind').textContent=(wd.wind_speed?Math.round(wd.wind_speed*3.6):wd.wind||'--')+' km/h';
      if(el('wxTime'))el('wxTime').textContent=hm;
      if(el('wxCard'))el('wxCard').style.display='';
      // Welcome screen mini badge
      var ww=document.getElementById('welcomeWx');
      if(ww){ww.innerHTML=icon+' '+wd.temp+'°C · '+(wd.city||city)+' · '+desc;ww.style.display='';}
    }
  }catch(e){var lod=document.getElementById('wxLoading');if(lod){lod.textContent='Weather unavailable';}}
}

// Langsearch removed
var _langsearchEnabled=false;
function _isNewsQuery(text){
  var tl=text.toLowerCase();
  var kws=['latest','news','today','current','2024','2025','2026','update','who won','what happened','recent','breaking','price of','stock','weather','score','result'];
  return kws.some(function(k){return tl.includes(k);});
}

// ══ Tool bar ══════════════════════════════════════════════════════════════════
var _activeTool='web';
var _explainLevel='eli5';
var _openTpanel=null;

function tbarToggle(tool){
  var wasOpen=(_openTpanel===tool);
  tbarClose();
  if(wasOpen) return;
  _openTpanel=tool;
  // Mark button active
  document.querySelectorAll('.tbar-btn').forEach(function(b){b.classList.remove('tbar-open');});
  var btn=document.getElementById('tbtn-'+tool); if(btn) btn.classList.add('tbar-open');
  // Show panel
  var pan=document.getElementById('tpanel-'+tool); if(pan){pan.style.display='block';}
  // Focus the first input in panel
  setTimeout(function(){
    var inp=document.querySelector('#tpanel-'+tool+' .tpanel-input');
    if(inp) inp.focus();
  },60);
}
function tbarClose(){
  if(_openTpanel){
    var pan=document.getElementById('tpanel-'+_openTpanel);
    if(pan) pan.style.display='none';
    var btn=document.getElementById('tbtn-'+_openTpanel);
    if(btn) btn.classList.remove('tbar-open');
    _openTpanel=null;
  }
}

// ── tpanel action handlers ────────────────────────────────────────────────────
function _tpVal(id){ var el=document.getElementById(id); return el?(el.value||'').trim():''; }
function _tpSend(tool){
  var q=_tpVal('tinput-'+tool); if(!q){return;}
  tbarClose();
  document.getElementById('msgIn').value=q;
  sendMsg();
}
function _tpSendPre(tool,prefix){
  var q=_tpVal('tinput-'+tool); if(!q) return;
  tbarClose();
  document.getElementById('msgIn').value=prefix+q;
  sendMsg();
}
function _tpDeep(){
  var q=_tpVal('tinput-deep'); if(!q) return;
  tbarClose();
  doDeepResearch(q);
}
function _tpMath(){
  var q=_tpVal('tinput-math'); if(!q) return;
  tbarClose();
  doSpecialResearch(q,'math');
}
function _tpCode(){
  var q=_tpVal('tinput-code'); if(!q) return;
  tbarClose();
  doSpecialResearch(q,'code');
}
function _tpTrans(){
  var q=_tpVal('tinput-trans'); if(!q) return;
  var from=(_tpVal('transFrom')||document.getElementById('transFrom').value)||'Auto-detect';
  var to=(_tpVal('transTo')||document.getElementById('transTo').value)||'English';
  tbarClose();
  document.getElementById('msgIn').value='Translate the following from '+from+' to '+to+'. Return only the translation:\n\n'+q;
  sendMsg();
}
function _tpExplain(){
  var q=_tpVal('tinput-explain'); if(!q) return;
  var lvlMap={eli5:'Explain this as if I am 5 years old, use simple words and analogies',normal:'Explain this clearly with good examples',expert:'Give a deep technical expert-level explanation with precise terminology',analogy:'Explain this using a creative real-world analogy or metaphor'};
  tbarClose();
  document.getElementById('msgIn').value=lvlMap[_explainLevel]+': '+q;
  sendMsg();
}
function _tpCompare(){
  var a=_tpVal('tinput-cmp1'), b=_tpVal('tinput-cmp2'), ctx=_tpVal('tinput-cmpctx');
  if(!a||!b){showToast('Enter both options to compare');return;}
  tbarClose();
  document.getElementById('msgIn').value='Compare '+a+' vs '+b+(ctx?' for: '+ctx:'')+'. Give a structured comparison with pros/cons and a recommendation.';
  sendMsg();
}
function _tpBrowser(){
  var q=_tpVal('tinput-browser'); if(!q) return;
  tbarClose();
  openSP('scraper');
  setTimeout(function(){fbSearch(q);},200);
}
function _setLevel(btn,lvl){
  _explainLevel=lvl;
  document.querySelectorAll('.tpanel-lvl').forEach(function(b){b.classList.remove('active');});
  btn.classList.add('active');
}

// ── Deep Research — live card grid, every available model ────────────────────
async function doDeepResearch(prompt){
  if(!prompt) prompt=document.getElementById('msgIn').value.trim();
  if(!prompt){showToast('Enter a question first');return;}
  document.getElementById('msgIn').value='';
  updateTokenCount('');
  hist.push({role:'user',content:prompt});
  addMsg('user',prompt);

  // Build the live result container in chat
  var bbl=addMsg('ai','',null,'Deep Research');
  var t0=Date.now();
  bbl.innerHTML=_drBuildShell(prompt,t0);
  _drAnimate(bbl,t0);

  try{
    var r=await apiFetch('/api/deep-research',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({prompt:prompt})});
    var d=await r.json();
    var elapsed=((Date.now()-t0)/1000).toFixed(1);
    if(!d.ok){bbl.innerHTML='<div style="color:var(--red);font-size:13px">Deep research failed: '+(d.error||'All models unavailable')+'</div>';return;}
    bbl.innerHTML=_drBuildResult(d,elapsed,prompt);
    hist.push({role:'assistant',content:d.synthesis||''});
    _renderMath(bbl);
    // Add copy button
    var copyBtn=bbl.querySelector('.dr-copy');
    if(copyBtn) copyBtn.onclick=function(){navigator.clipboard.writeText(d.synthesis||'');showToast('Copied synthesis');};
    // Expand card toggles
    bbl.querySelectorAll('.dr-card').forEach(function(card){
      card.onclick=function(){card.classList.toggle('expanded');};
    });
  }catch(e){
    var elapsed=((Date.now()-t0)/1000).toFixed(1);
    bbl.innerHTML='<div style="color:var(--red);font-size:13px">'+e.message+'</div>';
  }
}

function _drBuildShell(prompt,t0){
  return '<div class="dr-header"><span class="dr-title-ico">🔬</span><span class="dr-title-txt">Deep Research</span>'
    +'<span class="dr-live-badge">LIVE</span></div>'
    +'<div class="dr-prompt">'+escHtml(prompt.slice(0,120))+(prompt.length>120?'…':'')+'</div>'
    +'<div style="font-size:11px;color:var(--tx3);margin:8px 0 12px">Querying all available models simultaneously…</div>'
    +'<div class="dr-wrap" id="dr-cards-live">'
    +[1,2,3,4,5,6].map(function(i){return '<div class="dr-card pend"><div class="dr-card-hd"><div class="dr-card-dot"></div><div class="dr-card-name">MODEL '+i+'</div></div><div class="dr-card-txt" style="color:var(--tx3);font-size:11px">Waiting…</div></div>';}).join('')
    +'</div>';
}

function _drAnimate(bbl,t0){
  var timer=setInterval(function(){
    var el=bbl.querySelector('#dr-cards-live');
    if(!el||!bbl.parentNode){clearInterval(timer);return;}
    var elapsed=((Date.now()-t0)/1000).toFixed(1);
    var badge=bbl.querySelector('.dr-live-badge');
    if(badge) badge.textContent=elapsed+'s';
  },200);
}

function _drBuildResult(d,elapsed,prompt){
  var perspectives=d.perspectives||[];
  var cards=perspectives.map(function(p){
    var ms=p.ms?p.ms+'ms':'?';
    var speed=p.ms<1000?'fast':p.ms<3000?'normal':'slow';
    var speedCol=speed==='fast'?'#4ade80':speed==='normal'?'#60a5fa':'#f87171';
    return '<div class="dr-card done">'
      +'<div class="dr-card-hd">'
      +'<div class="dr-card-dot"></div>'
      +'<div class="dr-card-name">'+escHtml(p.model||p.prov)+'</div>'
      +'<div class="dr-card-ms" style="color:'+speedCol+'">'+ms+'</div>'
      +'</div>'
      +'<div class="dr-card-txt">'+escHtml(p.text||'')+'</div>'
      +'</div>';
  }).join('');

  var failed=Math.max(0,(d.models_queried||0)-(d.agents_used||0));
  return '<div class="dr-header">'
    +'<span class="dr-title-ico">🔬</span>'
    +'<span class="dr-title-txt">Deep Research</span>'
    +'<span class="dr-stat">'+d.agents_used+' AIs · '+elapsed+'s'+(failed?' · '+failed+' failed':'')+'</span>'
    +'<button class="dr-copy" style="margin-left:auto;background:rgba(255,255,255,.06);border:1px solid var(--glass-bdr);border-radius:7px;padding:3px 9px;font-size:10px;color:var(--tx3);cursor:pointer">📋 Copy</button>'
    +'</div>'
    +(d.web_ctx?'<div class="dr-web-ctx">🌐 Web-grounded · Bing context used</div>':'')
    +'<div class="dr-synth">'
    +'<div class="dr-synth-hd">✦ Synthesis</div>'
    +'<div class="dr-synth-txt">'+fmt(d.synthesis||'No synthesis available.')+'</div>'
    +'</div>'
    +'<details style="margin-top:12px"><summary style="font-size:11px;color:var(--tx3);cursor:pointer;user-select:none;list-style:none;display:flex;align-items:center;gap:6px"><span>▸</span> All '+perspectives.length+' model responses</summary>'
    +'<div class="dr-wrap" style="margin-top:10px">'+cards+'</div>'
    +'</details>';
}

// ── Special Research — math (all reasoning models) and code (all coding models)
async function doSpecialResearch(prompt,mode){
  hist.push({role:'user',content:(mode==='math'?'[Math] ':'[Code] ')+prompt});
  addMsg('user',(mode==='math'?'📐 ':'💻 ')+prompt);
  var bbl=addMsg('ai','',null,mode==='math'?'Math Solver':'Code Engine');
  var t0=Date.now();
  bbl.innerHTML='<div style="font-size:11px;color:var(--tx3)">Routing to '+(mode==='math'?'all reasoning models':'all coding models')+'…<div class="thinking-dots" style="display:inline-flex;margin-left:8px"><span></span><span></span><span></span></div></div>';
  try{
    var r=await apiFetch('/api/deep-research',{method:'POST',headers:{'Content-Type':'application/json'},
      body:JSON.stringify({prompt:prompt,mode:mode})});
    var d=await r.json();
    var elapsed=((Date.now()-t0)/1000).toFixed(1);
    if(!d.ok){bbl.innerHTML='<div style="color:var(--red)">Failed: '+(d.error||'No models responded')+'</div>';return;}
    bbl.innerHTML=_drBuildResult(d,elapsed,prompt);
    hist.push({role:'assistant',content:d.synthesis||''});
    _renderMath(bbl);
    bbl.querySelectorAll('.dr-card').forEach(function(c){c.onclick=function(){c.classList.toggle('expanded');};});
    var copyBtn=bbl.querySelector('.dr-copy');
    if(copyBtn) copyBtn.onclick=function(){navigator.clipboard.writeText(d.synthesis||'');showToast('Copied');};
  }catch(e){bbl.innerHTML='<div style="color:var(--red)">'+e.message+'</div>';}
}

// ── Fusion.Browser ────────────────────────────────────────────────────────────
var _fbHistory=[], _fbIdx=-1, _fbLastResults=[], _fbLastQuery='';
function fbNavigate(){
  var inp=document.getElementById('fbUrl'); if(!inp) return;
  var val=inp.value.trim(); if(!val) return;
  if(!val.includes('.')||val.includes(' ')) fbSearch(val);
  else fbLoadUrl(val.startsWith('http')?val:'https://'+val);
}
function fbSearch(q){var inp=document.getElementById('fbUrl');if(inp)inp.value=q;_fbQuery(q);}
function fbLoadUrl(url){var inp=document.getElementById('fbUrl');if(inp)inp.value=url;fbSearch('site:'+url.replace(/https?:\/\//,'').split('/')[0]);}
function fbBack(){if(_fbIdx>0){_fbIdx--;fbSearch(_fbHistory[_fbIdx]);}}
function fbForward(){if(_fbIdx<_fbHistory.length-1){_fbIdx++;fbSearch(_fbHistory[_fbIdx]);}}
function fbRefresh(){if(_fbLastQuery)fbSearch(_fbLastQuery);}
async function _fbQuery(q){
  _fbLastQuery=q;
  if(_fbHistory[_fbHistory.length-1]!==q){_fbHistory.push(q);_fbIdx=_fbHistory.length-1;}
  var welcome=document.getElementById('fbWelcome');
  var res=document.getElementById('fbResults');
  var loadBar=document.getElementById('fbLoadBar');
  var status=document.getElementById('fbStatus');
  var secure=document.getElementById('fbSecure');
  if(welcome)welcome.style.display='none';
  if(res){res.style.display='';res.innerHTML='<div style="padding:28px;text-align:center"><div class="thinking-dots" style="justify-content:center"><span></span><span></span><span></span></div><div style="font-size:11px;color:var(--tx3);margin-top:10px">Searching Bing…</div></div>';}
  if(loadBar)loadBar.style.width='40%';
  if(status){status.style.display='';status.textContent='Connecting…';}
  try{
    var r=await apiFetch('/api/search/web',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({query:q})});
    var d=await r.json();
    if(loadBar){loadBar.style.width='100%';setTimeout(function(){loadBar.style.width='0%';},350);}
    if(status){status.textContent=(d.results&&d.results.length?d.results.length+' results':'No results');setTimeout(function(){if(status)status.style.display='none';},2000);}
    if(secure)secure.textContent=d.ok?'🔒':'⚠️';
    _fbLastResults=d.results||[];
    var html='';
    if(d.instant){
      html+='<div class="fb-instant"><div style="font-size:9px;font-weight:700;color:var(--blue);text-transform:uppercase;letter-spacing:1px;margin-bottom:6px">⚡ Instant Answer</div>'
        +'<div style="font-size:13px;color:var(--tx);line-height:1.72">'+escHtml(d.instant)+'</div></div>';
    }
    if(_fbLastResults.length){
      html+='<div style="font-size:10px;font-weight:600;color:var(--tx3);text-transform:uppercase;letter-spacing:.8px;margin-bottom:9px">'+_fbLastResults.length+' results for "'+escHtml(q)+'"</div>';
      _fbLastResults.forEach(function(r2){
        html+='<div class="fb-result" onclick="window.open(\''+r2.url.replace(/'/g,"\\'")+'\')">'
          +'<div class="fb-result-title">'+escHtml(r2.title||'')+'</div>'
          +'<div class="fb-result-url">'+escHtml((r2.url||'').replace(/https?:\/\//,'').slice(0,70))+'</div>'
          +'<div class="fb-result-snip">'+escHtml(r2.snippet||'')+'</div></div>';
      });
      html+='<button onclick="fbSendToAI()" style="width:100%;margin-top:10px;background:var(--grad);border:none;border-radius:10px;padding:10px;color:#fff;font-size:13px;font-weight:600;cursor:pointer">💬 Ask AI about these results</button>';
    }else if(!d.instant){
      html='<div class="empty-state"><span class="ei">🔍</span>No results found.</div>';
    }
    if(res)res.innerHTML=html;
  }catch(e){
    if(loadBar)loadBar.style.width='0%';
    if(res)res.innerHTML='<div class="empty-state"><span class="ei">⚠️</span>'+escHtml(e.message)+'</div>';
  }
}
async function fbSendToAI(){
  if(!_fbLastResults.length&&!_fbLastQuery){showToast('Search first');return;}
  closeSP();
  var ctx='Search results for: "'+_fbLastQuery+'"\n\n';
  _fbLastResults.forEach(function(r,i){ctx+=(i+1)+'. '+r.title+'\n   '+r.url+'\n   '+r.snippet+'\n\n';});
  document.getElementById('msgIn').value='Based on these results, answer: '+_fbLastQuery;
  sendMsg();
}

// ── Auto web search — always on ──────────────────────────────────────────────
async function _autoWebSearch(text){
  if(!text||text.startsWith('/')||text.length<4) return null;
  try{
    var r=await apiFetch('/api/search/ddg',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({query:text.slice(0,200)})});
    var d=await r.json();
    if(!d.ok) return null;
    var ctx='';
    if(d.instant) ctx+='[Web: '+d.instant+']\n';
    (d.results||[]).slice(0,5).forEach(function(res){if(res.snippet) ctx+='['+res.title+': '+res.snippet+']\n';});
    return ctx.trim()||null;
  }catch(e){return null;}
}

// ── Token counter ────────────────────────────────────────────────────────────
var _TOK_MAX=4096;
function updateTokenCount(val){
  var n=Math.ceil((val||'').length/4);
  var el=document.getElementById('tokNum'),fill=document.getElementById('tokFill'),ctr=document.getElementById('tokCtr');
  if(!el) return;
  el.textContent=n>999?(n/1000).toFixed(1)+'k':String(n);
  var pct=Math.min(100,Math.round(n/_TOK_MAX*100));
  if(fill) fill.style.width=pct+'%';
  if(ctr) ctr.classList.toggle('warn',pct>=75);
}

// ── Floating particle system ──────────────────────────────────────────────────
(function(){
  var c=document.getElementById('particleCanvas');
  if(!c)return;
  var ctx=c.getContext('2d');
  var W=c.width=window.innerWidth,H=c.height=window.innerHeight;
  var particles=[];
  for(var i=0;i<55;i++){
    particles.push({x:Math.random()*W,y:Math.random()*H,r:Math.random()*1.4+0.3,
      vx:(Math.random()-0.5)*0.28,vy:(Math.random()-0.5)*0.28,
      opacity:Math.random()*0.5+0.1,color:Math.random()>0.5?'rgba(6,182,212,':'rgba(124,58,237,'});
  }
  function draw(){
    ctx.clearRect(0,0,W,H);
    // Draw connections
    for(var i=0;i<particles.length;i++){
      for(var j=i+1;j<particles.length;j++){
        var dx=particles[i].x-particles[j].x,dy=particles[i].y-particles[j].y;
        var dist=Math.sqrt(dx*dx+dy*dy);
        if(dist<130){
          ctx.beginPath();ctx.moveTo(particles[i].x,particles[i].y);ctx.lineTo(particles[j].x,particles[j].y);
          ctx.strokeStyle='rgba(80,120,200,'+(0.06*(1-dist/130))+')';ctx.lineWidth=0.5;ctx.stroke();
        }
      }
    }
    particles.forEach(function(p){
      p.x+=p.vx;p.y+=p.vy;
      if(p.x<0)p.x=W;if(p.x>W)p.x=0;if(p.y<0)p.y=H;if(p.y>H)p.y=0;
      ctx.beginPath();ctx.arc(p.x,p.y,p.r,0,Math.PI*2);
      ctx.fillStyle=p.color+p.opacity+')';ctx.fill();
    });
    requestAnimationFrame(draw);
  }
  draw();
  window.addEventListener('resize',function(){W=c.width=window.innerWidth;H=c.height=window.innerHeight;});
})();

// ── Bing AI Search Browser (via OpenSERP) ────────────────────────────────────
var _bingResults=[], _bingQuery='';
async function doBingSearch(){
  var q=(document.getElementById('ddgQuery')||{}).value||''; q=q.trim();
  if(!q){showToast('Enter a search query');return;}
  var st=document.getElementById('ddgStatus');var res=document.getElementById('ddgResults');
  var instant=document.getElementById('ddgInstant');var sendBtn=document.getElementById('ddgSendBtn');
  st.style.display='';st.textContent='🔍 Searching Bing…';
  res.innerHTML='';instant.style.display='none';sendBtn.style.display='none';
  try{
    var r=await apiFetch('/api/search/web',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({query:q})});
    var d=await r.json();
    st.textContent='';st.style.display='none';
    _bingQuery=q; _bingResults=d.results||[];
    if(d.instant){
      instant.style.display='';
      instant.innerHTML='<div style="font-size:10px;font-weight:700;color:var(--blue);text-transform:uppercase;letter-spacing:1px;margin-bottom:6px">⚡ Top Result</div>'
        +'<div style="font-size:13px;color:var(--tx);line-height:1.7">'+escHtml(d.instant)+'</div>';
    }
    if(_bingResults.length){
      _bingResults.forEach(function(r2){
        var card=document.createElement('div');card.className='ddg-result';
        card.innerHTML='<div class="ddg-result-title">'+escHtml(r2.title||'')+'</div>'
          +'<div class="ddg-result-url">'+escHtml((r2.url||'').replace(/https?:\/\//,'').substring(0,60))+'</div>'
          +'<div class="ddg-result-snip">'+escHtml(r2.snippet||'')+'</div>';
        card.onclick=function(){window.open(r2.url,'_blank');};
        res.appendChild(card);
      });
      sendBtn.style.display='';
    } else if(!d.instant){
      res.innerHTML='<div style="font-size:12px;color:var(--tx3);padding:16px;text-align:center">No results found. Try a different query.</div>';
    }
  }catch(e){st.textContent='Search failed: '+e.message;}
}
// keep old function name working
function doDDGSearch(){return doBingSearch();}

async function sendDDGToChat(){
  if(!_bingQuery&&!_bingResults.length){showToast('Search first');return;}
  closeSP();
  var ctx='Bing search results for: "'+_bingQuery+'"\n\n';
  _bingResults.forEach(function(r,i){ctx+=(i+1)+'. '+r.title+'\n   '+r.url+'\n   '+r.snippet+'\n\n';});
  var prompt='Based on these Bing search results, give me a clear, well-structured answer about "'+_bingQuery+'".\n\nSearch Results:\n'+ctx+'\nSummarise the key information in a helpful way. Use your knowledge to add context where useful.';
  document.getElementById('msgIn').value=prompt;
  sendMsg();
}

// ── Auto web search — always on ──────────────────────────────────────────────
async function _autoWebSearch(text){
  if(!text||text.startsWith('/')||text.length<4) return null;
  try{
    var r=await apiFetch('/api/search/web',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({query:text.slice(0,200)})});
    var d=await r.json();
    if(!d.ok) return null;
    var ctx='';
    if(d.instant) ctx+='[Web: '+d.instant+']\n';
    (d.results||[]).slice(0,5).forEach(function(res){if(res.snippet) ctx+='['+res.title+': '+res.snippet+']\n';});
    return ctx.trim()||null;
  }catch(e){return null;}
}


// ── Overseer uses cheapest Groq (llama-3.1-8b-instant) — handled server-side
// ── Image model auto-use: when an image model is selected, route image requests to it
function _getSelectedModelType(){
  var mkey=(document.getElementById('mOverride')||{value:'auto'}).value;
  if(mkey==='auto') return 'chat';
  var m=allModels.find(function(x){return x.key===mkey;});
  return m?m.type:'chat';
}

// ── setDevUI update to also show LSB dev btn ─────────────────────────────
var _origSetDevUI=typeof setDevUI==='function'?setDevUI:null;

function toggleTheme()   { 
  var newMode=isDark?'light':'dark';
  setThemeMode(newMode); 
  var mtc=document.getElementById('metaThemeColor');
  if(mtc) mtc.content=(newMode==='light'?'#ffffff':'#000000');
  closeDrop(); 
}

function updateThemeUI() {
  // Mode buttons
  ['dark','light'].forEach(function(m){
    var b=document.getElementById('tmb-'+m); if(!b) return;
    b.classList.toggle('active', (themePref.mode||'dark')===m);
  });
  // Accent swatches
  document.querySelectorAll('.accent-swatch').forEach(function(s){
    s.classList.toggle('active', s.dataset.accent===(themePref.accent||'red'));
  });
  // BG swatches
  document.querySelectorAll('.bg-swatch').forEach(function(s){
    s.classList.toggle('active', s.dataset.bg===(themePref.bg||'orbs'));
  });
  // Color pickers
  var c1=document.getElementById('customBgColor');
  var c2=document.getElementById('customBgColor2');
  if(c1&&themePref.customBg) c1.value=themePref.customBg;if(c2&&themePref.customBg2) c2.value=themePref.customBg2;
}

// Use system theme if no saved preference
// Default is 'system' — applyTheme resolves the actual dark/light
applyTheme(themePref);
// _applyMenuToggles called in launch()
// Track system theme changes — always follow OS when mode==='system'
if(window.matchMedia){
  window.matchMedia('(prefers-color-scheme:dark)').addEventListener('change',function(e){
    if(themePref.mode==='system') applyTheme(themePref);
  });
}

function switchTab(t) {
  ['login','register','forgot'].forEach(function(x){
    document.getElementById(x+'Form').style.display=x===t?'':'none';
    var id=x==='login'?'tabLogin':x==='register'?'tabReg':'tabFP';
    document.getElementById(id).classList.toggle('active',x===t);
  });
  document.getElementById('authErr').style.display='none';
  if(t==='register'){setTimeout(function(){try{if(window.hcaptcha)hcaptcha.reset();}catch(e){}},200);}
}
function showAuthErr(m) { var e=document.getElementById('authErr'); e.textContent=m; e.style.display='block'; }


// ── Math Puzzle (bot protection, no API key) ──────────────────────────────
// ── hCaptcha helpers ─────────────────────────────────────────────────────────

function onHcaptchaSuccess(token){ var e=document.getElementById('captchaErr'); if(e) e.style.display='none'; }
function onHcaptchaError(){ var e=document.getElementById('captchaErr'); if(e){e.style.display='';e.textContent='⚠ Captcha error — please try again.';} }
function _getHcaptchaToken(){try{return window.hcaptcha?hcaptcha.getResponse():'';}catch(e){return '';}}
function _resetHcaptcha(){try{if(window.hcaptcha)hcaptcha.reset();}catch(e){}}

async function doLogin() {
  var u=document.getElementById('lUser').value.trim(), p=document.getElementById('lPass').value;
  if(!u||!p) return showAuthErr('Enter username and password');
  var btn=document.getElementById('loginBtn'); btn.textContent='Signing in\u2026';btn.disabled=true;
  try{
    var r=await fetch('/api/login',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({username:u,password:p})});
    var d=await r.json();
    if(!r.ok){showAuthErr(d.error);btn.textContent='Sign In \u2192';btn.disabled=false;return;}
    authToken=d.token;_setCookie('fusion_token',d.token,365);localStorage.setItem('fusion_token',d.token);uName=d.username; currentMKey=d.model_key||'auto'; isDevUser=d.is_dev||false;
    if(d.theme) { themePref.mode=d.theme; saveTheme(); }
    launch();
  }catch(e){showAuthErr('Connection error');btn.textContent='Sign In \u2192';btn.disabled=false;}
}

async function doRegister() {
  var captchaToken=_getHcaptchaToken();
  var captchaEl=document.getElementById('captchaErr');
  if(captchaEl) captchaEl.style.display='none';
  var u=document.getElementById('rUser').value.trim(),p=document.getElementById('rPass').value,p2=document.getElementById('rPass2').value;
  if(!u||!p) return showAuthErr('Fill all fields');if(p!==p2) return showAuthErr('Passwords do not match');
  if(p.length<6) return showAuthErr('Password must be 6+ characters');
  var btn=document.getElementById('regBtn'); btn.textContent='Creating…';btn.disabled=true;
  try{
    var r=await fetch('/api/register',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({username:u,password:p,'h-captcha-response':captchaToken})});
    var d=await r.json();
    if(!r.ok){
      if(captchaEl) captchaEl.style.display='';
      _resetHcaptcha();
      showAuthErr(d.error);btn.textContent='Create Account →';btn.disabled=false;return;
    }
    document.getElementById('lUser').value=u; document.getElementById('lPass').value=p;switchTab('login'); await doLogin();
  }catch(e){showAuthErr('Connection error');btn.textContent='Create Account →';btn.disabled=false;}
}

async function doForgotPw() {
  var u=document.getElementById('fpUser').value.trim(),p=document.getElementById('fpPass').value,p2=document.getElementById('fpPass2').value;
  if(!u||!p) return showAuthErr('Fill all fields');if(p!==p2) return showAuthErr('Passwords do not match');
  if(p.length<6) return showAuthErr('Password must be 6+ characters');
  var btn=document.getElementById('fpBtn'); btn.textContent='Resetting\u2026';btn.disabled=true;
  try{
    var r=await fetch('/api/forgot-password',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({username:u,new_password:p})});
    var d=await r.json();
    if(!r.ok){showAuthErr(d.error);btn.textContent='Reset Password \u2192';btn.disabled=false;return;}
    showToast('\u2705 Password reset! Please sign in.'); document.getElementById('lUser').value=u; switchTab('login');
  }catch(e){showAuthErr('Connection error');}
  btn.textContent='Reset Password \u2192';btn.disabled=false;
}

async function guestLogin() {
  var btn=document.getElementById('guestBtn');btn.textContent='Loading\u2026';btn.disabled=true;
  try{
    var r=await fetch('/api/guest',{method:'POST'}); var d=await r.json();
    if(!r.ok){showAuthErr(d.error||'Guest login failed');btn.textContent='Continue as Guest';btn.disabled=false;return;}
    authToken=d.token; localStorage.setItem('fusion_token',d.token); uName='Guest'; currentMKey='auto'; isDevUser=false; launch();
  }catch(e){showAuthErr('Connection error');btn.textContent='Continue as Guest';btn.disabled=false;}
}

async function doLogout() {
  await apiFetch('/api/logout',{method:'POST'});authToken=''; localStorage.removeItem('fusion_token');_delCookie('fusion_token'); uName=''; hist=[]; isDevUser=false;
  document.getElementById('chatPage').classList.remove('active');document.getElementById('authPage').classList.add('active');
  resetChat(); closeDrop(); closeSP();document.getElementById('lPass').value='';
  document.getElementById('guestBtn').textContent='Continue as Guest'; document.getElementById('guestBtn').disabled=false;setDevUI(false);
}

function setDevUI(isdev) {
  var btn=document.getElementById('devHeaderBtn'),item=document.getElementById('devDropItem'),sep=document.getElementById('devDropSep');
  if(btn){btn.style.display=isdev?'':'none'; if(isdev)btn.classList.add('dev-btn');}
  if(item) item.style.display=isdev?'':'none';if(sep) sep.style.display=isdev?'':'none';
  var lsbDev=document.getElementById('lsbDevBtn');if(lsbDev)lsbDev.style.display=isdev?'':'none';
}

function launch() {
  document.getElementById('authPage').classList.remove('active');document.getElementById('chatPage').classList.add('active');
  document.getElementById('uName').textContent=uName;
  if(userAvatar){var av=document.getElementById('uAv');av.style.backgroundImage='url('+userAvatar+')';av.style.backgroundSize='cover';av.style.borderRadius='6px';av.textContent='';}
  else{document.getElementById('uAv').textContent=(uName[0]||'?').toUpperCase();}
  resetChat(); loadKeys(); loadModels(); loadSaved(); loadMemory(); setDevUI(isDevUser); loadConvs();
  _loadWeather();
  _applyMenuToggles();
  setTimeout(_loadCustomEndpointUI, 600);
  if(window.speechSynthesis){loadVoices();loadVCVoices();window.speechSynthesis.onvoiceschanged=function(){loadVoices();loadVCVoices();};}
  if(isDevUser){showToast('\u{1F6E0}\uFE0F Dev mode active');fetchDevStats();}
}

async function savePref() {
  if(!authToken) return;
  var sel=document.getElementById('mOverride');
  var mkey=sel?sel.value:'auto'; currentMKey=mkey;
  updateModelModeUI(mkey);
  try{await apiFetch('/api/prefs',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({theme:isDark?'dark':'light',model_key:mkey})});}catch(e){}
}

// Keys are server-side only — no client key management
async function loadKeys(){
  try{
    var r=await apiFetch('/api/keys'); var d=await r.json();
    var provs=d.providers||[];
    ['groq','openrouter','huggingface','github','stability','cloudflare','extra','custom_endpoint'].forEach(function(p){setStat(p,provs.includes(p));});
    // provider stats
  }catch(e){}
}

function setStat(p,set){var el=document.getElementById(p+'Stat');if(!el)return;el.textContent=set?'Set \u2713':'Not set';el.className='kstat '+(set?'set':'unset');}

async function saveKey(p) {
  var inp=document.getElementById(p+'In'); var k=inp.value.trim(); if(!k) return showToast('Enter a key first');
  try{
    var r=await apiFetch('/api/keys',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({provider:p,key:k})});
    var d=await r.json();
    if(!r.ok) return showToast('Error: '+d.error);setStat(p,true); inp.value=''; showToast('\u2705 '+p+' key saved');
    loadKeys(); // refresh banner
  }catch(e){showToast('Error saving key');}
}
async function delKey(p) {
  if(!confirm('Delete '+p+' key?')) return;
  await apiFetch('/api/keys',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({provider:p,key:''})});
  setStat(p,false); showToast('Key removed');
  loadKeys(); // refresh banner
}

async function loadModels() {
  try{
    var r=await apiFetch('/api/models'); var d=await r.json(); allModels=(d.models||[]).filter(function(m){return !m.disabled;});
    var sel=document.getElementById('mOverride');
    var byGroup={};
    allModels.forEach(function(m){
      var grp;
      grp=m.company;if(!byGroup[grp]) byGroup[grp]=[];byGroup[grp].push(m);
    });
    sel.innerHTML='<option value="auto">\u26A1 Auto \u2014 Overseer picks best model</option>';
    // Order: chat companies first, then Image, then Video, then 3D last
    var companies=Object.keys(byGroup).sort(function(a,b){
      if(a.includes('3D'))    return 1; if(b.includes('3D'))    return -1;if(a.includes('Video')) return 1; if(b.includes('Video')) return -1;
      if(a.includes('Image')) return 1; if(b.includes('Image')) return -1;
      return a.localeCompare(b);
    });
    companies.forEach(function(co){
      var ms=byGroup[co];
      var og=document.createElement('optgroup'); og.label=ms[0].emoji+' '+co;
      ms.forEach(function(m){var o=document.createElement('option');o.value=m.key;var sc=m.score?(' ['+Math.round(m.score*100)+']'):'';var tg=(m.thinking?' \uD83E\uDDE0':'')+((m.vision||m.provider==='image_or')?' \uD83D\uDDBC\uFE0F':'');o.textContent=m.label+tg+sc+' \u2014 '+m.desc;if(m.key===currentMKey)o.selected=true;og.appendChild(o);});
      sel.appendChild(og);
    });
    if(currentMKey!=='auto') sel.value=currentMKey;updateModelModeUI(currentMKey);
  }catch(e){}
}

function updateModelModeUI(mkey) {
  var note=document.getElementById('imgModeNote');
  var ta=document.getElementById('msgIn');
  var found=allModels.find(function(m){return m.key===mkey;});
  var provider=found?found.provider:'';
  if(note){
    if(provider==='image'||provider==='image_or') note.style.display='';else note.style.display='none';
  }
  if(ta){
    if(provider==='image')      ta.placeholder='Describe the image you want... Overseer will enhance your prompt';
    else if(provider==='video') ta.placeholder='Describe your video scene... Overseer will ask questions to enhance it';
    else if(provider==='video3d') ta.placeholder='Describe the 3D object, or just send to generate from your description...';
    else ta.placeholder='Ask anything... Select an \uD83D\uDDBC\uFE0F Image, \uD83C\uDFAC Video, or \uD83E\uDDCA 3D model for media generation';
  }
  // Update header pill label
  if(found){
    var parts=found.label.split('\u00B7');
    document.getElementById('mlabel').textContent=parts.length>1?parts[1].trim().split(' ').slice(0,3).join(' '):found.label.split(' ')[0];
  } else if(mkey==='auto'){
    document.getElementById('mlabel').textContent='Auto';
  }
}

var _genAbortCtrl=null,_genRunning=false,_lastGenResult=null;

function stopGeneration(btnId,stopId){
  if(_genAbortCtrl){try{_genAbortCtrl.abort();}catch(e){}_genAbortCtrl=null;}
  _genRunning=false;
  var btn=document.getElementById(btnId),stop=document.getElementById(stopId);
  if(btn){btn.disabled=false;btn.textContent=btn.dataset.orig||btn.textContent;}
  if(stop)stop.style.display='none';
}
function _genBtnStart(btnId,stopId,label){
  _genAbortCtrl=new AbortController();_genRunning=true;
  var btn=document.getElementById(btnId),stop=document.getElementById(stopId);
  if(btn){btn.dataset.orig=btn.textContent;btn.disabled=true;btn.textContent='\u23F3 '+label;}
  if(stop)stop.style.display='';
}
function _genBtnDone(btnId,stopId){
  _genRunning=false;_genAbortCtrl=null;
  var btn=document.getElementById(btnId),stop=document.getElementById(stopId);
  if(btn){btn.disabled=false;btn.textContent=btn.dataset.orig||btn.textContent;}
  if(stop)stop.style.display='none';
}
function _downloadB64(b64,mime,filename){
  var a=document.createElement('a');a.href='data:'+mime+';base64,'+b64;a.download=filename;a.click();
}
function downloadAudio(audioId,filename){
  var au=document.getElementById(audioId);
  if(!au||!au.src)return showToast('No audio to download');
  var a=document.createElement('a');a.href=au.src;a.download=filename;a.click();
}
function sendGenToChat(){
  if(!_lastGenResult)return showToast('Generate something first');
  var res=_lastGenResult;
  if(res.b64){
    pendingImageB64=res.b64;pendingImageMime=res.mime||'image/png';pendingImageName='generated.png';
    document.getElementById('imgPrev').src='data:'+res.mime+';base64,'+res.b64;document.getElementById('imgPrevName').textContent='Generated image';
    document.getElementById('imgPrevWrap').classList.add('show');closeSP();showToast('\uD83D\uDDBC\uFE0F Attached to chat');
  }else if(res.url){
    document.getElementById('msgIn').value='[Image: '+res.url+'] '+res.prompt;ar(document.getElementById('msgIn'));closeSP();
  }
}
function downloadImgFromChat(btn){
  var acts=btn.closest('.msg-actions');
  var b64=acts&&acts.dataset.b64?acts.dataset.b64:(_lastGenResult?_lastGenResult.b64:'');
  var mime=acts&&acts.dataset.mime?acts.dataset.mime:'image/png';
  if(b64)_downloadB64(b64,mime,'fusion-image.png');else showToast('No image data');
}
function attachImgFromChat(btn){
  var acts=btn.closest('.msg-actions');
  var b64=acts&&acts.dataset.b64?acts.dataset.b64:(_lastGenResult?_lastGenResult.b64:'');
  var mime=acts&&acts.dataset.mime?acts.dataset.mime:'image/png';
  if(!b64)return showToast('No image data');pendingImageB64=b64;pendingImageMime=mime;pendingImageName='generated.png';
  document.getElementById('imgPrev').src='data:'+mime+';base64,'+b64;document.getElementById('imgPrevName').textContent='Generated image';
  document.getElementById('imgPrevWrap').classList.add('show');showToast('\uD83D\uDDBC\uFE0F Attached \u2014 add message and send!');
}

var _panelImgAbort = null;
var _panelImgB64 = '', _panelImgMime = 'image/jpeg';

async function generateImage(){
  var prompt=document.getElementById('igPrompt').value.trim();
  if(!prompt)return showToast('Enter an image prompt');
  var modelId=(document.getElementById('igModel')||{value:'flux'}).value;
  var sz=((document.getElementById('igSize')||{value:'1024x1024'}).value).split('x');
  var w=parseInt(sz[0])||1024,h=parseInt(sz[1])||1024;
  _genBtnStart('igGenBtn','imgStopBtn','Generating…');
  var rw=document.getElementById('igResult'),img=document.getElementById('igImg');
  var tag=document.getElementById('igBackendTag');
  if(rw){rw.style.display='';rw.classList.remove('show');}
  try{
    // Use the AbortController from _genBtnStart
    var sig=_genAbortCtrl?_genAbortCtrl.signal:undefined;
    var isCfImg=modelId.startsWith('@cf/')||modelId.startsWith('@hf/');
    var isGhImg=modelId==='openai/gpt-image-1'||modelId==='openai/dall-e-3'||modelId==='gh_img_gpt'||modelId==='gh_img_dalle3';
    var isWorkerImg=modelId.startsWith('worker:');
    var imgApiPath=isWorkerImg?'/api/generate/image/worker':isCfImg?'/api/generate/image/cloudflare':isGhImg?'/api/generate/image/github':'/api/generate/image';
    var workerModelKey=isWorkerImg?modelId.slice(7):'sdxl';
    var imgBody=isWorkerImg?{prompt:prompt,model_key:workerModelKey,want_enrich:true}:isCfImg?{prompt:prompt,model:modelId}:isGhImg?{prompt:prompt,model:(modelId==='gh_img_gpt'?'openai/gpt-image-1':modelId==='gh_img_dalle3'?'openai/dall-e-3':modelId),size:w+'x'+h}:{prompt:prompt,model:modelId,width:w,height:h};
    var r=await apiFetch(imgApiPath,{method:'POST',
      signal:sig,
      headers:{'Content-Type':'application/json'},
      body:JSON.stringify(imgBody)});
    var rawText=await r.text();
    var d={};
    try{d=JSON.parse(rawText);}catch(_){d={error:'Bad response: '+rawText.slice(0,120)};}
    if(!r.ok||d.ok===false){
      var msg=d.error||('HTTP '+r.status);
      var devMsg=d.dev_error?('\n\nDetail: '+d.dev_error):'';
      showToast('❌ '+msg,8000);if(isDevUser&&d.dev_error) showErrOverlay('Image gen failed: '+d.dev_error);_genBtnDone('igGenBtn','imgStopBtn');
      return;
    }
    var src=d.b64?('data:'+(d.mime||'image/jpeg')+';base64,'+d.b64):(d.url||'');
    if(!src){showToast('❌ No image in response',5000);_genBtnDone('igGenBtn','imgStopBtn');return;}
    _panelImgB64=d.b64||''; _panelImgMime=d.mime||'image/jpeg'; _panelImgUrl=d.url||'';
    if(img){
      img.style.cssText='width:100%;border-radius:10px;border:1px solid var(--glass-bdr2);display:block;cursor:zoom-in;margin-bottom:8px';img.src=src;
      img.onerror=function(){
        if(!d.b64&&d.url){
          img.style.display='none';
          var lnk=document.getElementById('igImgFallback');
          if(!lnk){lnk=document.createElement('a');lnk.id='igImgFallback';lnk.target='_blank';lnk.style.cssText='display:block;padding:10px;background:rgba(255,255,255,.06);border:1px solid var(--glass-bdr);border-radius:10px;color:var(--blue);font-size:12px;text-align:center;margin-bottom:8px';rw.insertBefore(lnk,rw.firstChild);}
          lnk.href=d.url; lnk.textContent='🖼 Open image in new tab →';
        }
      };
    }
    if(tag)tag.textContent='⚡ '+(d.backend||'Pollinations')+(d.b64?' · '+Math.round(d.b64.length*0.75/1024)+'KB':'');if(rw)rw.classList.add('show');
    showToast('✅ Image ready!');
  }catch(e){
    if(e.name!=='AbortError')showToast('❌ '+e.message,6000);else showToast('⏹ Stopped');
  }
  _genBtnDone('igGenBtn','imgStopBtn');
}

var _panelImgUrl='';

function downloadGenImg(){
  if(_panelImgB64){_downloadB64(_panelImgB64,_panelImgMime,'fusion-image.png');}
  else if(_panelImgUrl){var a=document.createElement('a');a.href=_panelImgUrl;a.download='fusion-image.png';a.target='_blank';a.click();}
  else{var img=document.getElementById('igImg');if(img&&img.src&&img.src!==window.location.href){var a=document.createElement('a');a.href=img.src;a.download='fusion-image.png';a.target='_blank';a.click();}}
}

function sendGenImgToChat(){
  var img=document.getElementById('igImg');
  var src=_panelImgB64?('data:'+_panelImgMime+';base64,'+_panelImgB64):(_panelImgUrl||(img&&img.src)||'');
  if(!src||src===window.location.href){showToast('Generate an image first');return;}
  if(_panelImgB64){
    pendingImageB64=_panelImgB64; pendingImageMime=_panelImgMime; pendingImageName='generated.png';
    document.getElementById('imgPrev').src='data:'+_panelImgMime+';base64,'+_panelImgB64;document.getElementById('imgPrevName').textContent='Generated image';
    document.getElementById('imgPrevWrap').classList.add('show');closeSP(); showToast('📎 Attached! Add a message and send.');
  } else {
    // URL-only: embed directly in chat
    var bbl=addMsg('ai','',null,null);
    bbl.innerHTML='<img src="'+escHtml(src)+'" style="max-width:100%;border-radius:12px;border:1px solid var(--glass-bdr2);display:block;cursor:zoom-in" onclick="openImgFull(this.src)"/>'
      +'<br><span style="font-size:11px;color:var(--tx3)">Generated image — <a href="'+escHtml(src)+'" target="_blank" style="color:var(--blue)">open full size</a></span>';
    hist.push({role:'assistant',content:'[Image generated]'});closeSP(); scrollDown(); showToast('💬 Sent to chat!');
  }
}

var _videoAbort = null;
var _videoB64 = '', _videoMime = 'video/mp4', _videoUrl = '';

async function generateVideo(){
  var prompt=document.getElementById('videoPrompt').value.trim();
  if(!prompt)return showToast('Enter a video prompt');
  var model=(document.getElementById('videoModel')||{value:'animatediff'}).value;
  var sz=((document.getElementById('videoSize')||{value:'512x512'}).value).split('x');
  var w=parseInt(sz[0])||512,h=parseInt(sz[1])||512;
  _genBtnStart('videoGenBtn','videoStopBtn','Generating… (~60s)');
  var res=document.getElementById('videoResult');
  if(res)res.style.display='none';_videoAbort=new AbortController();
  try{
    var r=await apiFetch('/api/generate/video',{method:'POST',
      signal:_videoAbort.signal,
      headers:{'Content-Type':'application/json'},
      body:JSON.stringify({prompt:prompt,model:model,width:w,height:h})});
    if(!r.ok){var et=await r.text();var er={};try{er=JSON.parse(et);}catch(_){}showToast('❌ '+(er.error||'Video gen failed'),6000);_genBtnDone('videoGenBtn','videoStopBtn');return;}
    var d=await r.json();
    var vid=document.getElementById('videoEl');
    var tag=document.getElementById('videoBackendTag');
    if(d.b64){
      _videoB64=d.b64; _videoMime=d.mime||'video/mp4';
      var src='data:'+_videoMime+';base64,'+d.b64;
      if(vid){vid.src=src;vid.load();}
    } else if(d.url){
      _videoUrl=d.url; _videoB64=''; _videoMime='video/mp4';
      if(vid){vid.src=d.url;vid.load();}
    } else {showToast('❌ No video returned',5000);_genBtnDone('videoGenBtn','videoStopBtn');return;}
    if(tag)tag.textContent='⚡ '+escHtml(d.backend||'Pollinations');if(res)res.style.display='';showToast('✅ Video ready!');
  }catch(e){
    if(e.name!=='AbortError')showToast('❌ '+e.message,6000);else showToast('⏹ Stopped');
  }
  _videoAbort=null;_genBtnDone('videoGenBtn','videoStopBtn');
}

function stopVideoGen(){if(_videoAbort)_videoAbort.abort();}

async function _doGen3DFromChat(prompt){
  loading=true; document.getElementById('sbtn').disabled=true;
  var bbl=addMsg('ai','',null,null);
  bbl.innerHTML='<div style="display:flex;align-items:center;gap:10px"><div class="thinking-dots"><span></span><span></span><span></span></div><span style="color:var(--tx2);font-size:13px">🧊 Generating 3D render…</span></div>';
  scrollDown();
  var abort=new AbortController();
  try{
    var r=await apiFetch('/api/generate/3d',{method:'POST',signal:abort.signal,
      headers:{'Content-Type':'application/json'},
      body:JSON.stringify({prompt:prompt,style:'3d render',width:1024,height:1024})});
    if(!r.ok){var et=await r.text();var er={};try{er=JSON.parse(et);}catch(_){}
      bbl.innerHTML='<span style="color:var(--red)">❌ '+(er.error||'3D gen failed')+'</span>';
    }else{
      var d=await r.json();
      var src2=d.b64?'data:'+d.mime+';base64,'+d.b64:(d.url||'');
      if(src2){
        bbl.innerHTML='<img src="'+src2+'" style="max-width:100%;border-radius:12px;border:1px solid var(--glass-bdr2);display:block;cursor:zoom-in;margin-bottom:6px" onclick="openImgFull(this.src)"/>'
          +'<span style="font-size:10px;color:var(--tx3)">🧊 3D render · '+(d.backend||'Pollinations')+'</span>';
        hist.push({role:'assistant',content:'[3D render: '+prompt+']'});showToast('✅ 3D ready!');
      }else bbl.innerHTML='<span style="color:var(--red)">❌ No 3D image returned</span>';
    }
  }catch(e){bbl.innerHTML='<span style="color:var(--red)">❌ '+escHtml(e.message)+'</span>';}
  loading=false; document.getElementById('sbtn').disabled=false;
}

function downloadVideo(){
  if(_videoB64){_downloadB64(_videoB64,_videoMime,'fusion-video.mp4');}
  else if(_videoUrl){var a=document.createElement('a');a.href=_videoUrl;a.download='fusion-video.mp4';a.target='_blank';a.click();}
  else showToast('No video to download');
}

function sendVideoToChat(){
  var vid=document.getElementById('videoEl');
  if(!vid||!vid.src)return showToast('Generate a video first');
  var bbl=addMsg('ai','',null,null);
  bbl.innerHTML='<video controls style="width:100%;max-width:480px;border-radius:10px;border:1px solid var(--glass-bdr2)" src="'+vid.src+'"></video>'
    +'<br><span style="font-size:11px;color:var(--tx3)">Generated video</span>';
  hist.push({role:'assistant',content:'[Video generated]'});closeSP(); scrollDown(); showToast('💬 Sent to chat!');
}

var _threedAbort = null;
var _threedB64 = '', _threedMime = 'image/jpeg';

async function generate3D(){
  var prompt=document.getElementById('threedPrompt').value.trim();
  if(!prompt)return showToast('Describe your 3D object');
  var style=(document.getElementById('threedStyle')||{value:'3d render'}).value;
  var sz=((document.getElementById('threedSize')||{value:'1024x1024'}).value).split('x');
  var w=parseInt(sz[0])||1024,h=parseInt(sz[1])||1024;
  _genBtnStart('threedGenBtn','threedStopBtn','Generating 3D…');
  var res=document.getElementById('threedResult');
  if(res)res.style.display='none';_threedAbort=new AbortController();
  try{
    // Try real 3D model first, fallback to render
    var ep='/api/generate/3d/model';
    var r=await apiFetch(ep,{method:'POST',
      signal:_threedAbort.signal,
      headers:{'Content-Type':'application/json'},
      body:JSON.stringify({prompt:prompt,style:style,width:w,height:h})});
    if(!r.ok){var et=await r.text();var er={};try{er=JSON.parse(et);}catch(_){}showToast('❌ '+(er.error||'3D gen failed'),6000);_genBtnDone('threedGenBtn','threedStopBtn');return;}
    var d=await r.json();
    var img=document.getElementById('threedImg');
    var tag=document.getElementById('threedModelTag');
    var tdViewer=document.getElementById('threedModelViewer');
    // Handle real 3D GLB model vs rendered image
    if(d.format==='glb'&&d.b64){
      // Offer download as GLB
      var glbData='data:model/gltf-binary;base64,'+d.b64;
      if(img)img.style.display='none';
      if(tdViewer){
        tdViewer.style.display='';
        tdViewer.innerHTML='<div style="padding:16px;background:rgba(0,255,80,.08);border:1px solid rgba(0,255,80,.2);border-radius:10px;text-align:center"><div style="font-size:24px;margin-bottom:8px">🧊</div><div style="font-size:13px;color:var(--tx);font-weight:600">3D Model Ready (GLB)</div><div style="font-size:11px;color:var(--tx2);margin:6px 0">'+escHtml(d.backend||'HuggingFace')+'</div><a href="'+glbData+'" download="fusion-3d.glb" class="ig-btn" style="display:inline-block;margin-top:8px;padding:8px 18px;text-decoration:none">⬇️ Download .GLB</a></div>';
      }
      if(tag)tag.textContent='✅ Real 3D Model — '+escHtml(d.backend||'');
      showToast('✅ 3D model ready! Download as .GLB');
      _genBtnDone('threedGenBtn','threedStopBtn');
      if(res)res.style.display='';
      return;
    }
    if(tdViewer)tdViewer.style.display='none';
    if(img)img.style.display='';
    var src=d.b64?'data:'+d.mime+';base64,'+d.b64:(d.url||'');
    if(!src){showToast('❌ No result returned',5000);_genBtnDone('threedGenBtn','threedStopBtn');return;}
    _threedB64=d.b64||''; _threedMime=d.mime||'image/jpeg';if(img)img.src=src;
    if(tag)tag.textContent='⚡ '+escHtml(d.backend||'Pollinations')+(d.enhanced_prompt?' · '+escHtml(d.enhanced_prompt.slice(0,60))+'…':'');
    if(res)res.style.display='';showToast('✅ 3D render ready!');
  }catch(e){
    if(e.name!=='AbortError')showToast('❌ '+e.message,6000);else showToast('⏹ Stopped');
  }
  _threedAbort=null;_genBtnDone('threedGenBtn','threedStopBtn');
}

function stop3DGen(){if(_threedAbort)_threedAbort.abort();}

function download3D(){
  if(_threedB64){_downloadB64(_threedB64,_threedMime,'fusion-3d.png');}
  else{var img=document.getElementById('threedImg');if(img&&img.src){var a=document.createElement('a');a.href=img.src;a.download='fusion-3d.png';a.target='_blank';a.click();}}
}

function send3DToChat(){
  if(!_threedB64)return showToast('Generate a 3D render first');
  var src='data:'+_threedMime+';base64,'+_threedB64;
  var bbl=addMsg('ai','',null,null);
  bbl.innerHTML='<img src="'+src+'" style="max-width:100%;border-radius:12px;border:1px solid var(--glass-bdr2);display:block;cursor:zoom-in" onclick="openImgFull(this.src)"/>'
    +'<br><span style="font-size:11px;color:var(--tx3)">3D render</span>';
  hist.push({role:'assistant',content:'[3D render generated]'});closeSP(); scrollDown(); showToast('💬 Sent to chat!');
}

async function generateMusic() {
  var prompt=document.getElementById('musicPrompt').value.trim();
  if(!prompt) return showToast('Enter a music prompt');showToast('🎵 Music generation coming soon — use Browser TTS for speech',5000);
}

async function generateSpeech(){
  var text=document.getElementById('speechText').value.trim();
  if(!text)return showToast('Enter text to synthesise');
  var voice=(document.getElementById('speechVoice')||{value:'alloy'}).value;
  _genBtnStart('speechGenBtn','speechStopBtn','Synthesizing...');
  try{
    var r=await apiFetch('/api/generate/audio',{method:'POST',
      signal:_genAbortCtrl?_genAbortCtrl.signal:undefined,
      headers:{'Content-Type':'application/json'},
      body:JSON.stringify({text:text,voice:voice})});
    if(!r.ok){var er=await r.json();showToast('❌ '+(er.error||'TTS failed')+' — browser TTS used',5000);browserTTS();_genBtnDone('speechGenBtn','speechStopBtn');return;}
    var d=await r.json();
    if(d.b64){
      var au=document.getElementById('speechAudio');
      au.src='data:'+d.mime+';base64,'+d.b64;document.getElementById('speechResult').classList.add('show');
      var dl=document.getElementById('speechDlBtn');
      if(dl){dl.style.display='';dl.onclick=function(){_downloadB64(d.b64,d.mime,'fusion-speech.mp3');};}
      au.play();showToast('✅ Speech ready!');
    }
  }catch(e){if(e.name!=='AbortError'){showToast('❌ '+e.message,5000);browserTTS();}}
  _genBtnDone('speechGenBtn','speechStopBtn');
}

function browserTTS() {
  var text=document.getElementById('speechText').value.trim();
  if(!text){showToast('Enter text first');return;}
  if(!window.speechSynthesis){showToast('Browser TTS not supported');return;}
  var voices=window.speechSynthesis.getVoices();
  var utt=new SpeechSynthesisUtterance(text.slice(0,500));
  // Try to pick a natural-sounding English voice
  var preferred=voices.find(function(v){return v.lang.startsWith('en')&&v.localService;})
    ||voices.find(function(v){return v.lang.startsWith('en');})
    ||voices[0];
  if(preferred) utt.voice=preferred;utt.rate=1; utt.pitch=1;window.speechSynthesis.cancel();window.speechSynthesis.speak(utt);
  showToast('\u{1F50A} Speaking with browser TTS: '+(preferred?preferred.name:'default'));
}
function loadVCVoices() {
  var sel=document.getElementById('vcVoiceSelect'); if(!sel) return;
  var voices=window.speechSynthesis?window.speechSynthesis.getVoices():[];
  sel.innerHTML='<option value="">Default Voice</option>';
  voices.forEach(function(v,i){var o=document.createElement('option');o.value=i;o.textContent=v.name+' ('+v.lang+')';sel.appendChild(o);});
  if(vcPref.voiceIdx!==undefined) sel.value=vcPref.voiceIdx;
  if(vcPref.rate!==undefined){var r=document.getElementById('vcRate');if(r){r.value=vcPref.rate;document.getElementById('vcRateLbl').textContent=parseFloat(vcPref.rate).toFixed(1)+'x';}}
}
function saveVoiceChatPref(){
  var sel=document.getElementById('vcVoiceSelect'),rate=document.getElementById('vcRate');
  vcPref={voiceIdx:sel?sel.value:'',rate:rate?parseFloat(rate.value):1};localStorage.setItem('fusion_vc',JSON.stringify(vcPref));
}
function toggleVoiceChat() { if(vcActive) stopVoiceChat(); else startVoiceChat(); }
function startVoiceChat() {
  var SR=window.SpeechRecognition||window.webkitSpeechRecognition;
  if(!SR){showErrOverlay('Voice chat not supported. Try Chrome or Edge.');return;}
  vcActive=true;
  var btn=document.getElementById('vcBtn'); btn.classList.add('active'); btn.textContent='\u{1F534} Listening... (click to stop)';
  document.getElementById('vcStatus').textContent='Listening... Speak now';document.getElementById('vcTranscript').textContent='';
  vcRecog=new SR(); vcRecog.continuous=false; vcRecog.interimResults=true; vcRecog.lang='en-US';
  vcRecog.onresult=function(e){
    var final='',interim='';
    for(var i=e.resultIndex;i<e.results.length;i++){if(e.results[i].isFinal)final+=e.results[i][0].transcript;else interim+=e.results[i][0].transcript;}
    document.getElementById('vcTranscript').textContent=interim||final||'Listening...';
    if(final){document.getElementById('vcTranscript').textContent=final;sendVoiceChatMsg(final);}
  };
  vcRecog.onerror=function(e){stopVoiceChat();if(e.error!=='aborted')showErrOverlay('Mic error: '+e.error);};
  vcRecog.onend=function(){if(vcActive){setTimeout(function(){if(vcActive&&vcRecog)vcRecog.start();},800);}};vcRecog.start();
}
async function sendVoiceChatMsg(text) {
  if(!text.trim()||loading) return;document.getElementById('vcStatus').textContent='AI is thinking...';
  document.getElementById('msgIn').value=text;await sendMsg();
  setTimeout(function(){if(vcActive)document.getElementById('vcStatus').textContent='Listening... Speak now';},2000);
}
function stopVoiceChat() {
  vcActive=false;
  if(vcRecog){try{vcRecog.stop();}catch(e){}vcRecog=null;}
  var btn=document.getElementById('vcBtn'); btn.classList.remove('active'); btn.textContent='\u{1F3A4} Start Voice Chat';
  document.getElementById('vcStatus').textContent='Voice chat stopped. Click to restart.';
}

async function loadSaved(){try{var r=await apiFetch('/api/saved');var d=await r.json();renderSaved(d.items||[]);}catch(e){}}
function renderSaved(items){
  var el=document.getElementById('savedList');
  if(!items.length){el.innerHTML='<div class="empty-state"><span class="ei">&#x1F516;</span>Nothing saved yet.</div>';return;}
  el.innerHTML=items.map(function(it){
    return '<div class="saved-item" onclick="loadSavedItem(\''+encodeURIComponent(it.content)+'\')">'
      +'<button class="si-del" onclick="event.stopPropagation();deleteSaved('+it.id+')" title="Delete">&#x2715;</button>'
      +'<div class="si-title">'+escHtml(it.title)+'</div>'
      +'<div class="si-preview">'+escHtml(it.content.slice(0,80))+'</div>'
      +'<div class="si-ts">'+it.ts.slice(0,16).replace('T',' ')+'</div>'
      +'</div>';
  }).join('');
}
async function saveMessage(content){
  var title=content.slice(0,60)+(content.length>60?'\u2026':'');
  try{var r=await apiFetch('/api/saved',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({title:title,content:content})});var d=await r.json();if(r.ok){showToast('\u{1F516} Saved!');loadSaved();}else showToast('Error: '+d.error);}catch(e){showToast('Error saving');}
}
async function deleteSaved(id){await apiFetch('/api/saved/'+id,{method:'DELETE'});loadSaved();showToast('Deleted');}
function loadSavedItem(enc){var c=decodeURIComponent(enc);closeSP();document.getElementById('msgIn').value=c;ar(document.getElementById('msgIn'));document.getElementById('msgIn').focus();}
async function loadMemory(){try{var r=await apiFetch('/api/memory');var d=await r.json();renderMemory(d.memory||[]);}catch(e){}}
function renderMemory(items){
  var el=document.getElementById('memoryList');
  if(!items.length){el.innerHTML='<div class="empty-state"><span class="ei">&#x1F9E0;</span>No memories yet.</div>';return;}
  el.innerHTML=items.map(function(m){
    return '<div class="mem-item"><span class="mem-key">'+escHtml(m.key)+'</span><span class="mem-val">'+escHtml(m.value)+'</span><button class="mem-del" onclick="delMemory(\''+encodeURIComponent(m.key)+'\')">&#x2715;</button></div>';
  }).join('');
}
async function addMemory(){
  var k=document.getElementById('memKey').value.trim(),v=document.getElementById('memVal').value.trim();
  if(!k||!v) return showToast('Enter key and value');
  try{await apiFetch('/api/memory',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({key:k,value:v})});document.getElementById('memKey').value='';document.getElementById('memVal').value='';loadMemory();showToast('\u{1F9E0} Memory saved');}catch(e){showToast('Error');}
}
async function delMemory(enc){var key=decodeURIComponent(enc);await apiFetch('/api/memory/'+encodeURIComponent(key),{method:'DELETE'});loadMemory();showToast('Memory removed');}

var SP_TABS =['keys','model','imagine','video','threed','audio','voicechat','tts','convs','saved','memory','theme','scraper','extra'];
_initAvatar();
function openSP(tab){document.getElementById('sp').classList.add('open');document.getElementById('ovl').classList.add('open');closeDrop();switchSPTab(tab||'model');}
function closeSP(){document.getElementById('sp').classList.remove('open');document.getElementById('ovl').classList.remove('open');}
function switchSPTab(t){
  SP_TABS.forEach(function(x){
    var el=document.getElementById('spTab-'+x);
    var btn=document.getElementById('spt-'+x);
    if(el){
      // Browser tab needs flex, all others block
      var disp=(x===t)?(x==='scraper'?'flex':'block'):'none';
      el.style.display=disp;
    }
    if(btn) btn.classList.toggle('active',x===t);
  });
  if(t==='saved') loadSaved();if(t==='convs') loadConvs();if(t==='memory') loadMemory();
  if(t==='tts'||t==='voicechat'){loadVoices();loadVCVoices();}
  if(t==='theme') updateThemeUI();if(t==='threed') update3DUI();
  if(t==='extra') _loadCustomEndpointUI();
}

function update3DUI(){
  var model=document.getElementById('threedModel');
  if(!model) return;
  var v=model.value||'';
  var needsImg=['fal-ai/triposr','fal-ai/trellis','fal-ai/stable-zero123'].includes(v);
  var sec=document.getElementById('threed3dImageSection');
  if(sec) sec.style.display=needsImg?'':'none';
  var tp=document.getElementById('threedPrompt');
  if(tp) tp.placeholder=needsImg?'Optional: describe the object to guide 3D generation...':'Describe the 3D object in detail (required for text-to-3D)...';
}
function toggleDrop(){document.getElementById('uDrop').classList.toggle('open');}
function closeDrop(){document.getElementById('uDrop').classList.remove('open');}
document.addEventListener('click',function(e){if(document.getElementById('uChip')&&!document.getElementById('uChip').contains(e.target))closeDrop();});

async function fetchDevStats(){try{var r=await apiFetch('/api/dev/stats');if(!r.ok)return;devStats=await r.json();}catch(e){}}

// ── Dev checkbox helpers — avoid quote-escaping hell in HTML strings ──────────
function _devToggleFromEl(el){
  var key=el.closest('[data-key]')&&el.closest('[data-key]').dataset.key||el.dataset.key;
  if(key) devToggleModel(key, el.checked);
}
function _devMenuToggleFromEl(el){
  var key=el.closest('[data-key]')&&el.closest('[data-key]').dataset.key||el.dataset.key;
  if(key) devSetMenuToggle(key, el.checked);
}

// ── Dev: All Models Tab ────────────────────────────────────────────────────
function renderAllModelsTab(body){
  body.innerHTML='<div style="text-align:center;padding:20px;color:var(--tx3)">⏳ Loading models…</div>';
  apiFetch('/api/dev/disabled-models').then(function(r){return r.json();}).then(function(dd){
    var disabled=dd.disabled||[];
    var h='<h5 style="margin-bottom:12px">🤖 All Models ('+allModels.length+')</h5>';
    h+='<div style="font-size:11px;color:var(--tx2);margin-bottom:10px">Toggle models on/off for all users.</div>';
    h+='<div style="display:flex;gap:6px;flex-wrap:wrap;margin-bottom:14px">';
    h+='<button onclick="devToggleAllModels(true)" style="font-size:11px;padding:5px 12px;background:rgba(34,197,94,.15);border:1px solid rgba(34,197,94,.3);border-radius:6px;color:#4ade80;cursor:pointer">✅ Enable All</button>';
    h+='<button onclick="devToggleAllModels(false)" style="font-size:11px;padding:5px 12px;background:rgba(124,58,237,.12);border:1px solid rgba(124,58,237,.25);border-radius:6px;color:var(--red);cursor:pointer">🚫 Disable All</button>';
    h+='</div>';
    h+='<div class="dev-list" style="max-height:480px;overflow-y:auto">';
    // Group by provider
    var groups={};
    allModels.forEach(function(m){
      var p=m.provider||'other';
      if(!groups[p]) groups[p]=[];
      groups[p].push(m);
    });
    Object.entries(groups).forEach(function(gEntry){
      var prov=gEntry[0], models=gEntry[1];
      h+='<div style="font-size:9px;font-weight:700;color:var(--tx3);text-transform:uppercase;letter-spacing:1px;padding:8px 0 4px">'+escHtml(prov)+'</div>';
      models.forEach(function(m){
        var isDisabled=disabled.includes(m.key);
        var bg=isDisabled?'rgba(124,58,237,.06)':'rgba(34,197,94,.04)';
        var borderCol=isDisabled?'rgba(124,58,237,.2)':'rgba(34,197,94,.15)';
        h+='<div style="display:flex;align-items:center;gap:8px;padding:7px 10px;background:'+bg+';border:1px solid '+borderCol+';border-radius:8px;margin-bottom:4px">';
        h+='<span style="font-size:14px">'+(m.emoji||'🤖')+'</span>';
        h+='<div style="flex:1;min-width:0">';
        h+='<div style="font-size:12px;font-weight:600;color:var(--tx)">'+escHtml(m.label||m.key)+'</div>';
        h+='<div style="font-size:10px;color:var(--tx3)">'+escHtml(m.key)+'</div>';
        h+='</div>';
        h+='<label style="display:flex;align-items:center;gap:5px;cursor:pointer;user-select:none">';
        h+='<input type="checkbox" data-key="'+escHtml(m.key)+'" '+(isDisabled?'':'checked')+' onchange="_devToggleFromEl(this)" style="width:14px;height:14px;accent-color:var(--red)"/>';
        h+='<span style="font-size:10px;color:'+(isDisabled?'var(--red)':'#4ade80')+'">'+(isDisabled?'Off':'On')+'</span>';
        h+='</label></div>';
      });
    });
    h+='</div>';
    body.innerHTML=h;
  }).catch(function(e){body.innerHTML='<div style="color:var(--red);padding:20px">Error: '+escHtml(e.message)+'</div>';});
}
async function devToggleModel(key,enabled){
  try{
    await apiFetch('/api/dev/model-enable',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({key:key,enabled:enabled})});
    showToast((enabled?'✅ Enabled: ':'🚫 Disabled: ')+key);
  }catch(e){showToast('Error: '+e.message);}
}
async function devToggleAllModels(enabled){
  var keys=allModels.map(function(m){return m.key;});
  for(var i=0;i<keys.length;i++){
    try{await apiFetch('/api/dev/model-enable',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({key:keys[i],enabled:enabled})});}catch(e){}
  }
  showToast(enabled?'✅ All models enabled':'🚫 All models disabled');
  renderAllModelsTab(document.getElementById('devModalBody'));
}

// ── Dev: Menu Toggle Tab ────────────────────────────────────────────────────
// _devMenuToggles moved to top of script
function renderMenuToggleTab(body){
  var items=[
    {key:'rtb_imagegen',label:'Right Bar — Image Gen button'},
    {key:'rtb_video',label:'Right Bar — Video button'},
    {key:'rtb_threed',label:'Right Bar — 3D button'},
    {key:'rtb_audio',label:'Right Bar — Audio button'},
    {key:'rtb_voice',label:'Right Bar — Voice button'},
    {key:'sp_imagetab',label:'Settings — Image tab'},
    {key:'sp_videotab',label:'Settings — Video tab'},
    {key:'sp_threedtab',label:'Settings — 3D tab'},
    {key:'sp_audiotab',label:'Settings — Audio tab'},
    {key:'sp_voicetab',label:'Settings — Voice Chat tab'},
    {key:'sp_ttstab',label:'Settings — TTS tab'},
    {key:'sp_scrapetab',label:'Settings — Scraper tab'},
    {key:'header_voice',label:'Header — Voice Chat button'},
    {key:'lsb_weather',label:'Left Sidebar — Weather section'},
  ];
  var h='<h5 style="margin-bottom:8px">🚧 Menu Item Toggles</h5>';
  h+='<div style="font-size:11px;color:var(--tx2);margin-bottom:12px;line-height:1.6">Hide or show UI elements for all users. Changes are instant and saved locally.</div>';
  h+='<div class="dev-list">';
  items.forEach(function(item){
    var isOn=_devMenuToggles[item.key]!==false;
    h+='<div style="display:flex;align-items:center;gap:10px;padding:9px 12px;background:rgba(255,255,255,.03);border:1px solid var(--glass-bdr);border-radius:8px;margin-bottom:5px">';
    h+='<div style="flex:1;font-size:12px;color:var(--tx)">'+escHtml(item.label)+'</div>';
    h+='<label style="display:flex;align-items:center;gap:6px;cursor:pointer">';
    h+='<input type="checkbox" data-key="'+escHtml(item.key)+'" '+(isOn?'checked':'')+' onchange="_devMenuToggleFromEl(this)" style="width:15px;height:15px;accent-color:var(--red)"/>';
    h+='<span style="font-size:11px;color:'+(isOn?'#4ade80':'var(--red)')+'">'+(isOn?'Visible':'Hidden')+'</span>';
    h+='</label></div>';
  });
  h+='</div>';
  h+='<button onclick="devResetMenuToggles()" style="margin-top:12px;font-size:11px;padding:7px 14px;background:rgba(255,255,255,.06);border:1px solid var(--glass-bdr);border-radius:8px;color:var(--tx2);cursor:pointer">↺ Reset All to Visible</button>';
  body.innerHTML=h;
}
function devSetMenuToggle(key,visible){
  _devMenuToggles[key]=visible;
  localStorage.setItem('dev_menu_toggles',JSON.stringify(_devMenuToggles));
  _applyMenuToggles();
  showToast((visible?'👁 Shown: ':'🙈 Hidden: ')+key);
  // re-render to update colors
  renderMenuToggleTab(document.getElementById('devModalBody'));
}
function devResetMenuToggles(){
  _devMenuToggles={};localStorage.setItem('dev_menu_toggles','{}');
  _applyMenuToggles();showToast('↺ All menu items reset to visible');
  renderMenuToggleTab(document.getElementById('devModalBody'));
}

function openDevModal(){closeDrop();document.getElementById('devModal').classList.add('open');currentModalTab='overview';switchModalTab('overview');}
function closeDevModal(){document.getElementById('devModal').classList.remove('open');}
document.getElementById('devModal').addEventListener('click',function(e){if(e.target===document.getElementById('devModal'))closeDevModal();});
function switchModalTab(t){
  currentModalTab=t;
  ['overview','tokens','models','users','activity','visitors','logs','limits','flags','console','allmodels','menutoggle'].forEach(function(x){var el=document.getElementById('mdt-'+x);if(el)el.classList.toggle('active',x===t);});
  renderModalTab(t);
}
async function renderModalTab(t){
  var body=document.getElementById('devModalBody');
  if(t==='console'){renderConsoleTab(body);return;}
  if(t==='flags'){renderFlagsTab(body);return;}
  if(t==='limits'){renderLimitsTab(body);return;}
  if(t==='logs'){renderLogsTab(body);return;}
  if(t==='allmodels'){
    if(!allModels.length){
      body.innerHTML='<div style="text-align:center;padding:30px;color:var(--tx3)">⏳ Loading models…</div>';
      await loadModels();
    }
    renderAllModelsTab(body);return;
  }
  if(t==='menutoggle'){renderMenuToggleTab(body);return;}
  if(!devStats){body.innerHTML='<div style="text-align:center;padding:40px;color:var(--tx3)">\u23F3 Loading\u2026</div>';await fetchDevStats();if(!devStats){body.innerHTML='<div style="text-align:center;padding:40px;color:var(--red)">Failed to load stats.</div>';return;}}
  if(t==='overview')  renderOverviewTab(body);else if(t==='tokens')    renderTokensTab(body);
  else if(t==='models')   renderModelsTab(body);else if(t==='users')    renderUsersTab(body);
  else if(t==='activity') renderActivityTab(body);else if(t==='visitors') renderVisitorsTab(body);
}
function renderOverviewTab(body){
  var s=devStats; var h='<div class="stat-grid">';
  h+=statCard('Users',s.total_users,'');h+=statCard('Guests',s.total_guests,'all time');
  h+=statCard('Messages',s.total_messages,'');h+=statCard('Today',s.messages_today,'');
  h+=statCard('Visits 24h',s.visits_24h,'');h+=statCard('Total Visits',s.total_visits,'');h+=statCard('Online',s.online_now,'30min');
  h+='</div><div class="dev-list"><h5>\u{1F511} Server Env Keys</h5>';
  var keyLabels={GROQ_KEY:'⚡ Groq',OPENROUTER_KEY:'🔀 OpenRouter',HF_TOKEN:'🤗 HuggingFace',GITHUB_TOKEN:'💻 GitHub',FAL_KEY:'⚡ fal.ai',REPLICATE_KEY:'🧬 Replicate',STABILITY_KEY:'🎨 Stability AI',SECRET_KEY:'🔒 Secret Key',DEV_PASSWORD:'🔑 Dev Password',TP_USERNAME:'👤 Dev Username'};
  Object.entries(s.env_keys).forEach(function(e2){var label=keyLabels[e2[0]]||e2[0];var isset=e2[1];h+='<div class="dev-row" style="gap:10px"><span style="font-size:15px">'+(isset?'✅':'⬜')+'</span><span class="k" style="flex:1">'+escHtml(label)+'</span><code style="font-size:9px;color:var(--tx3);font-family:monospace">'+escHtml(e2[0])+'</code><span class="'+(isset?'ok':'no')+'" style="font-size:10px">'+(isset?'Set':'Not set')+'</span></div>';});
  h+='</div><div class="dev-list"><h5>\u{1F4CA} Top Models</h5>';
  if(s.top_models&&s.top_models.length){s.top_models.forEach(function(m){h+='<div class="dev-row"><span class="k">'+escHtml(m.model)+'</span><span class="v">'+m.count+' msgs</span></div>';});}
  else h+='<div style="font-size:12px;color:var(--tx3)">No messages yet</div>';h+='</div>'; body.innerHTML=h;
}

function renderUsersTab(body){
  body.innerHTML='<div style="text-align:center;padding:20px;color:var(--tx3)">⏳ Loading users…</div>';
  apiFetch('/api/dev/users').then(function(r){return r.json();}).then(function(d){
    var users=d.users||[];
    var h='<div style="display:flex;justify-content:space-between;align-items:center;margin-bottom:12px">';
    h+='<h5 style="margin:0">👥 All Users ('+users.length+')</h5>';
    h+='<button class="run-test-btn" style="margin:0;padding:6px 12px;font-size:11px" onclick="refreshDevStats()">🔄 Refresh</button></div>';
    h+='<div style="font-size:11px;color:var(--tx2);margin-bottom:10px">Click a username to see their message history.</div>';
    if(!users.length){h+='<div style="color:var(--tx3);font-size:12px">No users registered yet</div>';body.innerHTML=h;return;}
    h+='<div class="dev-list">';
    users.forEach(function(u){
      var isBanned=!!u.banned;
      var isGuest=u.username.startsWith('guest_');
      var rowBg=isBanned?'rgba(124,58,237,.08)':isGuest?'rgba(255,255,255,.02)':'';
      h+='<div class="dev-row" style="flex-wrap:wrap;gap:6px;background:'+rowBg+';border-radius:8px;padding:8px 10px;margin-bottom:6px;align-items:flex-start">';
      h+='<div style="flex:1;min-width:120px">';h+='<div style="font-weight:600;font-size:12px;color:var(--tx)">'+(isBanned?'🚫 ':isGuest?'👻 ':'👤 ')+escHtml(u.username)+'</div>';
      h+='<div style="font-size:10px;color:var(--tx3);margin-top:2px">ID: '+u.id+' · Joined: '+u.created.slice(0,10)+'</div>';
      if(u.last_active)h+='<div style="font-size:10px;color:var(--tx3)">Last active: '+u.last_active.slice(0,16).replace("T"," ")+'</div>';
      if(u.providers)h+='<div style="font-size:10px;color:var(--blue)">Keys: '+escHtml(u.providers)+'</div>';
      if(isBanned)h+='<div style="font-size:10px;color:var(--red);margin-top:2px">⛔ Ban reason: '+escHtml(u.banned)+'</div>';
      h+='</div>';h+='<div style="font-size:11px;font-weight:700;color:var(--tx2);min-width:50px;text-align:right">'+u.msg_count+' msgs</div>';
      h+='<div style="display:flex;flex-wrap:wrap;gap:4px;width:100%;margin-top:4px">';
      h+='<button onclick="devViewHistory('+u.id+',\''+escHtml(u.username)+'\')" style="font-size:10px;padding:3px 8px;background:rgba(6,182,212,.15);border:1px solid rgba(6,182,212,.3);border-radius:5px;color:var(--blue);cursor:pointer">📜 History</button>';
      if(!isBanned)h+='<button onclick="devBanUser('+u.id+',\''+escHtml(u.username)+'\')" style="font-size:10px;padding:3px 8px;background:rgba(124,58,237,.12);border:1px solid rgba(124,58,237,.25);border-radius:5px;color:var(--red);cursor:pointer">🚫 Ban</button>';
      else h+='<button onclick="devUnbanUser('+u.id+',\''+escHtml(u.username)+'\')" style="font-size:10px;padding:3px 8px;background:rgba(34,197,94,.12);border:1px solid rgba(34,197,94,.25);border-radius:5px;color:var(--green);cursor:pointer">✅ Unban</button>';
      h+='<button onclick="devResetKeys('+u.id+',\''+escHtml(u.username)+'\')" style="font-size:10px;padding:3px 8px;background:rgba(255,255,255,.06);border:1px solid var(--glass-bdr);border-radius:5px;color:var(--tx2);cursor:pointer">🔑 Reset Keys</button>';
      if(isGuest||isBanned)h+='<button onclick="devDeleteUser('+u.id+',\''+escHtml(u.username)+'\')" style="font-size:10px;padding:3px 8px;background:rgba(124,58,237,.18);border:1px solid rgba(124,58,237,.4);border-radius:5px;color:var(--red);cursor:pointer">🗑 Delete</button>';
      h+='</div></div>';
    });
    h+='</div>';body.innerHTML=h;
  }).catch(function(e){body.innerHTML='<div style="color:var(--red);padding:20px">Error loading users: '+escHtml(e.message)+'</div>';});
}
async function devViewHistory(uid,username){
  var body=document.getElementById('devModalBody');
  body.innerHTML='<div style="color:var(--tx3);padding:20px;text-align:center">⏳ Loading history for '+escHtml(username)+'…</div>';
  try{
    var r=await apiFetch('/api/dev/user/'+uid+'/history');
    var d=await r.json();
    var msgs=d.messages||[];
    var h='<button onclick="renderUsersTab(document.getElementById(\'devModalBody\'))" style="font-size:11px;padding:5px 12px;background:rgba(255,255,255,.06);border:1px solid var(--glass-bdr);border-radius:6px;color:var(--tx2);cursor:pointer;margin-bottom:12px">← Back</button>';
    h+='<h5 style="margin-bottom:10px">📜 Messages from @'+escHtml(username)+' (last '+msgs.length+')</h5>';
    if(!msgs.length){h+='<div style="color:var(--tx3);font-size:12px">No messages yet</div>';body.innerHTML=h;return;}
    h+='<div class="dev-list">';
    msgs.forEach(function(m){
      var rc=m.role==='user'?'var(--blue)':'var(--green)';
      h+='<div class="dev-row" style="flex-direction:column;align-items:flex-start;padding:8px;gap:4px">';
      h+='<div style="display:flex;gap:8px;width:100%"><span style="color:'+rc+';font-size:10px;font-weight:700;text-transform:uppercase">'+m.role+'</span>'+(m.model?'<span style="font-size:9px;color:var(--tx3);margin-left:auto">'+escHtml(m.model)+'</span>':'')+'<span style="font-size:9px;color:var(--tx3)">'+m.ts.slice(0,16).replace("T"," ")+'</span></div>';
      h+='<div style="font-size:12px;color:var(--tx);line-height:1.5;white-space:pre-wrap;word-break:break-word">'+escHtml(m.content.slice(0,400))+(m.content.length>400?'…':'')+'</div>';h+='</div>';
    });
    h+='</div>'; body.innerHTML=h;
  }catch(e){body.innerHTML='<div style="color:var(--red);padding:20px">Error: '+escHtml(e.message)+'</div>';}
}
async function devBanUser(uid,username){
  var reason=prompt('Ban reason for @'+username+' (shown to admin):','Violation of TOS');
  if(!reason)return;
  try{var r=await apiFetch('/api/dev/user/'+uid+'/ban',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({reason:reason})});var d=await r.json();showToast(d.message||'Banned');renderUsersTab(document.getElementById('devModalBody'));}
  catch(e){showToast('Error: '+e.message);}
}
async function devUnbanUser(uid,username){
  try{var r=await apiFetch('/api/dev/user/'+uid+'/unban',{method:'POST'});var d=await r.json();showToast(d.message||'Unbanned');renderUsersTab(document.getElementById('devModalBody'));}
  catch(e){showToast('Error: '+e.message);}
}
async function devResetKeys(uid,username){
  if(!confirm('Reset all API keys for @'+username+'?'))return;
  try{var r=await apiFetch('/api/dev/user/'+uid+'/reset_keys',{method:'POST'});var d=await r.json();showToast(d.message||'Keys reset');}
  catch(e){showToast('Error: '+e.message);}
}
async function devDeleteUser(uid,username){
  if(!confirm('Permanently delete @'+username+' and all their data?'))return;
  try{var r=await apiFetch('/api/dev/user/'+uid+'/delete',{method:'DELETE'});var d=await r.json();showToast(d.message||'Deleted');renderUsersTab(document.getElementById('devModalBody'));}
  catch(e){showToast('Error: '+e.message);}
}

async function renderLogsTab(body){
  body.innerHTML='<div style="color:var(--tx3);text-align:center;padding:20px">⏳ Loading logs…</div>';
  var h='<div style="display:flex;gap:8px;margin-bottom:12px;flex-wrap:wrap">';
  h+='<input id="logSearch" type="text" placeholder="Search logs…" style="flex:1;background:rgba(10,15,30,.7);border:1px solid var(--glass-bdr);border-radius:8px;padding:7px 10px;color:var(--tx);font-size:12px;min-width:150px" onkeydown="if(event.key===\'Enter\')devSearchLogs()"/>';
  h+='<button onclick="devSearchLogs()" class="run-test-btn" style="margin:0;padding:7px 14px;font-size:11px">🔍 Search</button>';
  h+='<button onclick="renderLogsTab(document.getElementById(\'devModalBody\'))" class="run-test-btn" style="margin:0;padding:7px 14px;font-size:11px;background:rgba(255,255,255,.06)">All Logs</button>';
  h+='<button onclick="devExportLogs()" class="run-test-btn" style="margin:0;padding:7px 14px;font-size:11px;background:rgba(34,197,94,.15)">⬇ Export JSON</button>';
  h+='</div><div id="devLogsContainer">⏳</div>';body.innerHTML=h;devLoadLogs('/api/dev/request-logs');
}
async function devLoadLogs(url){
  var c=document.getElementById('devLogsContainer');if(!c)return;
  try{
    var r=await apiFetch(url);var d=await r.json();
    var logs=d.logs||[];
    if(!logs.length){c.innerHTML='<div style="color:var(--tx3);font-size:12px;padding:10px">No logs found.</div>';return;}
    var h='<div style="font-size:10px;color:var(--tx3);margin-bottom:6px">Showing '+logs.length+' entries'+(d.query?' for "'+escHtml(d.query)+'"':'')+'</div>';
    h+='<div class="dev-list" style="max-height:420px;overflow-y:auto">';
    logs.forEach(function(m){
      var rc=m.role==='user'?'var(--blue)':'var(--green)';
      h+='<div class="dev-row" style="flex-direction:column;align-items:flex-start;padding:7px;gap:3px;border-left:2px solid '+rc+'20">';
      h+='<div style="display:flex;gap:8px;width:100%;font-size:9px;color:var(--tx3)">';h+='<span style="color:'+rc+';font-weight:700;text-transform:uppercase">'+m.role+'</span>';
      h+='<span style="color:var(--tx2)">@'+escHtml(m.username||'?')+'</span>';if(m.model)h+='<span>'+escHtml(m.model)+'</span>';
      h+='<span style="margin-left:auto">'+m.ts.slice(0,16).replace("T"," ")+'</span>';h+='<span style="color:var(--tx3)">#'+m.id+'</span></div>';
      h+='<div style="font-size:11px;color:var(--tx);line-height:1.5;word-break:break-word">'+escHtml(m.content.slice(0,200))+(m.content.length>200?'…':'')+'</div>';h+='</div>';
    });
    h+='</div>'; c.innerHTML=h;
  }catch(e){if(c)c.innerHTML='<div style="color:var(--red)">Error: '+escHtml(e.message)+'</div>';}
}
async function devSearchLogs(){
  var q=(document.getElementById('logSearch')||{value:''}).value.trim();
  if(!q){devLoadLogs('/api/dev/request-logs');return;}
  devLoadLogs('/api/dev/search-logs?q='+encodeURIComponent(q));
}
async function devExportLogs(){
  try{var r=await apiFetch('/api/dev/export-logs');var d=await r.json();var blob=new Blob([JSON.stringify(d.logs,null,2)],{type:'application/json'});var a=document.createElement('a');a.href=URL.createObjectURL(blob);a.download='fusion-logs-'+new Date().toISOString().slice(0,10)+'.json';a.click();showToast('✅ Exported '+d.exported+' log entries');}
  catch(e){showToast('Export error: '+e.message);}
}

async function renderFlagsTab(body){
  body.innerHTML='<div style="color:var(--tx3);text-align:center;padding:20px">⏳ Loading…</div>';
  try{
    var r=await apiFetch('/api/dev/feature-flags');var flags=await r.json();
    var flagDefs=[
      {k:'ai_requests_enabled',label:'💬 AI Requests',desc:'Allow all chat AI requests'},
      {k:'image_gen_enabled',label:'🖼 Image Generation',desc:'Allow image generation'},
      {k:'video_gen_enabled',label:'🎬 Video Generation',desc:'Allow video generation'},
      {k:'chat_streaming_enabled',label:'⚡ Streaming',desc:'Stream chat responses token by token'},
      {k:'worker_image_enabled',label:'🏭 Worker Image',desc:'Allow CF Worker image generation'},
      {k:'slow_mode',label:'🐢 Slow Mode',desc:'Throttle all AI requests (emergency)'},
    ];
    var h='<div style="font-size:11px;color:var(--tx2);margin-bottom:14px;line-height:1.6">⚑ Feature flags let you enable or disable functionality without redeploying. Changes take effect immediately.</div>';
    h+='<div class="dev-list">';
    flagDefs.forEach(function(fd){
      var val=flags[fd.k]!==undefined?flags[fd.k]:true;
      var isOn=!!val;
      h+='<div class="dev-row" style="padding:10px;margin-bottom:8px;border-radius:8px">';h+='<div style="flex:1"><div style="font-weight:600;font-size:12px">'+fd.label+'</div>';
      h+='<div style="font-size:10px;color:var(--tx3)">'+fd.desc+'</div></div>';
      h+='<button onclick="devToggleFlag(\''+fd.k+'\',this)" data-key="'+fd.k+'" data-val="'+(isOn?'1':'0')+'" style="padding:5px 14px;border-radius:20px;border:none;font-size:11px;font-weight:600;cursor:pointer;background:'+(isOn?'rgba(34,197,94,.2)':'rgba(124,58,237,.2)')+';color:'+(isOn?'var(--green)':'var(--red)')+'">'+(isOn?'ON':'OFF')+'</button>';
      h+='</div>';
    });
    h+='</div>';h+='<div style="margin-top:4px;font-size:10px;color:var(--tx3)">Flags are stored server-side and persist across restarts.</div>';body.innerHTML=h;
  }catch(e){body.innerHTML='<div style="color:var(--red);padding:20px">Error: '+escHtml(e.message)+'</div>';}
}
async function devToggleFlag(key,btn){
  var curVal=btn.dataset.val==='1';
  var newVal=!curVal;
  try{
    // Get current flags
    var r=await apiFetch('/api/dev/feature-flags');var flags=await r.json();
    flags[key]=newVal;
    var r2=await apiFetch('/api/dev/feature-flags',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify(flags)});
    if(!r2.ok){showToast('❌ Failed to save flag');return;}
    btn.dataset.val=newVal?'1':'0';btn.textContent=newVal?'ON':'OFF';
    btn.style.background=newVal?'rgba(34,197,94,.2)':'rgba(124,58,237,.2)';btn.style.color=newVal?'var(--green)':'var(--red)';
    showToast((newVal?'✅':'🚫')+' '+key+' = '+newVal);
  }catch(e){showToast('Error: '+e.message);}
}

async function renderLimitsTab(body){
  body.innerHTML='<div style="color:var(--tx3);text-align:center;padding:20px">⏳ Loading…</div>';
  try{
    var r=await apiFetch('/api/dev/rate-limits');var lim=await r.json();
    var h='<div style="font-size:11px;color:var(--tx2);margin-bottom:14px;line-height:1.6">🛡 Rate limits protect the app from abuse and high API costs.</div>';
    h+='<div class="dev-list">';
    var fields=[
      {k:'global_rpm',label:'Global RPM',desc:'Max requests per minute across all users',type:'number',min:1,max:10000},
      {k:'per_user_rpm',label:'Per-User RPM',desc:'Max requests per user per minute',type:'number',min:1,max:1000},
      {k:'image_per_user_hour',label:'Images/User/Hour',desc:'Max image generations per user per hour',type:'number',min:1,max:500},
    ];
    fields.forEach(function(f){
      h+='<div class="dev-row" style="padding:10px;margin-bottom:8px;border-radius:8px;flex-wrap:wrap;gap:8px">';h+='<div style="flex:1"><div style="font-weight:600;font-size:12px">'+f.label+'</div>';
      h+='<div style="font-size:10px;color:var(--tx3)">'+f.desc+'</div></div>';
      h+='<input id="rl_'+f.k+'" type="number" min="'+f.min+'" max="'+f.max+'" value="'+(lim[f.k]!==undefined?lim[f.k]:'?')+'" style="width:80px;background:rgba(10,15,30,.7);border:1px solid var(--glass-bdr);border-radius:6px;padding:4px 8px;color:var(--tx);font-size:12px;text-align:right"/>';
      h+='</div>';
    });
    var emergBg=lim.emergency_stop?'rgba(124,58,237,.2)':'rgba(34,197,94,.1)';
    var emergColor=lim.emergency_stop?'var(--red)':'var(--green)';
    h+='<div class="dev-row" style="padding:10px;margin-bottom:8px;border-radius:8px;background:'+emergBg+';border:1px solid '+emergColor+'30">';
    h+='<div style="flex:1"><div style="font-weight:700;font-size:13px;color:'+emergColor+'">'+(lim.emergency_stop?'🚨':'✅')+' Emergency Stop</div>';
    h+='<div style="font-size:10px;color:var(--tx3)">Immediately halt ALL AI requests sitewide</div></div>';
    h+='<button onclick="devToggleEmergency(this)" data-val="'+(lim.emergency_stop?'1':'0')+'" style="padding:6px 16px;border-radius:20px;border:none;font-weight:700;font-size:12px;cursor:pointer;background:'+(lim.emergency_stop?'var(--red)':'rgba(124,58,237,.2)')+';color:#fff">'+(lim.emergency_stop?'ACTIVE — Click to Deactivate':'Activate Emergency Stop')+'</button>';
    h+='</div>';h+='</div>';h+='<button onclick="devSaveLimits()" class="run-test-btn" style="margin:8px 0 0;padding:8px 20px">💾 Save Rate Limits</button>';
    body.innerHTML=h;
  }catch(e){body.innerHTML='<div style="color:var(--red);padding:20px">Error: '+escHtml(e.message)+'</div>';}
}
async function devSaveLimits(){
  try{
    var r=await apiFetch('/api/dev/rate-limits');var lim=await r.json();
    var fields=['global_rpm','per_user_rpm','image_per_user_hour'];
    fields.forEach(function(k){var el=document.getElementById('rl_'+k);if(el&&!isNaN(parseInt(el.value)))lim[k]=parseInt(el.value);});
    var r2=await apiFetch('/api/dev/rate-limits',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify(lim)});
    if(r2.ok)showToast('✅ Rate limits saved');else showToast('❌ Failed');
  }catch(e){showToast('Error: '+e.message);}
}
async function devToggleEmergency(btn){
  var cur=btn.dataset.val==='1';
  var newVal=!cur;
  if(newVal&&!confirm('⚠️ This will STOP all AI requests for every user. Confirm?'))return;
  try{
    var r=await apiFetch('/api/dev/rate-limits');var lim=await r.json();
    lim.emergency_stop=newVal;await apiFetch('/api/dev/rate-limits',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify(lim)});
    showToast(newVal?'🚨 Emergency stop ACTIVATED':'✅ Emergency stop deactivated');renderLimitsTab(document.getElementById('devModalBody'));
  }catch(e){showToast('Error: '+e.message);}
}

function renderConsoleTab(body){
  var modelOpts='';
  allModels.filter(function(m){return m.type==='chat';}).forEach(function(m){
    modelOpts+='<option value="'+m.key+'">'+escHtml(m.label)+' ('+escHtml(m.provider||'')+')</option>';
  });
  var h='<div style="font-size:11px;color:var(--tx2);margin-bottom:14px;line-height:1.6">🧪 Send test prompts directly to any model. Bypasses Overseer routing.</div>';
  h+='<div style="display:flex;flex-direction:column;gap:10px">';
  h+='<select id="consoleModel" style="background:rgba(10,15,30,.7);border:1px solid var(--glass-bdr);border-radius:8px;padding:7px 10px;color:var(--tx);font-size:12px">'+modelOpts+'</select>';
  h+='<textarea id="consolePrompt" rows="4" placeholder="Enter a test prompt…" style="background:rgba(10,15,30,.7);border:1px solid var(--glass-bdr);border-radius:8px;padding:8px 10px;color:var(--tx);font-size:12px;resize:vertical;font-family:\'DM Mono\',monospace"></textarea>';
  h+='<button class="run-test-btn" style="margin:0" onclick="devRunConsole()">▶ Send Test Prompt</button>';
  h+='<div id="consoleResult" style="display:none;background:rgba(0,0,0,.4);border:1px solid var(--glass-bdr);border-radius:8px;padding:12px;font-size:12px;font-family:\'DM Mono\',monospace;white-space:pre-wrap;max-height:300px;overflow:auto;color:var(--tx)"></div>';
  h+='</div>';h+='<div class="dev-list" style="margin-top:16px"><h5>💡 Prompt Templates</h5>';
  var templates=[
    ['Capability check','Say exactly "OK I am working" and nothing else.'],
    ['Reasoning test','What is 137 × 49? Show your work step by step.'],
    ['Coding test','Write a Python function that checks if a string is a palindrome. Include tests.'],
    ['Context test','My name is AlphaTest. What is my name? (Test memory/context)'],
    ['Long output','List 10 different programming languages with a one-sentence description each.'],
  ];
  templates.forEach(function(t){
    h+='<div class="dev-row" style="cursor:pointer;padding:7px 10px;border-radius:6px" onclick="document.getElementById(\'consolePrompt\').value='+JSON.stringify(t[1])+'"><span class="k">'+escHtml(t[0])+'</span><span class="v" style="font-size:10px;color:var(--tx3);overflow:hidden;text-overflow:ellipsis;white-space:nowrap;flex:1;min-width:0">'+escHtml(t[1])+'</span></div>';
  });
  h+='</div>';body.innerHTML=h;
}
async function devRunConsole(){
  var model=(document.getElementById('consoleModel')||{value:'groq_llama33_70b'}).value;
  var prompt=(document.getElementById('consolePrompt')||{value:''}).value.trim();
  if(!prompt){showToast('Enter a prompt');return;}
  var res=document.getElementById('consoleResult');
  if(res){res.style.display='block';res.textContent='⏳ Waiting for response…';res.style.color='var(--tx3)';}
  var t0=Date.now();
  try{
    var r=await apiFetch('/api/dev/test-prompt',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({prompt:prompt,model_key:model})});
    var d=await r.json();
    var ms=Date.now()-t0;
    if(d.ok){
      res.style.color='var(--tx)';res.textContent='Model: '+escHtml(d.model)+'\nLatency: '+ms+'ms\n\n'+d.reply;
    } else {
      res.style.color='var(--red)';res.textContent='❌ Error (HTTP '+d.status+')\n\n'+(d.error||'Unknown error')+(d.raw?'\n\nRaw:\n'+d.raw:'');
    }
  }catch(e){if(res){res.style.color='var(--red)';res.textContent='JS Error: '+e.message;}}
}

function renderTokensTab(body){
  var s=devStats;
  var keyInfo={
    GROQ_KEY:{label:'⚡ Groq',desc:'Fast LLM inference — chat & reasoning',hint:'gsk_...',link:'https://console.groq.com/keys'},
    OPENROUTER_KEY:{label:'🔀 OpenRouter',desc:'50+ free models + image gen',hint:'sk-or-...',link:'https://openrouter.ai/keys'},
    CF_ACCOUNT_ID:{label:'☁️ Cloudflare Account',desc:'Required for CF Workers AI image gen',hint:'abc123...'},
    CF_KEY:{label:'☁️ Cloudflare API Key',desc:'Required for CF Workers AI image gen',hint:'v1.0-...'},
    CF_ACCOUNT_ID2:{label:'☁️ Cloudflare Account 2',desc:'Second CF account for extra 10k neurons',hint:'abc123...'},
    CF_KEY2:{label:'☁️ Cloudflare Key 2',desc:'API key for second CF account',hint:'v1.0-...'},
    WORKER_KEY:{label:'🏭 Worker API Key',desc:'Auth key for custom CF Worker image endpoint',hint:'your-secret'},
    HF_TOKEN:{label:'🤗 HuggingFace',desc:'HF inference models access',hint:'hf_...',link:'https://huggingface.co/settings/tokens'},
    GITHUB_TOKEN:{label:'💻 GitHub AI',desc:'GPT-4o, o4-mini, DeepSeek via GitHub',hint:'ghp_...',link:'https://github.com/settings/tokens'},
    STABILITY_KEY:{label:'🎨 Stability AI',desc:'Stable Diffusion, Stable Image Ultra',hint:'sk-...',link:'https://platform.stability.ai/account/keys'},
    WORKER_URL2:{label:'🏭 Worker URL 2 (Image History)',desc:'Second CF Worker — AI Image Generator with History',hint:'https://fusionai.gamesdohas.workers.dev/'},
    OPENAI_KEY:{label:'🤖 OpenAI (Direct)',desc:'GPT-4o, o3 via your own OpenAI key',hint:'sk-...',link:'https://platform.openai.com/api-keys'},
    ANTHROPIC_KEY:{label:'🎭 Anthropic (Claude)',desc:'Claude Opus/Sonnet/Haiku via your own key',hint:'sk-ant-...',link:'https://console.anthropic.com/settings/keys'},
    GEMINI_KEY:{label:'💎 Google Gemini',desc:'Gemini 2.0/2.5 via your own key',hint:'AIza...',link:'https://aistudio.google.com/apikey'},
    DEEPSEEK_KEY:{label:'🔵 DeepSeek (Direct)',desc:'DeepSeek Chat & Reasoner via your own key',hint:'sk-...',link:'https://platform.deepseek.com/api_keys'},
    MOONSHOT_KEY:{label:'🌙 Kimi / Moonshot AI',desc:'Kimi K2 via your own Moonshot key',hint:'sk-...',link:'https://platform.moonshot.ai/console/api-keys'},
    MISTRAL_KEY:{label:'🌀 Mistral AI (Direct)',desc:'Mistral Large/Small via your own key',hint:'...',link:'https://console.mistral.ai/api-keys'},
    XAI_KEY:{label:'✖️ xAI (Grok)',desc:'Grok 4 via your own xAI key',hint:'xai-...',link:'https://console.x.ai'},
    COHERE_KEY:{label:'🔮 Cohere',desc:'Command R+ via your own key',hint:'...',link:'https://dashboard.cohere.com/api-keys'},
    TOGETHER_KEY:{label:'🦙 Together AI',desc:'Llama & other OSS models',hint:'...',link:'https://api.together.xyz/settings/api-keys'},
    PERPLEXITY_KEY:{label:'🟣 Perplexity',desc:'Sonar Pro — web-grounded answers',hint:'pplx-...',link:'https://www.perplexity.ai/settings/api'},
  };
  var h='<div style="font-size:11px;color:var(--tx2);margin-bottom:14px;line-height:1.6">Server-level env keys. Set these in your HF Space → Settings → Variables. Users can also add their own in Settings → Keys.</div>';
  Object.entries(keyInfo).forEach(function(e2){
    var k=e2[0],info=e2[1],isSet=(s&&s.env_keys)?s.env_keys[k]:false;
    h+='<div class="krow" style="margin-bottom:10px">';
    h+='<div class="krow-top"><span class="kprov">'+info.label+'</span><span class="kstat '+(isSet?'set':'unset')+'">'+(isSet?'✓ Set':'✗ Missing')+'</span></div>';
    h+='<div style="font-size:11px;color:var(--tx2);margin-bottom:6px;line-height:1.5">'+info.desc+'</div>';
    h+='<div style="font-size:10px;color:var(--tx3);margin-bottom:6px;font-family:\'DM Mono\',monospace">Env var: <strong style="color:var(--tx2)">'+k+'</strong>'+(info.hint?' &nbsp;·&nbsp; Format: <em>'+info.hint+'</em>':'')+'</div>';
    if(info.link)h+='<a href="'+info.link+'" target="_blank" style="font-size:11px;color:var(--blue);text-decoration:none">→ Get key at '+info.link.replace('https://','')+'</a>';
    h+='</div>';
  });
  h+='<div class="dev-list" style="margin-top:4px"><h5>🧪 Live Connectivity Tests</h5>';h+='<button class="run-test-btn" style="margin-bottom:10px" onclick="testImgGen()">🖼 Test CF Image Gen</button>';
  h+='<div id="imgTestResult" style="font-size:11px;font-family:\'DM Mono\',monospace;white-space:pre-wrap;color:var(--tx2);background:rgba(0,0,0,.3);border-radius:8px;padding:10px;display:none;max-height:220px;overflow:auto"></div>';
  h+='</div>';h+='<div class="dev-list" style="margin-top:4px"><h5>ℹ️ How to set env keys on HuggingFace Spaces</h5>';
  h+='<div style="font-size:11px;color:var(--tx2);line-height:1.8">1. Go to your HF Space → <strong>Settings</strong> tab<br>2. Scroll to <strong>Repository secrets</strong><br>3. Add each key as a secret variable<br>4. Restart the Space — keys load on startup</div></div>';
  body.innerHTML=h;
}
async function testImgGen(){
  var box=document.getElementById('imgTestResult');
  if(!box)return;box.style.display='block';box.textContent='⏳ Testing CF image gen…';
  try{
    var r=await apiFetch('/api/test/image');var d=await r.json();
    if(d.ok){box.style.color='var(--green)';box.textContent='✅ SUCCESS\nModel: '+d.model+'\nHTTP: '+d.status+'\nSize: '+d.size_bytes+' bytes\nLatency: '+d.ms+'ms';}
    else{box.style.color='var(--red)';box.textContent='❌ FAILED\nModel: '+d.model+'\nHTTP: '+d.status+'\nError: '+d.error+(d.body?'\nBody: '+d.body:'')+(d.traceback?'\n\nTraceback:\n'+d.traceback:'');}
  }catch(e){box.style.color='var(--red)';box.textContent='❌ JS Error: '+e.message;}
}
function statCard(label,val,sub){return '<div class="stat-card"><label>'+label+'</label><div class="val">'+val+'</div>'+(sub?'<div class="sub">'+sub+'</div>':'')+'</div>';}
function renderModelsTab(body){
  var h='<div style="display:flex;gap:8px;margin-bottom:10px;flex-wrap:wrap">';
  h+='<button class="run-test-btn" id="runTestBtn" onclick="runModelTest()" style="margin:0">\u{1F50D} Test All Models</button>';h+='</div>';
  if(devTestResults){
    var s=devTestResults.summary;
    h+='<div class="overseer-box"><h4>\u{1F9E0} Overseer Report</h4><p>'+escHtml(devTestResults.overseer_report)+'</p></div>';
    h+='<div style="display:flex;gap:12px;margin-bottom:14px;font-size:12px;font-weight:600"><span style="color:var(--green)">\u2705 '+s.ok+' working</span><span style="color:var(--red)">\u274c '+s.failed+' failed</span><span style="color:var(--tx3)">\uD83D\uDD11 '+s.no_key+' no key</span></div>';
    devTestResults.models.forEach(function(m){
      var det=m.status==='ok'?(m.reply||'OK'):m.status==='no_key'?'No API key':(m.error||'HTTP '+m.http||'failed');
      h+='<div class="model-test-row '+m.status+'"><span class="mtbadge '+m.status+'">'+(m.status==='ok'?'\u2713 OK':m.status==='no_key'?'NO KEY':'\u2717 FAIL')+'</span><span class="mtlabel">'+escHtml(m.label)+'</span><span class="mtdetail">'+escHtml(det)+'</span><span class="mtms">'+(m.ms?m.ms+'ms':'')+'</span></div>';
    });
  }else{h+='<div style="text-align:center;padding:32px;color:var(--tx3);font-size:13px">Click "Test All Models" to ping every model.</div>';}
  body.innerHTML=h;
}
async function runModelTest(){
  var btn=document.getElementById('runTestBtn');if(btn){btn.disabled=true;btn.textContent='\u23F3 Testing...';}
  var body=document.getElementById('devModalBody');
  body.innerHTML='<button class="run-test-btn" disabled>\u23F3 Testing (~20s)...</button><div style="color:var(--tx3);font-size:12px;margin-bottom:12px">Pinging every model simultaneously...</div>';
  try{var r=await apiFetch('/api/test');if(!r.ok){showToast('Test failed: '+r.status);return;}devTestResults=await r.json();renderModelsTab(body);}
  catch(e){body.innerHTML='<div style="color:var(--red);padding:20px">Test error: '+escHtml(e.message)+'</div>';}
}
function renderActivityTab(body){
  var s=devStats; var h='<div class="dev-list"><h5>\u{1F4AC} Recent Messages</h5>';
  if(s.recent_messages&&s.recent_messages.length){
    s.recent_messages.forEach(function(m){
      var rc=m.role==='user'?'var(--blue)':'var(--green)';
      h+='<div class="dev-row" style="flex-wrap:wrap;gap:4px"><span class="k" style="color:'+rc+'">'+escHtml(m.username)+' ['+m.role+']</span><span class="v" style="flex:1;min-width:0;overflow:hidden;text-overflow:ellipsis;white-space:nowrap">'+escHtml(m.content)+'</span>'+(m.model?'<span style="font-size:9px;color:var(--tx3);white-space:nowrap">'+escHtml(m.model)+'</span>':'')+'</div>';
    });
  }else h+='<div style="font-size:12px;color:var(--tx3)">No messages yet</div>';
  h+='</div>'; body.innerHTML=h;
}
function renderVisitorsTab(body){
  var s=devStats; var h='<div class="dev-list"><h5>\u{1F4CD} Visitor Log</h5>';
  if(s.visitor_log&&s.visitor_log.length){
    s.visitor_log.forEach(function(v){h+='<div class="dev-row" style="flex-wrap:wrap"><span class="k">'+escHtml(v.username||'?')+'</span><span style="color:var(--blue);font-size:11px;min-width:70px">'+escHtml(v.event)+'</span><span style="color:var(--tx3);font-size:10px;font-family:monospace">'+escHtml(v.ip||'')+'</span><span style="color:var(--tx3);font-size:10px;margin-left:auto">'+v.ts.slice(0,16).replace('T',' ')+'</span></div>';});
  }else h+='<div style="font-size:12px;color:var(--tx3)">No data yet</div>';
  h+='</div><button class="ghost-btn" style="font-size:12px;padding:8px;margin-top:8px" onclick="refreshDevStats()">\u{1F504} Refresh</button>';body.innerHTML=h;
}
async function refreshDevStats(){devStats=null;await fetchDevStats();renderModalTab(currentModalTab);showToast('\u2705 Refreshed');}

function openCamera(){document.getElementById('cameraInput').click();}
function handleImageFile(e){
  var file=e.target.files[0]; if(!file) return;
  var reader=new FileReader();
  reader.onload=function(ev){var data=ev.target.result,comma=data.indexOf(',');pendingImageB64=data.slice(comma+1);pendingImageMime=file.type||'image/jpeg';pendingImageName=file.name;document.getElementById('imgPrev').src=data;document.getElementById('imgPrevName').textContent=file.name;document.getElementById('imgPrevWrap').classList.add('show');clearFile();};
  reader.readAsDataURL(file); e.target.value='';
}
function clearImage(){pendingImageB64='';pendingImageMime='';pendingImageName='';document.getElementById('imgPrevWrap').classList.remove('show');document.getElementById('imgPrev').src='';}
function fmtSize(b){if(b<1024)return b+'B';if(b<1048576)return (b/1024).toFixed(1)+'KB';return (b/1048576).toFixed(1)+'MB';}
var FICONS={pdf:'\u{1F4D5}',doc:'\u{1F4D8}',docx:'\u{1F4D8}',txt:'\u{1F4C4}',md:'\u{1F4DD}',py:'\u{1F40D}',js:'\u{1F4DC}',ts:'\u{1F4DC}',json:'\u{1F4CB}',csv:'\u{1F4CA}',xls:'\u{1F4CA}',xlsx:'\u{1F4CA}',zip:'\u{1F4E6}'};
function getFileIcon(n){return FICONS[(n.split('.').pop()||'').toLowerCase()]||'\u{1F4C4}';}
async function handleFileUpload(e){
  var file=e.target.files[0]; if(!file) return; clearImage();
  document.getElementById('fileIcon').textContent=getFileIcon(file.name);document.getElementById('fileName').textContent=file.name;
  document.getElementById('fileSize').textContent=fmtSize(file.size);document.getElementById('filePrevWrap').classList.add('show');
  pendingFileName=file.name; pendingFileSize=file.size;
  var textExts=['txt','md','py','js','ts','jsx','tsx','json','xml','csv','html','css','yml','yaml','toml','sh','sql','rs','go','java','cpp','c','h','php','rb','swift','kt'];
  var ext=(file.name.split('.').pop()||'').toLowerCase();
  var isText=file.type.startsWith('text/')||file.type==='application/json'||textExts.includes(ext);
  if(isText){var reader=new FileReader();reader.onload=function(ev){
    pendingFileText=ev.target.result;
    window._lastAttachedFile={name:file.name,size:file.size,preview:ev.target.result.slice(0,400)};
    showToast('📄 '+file.name+' ready');
  };reader.readAsText(file);}
  else{pendingFileText='[Binary file: '+file.name+', size: '+fmtSize(file.size)+']';showToast('\u{1F4CE} '+file.name+' attached');}
  e.target.value='';
}
function clearFile(){pendingFileText='';pendingFileName='';pendingFileSize=0;document.getElementById('filePrevWrap').classList.remove('show');}

function toggleVoice(){if(isRecording)stopVoice();else startVoice();}
function startVoice(){
  var SR=window.SpeechRecognition||window.webkitSpeechRecognition; if(!SR){showErrOverlay('Voice not supported. Try Chrome or Edge.');return;}
  recognition=new SR();recognition.continuous=false;recognition.interimResults=true;recognition.lang='en-US';
  recognition.onstart=function(){isRecording=true;document.getElementById('micBtn').classList.add('active');document.getElementById('voiceBar').classList.add('show');document.getElementById('voiceTxt').textContent='Listening\u2026';};
  recognition.onresult=function(e){var final='',interim='';for(var i=e.resultIndex;i<e.results.length;i++){if(e.results[i].isFinal)final+=e.results[i][0].transcript;else interim+=e.results[i][0].transcript;}document.getElementById('voiceTxt').textContent=interim||'Listening\u2026';if(final){var ta=document.getElementById('msgIn');ta.value=(ta.value+' '+final).trim();ar(ta);}};
  recognition.onerror=function(e){stopVoice();if(e.error!=='aborted')showErrOverlay('Mic error: '+e.error);};
  recognition.onend=function(){stopVoice();};recognition.start();
}
function stopVoice(){isRecording=false;document.getElementById('micBtn').classList.remove('active');document.getElementById('voiceBar').classList.remove('show');if(recognition){try{recognition.stop();}catch(e){}recognition=null;}}

var _voicesLoaded=false;
function loadVoices(){
  var sel=document.getElementById('voiceSelect'); if(!sel) return;
  var voices=window.speechSynthesis?window.speechSynthesis.getVoices():[];
  if(!voices.length&&!_voicesLoaded){window.speechSynthesis.onvoiceschanged=function(){_voicesLoaded=true;loadVoices();};return;}
  sel.innerHTML='<option value="">\u{1F507} Off (no voice)</option>';
  voices.forEach(function(v,i){var opt=document.createElement('option');opt.value=i;opt.textContent=v.name+' ('+v.lang+')';sel.appendChild(opt);});
  if(voicePref.voiceIdx!==undefined) sel.value=voicePref.voiceIdx;
  if(voicePref.rate!==undefined){var r=document.getElementById('voiceRate');if(r){r.value=voicePref.rate;updateRateLabel(r);}}
  ttsEnabled=!!(voicePref.voiceIdx!==undefined&&voicePref.voiceIdx!=='');
}
function saveVoicePref(){
  var sel=document.getElementById('voiceSelect'),rate=document.getElementById('voiceRate');
  voicePref={voiceIdx:sel?sel.value:'',rate:rate?parseFloat(rate.value):1};ttsEnabled=!!(voicePref.voiceIdx!==''&&voicePref.voiceIdx!==undefined);
  localStorage.setItem('fusion_voice',JSON.stringify(voicePref));
  var st=document.getElementById('voiceStatus');if(st){st.textContent=ttsEnabled?'Enabled':'Off';st.style.display=ttsEnabled?'':'none';st.className='kstat '+(ttsEnabled?'set':'unset');}
}
function updateRateLabel(el){var lb=document.getElementById('voiceRateLabel');if(lb)lb.textContent=parseFloat(el.value).toFixed(1)+'x';}
function previewVoice(){
  var voices=window.speechSynthesis?window.speechSynthesis.getVoices():[];
  var sel=document.getElementById('voiceSelect');var idx=sel?parseInt(sel.value):NaN;
  if(isNaN(idx)||!sel||sel.value===''){showToast('Select a voice first');return;}
  var utt=new SpeechSynthesisUtterance('Hello! I am Fusion AI, ready to help.');
  utt.voice=voices[idx];utt.rate=voicePref.rate||1;window.speechSynthesis.cancel();window.speechSynthesis.speak(utt);
}
function stopSpeech(){if(window.speechSynthesis)window.speechSynthesis.cancel();}
function speakText(text){
  if(!window.speechSynthesis) return;
  var voices=window.speechSynthesis.getVoices();
  var idx=ttsEnabled?parseInt(voicePref.voiceIdx):(!isNaN(parseInt(vcPref.voiceIdx))?parseInt(vcPref.voiceIdx):-1);
  var rate=ttsEnabled?(voicePref.rate||1):(vcPref.rate||1);
  if(idx<0||isNaN(idx)) return;
  var clean=text.replace(/```[\s\S]*?```/g,'').replace(/`[^`]*`/g,'').replace(/\*\*([^*]*)\*\*/g,'$1').replace(/\*([^*]*)\*/g,'$1').replace(/<[^>]*>/g,'').trim();
  if(!clean) return;
  var utt=new SpeechSynthesisUtterance(clean.slice(0,1000));
  if(voices[idx]) utt.voice=voices[idx];utt.rate=rate; window.speechSynthesis.cancel(); window.speechSynthesis.speak(utt);
}
function speakForVC(text){
  if(!window.speechSynthesis||!vcActive) return;
  var voices=window.speechSynthesis.getVoices();
  var sel=document.getElementById('vcVoiceSelect');var idx=sel?parseInt(sel.value):NaN;
  var rate=vcPref.rate||1;
  var clean=text.replace(/```[\s\S]*?```/g,'').replace(/`[^`]*`/g,'').replace(/\*\*([^*]*)\*\*/g,'$1').replace(/\*([^*]*)\*/g,'$1').replace(/<[^>]*>/g,'').trim();
  if(!clean) return;
  var utt=new SpeechSynthesisUtterance(clean.slice(0,1200));
  if(!isNaN(idx)&&voices[idx]) utt.voice=voices[idx];utt.rate=rate; window.speechSynthesis.cancel(); window.speechSynthesis.speak(utt);
}

function showToast(msg,dur){dur=dur||2600;var t=document.getElementById('toast');t.textContent=msg;t.classList.add('show');setTimeout(function(){t.classList.remove('show');},dur);}
var errTimer=null;
function showErrOverlay(msg,dur){dur=dur===undefined?7000:dur;var el=document.getElementById('errOverlay');document.getElementById('errOverlayMsg').textContent=msg;el.classList.add('show');clearTimeout(errTimer);if(dur>0)errTimer=setTimeout(function(){el.classList.remove('show');},dur);}
function hideErrOverlay(){document.getElementById('errOverlay').classList.remove('show');}

// Dev sees technical error, regular users see a friendly message
function showUserError(techMsg, friendlyMsg, dur) {
  var msg = isDevUser ? (techMsg||friendlyMsg) : (friendlyMsg||'Something went wrong. Please try again.');
  showErrOverlay(msg, dur);
}
function userErrorToast(techMsg, friendlyMsg, dur) {
  var msg = isDevUser ? (techMsg||friendlyMsg) : (friendlyMsg||'Something went wrong. Please try again.');
  showToast(msg, dur||5000);
}
function ar(el){el.style.height='auto';el.style.height=Math.min(el.scrollHeight,140)+'px';}
function handleKey(e){if(e.key==='Enter'&&!e.shiftKey){e.preventDefault();sendMsg();}}
function suggest(t){document.getElementById('msgIn').value=t;sendMsg();}
function scrollDown(){document.getElementById('chatBody').scrollTop=document.getElementById('chatBody').scrollHeight;}

// ══ Fusion File Generation ═══════════════════════════════════════════════════
var _ffc_urls=[];
var _ffc_content_map={};  // cardId → raw content (avoids giant DOM attrs)

// Map extension → {mime, isBinary, serverEndpoint, previewType}
var _FMT={
  // Text
  md:   {mime:'text/markdown',       binary:false, prev:'md'},
  txt:  {mime:'text/plain',          binary:false, prev:'text'},
  html: {mime:'text/html',           binary:false, prev:'html'},
  htm:  {mime:'text/html',           binary:false, prev:'html'},
  py:   {mime:'text/x-python',       binary:false, prev:'code',lang:'python'},
  js:   {mime:'text/javascript',     binary:false, prev:'code',lang:'javascript'},
  ts:   {mime:'text/typescript',     binary:false, prev:'code',lang:'typescript'},
  jsx:  {mime:'text/javascript',     binary:false, prev:'code',lang:'jsx'},
  tsx:  {mime:'text/typescript',     binary:false, prev:'code',lang:'tsx'},
  css:  {mime:'text/css',            binary:false, prev:'code',lang:'css'},
  json: {mime:'application/json',    binary:false, prev:'json'},
  yaml: {mime:'text/yaml',           binary:false, prev:'code',lang:'yaml'},
  yml:  {mime:'text/yaml',           binary:false, prev:'code',lang:'yaml'},
  xml:  {mime:'application/xml',     binary:false, prev:'code',lang:'xml'},
  csv:  {mime:'text/csv',            binary:false, prev:'csv'},
  sql:  {mime:'text/plain',          binary:false, prev:'code',lang:'sql'},
  sh:   {mime:'text/plain',          binary:false, prev:'code',lang:'bash'},
  bash: {mime:'text/plain',          binary:false, prev:'code',lang:'bash'},
  rs:   {mime:'text/plain',          binary:false, prev:'code',lang:'rust'},
  go:   {mime:'text/plain',          binary:false, prev:'code',lang:'go'},
  java: {mime:'text/plain',          binary:false, prev:'code',lang:'java'},
  c:    {mime:'text/plain',          binary:false, prev:'code',lang:'c'},
  cpp:  {mime:'text/plain',          binary:false, prev:'code',lang:'cpp'},
  r:    {mime:'text/plain',          binary:false, prev:'code',lang:'r'},
  rb:   {mime:'text/plain',          binary:false, prev:'code',lang:'ruby'},
  php:  {mime:'text/plain',          binary:false, prev:'code',lang:'php'},
  swift:{mime:'text/plain',          binary:false, prev:'code',lang:'swift'},
  kt:   {mime:'text/plain',          binary:false, prev:'code',lang:'kotlin'},
  env:  {mime:'text/plain',          binary:false, prev:'text'},
  toml: {mime:'text/plain',          binary:false, prev:'code',lang:'toml'},
  ini:  {mime:'text/plain',          binary:false, prev:'text'},
  conf: {mime:'text/plain',          binary:false, prev:'text'},
  // Binary (server-generated)
  pptx: {mime:'application/vnd.openxmlformats-officedocument.presentationml.presentation', binary:true, server:'/api/gen_pptx', prev:'pptx'},
  xlsx: {mime:'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',         binary:true, server:'/api/gen_xlsx', prev:'csv'},
  docx: {mime:'application/vnd.openxmlformats-officedocument.wordprocessingml.document',   binary:true, server:'/api/gen_docx', prev:'md'},
};
// pptx.md is a special case: AI writes markdown, we convert to PPTX
var _PPTXMD='pptx.md';

function _extOf(name){
  if(name.endsWith('.pptx.md')) return 'pptx.md';
  return (name.split('.').pop()||'txt').toLowerCase();
}

function _prevHTML(ext,content){
  var info=_FMT[ext]||{prev:'text'};
  var prevType=info.prev;
  var trimmed=content.trim();
  var MAX_ROWS=10, MAX_CHARS=600, MAX_LINES=12;

  if(prevType==='csv'||ext==='csv'){
    var rows=trimmed.split('\n').slice(0,MAX_ROWS);
    var total=trimmed.split('\n').length;
    var tbl='<table class="ffc-table">';
    rows.forEach(function(row,ri){
      var cells=row.match(/("(?:[^"]|"")*"|[^,]*),?/g)||[];
      cells=cells.filter(function(c){return c!==undefined;});
      // Proper CSV parse (handles quoted)
      var parsed=[];
      var inq=false,cur='',i=0;
      var r2=row;
      for(;i<r2.length;i++){
        var ch=r2[i];
        if(ch==='"'){if(inq&&r2[i+1]==='"'){cur+='"';i++;}else{inq=!inq;}}
        else if(ch===','&&!inq){parsed.push(cur);cur='';}
        else{cur+=ch;}
      }
      parsed.push(cur);
      tbl+='<tr>';
      parsed.forEach(function(cell){
        var t=ri===0?'th':'td';
        tbl+='<'+t+'>'+escHtml(cell.replace(/^"|"$/g,'').trim())+'</'+t+'>';
      });
      tbl+='</tr>';
    });
    tbl+='</table>';
    return '<div class="ffc-preview">'+tbl+(total>MAX_ROWS?'<div class="ffc-more">+' +(total-MAX_ROWS)+' more rows</div>':'')+'</div>';
  }
  if(prevType==='html'){
    return '<div class="ffc-preview"><div class="ffc-html-prev"><iframe srcdoc="'+escHtml(trimmed).replace(/"/g,'&quot;')+'" sandbox="allow-scripts" style="width:100%;height:180px;border:none;border-radius:6px;background:#fff"></iframe></div></div>';
  }
  if(prevType==='json'){
    try{
      var parsed2=JSON.parse(trimmed);
      var pretty=JSON.stringify(parsed2,null,2);
      return '<div class="ffc-preview"><pre class="ffc-code-pre ffc-json">'+escHtml(pretty.slice(0,MAX_CHARS))+(pretty.length>MAX_CHARS?'\n…':'')+'</pre></div>';
    }catch(e){/*fall through*/}
  }
  if(prevType==='md'){
    // Render markdown basics
    var rendered=trimmed.slice(0,800);
    rendered=rendered.replace(/^#{1,3} (.+)$/mg,'<strong>$1</strong>');
    rendered=rendered.replace(/\*\*(.+?)\*\*/g,'<strong>$1</strong>');
    rendered=rendered.replace(/\*(.+?)\*/g,'<em>$1</em>');
    rendered=rendered.replace(/`([^`]+)`/g,'<code style="background:rgba(6,182,212,.15);padding:1px 5px;border-radius:4px;font-family:monospace;font-size:11px">$1</code>');
    rendered=rendered.replace(/\n/g,'<br>');
    return '<div class="ffc-preview"><div class="ffc-md-prev">'+rendered+(trimmed.length>800?'<br><em style="color:rgba(130,130,180,.6);font-size:10px">… more content</em>':'')+'</div></div>';
  }
  if(prevType==='pptx'){
    // Show slide count and first slide
    var slides=trimmed.split(/\n---+\n|^---+$/m);
    var firstSlide=slides[0]||'';
    var title=(firstSlide.match(/^#{1,2}\s+(.+)$/m)||[])[1]||'Slide 1';
    return '<div class="ffc-preview"><div class="ffc-pptx-prev">📊 <strong>'+escHtml(title)+'</strong><div class="ffc-more">'+slides.length+' slides · Real .pptx generated on download</div></div></div>';
  }
  // Default: code/text
  var lines=trimmed.split('\n').slice(0,MAX_LINES).join('\n');
  var total2=trimmed.split('\n').length;
  return '<div class="ffc-preview"><pre class="ffc-code-pre">'+escHtml(lines)+(total2>MAX_LINES?'\n…':'')+'</pre></div>';
}

function _genFileCard(name, content){
  var ext=_extOf(name);
  // For pptx.md: the output file is .pptx (binary), content is markdown
  var dlExt=(ext==='pptx.md')?'pptx':ext;
  var dlName=name.endsWith('.pptx.md')?name.replace(/\.pptx\.md$/,'.pptx'):name;

  var ICONS={csv:'📊',json:'📋',md:'📄',txt:'📄',html:'🌐',htm:'🌐',py:'🐍',js:'📜',ts:'📜',
    jsx:'📜',tsx:'📜',sql:'🗄',xml:'📰',yaml:'⚙️',yml:'⚙️',sh:'🖥',bash:'🖥',pptx:'🎯',
    xlsx:'📊',docx:'📘',css:'🎨',rs:'🦀',go:'🐹',java:'☕',c:'🔧',cpp:'🔧',r:'📈',
    rb:'💎',php:'🐘',swift:'🐦',kt:'📱',toml:'⚙️',ini:'⚙️',conf:'⚙️','pptx.md':'🎯'};
  var icon=ICONS[ext]||'📄';
  var info=_FMT[dlExt]||{binary:false,mime:'text/plain'};
  var isBinary=info.binary||(ext==='pptx.md');
  var lines=content.trim().split('\n').length;
  var sizeStr=content.length>1048576?(content.length/1048576).toFixed(1)+'MB':content.length>1024?(content.length/1024).toFixed(1)+'KB':content.length+'B';

  // Preview HTML
  var previewExt=(ext==='pptx.md')?'pptx':ext;
  var prevHtml=_prevHTML(previewExt, content);

  // Build card
  var cardId='ffc-'+Math.random().toString(36).slice(2,8);
  var badgeHtml=isBinary?'<span style="font-size:9px;background:rgba(6,182,212,.2);border:1px solid rgba(6,182,212,.4);border-radius:4px;padding:1px 5px;color:#80c0ff;margin-left:6px;vertical-align:middle">Binary</span>':'';

  var cardHtml='<div class="fusion-file-card" id="'+cardId+'">'
    +'<div class="ffc-header">'
    +'<span class="ffc-icon">'+icon+'</span>'
    +'<div class="ffc-info">'
    +'<div class="ffc-name">'+escHtml(dlName)+badgeHtml+'</div>'
    +'<div class="ffc-meta">'+dlExt.toUpperCase()+' · '+sizeStr+' · '+lines+' lines · AI-generated</div>'
    +'</div>';

  if(isBinary){
    // Binary: download via server endpoint
    var serverEp=dlExt==='pptx'?'/api/gen_pptx':dlExt==='xlsx'?'/api/gen_xlsx':'/api/gen_docx';
    _ffc_content_map[cardId]=content.trim();
  var themeExtra=(dlExt==='pptx'?', \'dark\'':'');
  cardHtml+='<button class="ffc-dl" onclick="_ffcBinaryDl(this,\''+cardId+'\',\''+serverEp+'\',\''+escHtml(dlName)+'\''+themeExtra+')"' +'>⬇ Download '+dlExt.toUpperCase()+'</button>';
  } else {
    // Text: blob URL
    _ffc_content_map[cardId]=content.trim();
    var blob=new Blob([content.trim()],{type:(info.mime||'text/plain')+';charset=utf-8'});
    var blobUrl=URL.createObjectURL(blob);
    _ffc_urls.push(blobUrl);
    cardHtml+='<a class="ffc-dl" href="'+blobUrl+'" download="'+escHtml(dlName)+'">⬇ Download</a>';
  }

  cardHtml+='<button class="ffc-copy" onclick="_ffc_copy(this)" data-c="'+encodeURIComponent(content.trim())+'">📋 Copy</button>'
    +'</div>'
    +prevHtml
    +'</div>';
  return cardHtml;
}

async function _ffcBinaryDl(btn,cardId,endpoint,filename,theme){
  var content=_ffc_content_map[cardId]||'';
  if(!content){showToast('⚠️ Content not found — try regenerating');return;}
  btn.textContent='⏳ Generating…'; btn.disabled=true;
  try{
    var resp=await fetch(endpoint,{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({content:content,filename:filename,theme:theme||'dark'})});
    if(!resp.ok){var err=await resp.json(); throw new Error(err.error||'Server error');}
    var blob=await resp.blob();
    var url=URL.createObjectURL(blob);
    var a=document.createElement('a'); a.href=url; a.download=filename; a.click();
    setTimeout(function(){URL.revokeObjectURL(url);},5000);
    btn.textContent='✅ Downloaded!';
    setTimeout(function(){btn.textContent='⬇ Download '+filename.split('.').pop().toUpperCase(); btn.disabled=false;},2500);
  }catch(ex){
    btn.textContent='❌ '+ex.message; btn.disabled=false;
    setTimeout(function(){btn.textContent='⬇ Download '+filename.split('.').pop().toUpperCase();btn.disabled=false;},3000);
    showToast('⚠️ '+ex.message);
  }
}

function _ffc_copy(btn){
  var cardEl=btn.closest('.fusion-file-card');
  var cid=cardEl?cardEl.id:'';
  var content=(cid&&_ffc_content_map[cid])||decodeURIComponent(btn.dataset.c||'');
  navigator.clipboard.writeText(content).then(function(){
    btn.textContent='✅ Copied!';
    setTimeout(function(){btn.textContent='📋 Copy';},2200);
  });
}

function _parseFusionFilesRaw(text){
  return text.replace(/<<<FUSIONFILE:([\w.\-]+)>>>([\s\S]*?)<<<END_FUSIONFILE>>>/g,function(_,name,content){
    return _genFileCard(name,content);
  });
}

function escHtml(t){return String(t).replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;').replace(/"/g,'&quot;');}
function toggleThought(hdr){var body=hdr.nextElementSibling;var tog=hdr.querySelector('.thought-toggle');var lbl=hdr.querySelector('.thought-label');var isOpen=body.classList.toggle('open');tog.classList.toggle('open',isOpen);lbl.textContent=isOpen?'Hide thought process':'View thought process';}

function _renderMath(el){
  if(!el) return;
  // Wait for KaTeX to be ready — retry up to 20 times
  if(typeof window.katex==='undefined'||typeof renderMathInElement==='undefined'){
    if(!(el._mathRetry)) el._mathRetry=0;
    if(el._mathRetry<20){
      el._mathRetry++;
      setTimeout(function(){_renderMath(el);},200);
    }
    return;
  }
  // Don't re-render elements that already have katex output (streaming safe)
  if(el.querySelector('.katex')) return;
  try{
    renderMathInElement(el,{
      delimiters:[
        {left:'$$',right:'$$',display:true},
        {left:'\\\\[',right:'\\\\]',display:true},
        {left:'\\\\(',right:'\\\\)',display:false},
        {left:'$',right:'$',display:false}
      ],
      ignoredTags:['script','noscript','style','textarea','pre','code','option'],
      ignoredClasses:['hlcode','inline-code','no-math'],
      throwOnError:false,
      strict:false,
      trust:true,
      output:'html'
    });
    el._mathRetry=0;
  }catch(err){}
}

// ── KaTeX ready poller ─────────────────────────────────────────────────────
window._katexCheckInterval = setInterval(function(){
  if(typeof window.katex !== 'undefined' && typeof renderMathInElement !== 'undefined'){
    clearInterval(window._katexCheckInterval);
    // Render any existing math in the page
    document.querySelectorAll('.bbl').forEach(function(el){
      if(!el.querySelector('.katex')) _renderMath(el);
    });
  }
}, 300);

// Force re-render math (used after streaming completes — clears katex blocks first)
function _rerenderMath(el){
  if(!el) return;
  el._mathRetry=0;
  // Remove existing katex to re-render cleanly
  el.querySelectorAll('.katex-html').forEach(function(k){
    var parent=k.closest('.katex');
    if(parent&&parent.parentNode) parent.outerHTML=parent.querySelector('.katex-mathml')&&parent.querySelector('.katex-mathml').textContent||parent.innerText;
  });
  _renderMath(el);
}

function fmt(text){
  var t=String(text);
  // ── Parse FUSIONFILE blocks FIRST (before any escaping) ──────────────
  if(t.indexOf('<<<FUSIONFILE:')>=0){
    t=t.replace(/<<<FUSIONFILE:([\w.\-]+)>>>([\s\S]*?)<<<END_FUSIONFILE>>>/g,function(_,name,content){
      return _genFileCard(name,content);
    });
  }
  // Strip stray leading hashes/lines
  t=t.replace(/^(?:[#\-=]{3,}\s*\n)+/,'');
  // Protect math blocks FIRST — before any other processing
  var _mb=[],_mi=[];
  // Display math: $$...$$ and \[...\]
  t=t.replace(/\$\$([\s\S]+?)\$\$/g,function(m){_mb.push(m);return '⁠MB'+((_mb.length-1))+'⁠';});
  t=t.replace(/\\\[([\s\S]+?)\\\]/g,function(m){_mb.push(m);return '⁠MB'+((_mb.length-1))+'⁠';});
  // Inline math: \(...\) and $...$  (not $$)
  t=t.replace(/\\\(([^)]{1,400}?)\\\)/g,function(m){_mi.push(m);return '⁠MI'+((_mi.length-1))+'⁠';});
  t=t.replace(/\$([^$\n`]{1,300}?)\$/g,function(m,inner){
    // Only treat as math if it looks like math (contains math chars)
    if(/[\\^_{}\[\]()=+\-*\/\\|<>]/.test(inner)||/[a-zA-Z]{1,3}\d/.test(inner)){
      _mi.push(m);return '⁠MI'+((_mi.length-1))+'⁠';
    }
    return m;
  });
  // Code blocks with syntax highlighting and collapsible for long code
  t=t.replace(/```([\w]*)\n?([\s\S]{0,50000}?)```/g,function(m,lang,code){
    var trimmed=code.trim();
    var lines=trimmed.split('\n').length;
    var langLabel=lang?lang.toUpperCase():'CODE';
    var highlighted=_syntaxHL(trimmed,lang);
    // ── Very large blocks: compact card + floating annotated window ──────
    if(lines>60){
      var wid='cw_'+Math.random().toString(36).slice(2,9);
      _codeWindowStore[wid]={code:trimmed,lang:lang||'text',label:langLabel};
      var preview=trimmed.split('\n').slice(0,3).join('\n');
      var outline=_codeOutline(trimmed,lang);
      return '<div class="code-card" onclick="_openCodeWindow(\''+wid+'\')">'
        +'<div class="code-card-top">'
        +'<span class="code-card-icon">📄</span>'
        +'<div class="code-card-info"><div class="code-card-title">'+escHtml(langLabel)+' — '+lines+' lines</div>'
        +'<div class="code-card-sub">'+(outline.length?escHtml(outline.slice(0,3).join(' · ')):'Click to open in a window')+'</div></div>'
        +'<span class="code-card-open">⤢ Open</span>'
        +'</div>'
        +'<div class="code-card-preview"><pre>'+escHtml(preview)+'\n…</pre></div>'
        +'</div>';
    }
    if(lines>20){
      var cid='cb_'+Math.random().toString(36).slice(2,8);
      return '<div class="code-block-wrap">'
        +'<div class="code-block-header">'
        +'<span class="code-lang-badge">'+escHtml(langLabel)+'</span>'
        +'<span class="code-lines">'+lines+' lines</span>'
        +'<button class="code-copy-btn" onclick="_copyCode(this)">📋 Copy</button>'
        +'<button class="code-expand-btn" onclick="_toggleCode(this,\''+cid+'\')">▼ Expand</button>'
        +'</div>'
        +'<div class="code-block-body code-collapsed" id="'+cid+'">'
        +'<pre><code class="hlcode lang-'+escHtml(lang||'text')+'">'+highlighted+'</code></pre>'
        +'</div></div>';
    }
    return '<div class="code-block-wrap">'
      +'<div class="code-block-header">'
      +'<span class="code-lang-badge">'+escHtml(langLabel)+'</span>'
      +'<span class="code-lines">'+lines+' line'+(lines===1?'':'s')+'</span>'
      +'<button class="code-copy-btn" onclick="_copyCode(this)">📋 Copy</button>'
      +'</div>'
      +'<pre><code class="hlcode lang-'+escHtml(lang||'text')+'">'+highlighted+'</code></pre>'
      +'</div>';
  });
  // Tables
  t=t.replace(/^\|(.+)\|\s*$/mg,function(row){
    var cells=row.split('|').filter(function(c,i,a){return i>0&&i<a.length-1;});
    if(/^[-:\s|]+$/.test(row)) return '<tr class="thr-sep"></tr>';
    return '<tr>'+cells.map(function(c){return '<td style="padding:5px 10px;border:1px solid var(--glass-bdr);font-size:13px">'+c.trim()+'</td>';}).join('')+'</tr>';
  });
  t=t.replace(/(<tr>.*<\/tr>)+/gs,function(m){return '<table style="border-collapse:collapse;margin:8px 0;width:100%">'+m+'</table>';});
  t=t.replace(/<tr class="thr-sep"><\/tr>/g,'');
  // Horizontal rule
  t=t.replace(/^---+$/mg,'<hr style="border:none;border-top:1px solid var(--glass-bdr);margin:10px 0">');
  // Headings
  t=t.replace(/^### (.+)$/mg,'<strong style="font-size:14px;display:block;margin:8px 0 4px">$1</strong>');
  t=t.replace(/^## (.+)$/mg,'<strong style="font-size:15px;display:block;margin:10px 0 5px">$1</strong>');
  t=t.replace(/^# (.+)$/mg,'<strong style="font-size:17px;display:block;margin:12px 0 6px">$1</strong>');
  // Bullet lists
  t=t.replace(/^[\*\-] (.+)$/mg,'<li style="margin:3px 0 3px 16px">$1</li>');
  t=t.replace(/(<li[^>]*>.*<\/li>\n?)+/gs,function(m){return '<ul style="margin:6px 0;padding:0;list-style:disc">'+m+'</ul>';});
  // Numbered lists
  t=t.replace(/^\d+\. (.+)$/mg,'<li style="margin:3px 0 3px 16px">$1</li>');
  t=t.replace(/(<li[^>]*>.*<\/li>\n?)+/gs,function(m){return '<ul style="margin:6px 0;padding:0">'+m+'</ul>';});
  t=t.replace(/`([^\`\n]{1,400})`/g,function(m,code){return '<code class="inline-code">'+escHtml(code)+'</code>';});
  t=t.replace(/\*\*([^\n*]{1,500})\*\*/g,'<strong>$1</strong>');t=t.replace(/\*([^\n*]{1,300})\*/g,'<em>$1</em>');
  // Convert newlines first
  t=t.replace(/\n/g,'<br>');
  // Restore math LAST — so it lands in the DOM unmolested for KaTeX
  t=t.replace(/⁠MB(\d+)⁠/g,function(_,i){return _mb[parseInt(i)]||'';});
  t=t.replace(/⁠MI(\d+)⁠/g,function(_,i){return _mi[parseInt(i)]||'';});
  return t;
}

function _syntaxHL(code, lang){
  // Escape HTML first — all replacements work on escaped text
  var e=escHtml(code);
  var l=(lang||'').toLowerCase();
  // Placeholder system: wrap spans with a unique marker so later regexes don't re-process them
  var _spans=[]; var _sc=0;
  function wrap(cls,m){ var ph='«'+(++_sc)+'»'; _spans.push('<span class="'+cls+'">'+m+'</span>'); return ph; }
  function restore(s){ return s.replace(/«(\d+)»/g,function(_,i){return _spans[parseInt(i)-1]||_;}); }

  if(['python','py'].includes(l)){
    e=e.replace(/(#[^\n]*)/g,function(_,m){return wrap('hl-cmt',m);});
    e=e.replace(/(&#39;&#39;&#39;[\s\S]*?&#39;&#39;&#39;|&quot;&quot;&quot;[\s\S]*?&quot;&quot;&quot;)/g,function(_,m){return wrap('hl-str',m);});
    e=e.replace(/(&quot;(?:[^&]|&(?!quot;))*&quot;|&#39;(?:[^&]|&(?!#39;))*&#39;)/g,function(_,m){return wrap('hl-str',m);});
    e=e.replace(/\b(def|class|import|from|return|if|elif|else|for|while|try|except|finally|with|as|in|not|and|or|True|False|None|pass|break|continue|raise|yield|lambda|async|await|del|global|nonlocal|assert)\b/g,function(_,m){return wrap('hl-kw',m);});
    e=e.replace(/\b(\d+\.?\d*)\b/g,function(_,m){return wrap('hl-num',m);});
  } else if(['js','javascript','ts','typescript'].includes(l)){
    e=e.replace(/(\/\/[^\n]*)/g,function(_,m){return wrap('hl-cmt',m);});
    e=e.replace(/(\/\*[\s\S]*?\*\/)/g,function(_,m){return wrap('hl-cmt',m);});
    e=e.replace(/(&quot;(?:[^&]|&(?!quot;))*&quot;|&#39;(?:[^&]|&(?!#39;))*&#39;|`[^`]*`)/g,function(_,m){return wrap('hl-str',m);});
    e=e.replace(/\b(const|let|var|function|return|if|else|for|while|do|switch|case|break|continue|class|extends|import|export|default|async|await|new|this|true|false|null|undefined|typeof|instanceof|void|throw|try|catch|finally)\b/g,function(_,m){return wrap('hl-kw',m);});
    e=e.replace(/\b(\d+\.?\d*)\b/g,function(_,m){return wrap('hl-num',m);});
  } else if(['html','xml'].includes(l)){
    e=e.replace(/(&lt;\/?[\w\-]+)/g,function(_,m){return wrap('hl-kw',m);});
    e=e.replace(/(&quot;(?:[^&]|&(?!quot;))*&quot;)/g,function(_,m){return wrap('hl-str',m);});
    e=e.replace(/([\w\-]+=)(?=&quot;)/g,function(_,m){return wrap('hl-attr',m);});
  } else if(['css','scss','less'].includes(l)){
    e=e.replace(/(\/\*[\s\S]*?\*\/)/g,function(_,m){return wrap('hl-cmt',m);});
    e=e.replace(/([\w\-]+\s*)(?=\s*:(?!:))/g,function(_,m){return wrap('hl-attr',m);});
    e=e.replace(/(&quot;(?:[^&]|&(?!quot;))*&quot;|&#39;(?:[^&]|&(?!#39;))*&#39;)/g,function(_,m){return wrap('hl-str',m);});
    e=e.replace(/#(?:[0-9a-fA-F]{3,8})\b/g,function(_,m){return wrap('hl-num',m);});
    e=e.replace(/\b(\d+\.?\d*(?:px|em|rem|vh|vw|%|s|ms|deg)?)\b/g,function(_,m){return wrap('hl-num',m);});
  } else if(['bash','sh','shell','zsh'].includes(l)){
    e=e.replace(/(#[^\n]*)/g,function(_,m){return wrap('hl-cmt',m);});
    e=e.replace(/(&quot;(?:[^&]|&(?!quot;))*&quot;|&#39;(?:[^&]|&(?!#39;))*&#39;)/g,function(_,m){return wrap('hl-str',m);});
    e=e.replace(/\b(echo|printf|cd|ls|pwd|mkdir|rm|cp|mv|grep|sed|awk|cut|sort|uniq|curl|wget|chmod|chown|export|source|if|then|else|elif|fi|for|do|done|while|until|case|esac|function|return|exit|set|unset)\b/g,function(_,m){return wrap('hl-kw',m);});
  } else if(['sql'].includes(l)){
    e=e.replace(/\b(SELECT|FROM|WHERE|JOIN|LEFT|RIGHT|INNER|OUTER|CROSS|FULL|ON|GROUP\s+BY|ORDER\s+BY|HAVING|INSERT|UPDATE|DELETE|CREATE|DROP|ALTER|TABLE|INDEX|AS|AND|OR|NOT|IN|IS|NULL|DISTINCT|LIMIT|OFFSET|SET|VALUES|INTO|BEGIN|COMMIT|ROLLBACK|TRANSACTION)\b/gi,function(_,m){return wrap('hl-kw',m);});
    e=e.replace(/(&quot;(?:[^&]|&(?!quot;))*&quot;|&#39;(?:[^&]|&(?!#39;))*&#39;)/g,function(_,m){return wrap('hl-str',m);});
    e=e.replace(/\b(\d+\.?\d*)\b/g,function(_,m){return wrap('hl-num',m);});
  } else if(['json'].includes(l)){
    e=e.replace(/(&quot;(?:[^&]|&(?!quot;))*&quot;)\s*:/g,function(_,m){return wrap('hl-attr',m.replace(/\s*$/,''))+':';});
    e=e.replace(/:\s*(&quot;(?:[^&]|&(?!quot;))*&quot;)/g,function(_,m){return ': '+wrap('hl-str',m.trim());});
    e=e.replace(/\b(true|false|null)\b/g,function(_,m){return wrap('hl-kw',m);});
    e=e.replace(/\b(\d+\.?\d*)\b/g,function(_,m){return wrap('hl-num',m);});
  }
  return restore(e);
}
function _copyCode(btn){
  var pre=btn.closest('.code-block-wrap').querySelector('pre code');
  if(!pre) return;
  navigator.clipboard.writeText(pre.innerText||pre.textContent).then(function(){
    btn.textContent='✅ Copied!'; setTimeout(function(){btn.textContent='📋 Copy';},2000);
  });
}

function _toggleCode(btn, cid){
  var body=document.getElementById(cid);
  if(!body) return;
  var collapsed=body.classList.toggle('code-collapsed');
  btn.textContent=collapsed?'▼ Expand':'▲ Collapse';
}

// ── Large code block → compact card → floating annotated window ─────────────
var _codeWindowStore={};
function _codeOutline(code,lang){
  var out=[];
  var lines=code.split('\n');
  var patterns=[
    /^\s*(?:export\s+)?(?:async\s+)?function\s+([A-Za-z0-9_$]+)/,
    /^\s*(?:export\s+)?class\s+([A-Za-z0-9_$]+)/,
    /^\s*def\s+([A-Za-z0-9_]+)\s*\(/,
    /^\s*class\s+([A-Za-z0-9_]+)\s*[:\(]/,
    /^\s*(?:public|private|protected|static)?\s*(?:async\s+)?function\s+([A-Za-z0-9_$]+)/,
    /^\s*const\s+([A-Za-z0-9_$]+)\s*=\s*(?:async\s*)?\(/,
    /^\s*func\s+([A-Za-z0-9_]+)\s*\(/,
  ];
  for(var i=0;i<lines.length&&out.length<12;i++){
    for(var j=0;j<patterns.length;j++){
      var m=lines[i].match(patterns[j]);
      if(m&&m[1]){ out.push(m[1]+'()'); break; }
    }
  }
  return out;
}
function _openCodeWindow(wid){
  var entry=_codeWindowStore[wid];
  if(!entry) return;
  var old=document.getElementById('codeWinOverlay'); if(old) old.remove();
  var outline=_codeOutline(entry.code,entry.lang);
  var numbered=entry.code.split('\n').map(function(l,i){return '<span class="cw-lnum">'+(i+1)+'</span>'+escHtml(l);}).join('\n');
  var ov=document.createElement('div');
  ov.id='codeWinOverlay'; ov.className='code-win-overlay';
  ov.innerHTML='<div class="code-win" id="codeWinBox">'
    +'<div class="code-win-head" id="codeWinDrag">'
    +'<span class="code-win-title">📄 '+escHtml(entry.label)+' <span class="code-win-lines">'+entry.code.split('\n').length+' lines</span></span>'
    +'<div class="code-win-actions">'
      +'<button class="cw-btn" onclick="_copyCodeWindow(\''+wid+'\')">📋 Copy</button>'
      +'<button class="cw-btn" onclick="_downloadCodeWindow(\''+wid+'\')">⬇ Save</button>'
      +'<button class="cw-btn cw-close" onclick="_closeCodeWindow()">✕</button>'
    +'</div></div>'
    +(outline.length?'<div class="code-win-outline">'+outline.map(function(o){return '<span class="cw-outline-chip">'+escHtml(o)+'</span>';}).join('')+'</div>':'')
    +'<div class="code-win-body"><pre class="cw-pre"><code>'+numbered+'</code></pre></div>'
  +'</div>';
  document.body.appendChild(ov);
  _makeCodeWinDraggable();
}
function _closeCodeWindow(){ var ov=document.getElementById('codeWinOverlay'); if(ov) ov.remove(); }
function _copyCodeWindow(wid){
  var entry=_codeWindowStore[wid]; if(!entry) return;
  navigator.clipboard.writeText(entry.code).then(function(){ showToast('✅ Code copied'); });
}
function _downloadCodeWindow(wid){
  var entry=_codeWindowStore[wid]; if(!entry) return;
  var extMap={python:'py',javascript:'js',typescript:'ts',bash:'sh',shell:'sh',json:'json',html:'html',css:'css',java:'java',cpp:'cpp',c:'c',go:'go',rust:'rs',ruby:'rb',php:'php',sql:'sql',yaml:'yml'};
  var ext=extMap[(entry.lang||'').toLowerCase()]||'txt';
  var blob=new Blob([entry.code],{type:'text/plain'});
  var a=document.createElement('a'); a.href=URL.createObjectURL(blob); a.download='code.'+ext; a.click();
}
function _makeCodeWinDraggable(){
  var box=document.getElementById('codeWinBox'); var handle=document.getElementById('codeWinDrag');
  if(!box||!handle) return;
  var dragging=false,ox=0,oy=0;
  handle.onmousedown=function(e){
    dragging=true; ox=e.clientX-box.offsetLeft; oy=e.clientY-box.offsetTop;
    e.preventDefault();
  };
  document.onmousemove=function(e){
    if(!dragging) return;
    box.style.left=(e.clientX-ox)+'px'; box.style.top=(e.clientY-oy)+'px';
    box.style.transform='none'; box.style.position='fixed';
  };
  document.onmouseup=function(){ dragging=false; };
}

var WELCOME_HTML='<div class="welcome" id="ws"><div class="wico">&#x26A1;</div><h1>Fusion.AI</h1>'
  +'<p>Your intelligent AI — chat, create documents, build datasets, generate images, video, 3D and more.</p>'
  +'<div id="welcomeWx" style="display:none;font-size:12px;color:var(--tx3);margin-bottom:4px"></div>'
  +'<div class="wsubt">Your AI Powerhouse</div>'
  +'<div class="chips">'
  +'<div class="chip" onclick="suggest(\'Write a Python web scraper for product prices\')">&#x1F40D; Python scraper</div>'
  +'<div class="chip" onclick="suggest(\'Solve: if 3x+7=22 what is x? Show steps\')">&#x1F9EE; Math problem</div>'
  +'<div class="chip" onclick="suggest(\'Write a scary short story in 200 words\')">&#x1F47B; Short story</div>'
  +'<div class="chip" onclick="suggest(\'Give me 5 profitable startup ideas for 2026\')">&#x1F4A1; Startup ideas</div>'
  +'<div class="chip" onclick="suggest(\'Explain quantum computing simply\')">&#x269B;&#xFE0F; Quantum</div>'
  +'<div class="chip" onclick="suggest(\'Compare Groq vs Gemini vs OpenRouter\')">&#x2696;&#xFE0F; Compare AIs</div>'
  +'<div class="chip" onclick="suggest(\'Debug this code: print(Hello World)\')">&#x1F41B; Debug code</div>'
  +'<div class="chip" onclick="suggest(\'What are the best free AI APIs in 2026?\')">&#x1F916; Free AI APIs</div>'
  +'<div class="chip doc-chip" onclick="suggest(\'/csv top 50 global companies with revenue employees and industry\')">&#x1F4CA; CSV Dataset</div>'
  +'<div class="chip doc-chip" onclick="suggest(\'/doc Complete guide to building REST APIs in Python\')">&#x1F4C4; Create Doc</div>'
  +'<div class="chip doc-chip" onclick="suggest(\'/ppt Introduction to Machine Learning for Beginners\')">&#x1F4CA; Make PPT</div>'
  +'</div></div>';

var _chatAbort = null;  // AbortController for current AI stream

function stopChat() {
  if(_chatAbort){ try{_chatAbort.abort();}catch(e){} _chatAbort=null; }
  loading=false;document.getElementById('sbtn').disabled=false;document.getElementById('stopChatBtn').style.display='none';removeThinking();
  showToast('⏹ Stopped');
}

function _chatStreamStart() {
  _chatAbort=new AbortController();document.getElementById('stopChatBtn').style.display='';document.getElementById('sbtn').disabled=true;
}

function _chatStreamEnd() {
  _chatAbort=null;document.getElementById('stopChatBtn').style.display='none';document.getElementById('sbtn').disabled=false;
}

var currentConvId = null;  // active conversation ID (null = no session yet)

function resetChat(){
  hist=[];document.getElementById('msgs').innerHTML=WELCOME_HTML;currentConvId=null;
}
function clearChat(){resetChat();closeDrop();clearImage();clearFile();}

async function loadConvs() {
  var list=document.getElementById('convList');
  if(!list||!authToken) return;
  try{
    var r=await apiFetch('/api/conversations');
    var d=await r.json();
    var convs=d.conversations||[];
    if(!convs.length){list.innerHTML='<div class="empty-state"><span class="ei">💬</span>No saved conversations yet.<br><small style="color:var(--tx3)">Every chat is auto-saved.</small></div>';return;}
    var h='';
    convs.forEach(function(c){
      var ts=c.updated?c.updated.slice(0,10):'';
      var active=c.id===currentConvId?' style="border-color:var(--blue);background:rgba(6,182,212,.1)"':'';
      h+='<div class="saved-item" onclick="loadConv('+c.id+')"'+active+'>'
        +'<div style="display:flex;align-items:center;justify-content:space-between">'
        +'<div class="si-title">'+escHtml(c.title)+'</div>'
        +'<div style="display:flex;gap:4px">'
        +'<button onclick="event.stopPropagation();renameConv('+c.id+',\''+escHtml(c.title.replace(/'/g,"\\'"))+'\')" style="background:none;border:none;color:var(--tx3);cursor:pointer;font-size:12px;padding:2px 4px" title="Rename">✏️</button>'
        +'<button onclick="event.stopPropagation();deleteConv('+c.id+')" style="background:none;border:none;color:var(--tx3);cursor:pointer;font-size:12px;padding:2px 4px" title="Delete">🗑️</button>'
        +'</div></div>'
        +'<div class="si-ts">'+ts+'</div>'
        +'</div>';
    });
    list.innerHTML=h;
  }catch(e){list.innerHTML='<div style="color:var(--tx3);font-size:12px;padding:8px">Failed to load</div>';}
}

async function newConv() {
  // Lazy creation: don't create server-side until first message
  currentConvId=null; hist=[];
  document.getElementById('msgs').innerHTML='';
  showWelcome(); closeSP(); showToast('✨ New conversation — start typing!');
  loadConvs();
}

async function loadConv(cid) {
  try{
    var r=await apiFetch('/api/conversations/'+cid);
    var d=await r.json();
    currentConvId=cid;hist=[];document.getElementById('msgs').innerHTML='';
    document.getElementById('ws') && (document.getElementById('ws').style.display='none');
    showToast('💬 Loading: '+escHtml((d.title||'Chat').slice(0,40)));
    // Render all messages and rebuild full AI context history
    (d.messages||[]).forEach(function(m){
      if(m.role==='user'){
        addMsg('user',m.content,null,null,'');hist.push({role:'user',content:m.content});
      } else if(m.role==='assistant'){
        var bbl=addMsg('ai','',m.model,null);
        bbl.innerHTML=fmt(m.content);_renderMath(bbl);hist.push({role:'assistant',content:m.content});
        // Add message actions
        var acts=document.createElement('div');acts.className='msg-actions';
        acts.innerHTML='<button class="mact" onclick="navigator.clipboard.writeText(this.closest(\'.mcont\').querySelector(\'.bbl\').innerText)">📋 Copy</button>'
          +'<button class="mact" onclick="speakText(this.closest(\'.mcont\').querySelector(\'.bbl\').innerText)">🔊 Speak</button>';
        bbl.closest('.mcont').appendChild(acts);
      }
    });
    loadConvs();closeSP();scrollDown();showToast('💬 Loaded: '+d.title);
  }catch(e){showToast('Failed to load conversation');}
}

async function deleteConv(cid) {
  if(!confirm('Delete this conversation?')) return;
  try{
    await apiFetch('/api/conversations/'+cid,{method:'DELETE'});
    if(currentConvId===cid){resetChat();}
    loadConvs();showToast('🗑️ Deleted');
  }catch(e){showToast('Delete failed');}
}

async function renameConv(cid, currentTitle) {
  var newTitle=prompt('Rename conversation:',currentTitle);
  if(!newTitle||newTitle===currentTitle) return;
  try{
    await apiFetch('/api/conversations/'+cid+'/rename',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({title:newTitle})});
    loadConvs();
  }catch(e){showToast('Rename failed');}
}

async function _ensureConvId(firstMsg) {
  if(currentConvId) return;
  try{
    var title=(firstMsg||'New Chat').slice(0,60).trim()||'New Chat';
    var r=await apiFetch('/api/conversations',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({title:title})});
    var d=await r.json();
    currentConvId=d.id;
    loadConvs();
  }catch(e){}
}

async function _saveConvMsg(role, content, model) {
  if(!currentConvId||!authToken||!content) return;
  try{
    await apiFetch('/api/conversations/'+currentConvId+'/message',{
      method:'POST',headers:{'Content-Type':'application/json'},
      body:JSON.stringify({role:role,content:content,model:model||''})
    });
    loadConvs();
  }catch(e){}
}

function addMsg(role,content,tag,reason,imgSrc){
  var ws=document.getElementById('ws');if(ws)ws.remove();
  var msgs=document.getElementById('msgs');var init=uName?uName[0].toUpperCase():'G';
  var div=document.createElement('div');div.className='msg '+role;
  var tagHtml=tag?'<span class="mtag">'+escHtml(tag)+'</span>':'';
  var imgHtml=imgSrc?'<img class="msg-img" src="'+imgSrc+'" alt="attached"/>':'';
  if(role==='ai'){
    div.innerHTML='<div class="mav">&#x26A1;</div><div class="mcont"><div class="mname">Fusion.AI '+tagHtml+'</div><div class="bbl">'+fmt(content)+'</div>'+(reason?'<div class="rtag">&#x21B3; '+escHtml(reason)+'</div>':'')+'<div class="msg-actions"><button class="mact" onclick="saveMessage(this.closest(\'.mcont\').querySelector(\'.bbl\').innerText)">&#x1F516; Save</button><button class="mact" onclick="navigator.clipboard.writeText(this.closest(\'.mcont\').querySelector(\'.bbl\').innerText)">&#x1F4CB; Copy</button><button class="mact" onclick="speakText(this.closest(\'.mcont\').querySelector(\'.bbl\').innerText)">&#x1F50A; Speak</button></div></div>';
  }else{
    // Build enhanced file preview if there's an attached file
  var filePreviewHtml='';
  if(window._lastAttachedFile&&window._lastAttachedFile.name){
    var af=window._lastAttachedFile;
    var fext=(af.name.split('.').pop()||'').toLowerCase();
    var ficons={pdf:'📕',doc:'📘',docx:'📘',txt:'📄',md:'📝',py:'🐍',js:'📜',ts:'📜',json:'📋',csv:'📊',xls:'📊',xlsx:'📊',zip:'📦',pptx:'📊',png:'🖼',jpg:'🖼',jpeg:'🖼',gif:'🖼',mp4:'🎬',mp3:'🎵'};
    var ficon=ficons[fext]||'📄';
    var fsize=af.size>1048576?(af.size/1048576).toFixed(1)+'MB':af.size>1024?(af.size/1024).toFixed(0)+'KB':af.size+'B';
    filePreviewHtml='<div class="upload-preview-box"><div style="font-size:24px;flex-shrink:0">'+ficon+'</div>'
      +'<div class="upload-preview-meta"><div class="upload-preview-name">'+escHtml(af.name)+'</div>'
      +'<div class="upload-preview-size">'+fsize+' · '+fext.toUpperCase()+'</div>'
      +(af.preview?'<div class="upload-preview-text">'+escHtml(af.preview.slice(0,200))+(af.preview.length>200?'...':'')+'</div>':'')
      +'</div></div>';
    window._lastAttachedFile=null;
  }
  div.innerHTML='<div class="mav">'+escHtml(init)+'</div><div class="mcont"><div class="mname">'+escHtml(uName||'You')+'</div><div class="bbl">'+imgHtml+filePreviewHtml+'<span>'+escHtml(content)+'</span></div></div>';
  }
  msgs.appendChild(div);scrollDown();
  var _bbl=div.querySelector('.bbl');if(_bbl)setTimeout(function(){_renderMath(_bbl);},50);
  return _bbl;
}
function addThinking(){
  var ws=document.getElementById('ws');if(ws)ws.remove();
  var div=document.createElement('div');div.className='msg ai';div.id='typ';
  // Hide/show quick-action bar based on input focus
document.addEventListener('DOMContentLoaded',function(){
  var msgIn=document.getElementById('msgIn');
  var qaBar=document.getElementById('qaBar');
  if(msgIn&&qaBar){
    msgIn.addEventListener('focus',function(){qaBar.style.opacity='0.4';qaBar.style.pointerEvents='none';});
    msgIn.addEventListener('blur',function(){if(!msgIn.value){qaBar.style.opacity='1';qaBar.style.pointerEvents='all';}});
    msgIn.addEventListener('input',function(){
      if(msgIn.value){qaBar.style.opacity='0';qaBar.style.maxHeight='0';qaBar.style.overflow='hidden';qaBar.style.padding='0';}
      else{qaBar.style.opacity='1';qaBar.style.maxHeight='';qaBar.style.overflow='';qaBar.style.padding='';}
    });
  }
});
var _thinkMsgs=['Thinking…','Figuring this out…','On it…','Let me check…','Working on it…','Just a sec…','Pulling that up…','Good question…','Hmm, let me think…','Give me a moment…'];
  var _thinkTxt=_thinkMsgs[Math.floor(Math.random()*_thinkMsgs.length)];
  var _modelLbl=window._lastSelectedModel||'Fusion.AI';
  div.innerHTML='<div class="mav" style="font-size:14px">⚡</div><div class="mcont"><div class="mname" id="aiMsgName">'+escHtml(_modelLbl)+'</div><div class="bbl"><div class="thinking-wrap"><div class="thinking-dots"><span></span><span></span><span></span></div><span class="thinking-txt">'+_thinkTxt+'</span></div></div></div>';
  document.getElementById('msgs').appendChild(div);scrollDown();
}
function removeThinking(){var el=document.getElementById('typ');if(el)el.remove();}

var _ovAnswers={}, _ovPrompt='', _ovMkey='', _ovMode='image';

function showOverseerModal(questions, prompt, mkey, mode) {
  _ovAnswers={}; _ovPrompt=prompt; _ovMkey=mkey; _ovMode=mode;
  var old=document.getElementById('ov-modal-bg'); if(old) old.remove();
  var modeLabel=mode==='video'?'\uD83C\uDFAC Video':mode==='video3d'?'\uD83E\uDDCA 3D':'\uD83D\uDDBC\uFE0F Image';
  var qs=questions.map(function(q,qi){
    var chips=q.options.map(function(opt){
      return '<div class="ov-chip" onclick="ovChip(this,'+qi+','+JSON.stringify(opt)+')">' + escHtml(opt) + '</div>';
    }).join('');
    return '<div class="ov-q">'
      +'<div class="ov-q-num">Question '+(qi+1)+'</div>'
      +'<div class="ov-q-text">'+escHtml(q.q)+'</div>'
      +'<div class="ov-chips" id="ov-c-'+qi+'">'+chips+'</div>'
      +'</div>';
  }).join('');
  var bg=document.createElement('div'); bg.id='ov-modal-bg'; bg.className='ov-modal-bg';
  bg.innerHTML='<div class="ov-modal" onclick="event.stopPropagation()">'
    +'<div class="ov-modal-hdr">'
    +'<div class="ov-modal-badge">\uD83E\uDDE0\uFE0F Overseer \u00B7 GPT-OSS Safety</div>'
    +'<div class="ov-modal-title">Enhance your '+modeLabel+' prompt</div>'
    +'<div class="ov-modal-sub">Answer a few quick questions and the Overseer will craft a detailed, optimised prompt.</div>'
    +'<div class="ov-modal-prompt">'+escHtml(prompt)+'</div>'
    +'</div>'
    +'<div class="ov-modal-body">'+qs+'</div>'
    +'<div class="ov-modal-footer">'
    +'<button class="ov-gen-btn" onclick="ovSubmit()">\uD83D\uDE80\uFE0F Generate with enhanced prompt</button>'
    +'<button class="ov-skip-btn" onclick="ovSkip()">Skip \u2014 use my original prompt as-is</button>'
    +'</div>'
    +'</div>';
  bg.addEventListener('click', function(e){ if(e.target===bg) ovSkip(); });document.body.appendChild(bg);
}

function ovChip(el, qi, val) {
  var row=document.getElementById('ov-c-'+qi);
  if(row) row.querySelectorAll('.ov-chip').forEach(function(c){c.classList.remove('sel');});el.classList.add('sel');
  var qText=el.closest('.ov-q').querySelector('.ov-q-text').innerText;
  _ovAnswers[qText]=val;
}

function closeOvModal() { var m=document.getElementById('ov-modal-bg'); if(m) m.remove(); }

async function ovSubmit() {
  closeOvModal(); addThinking();
  try{
    var r=await apiFetch('/api/overseer/enhance',{method:'POST',headers:{'Content-Type':'application/json'},
      body:JSON.stringify({prompt:_ovPrompt,mode:_ovMode,answers:_ovAnswers})});
    var d=await r.json(); removeThinking();
    var enhanced=d.enhanced_prompt||_ovPrompt;
    if(enhanced!==_ovPrompt){
      var eb=addMsg('ai','',null,'Overseer',null);
      if(eb) eb.innerHTML='<div style="font-size:10px;font-weight:700;color:rgba(180,150,255,.8);text-transform:uppercase;letter-spacing:.8px;margin-bottom:6px"> Enhanced Prompt</div>'
        +'<div style="font-size:13px;color:var(--tx);line-height:1.7;font-style:italic">'+escHtml(enhanced)+'</div>';
    }
    if(_ovMode==='image'){ await _doGenImgFromChat(enhanced,_ovAnswers); }
    else if(_ovMode==='video'||_ovMode==='video3d'){ await _doGenVideoFromChat(enhanced); }
    else { await _doGenerate(enhanced, _ovMkey); }
  }catch(e){
    removeThinking();
    if(_ovMode==='image'){ await _doGenImgFromChat(_ovPrompt,{}); }
    else if(_ovMode==='video'||_ovMode==='video3d'){ await _doGenVideoFromChat(_ovPrompt); }
    else { addMsg('ai','Something went wrong. Please try again.'); }
    loading=false; document.getElementById('sbtn').disabled=false;
  }
}

function ovSkip() {
  closeOvModal();
  if(_ovMode==='image'){ _doGenImgFromChat(_ovPrompt,{}); }
  else if(_ovMode==='video'||_ovMode==='video3d'){ _doGenVideoFromChat(_ovPrompt); }
  else { _doGenerate(_ovPrompt, _ovMkey); }
}

async function _doGenerate(prompt, mkey) {
  var found=allModels.find(function(m){return m.key===mkey;});
  var provider=found?found.provider:'';
  var modelId=found?found.model:'';
  // ── image_or provider ────────────────────────────────────────────────────
  if(provider==='image_or'){
    loading=true;document.getElementById('sbtn').disabled=true;
    var bbl=addMsg('ai','',null,null);
    bbl.innerHTML='<span style="font-style:italic;color:var(--tx2)"> Generating image…</span>';
    try{
      var ir=await apiFetch('/api/generate/image',{method:'POST',
        headers:{'Content-Type':'application/json'},
        body:JSON.stringify({prompt:prompt,model:modelId,width:1024,height:1024})});
      if(!ir.ok){var er=await ir.json();bbl.innerHTML='<span style="color:var(--red)">❌ '+(er.error||'Image gen failed')+'</span>';}
      else{
        var id=await ir.json();
        var isrc=id.b64?'data:'+id.mime+';base64,'+id.b64:(id.url||'');
        if(isrc){
          bbl.innerHTML='<div class="gen-img-card">'
            +'<img class="gen-img-preview" src="'+isrc+'" onclick="openImgFull(this.src)"/>'
            +'<div class="gen-img-actions">'
            +'<button class="gen-img-btn" onclick="_dlGenImg(this)">&#x2B07;&#xFE0F; Download</button>'
            +'<button class="gen-img-btn gen-img-btn-sec" onclick="_attachGenImg(this)">Use in Chat</button>'
            +'</div></div>';
          if(bbl.querySelector('.gen-img-card')){
            bbl.querySelector('.gen-img-card').dataset.b64=id.b64||'';bbl.querySelector('.gen-img-card').dataset.mime=id.mime||'image/jpeg';
            bbl.querySelector('.gen-img-card').dataset.url=id.url||'';
          }
          var acts=document.createElement('div');acts.className='msg-actions';
          acts.dataset.b64=id.b64||'';acts.dataset.mime=id.mime||'image/png';
          acts.innerHTML='<button class="mact" onclick="downloadImgFromChat(this)">⬇ Download</button><button class="mact" onclick="attachImgFromChat(this)">Attach to Chat</button>';
          bbl.closest('.mcont').appendChild(acts);_lastGenResult={type:'image',b64:id.b64||'',mime:id.mime||'image/png',url:id.url||'',prompt:prompt};
          hist.push({role:'assistant',content:'[Image generated: '+prompt+']'});showToast('✅ Image generated!');
        }else{bbl.innerHTML='<span style="color:var(--red)">❌ No image data returned</span>';}
      }
    }catch(e){bbl.innerHTML='<span style="color:var(--red)">❌ '+escHtml(e.message)+'</span>';}
    loading=false;document.getElementById('sbtn').disabled=false;document.getElementById('msgIn').focus();
    return;
  }
  // ── Chat fallback ─────────────────────────────────────────────────────
  loading=true; document.getElementById('sbtn').disabled=true; addThinking();
  try{
    var payload={message:prompt,history:hist.slice(0,-1),model_key:mkey==='auto'?null:mkey,
      image_b64:'',image_mime:'image/jpeg',file_text:''};
    var res=await apiFetch('/api/chat',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify(payload)});
    if(!res.ok){
      var errText2=await res.text(); var err2={};
      try{err2=JSON.parse(errText2);}catch(_){err2={error:'Server error ('+res.status+')'};}
      removeThinking();showUserError(err2.dev_error||err2.error||'Server error','Something went wrong. Please try again.');
      addMsg('ai','Something went wrong. Please try again.');loading=false;document.getElementById('sbtn').disabled=false;return;
    }
    removeThinking(); await _streamResponse(res);
  }catch(e){removeThinking();showErrOverlay('Network error: '+e.message);}
  loading=false; document.getElementById('sbtn').disabled=false; document.getElementById('msgIn').focus();
}

async function _streamResponse(res) {
  var bbl=addMsg('ai','',null,null); bbl.innerHTML='';
  var msgDiv=bbl.closest('.msg'); msgDiv.classList.add('streaming');
  var full='', reason='', usedModel='';
  var reader=res.body.getReader(); var dec=new TextDecoder();
  var stopped=false;
  _chatStreamStart();
  try{
  while(true){
    // Check if aborted
    if(_chatAbort&&_chatAbort.signal&&_chatAbort.signal.aborted){stopped=true;break;}
    var result=await reader.read(); if(result.done) break;
    var lines=dec.decode(result.value).split('\n').filter(function(l){return l.startsWith('data: ')&&!l.includes('[DONE]');});
    for(var li=0;li<lines.length;li++){
      try{
        var ev=JSON.parse(lines[li].slice(6));
        if(ev.type==='image'){
          var du='data:'+ev.mime+';base64,'+ev.b64;
          bbl.innerHTML='<img class="gen-img-msg" src="'+du+'"/><br><span style="font-size:11px;color:var(--tx3)">'+escHtml(ev.prompt||'')+'</span>';
          lastGenImgB64=ev.b64; lastGenImgMime=ev.mime; full='[Image]';
        }else if(ev.type==='video'){
          var vu='data:'+ev.mime+';base64,'+ev.b64;
          bbl.innerHTML='<video controls style="width:100%;max-width:480px;border-radius:10px;border:1px solid var(--glass-bdr2)" src="'+vu+'"></video>'
            +'<br><a href="'+vu+'" download="fusion-video.mp4" style="font-size:11px;color:var(--blue)">\u2B07 Download video</a>'
            +'<br><span style="font-size:11px;color:var(--tx3)">'+escHtml(ev.prompt||'')+'</span>';
          full='[Video]';
        }else if(ev.type==='retry'){
          if(isDevUser){var nb=document.createElement('div');nb.style.cssText='font-size:11px;color:var(--tx3);font-style:italic;padding:3px 0 5px;';nb.textContent='\u26a1 '+ev.failed+' \u2014 trying '+ev.next+'\u2026';var mc=bbl.closest('.mcont');if(mc)mc.insertBefore(nb,bbl);}
        }else if(ev.type==='meta'){
          reason=ev.reason; usedModel=ev.model;
          var nm=bbl.closest('.mcont')&&bbl.closest('.mcont').querySelector('.mname');
          if(nm)nm.innerHTML='Fusion.AI <span class="mtag">'+escHtml(ev.model)+'</span>';
          var pts=ev.model.split('\u00B7');document.getElementById('mlabel').textContent=pts.length>1?pts[1].trim().split(' ').slice(0,2).join(' '):ev.model.split(' ')[0];
        }else if(ev.type==='think_start'){
          var mc2=bbl.closest('.mcont');var tb=mc2.querySelector('.thought');
          if(!tb){tb=document.createElement('div');tb.className='thought thought-streaming';tb.innerHTML='<div class="thought-hdr" onclick="toggleThought(this)"><span class="thought-icon">\uD83E\uDDE0</span><span class="thought-label">Thinking<span class="thought-dots"><span>.</span><span>.</span><span>.</span></span></span><span class="thought-toggle">\u25BE</span></div><div class="thought-body thought-body-open"><div class="thought-txt" id="lt-'+Date.now()+'"></div></div>';mc2.insertBefore(tb,bbl);tb._liveId=tb.querySelector('.thought-txt').id;}
          scrollDown();
        }else if(ev.type==='think_delta'){
          var mc2=bbl.closest('.mcont');var tb2=mc2.querySelector('.thought');if(tb2&&tb2._liveId){var txt=document.getElementById(tb2._liveId);if(txt){txt.textContent+=ev.text;scrollDown();}}
          }else if(ev.type==='think_end'){
          var mc2=bbl.closest('.mcont');var tb3=mc2.querySelector('.thought');if(tb3){tb3.classList.remove('thought-streaming');var lb=tb3.querySelector('.thought-label');if(lb)lb.innerHTML='Hide thought process';var tog=tb3.querySelector('.thought-toggle');if(tog)tog.classList.add('open');var tbd=tb3.querySelector('.thought-body');if(tbd){tbd.classList.add('open');tbd.classList.remove('thought-body-open');}}
        }else if(ev.type==='delta'){
          full+=ev.text; bbl.innerHTML=fmt(full); scrollDown();
        }else if(ev.type==='error'){
          showUserError(ev.dev_message||ev.message, 'Something went wrong. Please try again.');
          bbl.innerHTML='<span style="color:#f87171">⚠️ '+(ev.message||'Something went wrong')+'</span>'
            +'<br><button onclick="_retryLast()" style="margin-top:8px;background:rgba(6,182,212,.15);border:1px solid rgba(6,182,212,.4);border-radius:9px;padding:6px 14px;font-size:11px;color:rgba(100,180,255,.9);cursor:pointer;font-family:DM Sans,sans-serif">↺ Retry</button>';
        }
      }catch(pe){}
    }
  }
  }catch(readErr){
    if(readErr.name!=='AbortError'&&!stopped) console.warn('Stream read error:',readErr);
  } finally {
    try{reader.cancel();}catch(e){}
    _chatStreamEnd();
  }
  if(stopped){
    if(full) bbl.innerHTML=fmt(full)+'<span style="display:block;font-size:10px;color:var(--tx3);margin-top:4px">⏹ Stopped</span>';
    else bbl.innerHTML='<span style="color:var(--tx3);font-style:italic">⏹ Generation stopped.</span>';
  }
  msgDiv.classList.remove('streaming');
  // Parse any FUSIONFILE blocks that arrived in the stream
  if(full && full.includes('<<<FUSIONFILE:')){
    var _parsed=_parseFusionFilesRaw(full);
    if(_parsed!==full){ bbl.innerHTML=_parsed; full=full; }
  }
  setTimeout(function(){_renderMath(bbl);},80);
  if(reason){var rt=document.createElement('div');rt.className='rtag';rt.textContent='\u21B3 '+reason;bbl.closest('.mcont').appendChild(rt);}
  if(full&&full!=='[Image]'&&full!=='[Video]'){
    var acts=document.createElement('div');acts.className='msg-actions';
    acts.innerHTML='<button class="mact" onclick="saveMessage(this.closest(\'.mcont\').querySelector(\'.bbl\').innerText)">\uD83D\uDD16 Save</button>'
      +'<button class="mact" onclick="navigator.clipboard.writeText(this.closest(\'.mcont\').querySelector(\'.bbl\').innerText)">\uD83D\uDCCB Copy</button>'
      +'<button class="mact" onclick="speakText(this.closest(\'.mcont\').querySelector(\'.bbl\').innerText)">\uD83D\uDD0A Speak</button>';
    bbl.closest('.mcont').appendChild(acts);hist.push({role:'assistant',content:full}); loadMemory();if(ttsEnabled) speakText(full);if(vcActive) speakForVC(full);
    // Save to conversation
    await _saveConvMsg('assistant', full, usedModel);
  }
}

function openImgFull(s){
  var ov=document.createElement('div');
  ov.style.cssText='position:fixed;inset:0;background:rgba(0,0,0,.92);z-index:9999;display:flex;align-items:center;justify-content:center;cursor:zoom-out';
  ov.onclick=function(){document.body.removeChild(ov);};
  var im=document.createElement('img');im.src=s;
  im.style.cssText='max-width:92vw;max-height:92vh;border-radius:12px;box-shadow:0 0 80px rgba(0,0,0,.8)';ov.appendChild(im);document.body.appendChild(ov);
}
function _dlGenImg(btn){
  var c=btn.closest('.gen-img-card');
  var b64=c&&c.dataset.b64,mime=(c&&c.dataset.mime)||'image/jpeg',url=(c&&c.dataset.url)||'';
  if(b64){_downloadB64(b64,mime,'fusion-image.png');}
  else if(url){var a=document.createElement('a');a.href=url;a.download='fusion-image.png';a.target='_blank';a.click();}
  else showToast('No image data');
}
function _attachGenImg(btn){
  var c=btn.closest('.gen-img-card');
  var b64=c&&c.dataset.b64,mime=(c&&c.dataset.mime)||'image/jpeg';
  if(!b64){showToast('URL-only images cannot be attached');return;}
  pendingImageB64=b64;pendingImageMime=mime;pendingImageName='generated.png';document.getElementById('imgPrev').src='data:'+mime+';base64,'+b64;
  document.getElementById('imgPrevName').textContent='Generated image';document.getElementById('imgPrevWrap').classList.add('show');
  showToast('Attached to chat!');
}

var _imgGenAbort = null;

async function _doGenImgFromChat(prompt, answers) {
  answers=answers||{};
  var w=1024,h=1024;
  var sa=answers['Aspect ratio?']||'';
  if(sa.includes('16:9')||sa.includes('Landscape')){w=1280;h=720;}
  else if(sa.includes('9:16')||sa.includes('Portrait')){w=720;h=1280;}
  var style=answers['Style?']||'';
  var full=prompt+(style?' — '+style+' style':'');
  loading=true; document.getElementById('sbtn').disabled=true;
  var bbl=addMsg('ai','',null,null);
  var stopId='imgStop_'+Date.now();
  bbl.innerHTML='<div style="display:flex;flex-direction:column;gap:10px">'
    +'<div style="display:flex;align-items:center;gap:10px">'
    +'<div class="thinking-dots"><span></span><span></span><span></span></div>'
    +'<span style="color:var(--tx2);font-size:13px" id="imgStatusTxt_'+stopId+'"> Generating image… (up to 90s)</span></div>'
    +'<div style="height:3px;background:var(--glass-bdr);border-radius:3px;overflow:hidden"><div id="imgProg_'+stopId+'" style="height:100%;width:0%;background:var(--grad);transition:width 0.5s;border-radius:3px"></div></div>'
    +'<button id="'+stopId+'" onclick="stopImgGen(\''+stopId+'\')" style="background:rgba(124,58,237,.15);border:1px solid rgba(124,58,237,.4);color:var(--red);border-radius:9px;padding:8px 14px;font-size:12px;font-weight:600;cursor:pointer;width:fit-content">'
    +'&#x23F9;&#xFE0F; Stop Generation</button></div>';
  scrollDown();
  // Animate progress bar
  var prog=document.getElementById('imgProg_'+stopId);
  var progPct=0; var progTimer=setInterval(function(){
    if(!prog)return; progPct=Math.min(progPct+1,92); prog.style.width=progPct+'%';
  },800);
  _imgGenAbort=new AbortController();
  try{
    var r=await apiFetch('/api/generate/image',{method:'POST',
      signal:_imgGenAbort.signal,
      headers:{'Content-Type':'application/json'},
      body:JSON.stringify({prompt:full,model:'flux',width:w,height:h})});    clearInterval(progTimer); if(prog)prog.style.width='100%';
    var stopBtn=document.getElementById(stopId);if(stopBtn)stopBtn.style.display='none';
    if(!r.ok){
      var errTxt=await r.text(); var er={};
      try{er=JSON.parse(errTxt);}catch(_){er={error:'Image gen failed ('+r.status+')'};}
      bbl.innerHTML='<span style="color:var(--red)">❌ '+(er.error||'Failed')+'</span>';
    } else {
      var d=await r.json();
      var isrc=d.b64?'data:'+d.mime+';base64,'+d.b64:(d.url||'');
      if(isrc){
        bbl.innerHTML='<div style="display:flex;flex-direction:column;gap:6px">'
          +'<img src="'+isrc+'" style="max-width:100%;border-radius:12px;border:1px solid var(--glass-bdr2);display:block;cursor:zoom-in" onclick="openImgFull(this.src)"/>'
          +'<div style="display:flex;align-items:center;justify-content:space-between;flex-wrap:wrap;gap:6px">'
          +'<span style="font-size:10px;color:var(--tx3)">⚡ '+escHtml(d.backend||'Pollinations')+'</span>'
          +'<div style="display:flex;gap:6px">'
          +(d.b64?'<button class="mact" onclick="(function(b,m){var a=document.createElement(\'a\');a.href=\'data:\'+m+\';base64,\'+b;a.download=\'fusion-image.png\';a.click();})(\''+d.b64+'\',\''+d.mime+'\')" >⬇ Download</button>':'')
          +'<button class="mact" onclick="_attachGenImgData(\''+escHtml(d.b64||'')+'\',\''+escHtml(d.mime||'image/jpeg')+'\')">📎 Use in Chat</button>'
          +'</div></div></div>';
        _lastGenResult={type:'image',b64:d.b64||'',mime:d.mime||'image/jpeg',url:d.url||'',prompt:full};hist.push({role:'assistant',content:'[Image: '+full+']'});
        showToast('✅ Image ready!');
      } else {
        bbl.innerHTML='<span style="color:var(--red)">❌ No image returned. Please try again.</span>';
      }
    }
  } catch(e) {
    clearInterval(progTimer);
    var stopBtn2=document.getElementById(stopId);if(stopBtn2)stopBtn2.style.display='none';
    if(e.name==='AbortError'){
      bbl.innerHTML='<span style="color:var(--tx3);font-style:italic">⏹ Image generation stopped.</span>';
    } else {
      bbl.innerHTML='<span style="color:var(--red)">❌ '+escHtml(e.message)+'</span>';
    }
  }
  _imgGenAbort=null;loading=false; document.getElementById('sbtn').disabled=false;document.getElementById('msgIn').focus();
}

function stopImgGen(stopId) {
  if(_imgGenAbort){ _imgGenAbort.abort(); }
  var btn=document.getElementById(stopId); if(btn) btn.style.display='none';
  var prog=document.getElementById('imgProg_'+stopId.replace('imgStop_','imgProg_'));
  // try to update status text
  var statTxt=document.getElementById('imgStatusTxt_'+stopId);
  if(statTxt) statTxt.textContent=' Stopping…';
}

function _attachGenImgData(b64, mime) {
  if(!b64){showToast('URL-only images cannot be attached');return;}
  pendingImageB64=b64; pendingImageMime=mime; pendingImageName='generated.png';document.getElementById('imgPrev').src='data:'+mime+';base64,'+b64;
  document.getElementById('imgPrevName').textContent='Generated image';document.getElementById('imgPrevWrap').classList.add('show');
  showToast('📎 Attached to chat!');
}

async function _doGenVideoFromChat(prompt) {
  var bbl=addMsg('ai','',null,null);
  bbl.innerHTML='<div style="background:rgba(6,182,212,.07);border:1px solid rgba(6,182,212,.18);border-radius:12px;padding:16px 18px">'
    +'<div style="font-size:15px;font-weight:700;margin-bottom:8px"> Video Generation</div>'
    +'<div style="font-size:12px;color:var(--tx2);line-height:1.8">'
    +'Video generation needs a paid API. <strong>Free options:</strong><br>'
    +'• <a href="https://runwayml.com" target="_blank" style="color:var(--blue)">RunwayML</a> — free trial<br>'
    +'• <a href="https://lumalabs.ai/dream-machine" target="_blank" style="color:var(--blue)">Luma Dream Machine</a> — free tier<br>'
    +'• <a href="https://pika.art" target="_blank" style="color:var(--blue)">Pika.art</a> — free tier</div>'
    +'<div style="font-size:11px;color:var(--tx3);margin-top:10px;border-top:1px solid var(--glass-bdr);padding-top:8px">'
    +'Your prompt: <em>'+escHtml(prompt)+'</em></div></div>';
  scrollDown();
}

function toggleChatOpts(){var p=document.getElementById('chatOptsPanel');if(!p)return;p.style.display=p.style.display==='none'?'':'none';}
function closeChatOpts(){var p=document.getElementById('chatOptsPanel');if(p)p.style.display='none';}
document.addEventListener('click',function(e){
  var btn=document.getElementById('chatOptBtn'),panel=document.getElementById('chatOptsPanel');
  if(btn&&panel&&!btn.contains(e.target)&&!panel.contains(e.target))panel.style.display='none';
});
function setChatFontSize(sz){
  var cb=document.getElementById('chatBody');if(!cb)return;
  cb.style.fontSize=sz==='small'?'13px':sz==='large'?'17px':'15px';closeChatOpts();showToast('Text: '+sz);
}
function exportChat(){
  var msgs=document.getElementById('msgs');if(!msgs)return showToast('No chat');
  var lines=[];msgs.querySelectorAll('.msg').forEach(function(m){
    var who=m.classList.contains('ai')?'Fusion.AI':'You';
    var t=m.querySelector('.bbl');if(t)lines.push(who+': '+t.innerText.trim());
  });
  if(!lines.length)return showToast('Nothing to export');
  var blob=new Blob([lines.join('\n\n')],{type:'text/plain'});
  var a=document.createElement('a');a.href=URL.createObjectURL(blob);
  a.download='fusion-chat-'+new Date().toISOString().slice(0,10)+'.txt';a.click();closeChatOpts();showToast(' Exported!');
}

async function _tryFreeAPI(text){
  var tl=text.toLowerCase().trim();
  var isJoke=/(tell|say|give|share).{0,15}joke|joke\s*$|make me laugh|funny/i.test(tl);
  var isQuote=/(give|share|tell|say|show).{0,15}quote|inspirational|motivat/i.test(tl);
  var isFact=/(give|share|tell|say|random|fun).{0,10}fact|did you know/i.test(tl);
  var ep=isJoke?'/api/freeapi/joke':isQuote?'/api/freeapi/quote':isFact?'/api/freeapi/fact':null;
  if(!ep) return null;
  try{
    var r=await apiFetch(ep);
    var d=await r.json();
    if(d.ok&&d.text) return d.text;
  }catch(e){}
  return null;
}

async function sendMsg(){
  if(loading) return;
  // Track selected model for display
  var mSel=document.getElementById('mOverride');
  window._lastSelectedModel=mSel&&mSel.value?mSel.options[mSel.selectedIndex].text:'Fusion.AI';
  var inp=document.getElementById('msgIn'); var text=inp.value.trim();
  if(!text&&!pendingImageB64&&!pendingFileText) return;
  // ── Document command detection (/doc, /ppt, /csv etc.) ────────────────
  var _docCheck=text.match(/^\/(doc|report|ppt|csv|dataset|json|xlsx|html|py|js|sql|txt|xml|yaml)\s+(.+)/i);
  if(_docCheck){
    var _ft=_docCheck[1].toLowerCase();
    var _topic=_docCheck[2].trim();
    var _fnBase=_topic.replace(/[^a-z0-9\s]/gi,'').trim().replace(/\s+/g,'_').toLowerCase().slice(0,45)||'file';
    var _extMap={doc:'md',report:'md',ppt:'pptx.md',dataset:'csv',xlsx:'csv'};
    var _ext=_extMap[_ft]||_ft;
    var _fname=_fnBase+'.'+_ext;
    var _tdescs={doc:'a comprehensive well-structured document (use rich markdown with headings, sections, tables)',report:'a detailed professional report with executive summary, sections, data, and conclusions',ppt:'presentation slides (use # Title for each slide title, --- between slides, bullet points for content)',csv:'CSV data with a descriptive header row and at least 20 rows of realistic data',dataset:'a rich CSV dataset with headers and at least 25 rows of diverse, realistic data',json:'well-structured valid JSON data',xlsx:'CSV data that can be opened in Excel (with proper headers)',html:'a complete self-contained styled HTML page',py:'complete working Python code with comments',js:'complete working JavaScript code with comments',sql:'SQL with CREATE TABLE statements and INSERT data',txt:'plain text formatted content',xml:'valid well-formed XML data',yaml:'valid YAML configuration or data file'};
    var _tdesc=_tdescs[_ft]||'a file';
    inp.value='';
    text='Create '+_tdesc+' about: '+_topic+'\n\n[SYSTEM: Generate the complete file content and wrap it EXACTLY in:\n<<<FUSIONFILE:'+_fname+'>>>\nYour complete content here\n<<<END_FUSIONFILE>>>\nDo NOT add anything outside the markers. Make it complete and production-ready.]';
  }
  var tl=text.toLowerCase();
  // Slash commands — bypass normal flow entirely
  if(tl.startsWith('/imagine ')||tl.startsWith('/img ')){
    var imgPrompt=text.slice(tl.startsWith('/imagine ')?9:5).trim();
    if(!imgPrompt){showToast('Usage: /imagine [your prompt]');return;}
    inp.value=''; inp.style.height='auto';addMsg('user',text,null,null,''); hist.push({role:'user',content:text});
    loading=true; document.getElementById('sbtn').disabled=true;
    await _doGenImgFromChat(imgPrompt,{}); loading=false; document.getElementById('sbtn').disabled=false; return;
  }
  if(tl.startsWith('/video ')){
    var vidPrompt=text.slice(7).trim(); if(!vidPrompt){showToast('Usage: /video [prompt]');return;}
    inp.value=''; inp.style.height='auto';addMsg('user',text,null,null,''); hist.push({role:'user',content:text});
    loading=true; document.getElementById('sbtn').disabled=true;
    await _doGenVideoFromChat(vidPrompt); loading=false; document.getElementById('sbtn').disabled=false; return;
  }
  if(tl.startsWith('/3d ')){
    var tdPrompt=text.slice(4).trim(); if(!tdPrompt){showToast('Usage: /3d [prompt]');return;}
    inp.value=''; inp.style.height='auto';addMsg('user',text,null,null,''); hist.push({role:'user',content:text});
    loading=true; document.getElementById('sbtn').disabled=true;
    await _doGen3DFromChat(tdPrompt); loading=false; document.getElementById('sbtn').disabled=false; return;
  }
  if(tl.startsWith('/scrape ')){
    var scrapeTarget=text.slice(8).trim(); if(!scrapeTarget){showToast('Usage: /scrape [url]');return;}
    inp.value=''; inp.style.height='auto'; addMsg('user',text,null,null,''); hist.push({role:'user',content:text});
    loading=true; document.getElementById('sbtn').disabled=true;
    try{
      var sr=await apiFetch('/api/scrape',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({url:scrapeTarget})});
      var sd=await sr.json();
      if(sd.ok){
        var scrapeMsg='**Scraped: '+escHtml(sd.title||scrapeTarget)+'**\n\n'+sd.text.slice(0,2000)+(sd.text.length>2000?'\n\n*(truncated)*':'');
        var sb=addMsg('ai',scrapeMsg,null,'Web Scraper');
        hist.push({role:'assistant',content:scrapeMsg});
        await _saveConvMsg('assistant',scrapeMsg,'Scraper');
      } else {
        addMsg('ai','⚠️ Could not scrape that URL: '+(sd.error||'Unknown error'),null,'Scraper');
      }
    }catch(e){addMsg('ai','❌ Scrape failed: '+e.message,null,'Scraper');}
    loading=false; document.getElementById('sbtn').disabled=false; return;
  }
  inp.value=''; inp.style.height='auto';
  var mkey=(document.getElementById('mOverride')||{value:'auto'}).value;
  var found=allModels.find(function(m){return m.key===mkey;});
  var isMedia=false;
  var wc=text?text.split(/\s+/).filter(Boolean).length:0;
  loading=true; document.getElementById('sbtn').disabled=true;
  var imgSrc=pendingImageB64?'data:'+pendingImageMime+';base64,'+pendingImageB64:'';
  var displayText=text||(pendingImageB64?'[Image attached]':'')+(pendingFileName?'[File: '+pendingFileName+']':'');
  addMsg('user',displayText||text,null,null,imgSrc);
  hist.push({role:'user',content:pendingFileText?(text+'\n\n'+pendingFileText).trim():(text||'Analyse this image.')});
  var savedFile=pendingFileText, savedImg=pendingImageB64, savedMime=pendingImageMime;
  clearImage(); clearFile();
  // ── News/current queries → groq_compound (has web search) ──────────────────
  if(text&&!savedImg&&!savedFile&&_isNewsQuery(text)){
    var _mkey2=(document.getElementById('mOverride')||{value:'auto'}).value;
    if(_mkey2==='auto'){
      // override to compound for this request
      document.getElementById('mOverride')&&(document.getElementById('mOverride').value='groq_compound');
      setTimeout(function(){document.getElementById('mOverride')&&(document.getElementById('mOverride').value='auto');},100);
    }
  }
  // ── FreeAPI shortcut: try free APIs before burning LLM tokens ───────────────
  if(text&&!savedImg&&!savedFile){
    try{
      var _far=await apiFetch('/api/freeapi/query',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({text:text})});
      var _fad=await _far.json();
      if(_fad.ok&&_fad.result){
        var _fab=addMsg('ai',_fad.result,null,'FreeAPI · '+(_fad.api||''));
        hist.push({role:'assistant',content:_fad.result});
        await _saveConvMsg('assistant',_fad.result,'FreeAPI');
        loading=false; document.getElementById('sbtn').disabled=false; return;
      }
    }catch(_){}
    // Legacy joke/quote/fact
    var _fa=await _tryFreeAPI(text);
    if(_fa){
      var _fab2=addMsg('ai',_fa,null,'FreeAPI');
      hist.push({role:'assistant',content:_fa});
      await _saveConvMsg('assistant',_fa,'FreeAPI');
      loading=false; document.getElementById('sbtn').disabled=false; return;
    }
  }
  // ── Overseer: Image/Video/3D — modal Q&A for short prompts ──────────────────
  if(isMedia && text && wc < 12){
    loading=false; document.getElementById('sbtn').disabled=false;
    try{
      var er=await apiFetch('/api/overseer/enhance',{method:'POST',headers:{'Content-Type':'application/json'},
        body:JSON.stringify({prompt:text,mode:found.provider,model_label:found.label,answers:{}})});
      var ed=await er.json();
      if(ed.action==='ask'&&ed.questions&&ed.questions.length){
        showOverseerModal(ed.questions,text,mkey,found.provider); return;
      }
      if(ed.enhanced_prompt&&ed.enhanced_prompt!==text) text=ed.enhanced_prompt;
    }catch(e){}
    loading=true; document.getElementById('sbtn').disabled=true;await _doGenerate(text, mkey); return;
  }
  // ── Image model auto-use: if user picks an image model, route directly to image gen ──
  if(text && !savedImg && !savedFile){
    var _selMKey=(document.getElementById('mOverride')||{value:'auto'}).value;
    var _selM=allModels.find(function(m){return m.key===_selMKey;});
    if(_selM && _selM.type==='image'){
      loading=false; document.getElementById('sbtn').disabled=false;
      await _doGenImgFromChat(text,{model:_selM.key});
      return;
    }
  }
  // ── Intent detection — image/video/3d from ANY message ────────────────────
  if(text && !savedImg){
    try{
      var ir=await apiFetch('/api/detect-intent',{method:'POST',
        headers:{'Content-Type':'application/json'},body:JSON.stringify({text:text})});
      var id2=await ir.json();
      if(id2.intent && id2.intent!=='chat'){
        loading=false; document.getElementById('sbtn').disabled=false;
        var qr=await apiFetch('/api/overseer/enhance',{method:'POST',
          headers:{'Content-Type':'application/json'},
          body:JSON.stringify({prompt:text,mode:id2.intent,answers:{},word_count:99})});
        var qd=await qr.json();
        if(qd.action==='ask'&&qd.questions&&qd.questions.length){
          showOverseerModal(qd.questions,text,'auto',id2.intent); return;
        }
        loading=true; document.getElementById('sbtn').disabled=true;if(id2.intent==='image') await _doGenImgFromChat(qd.enhanced_prompt||text,{});
        else await _doGenVideoFromChat(qd.enhanced_prompt||text);loading=false; document.getElementById('sbtn').disabled=false; return;
      }
    }catch(e){}
  }
  // Web search via Langsearch removed
  // ── Overseer: silently enhance short vague chat prompts ──────────────────
  if(!isMedia && mkey==='auto' && text && wc>0 && wc<6 && !savedImg){
    try{
      var cr=await apiFetch('/api/overseer/enhance',{method:'POST',headers:{'Content-Type':'application/json'},
        body:JSON.stringify({prompt:text,mode:'chat',answers:{}})});
      var cd=await cr.json();
      if(cd.enhanced_prompt&&cd.enhanced_prompt!==text) text=cd.enhanced_prompt;
    }catch(e){}
  }
  // ── Tool routing ─────────────────────────────────────────────────────────
  if(text && !text.startsWith('/')){
    if(_activeTool==='deep'){
      document.getElementById('msgIn').value=''; updateTokenCount('');
      doDeepResearch(text); return;
    }
    if(_activeTool==='math'||text.startsWith('[Math] ')){
      document.getElementById('msgIn').value=''; updateTokenCount('');
      doSpecialResearch(text.replace(/^\[Math\] /,''),'math'); return;
    }
    if(_activeTool==='code'||text.startsWith('[Code] ')){
      document.getElementById('msgIn').value=''; updateTokenCount('');
      doSpecialResearch(text.replace(/^\[Code\] /,''),'code'); return;
    }
  }
  // ── Auto web-ground every message — inject DDG results ───────────────────
  var _webCtx='';
  if(text && !text.startsWith('/') && text.length>=4){
    try{_webCtx=await _autoWebSearch(text)||'';}catch(e){}
  }
  // ── Send to chat ──────────────────────────────────────────────────────────
  // Auto-create conversation session if needed
  await _ensureConvId(text);
  // Save user message to conversation
  await _saveConvMsg('user', text, '');
  var augmentedMsg=_webCtx?text+'\n\n<web_context>\n'+_webCtx+'\n</web_context>':text;
  var payload={message:augmentedMsg,history:hist.slice(0,-1),model_key:mkey==='auto'?null:mkey,
    image_b64:savedImg||'',image_mime:savedMime||'image/jpeg',file_text:savedFile||''};
  addThinking();
  try{
    _chatStreamStart();
    var chatSig=_chatAbort?_chatAbort.signal:undefined;
    var res=await apiFetch('/api/chat',{method:'POST',headers:{'Content-Type':'application/json'},
      signal:chatSig,body:JSON.stringify(payload)});
    if(!res.ok){
      var errText=await res.text(); var err={};
      try{err=JSON.parse(errText);}catch(_){err={error:'Server error ('+res.status+')'};}
      removeThinking();showUserError(err.dev_error||err.error||'Server error','Something went wrong. Please try again.');
      addMsg('ai','Something went wrong. Please try again.');loading=false;_chatStreamEnd();return;
    }
    removeThinking(); await _streamResponse(res);
  }catch(e){
    if(e.name!=='AbortError'){removeThinking();showErrOverlay('Network error: '+e.message);addMsg('ai','Network error. Check your connection.');}
    _chatStreamEnd();
  }
  loading=false; document.getElementById('sbtn').disabled=false; document.getElementById('msgIn').focus();
}


// ── Custom Endpoint ────────────────────────────────────────────────────────
var _customEndpointData = JSON.parse(localStorage.getItem('fusion_custom_ep')||'{}');
function saveCustomEndpoint(){
  var url=(document.getElementById('customEndpointUrl')||{}).value.trim();
  var model=(document.getElementById('customEndpointModel')||{}).value.trim();
  var key=(document.getElementById('customEndpointKey')||{}).value.trim();
  if(!url){showToast('Enter endpoint URL');return;}
  _customEndpointData={url:url,model:model,key:key};
  localStorage.setItem('fusion_custom_ep',JSON.stringify(_customEndpointData));
  // Also save key to server
  if(key) apiFetch('/api/keys',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({provider:'custom_endpoint',key:key})});
  if(url) apiFetch('/api/keys',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({provider:'extra',key:url+(model?'||'+model:'')})});
  showToast('\u2705 Custom endpoint saved');
  var st=document.getElementById('customEndpointStat');if(st)st.textContent='Saved';
  // Add to model selector
  _addCustomModelOption(url,model);
}
function _addCustomModelOption(url,model){
  var sel=document.getElementById('mOverride');if(!sel)return;
  var existing=sel.querySelector('option[value="custom_ep"]');
  if(existing)sel.removeChild(existing);
  if(!url)return;
  var opt=document.createElement('option');opt.value='custom_ep';opt.textContent='\u{1F517} Custom — '+(model||url.split('/')[2]||'Your API');
  sel.appendChild(opt);
}
function _loadCustomEndpointUI(){
  var d=_customEndpointData;
  var urlEl=document.getElementById('customEndpointUrl');
  var modEl=document.getElementById('customEndpointModel');
  if(urlEl&&d.url){urlEl.value=d.url;}
  if(modEl&&d.model){modEl.value=d.model;}
  var st=document.getElementById('customEndpointStat');
  if(st) st.textContent=d.url?'Saved':'Not set';
  if(d.url)_addCustomModelOption(d.url,d.model);
}

// ── Scraper JS ─────────────────────────────────────────────────────────────
var _scraperData = null;
async function doScrape(){
  var url=(document.getElementById('scraperUrl')||{}).value.trim();
  if(!url){showToast('Enter a URL to scrape');return;}
  var st=document.getElementById('scraperStatus');
  var res=document.getElementById('scraperResult');
  if(st){st.style.display='';st.textContent='\u23F3 Scraping...';}
  if(res) res.style.display='none';
  try{
    var r=await apiFetch('/api/scrape',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({url:url})});
    var d=await r.json();
    if(d.ok){
      _scraperData=d;
      document.getElementById('scraperTitle').textContent=d.title||url;
      document.getElementById('scraperDesc').textContent=d.description||'';
      document.getElementById('scraperMeta').textContent=(d.word_count||0)+' words · '+d.content_type+' · HTTP '+d.status;
      document.getElementById('scraperText').value=d.text||d.data&&JSON.stringify(d.data,null,2)||'';
      var linksWrap=document.getElementById('scraperLinksWrap');
      var linksEl=document.getElementById('scraperLinks');
      if(d.links&&d.links.length&&linksWrap&&linksEl){
        document.getElementById('scraperLinkCount').textContent=d.links.length;
        linksEl.innerHTML=d.links.map(function(l){return '<a href="'+l+'" target="_blank" style="display:block;overflow:hidden;text-overflow:ellipsis;white-space:nowrap">'+l+'</a>';}).join('');
        linksWrap.style.display='';
      }
      if(st)st.textContent='\u2705 Scraped successfully — '+d.word_count+' words';
      if(res)res.style.display='';
    } else {
      if(st)st.textContent='\u274C '+d.error;
    }
  }catch(e){if(st)st.textContent='\u274C '+e.message;}
}
function sendScrapedToChat(){
  if(!_scraperData||!_scraperData.text){showToast('Scrape a URL first');return;}
  var txt='[Scraped: '+(_scraperData.title||_scraperData.url)+']\n\n'+_scraperData.text.slice(0,3000);
  document.getElementById('msgIn').value=txt;ar(document.getElementById('msgIn'));
  closeSP();document.getElementById('msgIn').focus();showToast('\u{1F4CB} Scraped content added to input');
}

(async function(){
  // Handle Google OAuth callback token in URL
  var _up=new URLSearchParams(window.location.search);
  var _gt=_up.get('g_token'), _gn=_up.get('g_name'), _ga=_up.get('g_avatar'), _ge=_up.get('auth_error');
  if(_ge){document.getElementById('authErr').textContent='Sign-in error: '+_ge;document.getElementById('authErr').style.display='block';}
  if(_gt){
    authToken=_gt; _setCookie('fusion_token',_gt,365); localStorage.setItem('fusion_token',_gt);
    if(_gn) uName=decodeURIComponent(_gn);
    if(_ga) userAvatar=decodeURIComponent(_ga);
    window.history.replaceState({},'','/');
  }
  if(!authToken) return;
  try{
    var r=await apiFetch('/api/me');var d=await r.json();
    if(d.logged_in){
      uName=d.username; currentMKey=d.model_key||'auto'; isDevUser=d.is_dev||false;
      if(d.avatar) userAvatar=d.avatar;
      if(d.theme){themePref.mode=d.theme;saveTheme();}
      launch();
    }else{authToken='';localStorage.removeItem('fusion_token');_delCookie('fusion_token');}
  }
  catch(e){authToken='';localStorage.removeItem('fusion_token');_delCookie('fusion_token');}
})();

// ══ Toolbox dropdown ══════════════════════════════════════════════════════
function toggleToolbox(e){
  e.stopPropagation();
  var btn=document.getElementById('toolboxBtn');
  var menu=document.getElementById('toolboxMenu');
  var open=!menu.classList.contains('hidden');
  if(open){closeToolbox();}
  else{menu.classList.remove('hidden');btn.classList.add('open');}
}
function closeToolbox(){
  var menu=document.getElementById('toolboxMenu');
  var btn=document.getElementById('toolboxBtn');
  if(menu){menu.classList.add('hidden');}
  if(btn){btn.classList.remove('open');}
}
document.addEventListener('click',function(e){
  var wrap=document.getElementById('toolboxWrap');
  if(wrap&&!wrap.contains(e.target)){closeToolbox();}
});

// ══ File Creator Modal ═════════════════════════════════════════════════════
var _FILE_TYPES=[
  // Documents
  {icon:'📄',name:'Markdown Doc', ext:'md',      cat:'Docs',   binary:false, prompt:'a comprehensive, well-structured markdown document with headings, subheadings, tables, and rich content'},
  {icon:'📑',name:'Word Doc',     ext:'docx',    cat:'Docs',   binary:true,  prompt:'a professional document with structured headings, paragraphs, bullet lists, and tables (in markdown format for conversion)'},
  {icon:'📋',name:'Report',       ext:'md',      cat:'Docs',   binary:false, prompt:'a detailed professional report with executive summary, methodology, findings, data analysis, conclusions, and recommendations'},
  {icon:'📝',name:'Plain Text',   ext:'txt',     cat:'Docs',   binary:false, prompt:'plain text formatted content, clearly organized with spacing and separators'},
  // Spreadsheets & Data
  {icon:'📊',name:'Excel (.xlsx)',ext:'xlsx',    cat:'Data',   binary:true,  prompt:'CSV data with a descriptive header row and at least 25 rows of realistic, diverse data'},
  {icon:'📊',name:'CSV Dataset',  ext:'csv',     cat:'Data',   binary:false, prompt:'CSV data with a clear header row and at least 25 rows of realistic, varied, detailed data'},
  {icon:'📋',name:'JSON Data',    ext:'json',    cat:'Data',   binary:false, prompt:'well-structured valid JSON data, properly formatted and indented with realistic values'},
  {icon:'⚙️',name:'YAML Config', ext:'yaml',    cat:'Data',   binary:false, prompt:'valid, well-structured YAML configuration or dataset'},
  {icon:'📰',name:'XML',          ext:'xml',     cat:'Data',   binary:false, prompt:'valid, well-formed XML data with proper nesting and realistic values'},
  // Presentations
  {icon:'🎯',name:'PowerPoint',   ext:'pptx.md', cat:'Present',binary:true,  prompt:'presentation slides — use # Title for each slide title, --- to separate slides, and bullet points for content. Include at least 8 slides.'},
  // Web & UI
  {icon:'🌐',name:'HTML Page',    ext:'html',    cat:'Web',    binary:false, prompt:'a complete, self-contained, beautifully styled HTML page with inline CSS and JavaScript, responsive design'},
  {icon:'🎨',name:'CSS Styles',   ext:'css',     cat:'Web',    binary:false, prompt:'complete, well-organized CSS stylesheet with variables, responsive breakpoints, and modern design'},
  {icon:'📜',name:'JavaScript',   ext:'js',      cat:'Web',    binary:false, prompt:'complete, well-commented, modern JavaScript code with error handling'},
  {icon:'📜',name:'TypeScript',   ext:'ts',      cat:'Web',    binary:false, prompt:'complete, typed TypeScript code with interfaces, types, and proper annotations'},
  // Backend Code
  {icon:'🐍',name:'Python',       ext:'py',      cat:'Code',   binary:false, prompt:'complete, well-commented, runnable Python code with docstrings and error handling'},
  {icon:'🗄',name:'SQL',          ext:'sql',     cat:'Code',   binary:false, prompt:'SQL script with CREATE TABLE statements, indexes, constraints, and INSERT data rows'},
  {icon:'🦀',name:'Bash Script',  ext:'sh',      cat:'Code',   binary:false, prompt:'complete bash shell script with shebang, comments, and error handling'},
  {icon:'🐹',name:'Go',           ext:'go',      cat:'Code',   binary:false, prompt:'complete, idiomatic Go code with proper package structure and error handling'},
  {icon:'☕',name:'Java',         ext:'java',    cat:'Code',   binary:false, prompt:'complete Java class with proper structure, comments, and best practices'},
  {icon:'🔧',name:'C/C++',        ext:'cpp',     cat:'Code',   binary:false, prompt:'complete, well-commented C++ code with proper includes and main function'},
  {icon:'💎',name:'Ruby',         ext:'rb',      cat:'Code',   binary:false, prompt:'complete, idiomatic Ruby code with comments'},
  {icon:'🐘',name:'PHP',          ext:'php',     cat:'Code',   binary:false, prompt:'complete PHP script with proper structure and comments'},
  // Config
  {icon:'⚙️',name:'TOML Config', ext:'toml',    cat:'Config', binary:false, prompt:'valid TOML configuration file with sections, comments, and realistic values'},
  {icon:'⚙️',name:'ENV File',    ext:'env',     cat:'Config', binary:false, prompt:'environment variables .env file with comments explaining each variable'},
];
// Group by category
var _FILE_CATS=['Docs','Data','Present','Web','Code','Config'];

var _selectedFileType=null;

function openFileCreator(){
  var grid=document.getElementById('fileTypeGrid');
  if(grid&&grid.children.length===0){
    // Build categorized grid
    _FILE_CATS.forEach(function(cat){
      var catDiv=document.createElement('div');
      catDiv.style.cssText='grid-column:1/-1;font-size:9px;font-weight:700;letter-spacing:1.2px;color:rgba(80,120,200,.6);text-transform:uppercase;padding:8px 4px 2px;margin-top:4px';
      catDiv.textContent=cat;
      grid.appendChild(catDiv);
      _FILE_TYPES.forEach(function(ft,i){
        if(ft.cat!==cat) return;
        var card=document.createElement('div');
        card.className='ftype-card';
        card.dataset.idx=i;
        var badge=ft.binary?'<span style="font-size:8px;background:rgba(6,182,212,.2);border:1px solid rgba(6,182,212,.3);border-radius:3px;padding:1px 4px;color:#80c0ff;display:block;margin-top:2px">Real '+ft.ext.replace('pptx.md','pptx').toUpperCase()+'</span>':'';
        card.innerHTML='<div class="ftype-card-icon">'+ft.icon+'</div><div class="ftype-card-name">'+ft.name+'</div><div class="ftype-card-ext">.'+ft.ext.replace('pptx.md','pptx')+'</div>'+badge;
        card.onclick=function(){
          document.querySelectorAll('.ftype-card').forEach(function(c){c.classList.remove('selected');});
          card.classList.add('selected');
          _selectedFileType=i;
          document.getElementById('fileTopicInput').placeholder='Describe your '+ft.name.toLowerCase()+'…';
        };
        grid.appendChild(card);
      });
    });
  }
  _selectedFileType=null;
  document.querySelectorAll('.ftype-card').forEach(function(c){c.classList.remove('selected');});
  document.getElementById('fileTopicInput').value='';
  document.getElementById('fileModalOverlay').classList.remove('hidden');
  setTimeout(function(){document.getElementById('fileTopicInput').focus();},200);
}
function closeFileCreator(){
  document.getElementById('fileModalOverlay').classList.add('hidden');
}
function fileModalClickOut(e){
  if(e.target===document.getElementById('fileModalOverlay')){closeFileCreator();}
}
function submitFileCreator(){
  if(_selectedFileType===null){showToast('Pick a file type first');return;}
  var topic=document.getElementById('fileTopicInput').value.trim();
  if(!topic){showToast('Enter a topic or description');return;}
  var ft=_FILE_TYPES[_selectedFileType];
  var fnBase=topic.replace(/[^a-z0-9 ]/gi,'').trim().replace(/ +/g,'_').toLowerCase().slice(0,45)||'file';
  var fname=fnBase+'.'+ft.ext;
  closeFileCreator();
  var prompt='Create '+ft.prompt+' about: '+topic
    +'\n\n[SYSTEM: Wrap ENTIRE content EXACTLY as:\n<<<FUSIONFILE:'+fname+'>>>\n[content]\n<<<END_FUSIONFILE>>>\nRules: (1) Nothing outside the markers. (2) Complete and production-ready — zero placeholders. (3) For CSV/xlsx: 25+ rows. (4) For pptx.md slides: 8+ slides with --- separator. (5) For code: fully runnable.]';
  var inp=document.getElementById('msgIn');
  inp.value=prompt;
  sendMsg();
}

// ══ Model Arena ════════════════════════════════════════════════════════════
var _ARENA_MODELS=[
  {key:'gh_gpt4o',label:'GPT-4o',emoji:'🤖',company:'OpenAI'},
  {key:'gh_o4_mini',label:'o4 Mini',emoji:'🧮',company:'OpenAI'},
  {key:'or_deepseek_r1',label:'DeepSeek R1',emoji:'🧊',company:'DeepSeek'},
  {key:'or_deepseek_v3',label:'DeepSeek V3',emoji:'🧊',company:'DeepSeek'},
  {key:'groq_llama33_70b',label:'Llama 3.3 70B',emoji:'🦙',company:'Meta'},
  {key:'groq_compound',label:'Compound Beta',emoji:'⚗️',company:'Groq'},
  {key:'or_qwen3_235b',label:'Qwen3 235B',emoji:'🌸',company:'Alibaba'},
  {key:'cf_qwen3_30b',label:'Qwen3 30B',emoji:'🌸',company:'Alibaba'},
  {key:'or_nemotron_super',label:'Nemotron 253B',emoji:'⚡',company:'NVIDIA'},
  {key:'gh_deepseek_r1',label:'DeepSeek R1',emoji:'🧊',company:'DeepSeek'},
  {key:'gh_phi4',label:'Phi-4',emoji:'Φ',company:'Microsoft'},
  {key:'or_mistral_small',label:'Mistral Small',emoji:'🌊',company:'Mistral'},
  {key:'cf_nemotron',label:'Nemotron 120B',emoji:'⚡',company:'NVIDIA'},
  {key:'groq_qwen3_32b',label:'Qwen3 32B',emoji:'🌸',company:'Alibaba'},
];
var _arenaSelKeys=new Set();

function openArena(){
  var sel=document.getElementById('arenaModelSel');
  if(sel&&sel.children.length===0){
    _ARENA_MODELS.forEach(function(m){
      var chip=document.createElement('span');
      chip.className='arena-model-chip';
      chip.dataset.key=m.key;
      chip.innerHTML=m.emoji+' '+m.label;
      chip.onclick=function(){
        if(_arenaSelKeys.has(m.key)){
          _arenaSelKeys.delete(m.key); chip.classList.remove('sel');
        } else if(_arenaSelKeys.size<5){
          _arenaSelKeys.add(m.key); chip.classList.add('sel');
        } else {showToast('Max 5 models');}
        var cnt=document.getElementById('arenaSelCount');
        if(cnt) cnt.textContent=_arenaSelKeys.size>0?_arenaSelKeys.size+' selected':'0 selected · auto-pick 5';
      };
      sel.appendChild(chip);
    });
  }
  _arenaSelKeys.clear();
  document.querySelectorAll('.arena-model-chip').forEach(function(c){c.classList.remove('sel');});
  document.getElementById('arenaPromptInput').value='';
  document.getElementById('arenaSelCount').textContent='0 selected · auto-pick 5';
  document.getElementById('arenaModalOverlay').classList.remove('hidden');
  setTimeout(function(){document.getElementById('arenaPromptInput').focus();},150);
}
function closeArena(){document.getElementById('arenaModalOverlay').classList.add('hidden');}
function arenaClickOut(e){if(e.target.id==='arenaModalOverlay')closeArena();}

async function runArena(){
  var prompt=document.getElementById('arenaPromptInput').value.trim();
  if(!prompt){showToast('Enter a question');return;}
  var btn=document.getElementById('arenaRunBtn');
  btn.disabled=true; btn.textContent='⏳ Running…';
  closeArena();

  addMsg('user',prompt);
  var bbl=addMsg('ai','',null,'Arena');
  var t0=Date.now();
  bbl.innerHTML='<div class="arena-wrap"><div class="arena-header"><span class="arena-badge">ARENA</span><span class="arena-prompt">'+escHtml(prompt.slice(0,80))+'</span></div>'
    +'<div class="arena-loading"><div class="arena-spinner"></div><span>Running '+(Array.from(_arenaSelKeys).length||5)+' models simultaneously…</span></div></div>';

  try{
    var body={prompt:prompt};
    if(_arenaSelKeys.size>0) body.models=Array.from(_arenaSelKeys);
    var r=await apiFetch('/api/arena',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify(body)});
    var d=await r.json();
    var elapsed=((Date.now()-t0)/1000).toFixed(1);
    if(!d.ok){bbl.innerHTML='<div style="color:#f87171">Arena failed.</div>';btn.disabled=false;btn.textContent='🏆 Run Arena';return;}
    _renderArenaResult(bbl,d,elapsed);
  }catch(ex){
    bbl.innerHTML='<div style="color:#f87171">'+escHtml(ex.message)+'</div>';
  }
  btn.disabled=false; btn.textContent='🏆 Run Arena';
}

var _arenaChosenText='';
function _renderArenaResult(bbl,d,elapsed){
  var results=d.results||[];
  var cards=results.map(function(r,i){
    return '<div class="arena-card" id="acard-'+i+'" onclick="_arenaSelect('+i+','+JSON.stringify(r.text).replace(/</g,'&lt;')+')">'+
      '<div class="arena-card-head"><span class="arena-card-emoji">'+escHtml(r.emoji||'🤖')+'</span>'+
      '<span class="arena-card-name">'+escHtml(r.model)+'</span>'+
      '<span class="arena-card-ms">'+r.ms+'ms</span></div>'+
      '<div class="arena-card-body">'+fmt(r.text)+'</div></div>';
  }).join('');
  bbl.innerHTML='<div class="arena-wrap">'+
    '<div class="arena-header"><span class="arena-badge">ARENA</span>'+
    '<span class="arena-prompt">'+escHtml((d.prompt||'').slice(0,80))+'</span>'+
    '<span style="font-size:10px;color:rgba(130,130,180,.5);margin-left:auto">'+results.length+' models · '+elapsed+'s</span></div>'+
    '<div style="font-size:11px;color:rgba(100,160,255,.6);margin:4px 0 8px">Click a card to select it, then use that answer.</div>'+
    '<div class="arena-grid">'+cards+'</div>'+
    '<div class="arena-actions"><button class="arena-use-btn" id="arenaUseBtn" onclick="_arenaUse()" disabled>Use Selected Answer</button></div></div>';
  _renderMath(bbl);
}

function _arenaSelect(i,text){
  document.querySelectorAll('.arena-card').forEach(function(c){c.classList.remove('selected');});
  var card=document.getElementById('acard-'+i);
  if(card) card.classList.add('selected');
  _arenaChosenText=text;
  var btn=document.getElementById('arenaUseBtn');
  if(btn){btn.disabled=false; btn.textContent='✓ Use This Answer';}
}
function _arenaUse(){
  if(!_arenaChosenText) return;
  // Add the chosen answer as an AI message
  var bbl2=addMsg('ai','',null,'Arena · Selected');
  bbl2.innerHTML=fmt(_arenaChosenText);
  hist.push({role:'assistant',content:_arenaChosenText});
  _renderMath(bbl2);
  showToast('✅ Answer added to chat');
}

// ══ Extreme Deep Think ═════════════════════════════════════════════════════
function openExtremeThink(){
  document.getElementById('edtPromptInput').value='';
  document.getElementById('edtModalOverlay').classList.remove('hidden');
  setTimeout(function(){document.getElementById('edtPromptInput').focus();},150);
}
function closeEDT(){document.getElementById('edtModalOverlay').classList.add('hidden');}
function edtClickOut(e){if(e.target.id==='edtModalOverlay')closeEDT();}

var _edtPhases=['🔍 Web Search (12×)','🤖 Round 1: Problem Analysis','🔬 Round 2: Hypotheses','⚡ Round 3: Challenge','🎯 Round 4: Integration','✍️ Round 5: Final Answer','🧬 Master Synthesis'];

async function runEDT(){
  var prompt=document.getElementById('edtPromptInput').value.trim();
  if(!prompt){showToast('Enter a question');return;}
  var btn=document.getElementById('edtRunBtn');
  btn.disabled=true; btn.textContent='⏳ Thinking…';
  closeEDT();
  addMsg('user',prompt);
  var bbl=addMsg('ai','',null,'Extreme Deep Think');
  var t0=Date.now();

  // Build loading UI
  var phaseRows=_edtPhases.map(function(ph,i){
    return '<div class="edt-loading-row" id="edt-ph-'+i+'"><div class="arena-spinner"></div><span>'+escHtml(ph)+'</span></div>';
  }).join('');
  bbl.innerHTML='<div class="edt-wrap">'+
    '<div class="edt-header"><span class="edt-badge">EXTREME</span>'+
    '<span class="edt-meta">Deep thinking in progress…</span></div>'+
    '<div class="edt-loading-phases">'+phaseRows+'</div></div>';

  // Animate phases while waiting
  var phaseTimer=0;
  _edtPhases.forEach(function(_,i){
    setTimeout(function(){
      var row=document.getElementById('edt-ph-'+i);
      if(row) row.classList.add('active');
      if(i>0){
        var prev=document.getElementById('edt-ph-'+(i-1));
        if(prev) prev.classList.add('done');
      }
    },i*8000);
  });

  try{
    var r=await apiFetch('/api/extreme-think',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({prompt:prompt})});
    var d=await r.json();
    var elapsed=((Date.now()-t0)/1000).toFixed(1);
    if(!d.ok){bbl.innerHTML='<div style="color:#f87171">Deep think failed.</div>';btn.disabled=false;btn.textContent='🧠 Start Deep Think';return;}
    _renderEDTResult(bbl,d,elapsed);
    hist.push({role:'assistant',content:d.synthesis||''});
  }catch(ex){
    bbl.innerHTML='<div style="color:#f87171">'+escHtml(ex.message)+'</div>';
  }
  btn.disabled=false; btn.textContent='🧠 Start Deep Think';
}

function _renderEDTResult(bbl,d,elapsed){
  var persp=(d.perspectives||[]).map(function(p){
    return '<div class="edt-persp-card">'+
      '<div class="edt-persp-head"><span style="font-size:14px">'+escHtml(p.emoji||'🤖')+'</span>'+
      '<span class="edt-persp-name">'+escHtml(p.model)+'</span>'+
      '<span class="edt-persp-rounds">'+p.rounds+' rounds</span></div>'+
      '<div class="edt-persp-body">'+escHtml(p.text.slice(0,300))+(p.text.length>300?'…':'')+'</div></div>';
  }).join('');

  bbl.innerHTML='<div class="edt-wrap">'+
    '<div class="edt-header"><span class="edt-badge">EXTREME</span>'+
    '<span class="edt-meta">'+d.models_used+' models · '+d.rounds+' rounds · '+d.web_searches+' web searches · '+elapsed+'s</span></div>'+
    '<div class="edt-stats">'+
      '<div class="edt-stat"><div class="edt-stat-val">'+d.web_searches+'</div><div class="edt-stat-lbl">Web Searches</div></div>'+
      '<div class="edt-stat"><div class="edt-stat-val">'+d.models_used+'</div><div class="edt-stat-lbl">Models</div></div>'+
      '<div class="edt-stat"><div class="edt-stat-val">'+d.rounds+'</div><div class="edt-stat-lbl">Think Rounds</div></div>'+
      '<div class="edt-stat"><div class="edt-stat-val">'+elapsed+'s</div><div class="edt-stat-lbl">Total Time</div></div>'+
    '</div>'+
    '<div style="font-size:11px;font-weight:700;color:rgba(100,160,255,.7);text-transform:uppercase;letter-spacing:.8px;margin:4px 0 6px">Master Synthesis</div>'+
    '<div class="edt-synthesis">'+fmt(d.synthesis||'No synthesis generated.')+'</div>'+
    '<div style="font-size:11px;font-weight:700;color:rgba(100,160,255,.7);text-transform:uppercase;letter-spacing:.8px;margin:10px 0 6px">Model Perspectives ('+persp.length+')</div>'+
    '<div class="edt-perspectives">'+persp+'</div>'+
    '<div style="display:flex;gap:8px;margin-top:10px;justify-content:flex-end">'+
      '<button onclick="navigator.clipboard.writeText('+JSON.stringify(d.synthesis||'').replace(/</g,'&lt;')+');showToast(\'Copied\')" style="background:rgba(6,182,212,.12);border:1px solid rgba(6,182,212,.3);border-radius:9px;padding:6px 14px;font-size:11px;color:rgba(100,180,255,.9);cursor:pointer;font-family:\'DM Sans\',sans-serif">📋 Copy Synthesis</button>'+
    '</div></div>';
  _renderMath(bbl);
}


// ══ Conversations Sidebar ══════════════════════════════════════════════════
var _convSbData=[];
function openConvSidebar(){
  document.getElementById('convSidebar').classList.remove('closed');
  document.getElementById('convSidebar').classList.add('open');
  document.getElementById('convSbOverlay').classList.add('open');
  var cp=document.getElementById('chatPage'); if(cp) cp.classList.remove('sidebar-closed');
  _loadConvSidebar();
}
function closeConvSidebar(){
  document.getElementById('convSidebar').classList.add('closed');
  document.getElementById('convSidebar').classList.remove('open');
  document.getElementById('convSbOverlay').classList.remove('open');
  var cp=document.getElementById('chatPage'); if(cp) cp.classList.add('sidebar-closed');
}
function toggleConvSidebar(){
  var sb=document.getElementById('convSidebar');
  var isMobile=window.innerWidth<=768;
  var isVisible=isMobile?sb.classList.contains('open'):!sb.classList.contains('closed');
  if(isVisible) closeConvSidebar(); else openConvSidebar();
}
async function _loadConvSidebar(){
  var list=document.getElementById('convSbList');
  if(!authToken){list.innerHTML='<div style="padding:20px;text-align:center;font-size:11px;color:rgba(80,130,200,.5)">Sign in to see conversations</div>';return;}
  try{
    var r=await apiFetch('/api/conversations');
    var d=await r.json();
    _convSbData=d.conversations||[];
    _renderConvSidebar(_convSbData);
  }catch(e){
    list.innerHTML='<div style="padding:20px;text-align:center;font-size:11px;color:rgba(200,80,80,.6)">Failed to load</div>';
  }
}
function _renderConvSidebar(convs){
  var list=document.getElementById('convSbList');
  if(!convs.length){
    list.innerHTML='<div style="padding:30px 20px;text-align:center;font-size:12px;color:rgba(60,100,160,.5)">No conversations yet.<br><small>Start chatting to create one!</small></div>';
    return;
  }
  // Group by day
  var groups={};
  var now=new Date(); var today=now.toDateString();
  var yesterday=new Date(now-86400000).toDateString();
  convs.forEach(function(c){
    var d=new Date(c.updated||c.created); var ds=d.toDateString();
    var label=ds===today?'Today':ds===yesterday?'Yesterday':d.toLocaleDateString('en',{month:'short',day:'numeric'});
    if(!groups[label]) groups[label]=[];
    groups[label].push(c);
  });
  var html='';
  Object.keys(groups).forEach(function(day){
    html+='<div class="conv-sb-day">'+escHtml(day)+'</div>';
    groups[day].forEach(function(c){
      var icon=c.title.match(/code|python|js|script/i)?'💻':c.title.match(/image|photo|picture/i)?'🖼':c.title.match(/data|csv|excel/i)?'📊':c.title.match(/write|essay|doc/i)?'📝':'💬';
      html+='<div class="conv-sb-item'+(currentConvId===c.id?' active':'')+'" onclick="loadConvFromSidebar('+c.id+')">'
        +'<span class="conv-sb-item-icon">'+icon+'</span>'
        +'<div class="conv-sb-item-body">'
        +'<div class="conv-sb-item-title">'+escHtml((c.title||'Chat').slice(0,40))+'</div>'
        +'<div class="conv-sb-item-preview">'+escHtml((c.preview||'').slice(0,45))+'</div>'
        +'</div>'
        +'<button class="conv-sb-item-del" onclick="event.stopPropagation();deleteConvSb('+c.id+')" title="Delete">🗑</button>'
        +'</div>';
    });
  });
  list.innerHTML=html;
}
function filterConvs(q){
  if(!q) return _renderConvSidebar(_convSbData);
  var filtered=_convSbData.filter(function(c){return (c.title||'').toLowerCase().includes(q.toLowerCase());});
  _renderConvSidebar(filtered);
}
async function loadConvFromSidebar(id){
  if(window.innerWidth<=768) closeConvSidebar();
  // Load conversation into chat
  currentConvId=id;
  hist=[];
  var cb=document.getElementById('chatBody'); if(cb) cb.innerHTML='';
  try{
    var r=await apiFetch('/api/conversations/'+id+'/messages');
    var d=await r.json();
    (d.messages||[]).forEach(function(m){
      addMsg(m.role,m.content,null,m.model);
      hist.push({role:m.role,content:m.content});
    });
  }catch(e){showToast('Failed to load conversation');}
}
async function deleteConvSb(id){
  try{
    await apiFetch('/api/conversations/'+id,{method:'DELETE'});
    _convSbData=_convSbData.filter(function(c){return c.id!==id;});
    _renderConvSidebar(_convSbData);
    if(currentConvId===id){newConv();currentConvId=null;}
  }catch(e){showToast('Delete failed');}
}

// ══ Profile Panel ══════════════════════════════════════════════════════════
var _profileEmoji='😊';
var _AVATAR_EMOJIS=['🦊','🐺','🦁','🐯','🦝','🐻','🦄','🐸','🦋','🌊','⚡','🌙','🔥','💎','🌸','🎯','🚀','🎸','🌈','🏔','🌺','🦅','🐉','🦉','🎭','🧊','🌴','🦈','🦂','🎪','👾','🤖','🧬','🔮','🎲','🦚','🌋','🐋','🦊','🍄'];

function openProfile(){
  document.getElementById('profilePanel').classList.remove('hidden');
  _loadProfileData();
}
function closeProfile(){
  document.getElementById('profilePanel').classList.add('hidden');
  document.getElementById('emojiPickerWrap').style.display='none';
}
async function _loadProfileData(){
  try{
    var r=await apiFetch('/api/me');
    var d=await r.json();
    if(!d.logged_in){closeProfile();return;}
    _profileEmoji=d.avatar_emoji||d.username[0].toUpperCase();
    document.getElementById('profileName').textContent=d.display_name||d.username;
    document.getElementById('profileGid').textContent=d.guest_id||'—';
    document.getElementById('profileAvatarBig').textContent=_profileEmoji;
    // Update header avatar too
    var av=document.getElementById('uAv');
    if(av){av.textContent=_profileEmoji;av.style.fontSize='16px';av.style.lineHeight='1';}
    // Load stats
    try{
      var r2=await apiFetch('/api/conversations');
      var d2=await r2.json();
      document.getElementById('pStatConvs').textContent=(d2.conversations||[]).length;
    }catch(e){}
    try{
      var r3=await apiFetch('/api/memory');
      var d3=await r3.json();
      document.getElementById('pStatMem').textContent=(d3.memory||[]).length;
    }catch(e){}
    document.getElementById('pStatMsgs').textContent=hist.length;
  }catch(e){}
}
function toggleEmojiPicker(){
  var wrap=document.getElementById('emojiPickerWrap');
  var grid=document.getElementById('emojiGrid');
  if(wrap.style.display==='none'||!wrap.style.display){
    // Build grid
    if(grid.children.length===0){
      _AVATAR_EMOJIS.forEach(function(em){
        var btn=document.createElement('button');
        btn.className='emoji-opt'+(em===_profileEmoji?' selected':'');
        btn.textContent=em;
        btn.onclick=function(){
          _setAvatar(em);
          document.querySelectorAll('.emoji-opt').forEach(function(b){b.classList.remove('selected');});
          btn.classList.add('selected');
        };
        grid.appendChild(btn);
      });
    }
    wrap.style.display='block';
  } else {
    wrap.style.display='none';
  }
}
async function _setAvatar(emoji){
  _profileEmoji=emoji;
  document.getElementById('profileAvatarBig').textContent=emoji;
  var av=document.getElementById('uAv');
  if(av){av.textContent=emoji;av.style.fontSize='16px';}
  showToast('Avatar updated!');
  try{
    await apiFetch('/api/update_avatar',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({emoji:emoji})});
  }catch(e){}
}

// ══ AI Tools Hub ════════════════════════════════════════════════════════════
var _TOOLS=[
  {id:'translate',icon:'🌍',name:'Translator',desc:'Detect & translate any language instantly'},
  {id:'summarise',icon:'📋',name:'Summariser',desc:'Condense text into key points or TL;DR'},
  {id:'rewrite',  icon:'✍️',name:'Rewriter',  desc:'Change tone: formal, casual, concise'},
  {id:'quiz',     icon:'🧠',name:'Quiz Maker', desc:'Generate interactive quizzes on any topic'},
  {id:'grammar',  icon:'✅',name:'Grammar Fix', desc:'Fix spelling, grammar and punctuation'},
  {id:'hashtags', icon:'#️⃣',name:'Hashtags',   desc:'Generate hashtags for social media posts'},
  {id:'email',    icon:'📧',name:'Email Draft', desc:'Write professional emails from bullet points'},
  {id:'code_explain',icon:'💻',name:'Code Explainer',desc:'Paste code, get plain English explanation'},
];
var _activeToolId=null;

function openToolsHub(){
  document.getElementById('toolsHubOverlay').classList.remove('hidden');
  document.getElementById('toolsHubTitle').textContent='🛠 AI Tools';
  _showToolGrid();
}
function closeToolsHub(){document.getElementById('toolsHubOverlay').classList.add('hidden');_activeToolId=null;}
function toolsHubClickOut(e){if(e.target.id==='toolsHubOverlay')closeToolsHub();}

function _showToolGrid(){
  _activeToolId=null;
  document.getElementById('toolsHubTitle').textContent='🛠 AI Tools';
  var body=document.getElementById('toolsHubBody');
  body.innerHTML='<div class="tools-hub-grid">'
    +_TOOLS.map(function(t){
      return '<div class="tool-card" onclick="openTool(\''+t.id+'\')">'
        +'<div class="tool-card-icon">'+t.icon+'</div>'
        +'<div class="tool-card-name">'+t.name+'</div>'
        +'<div class="tool-card-desc">'+t.desc+'</div>'
        +'</div>';
    }).join('')+'</div>';
}

function openTool(id){
  _activeToolId=id;
  var t=_TOOLS.find(function(x){return x.id===id;});
  if(!t) return;
  document.getElementById('toolsHubTitle').textContent=t.icon+' '+t.name;
  var body=document.getElementById('toolsHubBody');
  var backBtn='<div class="tool-panel-back" onclick="_showToolGrid()">← All Tools</div>';

  if(id==='translate'){
    body.innerHTML=backBtn+
      '<div class="tool-panel active">'+
      '<textarea class="tool-textarea" id="toolTxt" placeholder="Paste text to translate…" rows="5"></textarea>'+
      '<div style="display:flex;gap:8px;align-items:center">'+
      '<label style="font-size:11px;color:rgba(80,130,200,.7)">To:</label>'+
      '<select class="tool-select" id="toolLang">'+
      ['English','Arabic','French','Spanish','German','Chinese','Japanese','Italian','Portuguese','Hindi','Russian','Korean','Dutch','Turkish','Swedish'].map(function(l){return '<option>'+l+'</option>';}).join('')+
      '</select>'+
      '<button class="tool-run-btn" id="toolRunBtn" onclick="_runTranslate()">Translate</button></div>'+
      '<div id="toolResult" style="display:none" class="tool-result-box"><button class="tool-result-copy" onclick="_copyToolResult()">📋</button><span id="toolResultTxt"></span></div>'+
      '</div>';
  } else if(id==='summarise'){
    body.innerHTML=backBtn+
      '<div class="tool-panel active">'+
      '<textarea class="tool-textarea" id="toolTxt" placeholder="Paste text to summarise…" rows="6"></textarea>'+
      '<div style="display:flex;gap:8px;align-items:center">'+
      '<label style="font-size:11px;color:rgba(80,130,200,.7)">Style:</label>'+
      '<select class="tool-select" id="toolStyle">'+
      '<option value="bullet">Bullet Points</option><option value="paragraph">Paragraphs</option><option value="eli5">ELI5</option><option value="tldr">TL;DR</option>'+
      '</select>'+
      '<button class="tool-run-btn" id="toolRunBtn" onclick="_runSummarise()">Summarise</button></div>'+
      '<div id="toolResult" style="display:none" class="tool-result-box"><button class="tool-result-copy" onclick="_copyToolResult()">📋</button><span id="toolResultTxt"></span></div>'+
      '</div>';
  } else if(id==='rewrite'){
    body.innerHTML=backBtn+
      '<div class="tool-panel active">'+
      '<textarea class="tool-textarea" id="toolTxt" placeholder="Paste text to rewrite…" rows="5"></textarea>'+
      '<div style="display:flex;gap:8px;align-items:center">'+
      '<label style="font-size:11px;color:rgba(80,130,200,.7)">Tone:</label>'+
      '<select class="tool-select" id="toolTone">'+
      '<option value="professional">Professional</option><option value="casual">Casual</option><option value="concise">Concise</option><option value="creative">Creative</option><option value="academic">Academic</option><option value="persuasive">Persuasive</option>'+
      '</select>'+
      '<button class="tool-run-btn" id="toolRunBtn" onclick="_runRewrite()">Rewrite</button></div>'+
      '<div id="toolResult" style="display:none" class="tool-result-box"><button class="tool-result-copy" onclick="_copyToolResult()">📋</button><span id="toolResultTxt"></span></div>'+
      '</div>';
  } else if(id==='quiz'){
    body.innerHTML=backBtn+
      '<div class="tool-panel active">'+
      '<div style="font-size:12px;color:rgba(100,160,220,.7);margin-bottom:8px">Generate an interactive quiz on any topic.</div>'+
      '<input type="text" class="tool-textarea" id="toolTxt" placeholder="Topic — e.g. World History, Python, Astronomy…" style="min-height:unset;padding:10px 13px;border-radius:12px">'+
      '<div style="display:flex;gap:8px;align-items:center">'+
      '<label style="font-size:11px;color:rgba(80,130,200,.7)">Questions:</label>'+
      '<select class="tool-select" id="toolCount"><option value="5">5</option><option value="8">8</option><option value="10">10</option></select>'+
      '<button class="tool-run-btn" id="toolRunBtn" onclick="_runQuiz()">Generate Quiz</button></div>'+
      '<div id="quizArea" style="display:none"></div>'+
      '</div>';
  } else if(id==='grammar'){
    body.innerHTML=backBtn+
      '<div class="tool-panel active">'+
      '<textarea class="tool-textarea" id="toolTxt" placeholder="Paste text to fix grammar, spelling and punctuation…" rows="6"></textarea>'+
      '<div style="text-align:right"><button class="tool-run-btn" id="toolRunBtn" onclick="_runGrammar()">Fix Grammar</button></div>'+
      '<div id="toolResult" style="display:none" class="tool-result-box"><button class="tool-result-copy" onclick="_copyToolResult()">📋</button><span id="toolResultTxt"></span></div>'+
      '</div>';
  } else if(id==='hashtags'){
    body.innerHTML=backBtn+
      '<div class="tool-panel active">'+
      '<textarea class="tool-textarea" id="toolTxt" placeholder="Describe your post or paste its text…" rows="4"></textarea>'+
      '<div style="display:flex;gap:8px;align-items:center">'+
      '<label style="font-size:11px;color:rgba(80,130,200,.7)">Platform:</label>'+
      '<select class="tool-select" id="toolPlatform"><option>Instagram</option><option>Twitter/X</option><option>LinkedIn</option><option>TikTok</option></select>'+
      '<button class="tool-run-btn" id="toolRunBtn" onclick="_runHashtags()">Generate</button></div>'+
      '<div id="toolResult" style="display:none" class="tool-result-box"><button class="tool-result-copy" onclick="_copyToolResult()">📋</button><span id="toolResultTxt"></span></div>'+
      '</div>';
  } else if(id==='email'){
    body.innerHTML=backBtn+
      '<div class="tool-panel active">'+
      '<textarea class="tool-textarea" id="toolTxt" placeholder="Bullet points of what the email should say…\n• Recipient: client\n• Topic: project delay\n• Tone: apologetic but professional" rows="5"></textarea>'+
      '<div style="text-align:right"><button class="tool-run-btn" id="toolRunBtn" onclick="_runEmail()">Write Email</button></div>'+
      '<div id="toolResult" style="display:none" class="tool-result-box"><button class="tool-result-copy" onclick="_copyToolResult()">📋</button><span id="toolResultTxt"></span></div>'+
      '</div>';
  } else if(id==='code_explain'){
    body.innerHTML=backBtn+
      '<div class="tool-panel active">'+
      '<textarea class="tool-textarea" id="toolTxt" placeholder="Paste your code here…" rows="7" style="font-family:\'DM Mono\',monospace;font-size:12px"></textarea>'+
      '<div style="text-align:right"><button class="tool-run-btn" id="toolRunBtn" onclick="_runCodeExplain()">Explain Code</button></div>'+
      '<div id="toolResult" style="display:none" class="tool-result-box"><button class="tool-result-copy" onclick="_copyToolResult()">📋</button><span id="toolResultTxt"></span></div>'+
      '</div>';
  }
}

function _showToolLoading(){
  var btn=document.getElementById('toolRunBtn');
  if(btn){btn.disabled=true;btn.textContent='⏳ Working…';}
}
function _showToolDone(btnTxt){
  var btn=document.getElementById('toolRunBtn');
  if(btn){btn.disabled=false;btn.textContent=btnTxt||'Run';}
}
function _showToolResult(text){
  var box=document.getElementById('toolResult');
  var txt=document.getElementById('toolResultTxt');
  if(box&&txt){box.style.display='block';txt.textContent=text;}
}
function _copyToolResult(){
  var txt=document.getElementById('toolResultTxt');
  if(txt) navigator.clipboard.writeText(txt.textContent).then(function(){showToast('Copied!');});
}

async function _runTranslate(){
  var text=document.getElementById('toolTxt').value.trim();
  var lang=document.getElementById('toolLang').value;
  if(!text){showToast('Paste some text first');return;}
  _showToolLoading();
  try{
    var r=await apiFetch('/api/ai_tools/translate',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({text:text,target:lang})});
    var d=await r.json();
    _showToolResult(d.translation||d.error||'Error');
  }catch(e){_showToolResult('Error: '+e.message);}
  _showToolDone('Translate');
}
async function _runSummarise(){
  var text=document.getElementById('toolTxt').value.trim();
  var style=document.getElementById('toolStyle').value;
  if(!text){showToast('Paste some text first');return;}
  _showToolLoading();
  try{
    var r=await apiFetch('/api/ai_tools/summarise',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({text:text,style:style})});
    var d=await r.json();
    _showToolResult(d.summary||d.error||'Error');
  }catch(e){_showToolResult('Error: '+e.message);}
  _showToolDone('Summarise');
}
async function _runRewrite(){
  var text=document.getElementById('toolTxt').value.trim();
  var tone=document.getElementById('toolTone').value;
  if(!text){showToast('Paste some text first');return;}
  _showToolLoading();
  try{
    var r=await apiFetch('/api/ai_tools/rewrite',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({text:text,tone:tone})});
    var d=await r.json();
    _showToolResult(d.result||d.error||'Error');
  }catch(e){_showToolResult('Error: '+e.message);}
  _showToolDone('Rewrite');
}

// Quiz
var _quizData=[];var _quizIdx=0;var _quizScore=0;var _quizAnswered=false;
async function _runQuiz(){
  var topic=document.getElementById('toolTxt').value.trim()||'General Knowledge';
  var count=parseInt(document.getElementById('toolCount').value)||5;
  _showToolLoading();
  try{
    var r=await apiFetch('/api/ai_tools/quiz',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({topic:topic,count:count})});
    var d=await r.json();
    if(!d.ok||!d.questions){_showToolResult('Quiz failed. Try a different topic.');_showToolDone('Generate Quiz');return;}
    _quizData=d.questions;_quizIdx=0;_quizScore=0;_quizAnswered=false;
    document.getElementById('quizArea').style.display='block';
    document.getElementById('toolRunBtn').style.display='none';
    _renderQuizQ();
  }catch(e){_showToolResult('Error: '+e.message);}
  _showToolDone('Generate Quiz');
}
function _renderQuizQ(){
  var area=document.getElementById('quizArea');
  if(_quizIdx>=_quizData.length){
    area.innerHTML='<div class="quiz-score">🏆 '+_quizScore+' / '+_quizData.length+'<br><small style="font-size:12px;font-weight:400;color:rgba(150,200,255,.7)">Quiz complete!</small></div>'+
      '<div style="text-align:center"><button class="tool-run-btn" onclick="_showToolGrid()">Back to Tools</button></div>';
    return;
  }
  var q=_quizData[_quizIdx];
  area.innerHTML='<div class="quiz-progress">Question '+(_quizIdx+1)+' of '+_quizData.length+' · Score: '+_quizScore+'</div>'+
    '<div class="quiz-q">'+escHtml(q.q)+'</div>'+
    '<div class="quiz-opts">'+
    (q.opts||[]).map(function(opt,i){
      return '<button class="quiz-opt-btn" onclick="_quizAnswer('+i+')">'+escHtml(opt)+'</button>';
    }).join('')+'</div>';
}
function _quizAnswer(idx){
  var q=_quizData[_quizIdx];
  var btns=document.querySelectorAll('.quiz-opt-btn');
  btns.forEach(function(b,i){b.disabled=true;if(i===q.ans)b.classList.add('correct');else if(i===idx)b.classList.add('wrong');});
  if(idx===q.ans)_quizScore++;
  setTimeout(function(){_quizIdx++;_renderQuizQ();},1000);
}

async function _runGrammar(){
  var text=document.getElementById('toolTxt').value.trim();
  if(!text){showToast('Paste some text first');return;}
  _showToolLoading();
  try{
    var r=await apiFetch('/api/chat',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({messages:[{role:'user',content:'Fix all grammar, spelling and punctuation errors in the following text. Return ONLY the corrected text:\n\n'+text}],stream:false})});
    var d=await r.json();
    _showToolResult(d.content||d.error||'Error');
  }catch(e){_showToolResult('Error: '+e.message);}
  _showToolDone('Fix Grammar');
}
async function _runHashtags(){
  var text=document.getElementById('toolTxt').value.trim();
  var plat=document.getElementById('toolPlatform').value;
  if(!text){showToast('Describe your post first');return;}
  _showToolLoading();
  try{
    var prompt='Generate 15-20 relevant hashtags for '+plat+' for this post: '+text+'. Return ONLY hashtags, space-separated, starting with #.';
    var r=await apiFetch('/api/chat',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({messages:[{role:'user',content:prompt}],stream:false})});
    var d=await r.json();
    _showToolResult(d.content||d.error||'Error');
  }catch(e){_showToolResult('Error: '+e.message);}
  _showToolDone('Generate');
}
async function _runEmail(){
  var text=document.getElementById('toolTxt').value.trim();
  if(!text){showToast('Add some bullet points');return;}
  _showToolLoading();
  try{
    var prompt='Write a professional email based on these notes:\n'+text+'\nReturn ONLY the complete email with Subject line.';
    var r=await apiFetch('/api/chat',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({messages:[{role:'user',content:prompt}],stream:false})});
    var d=await r.json();
    _showToolResult(d.content||d.error||'Error');
  }catch(e){_showToolResult('Error: '+e.message);}
  _showToolDone('Write Email');
}
async function _runCodeExplain(){
  var text=document.getElementById('toolTxt').value.trim();
  if(!text){showToast('Paste some code first');return;}
  _showToolLoading();
  try{
    var prompt='Explain this code in plain English. Describe what it does, how it works, and any potential issues:\n\n'+text;
    var r=await apiFetch('/api/chat',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({messages:[{role:'user',content:prompt}],stream:false})});
    var d=await r.json();
    _showToolResult(d.content||d.error||'Error');
  }catch(e){_showToolResult('Error: '+e.message);}
  _showToolDone('Explain Code');
}

// ══ Quick send helper ══════════════════════════════════════════════════════
function _retryLast(){
  if(hist.length>0){
    var lastUser=null;
    for(var i=hist.length-1;i>=0;i--){if(hist[i].role==='user'){lastUser=hist[i].content;break;}}
    if(lastUser){
      // Remove the last AI message from DOM and hist
      var msgs=document.querySelectorAll('.msg.ai');
      if(msgs.length) msgs[msgs.length-1].remove();
      hist=hist.filter(function(h){return h.role!=='assistant'||(hist.indexOf(h)<hist.length-1);});
      var inp=document.getElementById('msgIn');
      if(inp){inp.value=lastUser;sendMsg();}
    }
  }
}
function sendQuickMsg(prefix){
  var inp=document.getElementById('msgIn');
  if(!inp) return;
  var existing=inp.value.trim();
  if(existing) inp.value=prefix+': '+existing;
  else inp.value=prefix+': ';
  inp.focus();
  inp.setSelectionRange(inp.value.length,inp.value.length);
}

// ══ Load avatar on init ═══════════════════════════════════════════════════
function _initAvatar(){
  if(!authToken) return;
  apiFetch('/api/me').then(function(r){return r.json();}).then(function(d){
    if(d.logged_in&&d.avatar_emoji){
      var av=document.getElementById('uAv');
      if(av){av.textContent=d.avatar_emoji;av.style.fontSize='16px';av.style.lineHeight='1';}
      _profileEmoji=d.avatar_emoji;
    }
  }).catch(function(){});
}

// ESC key closes modals
document.addEventListener('keydown',function(e){
  if(e.key==='Escape'){
    closeFileCreator();
    closeToolbox();
    closeArena();
    closeEDT();
    closeProfile();
    closeToolsHub();
    closeConvSidebar();
  }
});

</script>

<style>
/* ═══ FusionOS + Computer + SVG modal shared ═══════════════════════════════ */
.modal-overlay{position:fixed;inset:0;z-index:9000;background:rgba(2,4,12,.82);backdrop-filter:blur(12px);display:flex;align-items:center;justify-content:center;padding:16px;box-sizing:border-box}
.modal-overlay.hidden{display:none}
.comp-chip{background:rgba(6,182,212,.12);border:1px solid rgba(6,182,212,.25);border-radius:20px;padding:4px 11px;font-size:11px;color:var(--tx2);cursor:pointer;transition:all .14s;white-space:nowrap}
.comp-chip:hover{background:rgba(6,182,212,.22);color:var(--tx)}
/* ═══ FusionOS desktop ══════════════════════════════════════════════════════ */
#fos-overlay{position:fixed;inset:0;z-index:9500;display:none;flex-direction:column;background:#050710;font-family:'Segoe UI',system-ui,sans-serif}
#fos-overlay.fos-on{display:flex}
#fos-bar{height:26px;background:rgba(6,8,18,.96);backdrop-filter:blur(20px);border-bottom:1px solid rgba(255,255,255,.055);display:flex;align-items:center;padding:0 10px;gap:0;flex-shrink:0;user-select:none}
.fosb-logo{font-size:12px;font-weight:800;background:linear-gradient(120deg,#4a9eff,#a855f7);-webkit-background-clip:text;-webkit-text-fill-color:transparent;margin-right:14px;cursor:pointer;letter-spacing:.3px}
.fosb-item{font-size:11px;color:rgba(255,255,255,.65);padding:2px 9px;border-radius:5px;cursor:pointer;transition:background .1s}
.fosb-item:hover{background:rgba(255,255,255,.1);color:#fff}
.fosb-right{margin-left:auto;display:flex;align-items:center;gap:12px}
.fosb-clock{font-size:11px;color:rgba(255,255,255,.45);font-variant-numeric:tabular-nums}
#fos-desktop{flex:1;position:relative;overflow:hidden;background:radial-gradient(ellipse at 28% 22%,#0c1840 0%,#04060f 60%,#000 100%)}
#fos-desktop::before{content:'';position:absolute;inset:0;background-image:radial-gradient(circle,rgba(255,255,255,.025) 1px,transparent 1px);background-size:28px 28px;pointer-events:none}
#fos-dock{height:62px;background:rgba(6,8,18,.9);backdrop-filter:blur(28px);border-top:1px solid rgba(255,255,255,.055);display:flex;align-items:center;justify-content:center;gap:4px;padding:0 16px;flex-shrink:0}
.fos-dck{display:flex;flex-direction:column;align-items:center;gap:2px;cursor:pointer;padding:5px 9px;border-radius:11px;transition:all .15s;border:none;background:none;position:relative}
.fos-dck:hover{background:rgba(255,255,255,.08);transform:translateY(-4px)}
.fos-dck-ico{font-size:24px;transition:transform .15s}
.fos-dck:hover .fos-dck-ico{transform:scale(1.18)}
.fos-dck-lbl{font-size:9px;color:rgba(255,255,255,.4);white-space:nowrap}
.fos-dck-dot{position:absolute;bottom:2px;left:50%;transform:translateX(-50%);width:4px;height:4px;border-radius:50%;background:#4a9eff;display:none}
.fos-dck.fos-running .fos-dck-dot{display:block}
/* Windows */
.fos-win{position:absolute;background:rgba(10,14,26,.97);border:1px solid rgba(255,255,255,.09);border-radius:12px;box-shadow:0 20px 55px rgba(0,0,0,.85);display:flex;flex-direction:column;overflow:hidden;min-width:280px;min-height:180px;animation:fosIn .16s ease-out}
@keyframes fosIn{from{opacity:0;transform:scale(.95) translateY(6px)}to{opacity:1;transform:none}}
.fos-win.fos-focused{border-color:rgba(74,158,255,.25);box-shadow:0 24px 70px rgba(0,0,0,.9),0 0 0 1px rgba(74,158,255,.1)}
.fos-win.fos-mini{display:none}
.fos-titlebar{height:34px;background:rgba(12,16,30,.99);display:flex;align-items:center;padding:0 11px;gap:8px;cursor:move;flex-shrink:0;border-bottom:1px solid rgba(255,255,255,.046);user-select:none}
.fos-traf{display:flex;gap:5px}.fos-tb{width:11px;height:11px;border-radius:50%;border:none;cursor:pointer;flex-shrink:0;outline:none;transition:filter .1s}.fos-tb:hover{filter:brightness(1.5)}
.fos-tbc{background:#ff5f56}.fos-tbm{background:#ffbd2e}.fos-tbx{background:#27c93f}
.fos-wico{font-size:13px;flex-shrink:0}.fos-wtitle{font-size:11.5px;color:rgba(255,255,255,.68);font-weight:600;flex:1;white-space:nowrap;overflow:hidden;text-overflow:ellipsis}
.fos-body{flex:1;overflow:hidden;display:flex;flex-direction:column;position:relative}
/* Resize handles */
.fos-rz{position:absolute;z-index:5}
.fos-rz-n{top:-3px;left:8px;right:8px;height:6px;cursor:n-resize}
.fos-rz-s{bottom:-3px;left:8px;right:8px;height:6px;cursor:s-resize}
.fos-rz-e{top:8px;right:-3px;bottom:8px;width:6px;cursor:e-resize}
.fos-rz-w{top:8px;left:-3px;bottom:8px;width:6px;cursor:w-resize}
.fos-rz-ne{top:-3px;right:-3px;width:12px;height:12px;cursor:ne-resize}
.fos-rz-nw{top:-3px;left:-3px;width:12px;height:12px;cursor:nw-resize}
.fos-rz-sw{bottom:-3px;left:-3px;width:12px;height:12px;cursor:sw-resize}
.fos-rz-se{bottom:-3px;right:-3px;width:12px;height:12px;cursor:se-resize;opacity:.35;display:flex;align-items:flex-end;justify-content:flex-end;padding:2px;color:rgba(255,255,255,.5);font-size:11px}
/* Terminal */
.fos-term{background:#06070e;height:100%;display:flex;flex-direction:column}
.fos-term-out{flex:1;overflow-y:auto;padding:10px 13px;line-height:1.65;white-space:pre-wrap;word-break:break-all;font-family:'Cascadia Code','Fira Code','DM Mono',monospace;font-size:12.5px;color:#c8ddf0}
.fos-term-out::-webkit-scrollbar{width:3px}.fos-term-out::-webkit-scrollbar-thumb{background:rgba(255,255,255,.1);border-radius:2px}
.fos-term-row{display:flex;align-items:center;padding:5px 11px;border-top:1px solid rgba(255,255,255,.04);flex-shrink:0;gap:6px}
.fos-prompt{color:#4a9eff;font-weight:700;font-size:12.5px;white-space:nowrap;font-family:inherit}
.fos-inp{flex:1;background:none;border:none;outline:none;color:#c8ddf0;font-family:inherit;font-size:12.5px;caret-color:#4a9eff}
.fos-cmd-echo{color:#7ec8e3}.fos-ok-txt{color:#7ee8a2}.fos-err-txt{color:#ff8484}.fos-ai-txt{color:#c084fc;font-style:italic}.fos-dim{color:rgba(180,200,240,.38)}
/* Files */
.fos-files{display:flex;height:100%}
.fos-fsb{width:138px;background:rgba(5,7,14,.65);border-right:1px solid rgba(255,255,255,.046);padding:7px 5px;overflow-y:auto;flex-shrink:0;display:flex;flex-direction:column;gap:1px}
.fos-fmain{flex:1;display:flex;flex-direction:column;overflow:hidden}
.fos-fbar{padding:6px 9px;border-bottom:1px solid rgba(255,255,255,.046);display:flex;align-items:center;gap:6px;flex-shrink:0}
.fos-flist{flex:1;overflow-y:auto;padding:5px}
.fos-fitem{display:flex;align-items:center;gap:7px;padding:5px 7px;border-radius:7px;cursor:pointer;transition:background .1s}
.fos-fitem:hover{background:rgba(255,255,255,.05)}.fos-fitem.fos-sel{background:rgba(74,158,255,.14)}
.fos-fname{font-size:11.5px;color:#b8cce8;flex:1;overflow:hidden;text-overflow:ellipsis;white-space:nowrap}
.fos-fmeta{font-size:9px;color:rgba(255,255,255,.22);white-space:nowrap}
.fos-sbtn{display:flex;align-items:center;gap:6px;padding:5px 7px;border-radius:7px;font-size:11px;color:rgba(255,255,255,.48);cursor:pointer;border:none;background:none;width:100%;text-align:left;transition:all .1s}
.fos-sbtn:hover{background:rgba(255,255,255,.07);color:#fff}
/* Editor */
.fos-editor{display:flex;flex-direction:column;height:100%}
.fos-ebar{display:flex;align-items:center;padding:5px 8px;gap:6px;border-bottom:1px solid rgba(255,255,255,.046);flex-shrink:0;background:rgba(7,9,17,.85);flex-wrap:wrap}
.fos-eta{flex:1;background:#05060c;color:#c8ddf0;font-family:'Cascadia Code','Fira Code','DM Mono',monospace;font-size:12.5px;border:none;outline:none;resize:none;padding:11px 13px;line-height:1.65;tab-size:2}
.fos-sbar{height:21px;background:rgba(5,6,12,.9);border-top:1px solid rgba(255,255,255,.04);display:flex;align-items:center;padding:0 11px;flex-shrink:0}
/* Agent */
.fos-agent{display:flex;flex-direction:column;height:100%;background:#060810}
.fos-alog{flex:1;overflow-y:auto;padding:10px;display:flex;flex-direction:column;gap:7px}
.fos-alog::-webkit-scrollbar{width:3px}.fos-alog::-webkit-scrollbar-thumb{background:rgba(255,255,255,.08);border-radius:2px}
.fos-astep{background:rgba(255,255,255,.028);border:1px solid rgba(255,255,255,.065);border-radius:8px;padding:8px 10px;font-size:11.5px;animation:fosIn .15s}
.fos-astep.fos-ok{border-color:rgba(78,230,120,.18)}.fos-astep.fos-err{border-color:rgba(255,100,100,.2)}
.fos-acmd{font-family:'Cascadia Code',monospace;font-size:10.5px;color:#7ec8e3;margin-bottom:3px;word-break:break-all}
.fos-aout{color:rgba(188,210,248,.52);white-space:pre-wrap;max-height:110px;overflow-y:auto;font-size:10.5px;font-family:'Cascadia Code',monospace;line-height:1.5}
.fos-afooter{padding:9px 10px;border-top:1px solid rgba(255,255,255,.046);flex-shrink:0;display:flex;flex-direction:column;gap:6px}
/* Monitor */
.fos-mon{padding:12px;display:flex;flex-direction:column;gap:8px;overflow-y:auto;height:100%;background:#050710}
.fos-mcard{background:rgba(255,255,255,.028);border:1px solid rgba(255,255,255,.065);border-radius:8px;padding:9px 11px}
.fos-mlbl{font-size:9px;color:rgba(255,255,255,.32);font-weight:700;letter-spacing:.7px;text-transform:uppercase;margin-bottom:4px}
.fos-mval{font-size:11px;color:#b8ccff;font-family:'Cascadia Code',monospace;white-space:pre;overflow:hidden}
/* Browser */
.fos-browser{display:flex;flex-direction:column;height:100%}
.fos-bbar{display:flex;align-items:center;gap:6px;padding:6px 8px;border-bottom:1px solid rgba(255,255,255,.046);flex-shrink:0;background:rgba(7,9,17,.85)}
.fos-bfavs{display:flex;gap:5px;padding:4px 8px;border-bottom:1px solid rgba(255,255,255,.036);flex-shrink:0;flex-wrap:wrap}
.fos-frame{flex:1;border:none;background:#fff}
/* Context menu */
.fos-ctx{position:fixed;background:rgba(12,16,28,.97);backdrop-filter:blur(18px);border:1px solid rgba(255,255,255,.09);border-radius:9px;padding:4px;z-index:99999;min-width:148px;box-shadow:0 10px 36px rgba(0,0,0,.8)}
.fos-ctx-item{padding:6px 12px;font-size:11.5px;color:rgba(255,255,255,.75);border-radius:6px;cursor:pointer;transition:background .1s}
.fos-ctx-item:hover{background:rgba(74,158,255,.18);color:#fff}
.fos-ctx-sep{height:1px;background:rgba(255,255,255,.07);margin:3px 0}
/* ═══ AI Computer modal ═════════════════════════════════════════════════════ */
.comp-src-card{background:rgba(255,255,255,.036);border:1px solid var(--glass-bdr);border-radius:9px;padding:9px 11px;margin-bottom:7px;transition:all .14s;cursor:pointer}
.comp-src-card:hover{border-color:rgba(6,182,212,.32);background:rgba(6,182,212,.05)}
.comp-cite{display:inline-block;background:rgba(6,182,212,.2);border:1px solid rgba(6,182,212,.35);border-radius:4px;padding:1px 5px;font-size:9.5px;font-weight:700;color:var(--blue);cursor:pointer;margin:0 2px;vertical-align:super}
</style>

<!-- ═════════ FusionOS overlay ═══════════════════════════════════════════ -->
<div id="fos-overlay">
  <div id="fos-bar">
    <div class="fosb-logo" onclick="closeFusionOS()">⬛ FusionOS</div>
    <div class="fosb-item" onclick="fosOpen('terminal')">Terminal</div>
    <div class="fosb-item" onclick="fosOpen('files')">Files</div>
    <div class="fosb-item" onclick="fosOpen('editor')">Editor</div>
    <div class="fosb-item" onclick="fosOpen('agent')">Agent</div>
    <div class="fosb-item" onclick="fosOpen('browser')">Browser</div>
    <div class="fosb-item" onclick="fosOpen('monitor')">Monitor</div>
    <div class="fosb-right">
      <span style="font-size:9px;background:rgba(45,164,78,.18);border:1px solid rgba(45,164,78,.35);border-radius:10px;padding:2px 8px;color:#6ee77a">LIVE SANDBOX</span>
      <span class="fosb-clock" id="fos-clock">00:00</span>
      <button class="fosb-item" onclick="closeFusionOS()" style="border:none;background:none;color:rgba(255,255,255,.5);cursor:pointer">✕ Exit</button>
    </div>
  </div>
  <div id="fos-desktop"></div>
  <div id="fos-dock">
    <button class="fos-dck" data-app="terminal" onclick="fosOpen('terminal')"><span class="fos-dck-ico">🖳</span><span class="fos-dck-lbl">Terminal</span><span class="fos-dck-dot"></span></button>
    <button class="fos-dck" data-app="files" onclick="fosOpen('files')"><span class="fos-dck-ico">🗂</span><span class="fos-dck-lbl">Files</span><span class="fos-dck-dot"></span></button>
    <button class="fos-dck" data-app="editor" onclick="fosOpen('editor')"><span class="fos-dck-ico">📝</span><span class="fos-dck-lbl">Editor</span><span class="fos-dck-dot"></span></button>
    <button class="fos-dck" data-app="agent" onclick="fosOpen('agent')"><span class="fos-dck-ico">🤖</span><span class="fos-dck-lbl">Agent</span><span class="fos-dck-dot"></span></button>
    <button class="fos-dck" data-app="browser" onclick="fosOpen('browser')"><span class="fos-dck-ico">🌐</span><span class="fos-dck-lbl">Browser</span><span class="fos-dck-dot"></span></button>
    <button class="fos-dck" data-app="monitor" onclick="fosOpen('monitor')"><span class="fos-dck-ico">📊</span><span class="fos-dck-lbl">Monitor</span><span class="fos-dck-dot"></span></button>
  </div>
</div>

<!-- ═════════ AI Computer modal ══════════════════════════════════════════ -->
<div class="modal-overlay hidden" id="computerOverlay" onclick="if(event.target===this)closeComputer()">
  <div style="background:var(--bg2);border:1px solid var(--glass-bdr2);border-radius:18px;width:min(96vw,780px);max-height:88vh;display:flex;flex-direction:column;overflow:hidden;box-shadow:0 24px 80px rgba(0,0,0,.7)">
    <div style="display:flex;align-items:center;gap:10px;padding:16px 18px 12px;border-bottom:1px solid var(--glass-bdr);flex-shrink:0">
      <span style="font-size:20px">🖥</span>
      <div style="flex:1"><div style="font-size:15px;font-weight:700;color:var(--tx)">AI Computer</div><div style="font-size:10.5px;color:var(--tx3)">Searches the web · reads pages · gives cited answers</div></div>
      <button onclick="closeComputer()" style="background:none;border:none;color:var(--tx3);font-size:17px;cursor:pointer;padding:3px 7px;border-radius:7px">✕</button>
    </div>
    <div style="padding:14px 16px;flex-shrink:0;border-bottom:1px solid var(--glass-bdr)">
      <div style="display:flex;gap:7px">
        <input id="compQuery" placeholder="Ask anything — news, research, comparisons…" onkeydown="if(event.key==='Enter')runComputer()"
          style="flex:1;background:rgba(255,255,255,.055);border:1px solid var(--glass-bdr2);border-radius:9px;padding:9px 13px;color:var(--tx);font-size:12.5px;outline:none"/>
        <select id="compDepth" style="background:rgba(255,255,255,.055);border:1px solid var(--glass-bdr2);border-radius:9px;padding:8px 9px;color:var(--tx);font-size:11.5px;cursor:pointer">
          <option value="fast">⚡ Fast</option><option value="normal" selected>🔍 Normal</option><option value="deep">🔬 Deep</option>
        </select>
        <button onclick="runComputer()" style="background:var(--grad);border:none;border-radius:9px;padding:9px 16px;color:#fff;font-size:12.5px;font-weight:700;cursor:pointer">Go ↵</button>
      </div>
      <div style="display:flex;gap:5px;margin-top:9px;flex-wrap:wrap">
        <button class="comp-chip" onclick="setComp('Latest AI news today')">AI news</button>
        <button class="comp-chip" onclick="setComp('How does quantum computing work')">Quantum computing</button>
        <button class="comp-chip" onclick="setComp('Top free AI tools 2025')">Free AI tools</button>
      </div>
    </div>
    <div id="compBody" style="flex:1;overflow-y:auto">
      <div id="compPlaceholder" style="display:flex;flex-direction:column;align-items:center;justify-content:center;padding:50px 24px;gap:12px;text-align:center">
        <div style="font-size:40px;opacity:.28">🖥</div>
        <div style="font-size:12px;color:var(--tx3);max-width:320px;line-height:1.6">Ask a question — searches the web, reads pages, returns a cited answer.</div>
      </div>
      <div id="compResult" style="display:none;padding:18px"></div>
    </div>
    <div style="padding:10px 16px;border-top:1px solid var(--glass-bdr);display:flex;justify-content:space-between;align-items:center;flex-shrink:0">
      <div id="compStatus" style="font-size:10.5px;color:var(--tx3)"></div>
      <button onclick="sendCompToChat()" id="compSendBtn" style="display:none;background:rgba(6,182,212,.16);border:1px solid rgba(6,182,212,.38);border-radius:7px;padding:6px 13px;color:var(--blue);font-size:11.5px;font-weight:600;cursor:pointer">📤 Send to Chat</button>
    </div>
  </div>
</div>

<!-- ═════════ SVG Art (Free) modal ════════════════════════════════════════ -->
<div class="modal-overlay hidden" id="svgOverlay" onclick="if(event.target===this)closeSvgGen()">
  <div style="background:var(--bg2);border:1px solid var(--glass-bdr2);border-radius:18px;width:min(96vw,560px);display:flex;flex-direction:column;overflow:hidden;box-shadow:0 24px 80px rgba(0,0,0,.7)">
    <div style="display:flex;align-items:center;gap:10px;padding:16px 18px 13px;border-bottom:1px solid var(--glass-bdr)">
      <span style="font-size:20px">🖼</span>
      <div style="flex:1"><div style="font-size:15px;font-weight:700;color:var(--tx)">AI Art — Free &amp; Fast</div><div style="font-size:10.5px;color:var(--tx3)">LLM generates SVG → convert to PNG. No image API needed.</div></div>
      <button onclick="closeSvgGen()" style="background:none;border:none;color:var(--tx3);font-size:17px;cursor:pointer;padding:3px 7px;border-radius:7px">✕</button>
    </div>
    <div style="padding:16px 18px">
      <textarea id="svgPrompt" rows="3" placeholder="e.g. a futuristic city at night with neon lights…"
        style="width:100%;background:rgba(255,255,255,.055);border:1px solid var(--glass-bdr2);border-radius:9px;padding:10px 13px;color:var(--tx);font-size:12.5px;resize:vertical;outline:none;box-sizing:border-box"></textarea>
      <div style="display:flex;gap:7px;margin-top:9px;flex-wrap:wrap">
        <select id="svgStyle" style="flex:1;background:rgba(255,255,255,.055);border:1px solid var(--glass-bdr2);border-radius:8px;padding:7px 9px;color:var(--tx);font-size:11.5px;cursor:pointer">
          <option>detailed illustration</option><option>neon cyberpunk art</option><option>watercolor painting</option>
          <option>minimalist vector art</option><option>abstract geometric art</option><option>flat design icon</option>
          <option>dark fantasy art</option><option>cartoon comic style</option>
        </select>
        <select id="svgSize" style="background:rgba(255,255,255,.055);border:1px solid var(--glass-bdr2);border-radius:8px;padding:7px 9px;color:var(--tx);font-size:11.5px;cursor:pointer">
          <option value="800 600">Landscape 800×600</option><option value="600 800">Portrait 600×800</option><option value="800 800">Square 800×800</option>
        </select>
      </div>
      <div id="svgPreviewBox" style="display:none;margin-top:12px;border:1px solid var(--glass-bdr);border-radius:10px;overflow:hidden;text-align:center"><canvas id="svgCanvas" style="max-width:100%;display:block;margin:0 auto"></canvas></div>
    </div>
    <div style="padding:12px 18px;border-top:1px solid var(--glass-bdr);display:flex;gap:7px">
      <button onclick="runSvgGen()" id="svgGenBtn" style="flex:1;background:var(--grad);border:none;border-radius:9px;padding:10px;color:#fff;font-size:12.5px;font-weight:700;cursor:pointer">✨ Generate Art</button>
      <button onclick="svgDownload()" id="svgDlBtn" style="display:none;background:rgba(6,182,212,.16);border:1px solid rgba(6,182,212,.38);border-radius:9px;padding:10px 14px;color:var(--blue);font-size:12px;font-weight:600;cursor:pointer">⬇ PNG</button>
      <button onclick="svgToChat()" id="svgChatBtn" style="display:none;background:rgba(78,200,100,.14);border:1px solid rgba(78,200,100,.32);border-radius:9px;padding:10px 14px;color:#7ee8a2;font-size:12px;font-weight:600;cursor:pointer">📤 Chat</button>
    </div>
  </div>
</div>

<script>
// ═══════════════════════════════════════════════════════════════════════════
// FusionOS — complete window manager + 6 apps
// ═══════════════════════════════════════════════════════════════════════════
var _fos={wins:{},zTop:100,nid:1,cwd:'',clkTimer:null};
var _FOS_APP={
  terminal:{ico:'🖳',title:'Terminal',w:620,h:420},
  files:   {ico:'🗂',title:'Files',   w:660,h:460},
  editor:  {ico:'📝',title:'Editor',  w:700,h:480},
  agent:   {ico:'🤖',title:'AI Agent',w:420,h:540},
  browser: {ico:'🌐',title:'Browser', w:860,h:580},
  monitor: {ico:'📊',title:'Monitor', w:450,h:440},
};

function openFusionOS(){
  document.getElementById('fos-overlay').classList.add('fos-on');
  if(!_fos.clkTimer){ _fosTick(); _fos.clkTimer=setInterval(_fosTick,1000); }
  if(!Object.keys(_fos.wins).length){ fosOpen('terminal'); setTimeout(function(){fosOpen('agent',{x:30,y:50});},80); }
}
function closeFusionOS(){ document.getElementById('fos-overlay').classList.remove('fos-on'); }
function _fosTick(){ var el=document.getElementById('fos-clock'); if(el){ var d=new Date(); el.textContent=String(d.getHours()).padStart(2,'0')+':'+String(d.getMinutes()).padStart(2,'0'); } }

function fosOpen(app,opts){
  opts=opts||{};
  if(!opts.forceNew){
    for(var id in _fos.wins){ if(_fos.wins[id].app===app&&!opts.path){ fosFocus(id); if(_fos.wins[id].mini){ _fos.wins[id].el.classList.remove('fos-mini'); _fos.wins[id].mini=false; } return id; } }
  }
  var cfg=_FOS_APP[app]||{ico:'⬛',title:app,w:560,h:400};
  var desk=document.getElementById('fos-desktop'); var dr=desk.getBoundingClientRect();
  var nc=Object.keys(_fos.wins).length;
  var x=opts.x!=null?opts.x:Math.max(18,Math.min(dr.width-cfg.w-18,60+nc*26));
  var y=opts.y!=null?opts.y:Math.max(12,Math.min(dr.height-cfg.h-18,28+nc*22));
  var id='fw'+(_fos.nid++);
  var win=document.createElement('div'); win.className='fos-win'; win.id=id;
  win.style.cssText='left:'+x+'px;top:'+y+'px;width:'+cfg.w+'px;height:'+cfg.h+'px;z-index:'+(++_fos.zTop);
  win.innerHTML='<div class="fos-titlebar">'
    +'<div class="fos-traf"><button class="fos-tb fos-tbc" onclick="fosClose(\''+id+'\')"></button><button class="fos-tb fos-tbm" onclick="fosMini(\''+id+'\')"></button><button class="fos-tb fos-tbx" onclick="fosMax(\''+id+'\')"></button></div>'
    +'<span class="fos-wico">'+cfg.ico+'</span><span class="fos-wtitle">'+(opts.title||cfg.title)+'</span></div>'
    +'<div class="fos-body" id="'+id+'-b"></div>'
    +'<div class="fos-rz fos-rz-n" data-d="n"></div><div class="fos-rz fos-rz-s" data-d="s"></div>'
    +'<div class="fos-rz fos-rz-e" data-d="e"></div><div class="fos-rz fos-rz-w" data-d="w"></div>'
    +'<div class="fos-rz fos-rz-ne" data-d="ne"></div><div class="fos-rz fos-rz-nw" data-d="nw"></div>'
    +'<div class="fos-rz fos-rz-sw" data-d="sw"></div><div class="fos-rz fos-rz-se" data-d="se">⌟</div>';
  desk.appendChild(win);
  _fos.wins[id]={el:win,app:app,mini:false,max:false,prev:null};
  _fosDrag(win,id); _fosResize(win,id);
  win.addEventListener('mousedown',function(){ fosFocus(id); });
  fosFocus(id); _fosUpdateDock();
  var b=document.getElementById(id+'-b');
  if(app==='terminal') _fosTerminal(b,id);
  else if(app==='files') _fosFiles(b,id,opts.path||'');
  else if(app==='editor') _fosEditor(b,id,opts.path||'',opts.content);
  else if(app==='agent') _fosAgent(b,id);
  else if(app==='browser') _fosBrowser(b,id,opts.url||'');
  else if(app==='monitor') _fosMon(b,id);
  return id;
}
function fosFocus(id){ document.querySelectorAll('.fos-win').forEach(function(w){w.classList.remove('fos-focused');}); var w=_fos.wins[id]; if(w){w.el.classList.add('fos-focused');w.el.style.zIndex=++_fos.zTop;} }
function fosClose(id){ var w=_fos.wins[id]; if(!w)return; w.el.remove(); delete _fos.wins[id]; _fosUpdateDock(); if(id==='_monid') clearInterval(_fos._monTimer); }
function fosMini(id){ var w=_fos.wins[id]; if(!w)return; w.mini=!w.mini; w.el.classList.toggle('fos-mini',w.mini); _fosUpdateDock(); }
function fosMax(id){
  var w=_fos.wins[id]; if(!w)return;
  var desk=document.getElementById('fos-desktop');
  if(!w.max){ w.prev={l:w.el.style.left,t:w.el.style.top,ww:w.el.style.width,hh:w.el.style.height}; w.el.style.cssText='left:0;top:0;width:'+desk.clientWidth+'px;height:'+desk.clientHeight+'px;z-index:'+(++_fos.zTop); w.max=true; }
  else{ var p=w.prev; w.el.style.left=p.l; w.el.style.top=p.t; w.el.style.width=p.ww; w.el.style.height=p.hh; w.max=false; }
  fosFocus(id);
}
function _fosUpdateDock(){ var r={}; for(var id in _fos.wins)r[_fos.wins[id].app]=true; document.querySelectorAll('.fos-dck').forEach(function(b){b.classList.toggle('fos-running',!!r[b.dataset.app]);}); }
function _fosDrag(win,id){
  var bar=win.querySelector('.fos-titlebar'); var sx,sy,ox,oy,on=false;
  bar.addEventListener('mousedown',function(e){ if(e.target.classList.contains('fos-tb'))return; on=true; sx=e.clientX; sy=e.clientY; ox=win.offsetLeft; oy=win.offsetTop; document.body.style.userSelect='none'; fosFocus(id); });
  document.addEventListener('mousemove',function(e){ if(!on)return; var desk=document.getElementById('fos-desktop'); var nx=Math.max(0,Math.min(desk.clientWidth-50,ox+e.clientX-sx)); var ny=Math.max(0,Math.min(desk.clientHeight-30,oy+e.clientY-sy)); win.style.left=nx+'px'; win.style.top=ny+'px'; });
  document.addEventListener('mouseup',function(){ on=false; document.body.style.userSelect=''; });
}
function _fosResize(win,id){
  win.querySelectorAll('.fos-rz').forEach(function(h){
    var dir=h.dataset.d||'se',sx,sy,ox,oy,ow,oh,on=false;
    h.addEventListener('mousedown',function(e){ on=true; sx=e.clientX; sy=e.clientY; ox=win.offsetLeft; oy=win.offsetTop; ow=win.offsetWidth; oh=win.offsetHeight; e.stopPropagation(); e.preventDefault(); document.body.style.userSelect='none'; fosFocus(id); });
    document.addEventListener('mousemove',function(e){ if(!on)return; var dx=e.clientX-sx,dy=e.clientY-sy,nx=ox,ny=oy,nw=ow,nh=oh; if(dir.indexOf('e')>=0)nw=Math.max(280,ow+dx); if(dir.indexOf('s')>=0)nh=Math.max(180,oh+dy); if(dir.indexOf('w')>=0){nw=Math.max(280,ow-dx);nx=ox+(ow-nw);} if(dir.indexOf('n')>=0){nh=Math.max(180,oh-dy);ny=oy+(oh-nh);} win.style.left=nx+'px';win.style.top=ny+'px';win.style.width=nw+'px';win.style.height=nh+'px'; });
    document.addEventListener('mouseup',function(){ on=false; document.body.style.userSelect=''; });
  });
}

async function _fosPost(ep,body){ var r=await apiFetch(ep,{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify(body||{})}); return await r.json(); }
function _fosIcon(name){ var e=(name.split('.').pop()||'').toLowerCase(); return {py:'🐍',js:'📜',ts:'📘',html:'🌐',css:'🎨',json:'📋',md:'📄',txt:'📝',sh:'🦀',csv:'📊',yml:'⚙️',yaml:'⚙️',png:'🖼',jpg:'🖼',svg:'🖼',pdf:'📕',zip:'🗜',go:'🐹',rs:'🦀',cpp:'🔧',c:'🔧',cs:'🔷',rb:'💎'}[e]||'📄'; }
function _fosFmt(b){ if(b<1024)return b+'B'; if(b<1048576)return (b/1024).toFixed(1)+'K'; return (b/1048576).toFixed(1)+'M'; }

// ── Terminal ─────────────────────────────────────────────────────────────────
function _fosTerminal(body,id){
  body.innerHTML='<div class="fos-term"><div class="fos-term-out" id="'+id+'-o"><span class="fos-dim">FusionOS Terminal — real Linux shell. Type help for tips.\n\n</span></div><div class="fos-term-row"><span class="fos-prompt" id="'+id+'-pr">~ $</span><input class="fos-inp" id="'+id+'-i" autocomplete="off" spellcheck="false" placeholder="type a command…"/></div></div>';
  var inp=document.getElementById(id+'-i'); var hist=[],hi=0;
  inp.addEventListener('keydown',async function(e){
    if(e.key==='Enter'){ var cmd=inp.value; inp.value=''; if(!cmd.trim())return; hist.push(cmd); hi=hist.length; await _fosRunCmd(id,cmd); }
    else if(e.key==='ArrowUp'){ if(hi>0){hi--;inp.value=hist[hi];setTimeout(function(){inp.setSelectionRange(99,99);},0);} e.preventDefault(); }
    else if(e.key==='ArrowDown'){ if(hi<hist.length-1){hi++;inp.value=hist[hi];}else{hi=hist.length;inp.value='';} e.preventDefault(); }
  });
  setTimeout(function(){inp.focus();},100);
}
async function _fosRunCmd(id,cmd){
  var out=document.getElementById(id+'-o'); if(!out)return;
  var pr=document.getElementById(id+'-pr');
  var cwd=(pr.textContent.split(' $')[0]||'~').replace(/^~\//,'').replace(/^~$/,'');
  if(cmd.trim()==='help'){ out.innerHTML+='<span class="fos-cmd-echo">$ help</span>\n<span class="fos-dim">Commands: any bash, python3, node, curl, git, gcc, npm, pip\nSpecials: clear, cd, help\nFor multi-step tasks use the Agent app.</span>\n\n'; out.scrollTop=out.scrollHeight; return; }
  out.innerHTML+='<span class="fos-cmd-echo">$ '+escHtml(cmd)+'</span>\n';
  out.scrollTop=out.scrollHeight;
  try{
    var d=await _fosPost('/api/fos/exec',{cmd:cmd,cwd:cwd});
    if(d.stdout==='__CLEAR__'){ out.innerHTML=''; return; }
    if(d.stdout) out.innerHTML+='<span class="fos-ok-txt">'+escHtml(d.stdout)+(d.stdout.endsWith('\n')?'':'\n')+'</span>';
    if(d.stderr) out.innerHTML+='<span class="fos-err-txt">'+escHtml(d.stderr)+(d.stderr.endsWith('\n')?'':'\n')+'</span>';
    if(!d.stdout&&!d.stderr) out.innerHTML+='\n';
    pr.textContent='~'+(d.cwd?'/'+d.cwd:'')+' $';
  }catch(e){ out.innerHTML+='<span class="fos-err-txt">Error: '+escHtml(e.message)+'</span>\n'; }
  out.innerHTML+='\n'; out.scrollTop=out.scrollHeight;
}

// ── Files ─────────────────────────────────────────────────────────────────────
async function _fosFiles(body,id,startPath){
  body.innerHTML='<div class="fos-files"><div class="fos-fsb"><button class="fos-sbtn" onclick="_fosFNav(\''+id+'\',\'\')">🏠 Home</button><button class="fos-sbtn" onclick="_fosFNew(\''+id+'\',\'file\')">📄 New File</button><button class="fos-sbtn" onclick="_fosFNew(\''+id+'\',\'folder\')">📁 New Folder</button><button class="fos-sbtn" onclick="_fosFNav(\''+id+'\',_fos.wins[\''+id+'\']&&_fos.wins[\''+id+'\'].curPath||\'\')">🔄 Refresh</button></div><div class="fos-fmain"><div class="fos-fbar"><span style="font-size:11px;color:rgba(255,255,255,.35)">📍</span><span id="'+id+'-fp" style="font-size:11px;color:rgba(255,255,255,.5);font-family:monospace">~</span></div><div class="fos-flist" id="'+id+'-fl"></div></div></div>';
  _fos.wins[id].curPath=startPath||'';
  await _fosFNav(id,startPath||'');
}
async function _fosFNav(id,path){
  var w=_fos.wins[id]; if(!w)return; w.curPath=path;
  var d=await _fosPost('/api/fos/files',{action:'list',path:path});
  var el=document.getElementById(id+'-fl'); var pe=document.getElementById(id+'-fp');
  if(!el)return;
  if(d.error){el.innerHTML='<div style="padding:16px;color:#ff8484;font-size:11.5px">'+escHtml(d.error)+'</div>';return;}
  pe.textContent='~/'+(d.path||'');
  var h='';
  if(path){ var par=path.split('/').slice(0,-1).join('/'); h+='<div class="fos-fitem" onclick="_fosFNav(\''+id+'\',\''+escHtml(par)+'\')"><span>⬅</span><span class="fos-fname">.. (up)</span></div>'; }
  if(!d.items.length) h+='<div style="padding:14px;color:rgba(255,255,255,.22);font-size:11px;text-align:center">Empty</div>';
  d.items.forEach(function(it){
    var full=(path?path+'/':'')+it.name; var q=full.replace(/'/g,"\\'");
    h+='<div class="fos-fitem" oncontextmenu="_fosFCtx(event,\''+id+'\',\''+q+'\','+it.is_dir+');return false;" ondblclick="'+(it.is_dir?'_fosFNav(\''+id+'\',\''+q+'\')':'fosOpen(\'editor\',{forceNew:false,path:\''+q+'\'})') +'">'+'<span>'+( it.is_dir?'📁':_fosIcon(it.name))+'</span><span class="fos-fname">'+escHtml(it.name)+'</span><span class="fos-fmeta">'+(!it.is_dir?_fosFmt(it.size):' ')+'</span></div>';
  });
  el.innerHTML=h;
}
async function _fosFNew(id,type){
  var n=prompt(type==='file'?'File name:':'Folder name:',type==='file'?'untitled.py':'new_folder'); if(!n)return;
  var w=_fos.wins[id]; var path=(w.curPath?w.curPath+'/':'')+n;
  await _fosPost('/api/fos/files',{action:type==='file'?'new_file':'mkdir',path:path});
  _fosFNav(id,w.curPath||'');
}
function _fosFCtx(e,id,path,isDir){
  document.querySelectorAll('.fos-ctx').forEach(function(c){c.remove();});
  var m=document.createElement('div'); m.className='fos-ctx';
  m.style.cssText='left:'+e.clientX+'px;top:'+e.clientY+'px';
  var items=[];
  if(!isDir) items.push(['✏️ Open in Editor','fosOpen(\'editor\',{forceNew:false,path:\''+path.replace(/'/g,"\\'")+'\'});']);
  items.push(['✏️ Rename','_fosFRename(\''+id+'\',\''+path.replace(/'/g,"\\'")+'\');']);
  items.push(['🗑 Delete','_fosFDel(\''+id+'\',\''+path.replace(/'/g,"\\'")+'\');']);
  items.forEach(function(it){ var d=document.createElement('div'); d.className='fos-ctx-item'; d.textContent=it[0]; d.onclick=function(){m.remove();eval(it[1]);}; m.appendChild(d); });
  document.body.appendChild(m);
  setTimeout(function(){document.addEventListener('click',function h(){m.remove();document.removeEventListener('click',h);},{once:true});},10);
}
async function _fosFRename(id,path){ var n=prompt('Rename to:',path.split('/').pop()); if(!n)return; var par=path.split('/').slice(0,-1).join('/'); await _fosPost('/api/fos/files',{action:'rename',path:path,new_path:(par?par+'/':'')+n}); _fosFNav(id,_fos.wins[id].curPath||''); }
async function _fosFDel(id,path){ if(!confirm('Delete "'+path+'"?'))return; await _fosPost('/api/fos/files',{action:'delete',path:path}); _fosFNav(id,_fos.wins[id].curPath||''); }

// ── Editor ────────────────────────────────────────────────────────────────────
async function _fosEditor(body,id,path,initContent){
  body.innerHTML='<div class="fos-editor"><div class="fos-ebar"><input id="'+id+'-fn" placeholder="filename.py" style="background:rgba(255,255,255,.055);border:1px solid rgba(255,255,255,.09);border-radius:6px;padding:4px 8px;color:#c8ddf0;font-size:11px;font-family:monospace;width:180px" value="'+escHtml(path||'')+'"/><button onclick="_fosEdSave(\''+id+'\')" style="background:rgba(74,158,255,.16);border:1px solid rgba(74,158,255,.32);border-radius:6px;padding:4px 11px;color:#7ec3ff;font-size:11px;font-weight:600;cursor:pointer">💾 Save</button><button onclick="_fosEdRun(\''+id+'\')" style="background:rgba(78,230,120,.13);border:1px solid rgba(78,230,120,.3);border-radius:6px;padding:4px 11px;color:#7ee8a2;font-size:11px;font-weight:600;cursor:pointer">▶ Run</button><span id="'+id+'-est" style="font-size:10px;color:rgba(255,255,255,.28);margin-left:auto"></span></div><textarea class="fos-eta" id="'+id+'-ta" spellcheck="false" placeholder="Start typing or open a file from Files…"></textarea><div class="fos-sbar"><span style="font-size:9.5px;color:rgba(255,255,255,.28)" id="'+id+'-sl">Ready</span></div></div>';
  var ta=document.getElementById(id+'-ta');
  ta.addEventListener('keydown',function(e){ if(e.key==='Tab'){e.preventDefault();var s=ta.selectionStart;ta.value=ta.value.substring(0,s)+'  '+ta.value.substring(ta.selectionEnd);ta.selectionStart=ta.selectionEnd=s+2;} if(e.key==='s'&&(e.ctrlKey||e.metaKey)){e.preventDefault();_fosEdSave(id);} });
  if(path){ var rd=await _fosPost('/api/fos/files',{action:'read',path:path}); if(rd.content!=null) ta.value=rd.content; document.getElementById(id+'-sl').textContent=path; }
  else if(initContent!=null) ta.value=initContent;
}
async function _fosEdSave(id){
  var path=document.getElementById(id+'-fn').value.trim(); if(!path){showToast('Enter filename first');return;}
  var st=document.getElementById(id+'-est'); st.textContent='Saving…';
  await _fosPost('/api/fos/files',{action:'write',path:path,content:document.getElementById(id+'-ta').value});
  st.textContent='✅ '+new Date().toLocaleTimeString(); document.getElementById(id+'-sl').textContent=path;
}
async function _fosEdRun(id){
  var path=document.getElementById(id+'-fn').value.trim(); if(!path){showToast('Save with a filename first');return;}
  await _fosEdSave(id);
  var runners={py:'python3',js:'node',sh:'bash',rb:'ruby',pl:'perl',go:'go run',c:'gcc -o /tmp/_fos_out "$f" && /tmp/_fos_out',cpp:'g++ -o /tmp/_fos_out "$f" && /tmp/_fos_out'};
  var ext=path.split('.').pop(); var runner=runners[ext]; if(!runner){showToast('Don\'t know how to run .'+ext);return;}
  var termId=fosOpen('terminal');
  setTimeout(async function(){ var cmd=runner.indexOf('"$f"')>=0?runner.replace(/"?\$f"?/g,'"'+path+'"'):(runner+' "'+path+'"'); await _fosRunCmd(termId,cmd); },280);
}

// ── Browser ───────────────────────────────────────────────────────────────────
function _fosBrowser(body,id,startUrl){
  body.innerHTML='<div class="fos-browser">'
    +'<div class="fos-bbar">'
      +'<button onclick="_fosBNav(\''+id+'\',\'back\')" style="background:rgba(255,255,255,.07);border:none;border-radius:6px;width:28px;height:26px;color:#fff;cursor:pointer;font-size:14px">←</button>'
      +'<button onclick="_fosBNav(\''+id+'\',\'fwd\')" style="background:rgba(255,255,255,.07);border:none;border-radius:6px;width:28px;height:26px;color:#fff;cursor:pointer;font-size:14px">→</button>'
      +'<button onclick="_fosBRld(\''+id+'\')" style="background:rgba(255,255,255,.07);border:none;border-radius:6px;width:28px;height:26px;color:#fff;cursor:pointer">↻</button>'
      +'<input id="'+id+'-url" placeholder="Enter URL or search…" style="flex:1;background:rgba(255,255,255,.055);border:1px solid rgba(255,255,255,.09);border-radius:7px;padding:5px 10px;color:#c8ddf0;font-size:11.5px;outline:none" onkeydown="if(event.key===\'Enter\')_fosBGo(\''+id+'\')"/>'
      +'<button onclick="_fosBGo(\''+id+'\')" style="background:rgba(124,58,237,.16);border:1px solid rgba(124,58,237,.3);border-radius:7px;padding:5px 11px;color:#c4a3ff;font-size:11px;font-weight:600;cursor:pointer">Go</button>'
      +'<button id="'+id+'-chromeBtn" onclick="_fosBToggleChromium(\''+id+'\')" title="Use a real headless Chromium instead of the proxy iframe" style="background:rgba(6,182,212,.14);border:1px solid rgba(6,182,212,.3);border-radius:7px;padding:5px 10px;color:#5fe8fb;font-size:11px;font-weight:600;cursor:pointer;white-space:nowrap">🖥 Chromium</button>'
    +'</div>'
    +'<div class="fos-bfavs">'
      +'<button class="comp-chip" onclick="_fosBLoad(\''+id+'\',\'https://www.google.com\')">Google</button>'
      +'<button class="comp-chip" onclick="_fosBLoad(\''+id+'\',\'https://github.com\')">GitHub</button>'
      +'<button class="comp-chip" onclick="_fosBLoad(\''+id+'\',\'https://wikipedia.org\')">Wikipedia</button>'
      +'<button class="comp-chip" onclick="_fosBLoad(\''+id+'\',\'https://developer.mozilla.org\')">MDN</button>'
      +'<button class="comp-chip" onclick="_fosBLoad(\''+id+'\',\'https://pypi.org\')">PyPI</button>'
      +'<button class="comp-chip" onclick="_fosBLoad(\''+id+'\',\'https://npmjs.com\')">npm</button>'
      +'<button class="comp-chip" onclick="_fosBLoad(\''+id+'\',\'https://news.ycombinator.com\')">HN</button>'
    +'</div>'
    +'<div id="'+id+'-info" style="display:none;padding:6px 10px;font-size:10px;color:rgba(255,255,255,.35);border-bottom:1px solid rgba(255,255,255,.04);flex-shrink:0">Loading via proxy…</div>'
    +'<iframe id="'+id+'-fr" class="fos-frame" src="about:blank" referrerpolicy="no-referrer"></iframe>'
    +'<div id="'+id+'-chromeView" style="display:none;flex:1;overflow:auto;background:#0a0a12;text-align:center;padding:10px">'
      +'<img id="'+id+'-chromeImg" style="max-width:100%;border-radius:6px;box-shadow:0 4px 24px rgba(0,0,0,.5)"/>'
      +'<div id="'+id+'-chromeMsg" style="color:rgba(255,255,255,.4);font-size:12px;padding:20px"></div>'
    +'</div>'
  +'</div>';
  if(startUrl) _fosBLoad(id, startUrl);
}
function _fosBGo(id){
  var u=(document.getElementById(id+'-url')||{}).value||''; u=u.trim();
  if(!u) return;
  if(!u.startsWith('http') && !u.includes(' ') && u.includes('.')) u='https://'+u;
  else if(!u.startsWith('http')) u='https://www.google.com/search?q='+encodeURIComponent(u);
  _fosBLoad(id,u);
}
var _fosBChromiumMode={};
function _fosBToggleChromium(id){
  _fosBChromiumMode[id]=!_fosBChromiumMode[id];
  var btn=document.getElementById(id+'-chromeBtn');
  if(btn){
    btn.style.background=_fosBChromiumMode[id]?'rgba(6,182,212,.35)':'rgba(6,182,212,.14)';
    btn.textContent=_fosBChromiumMode[id]?'🖥 Chromium ON':'🖥 Chromium';
  }
  var cur=(document.getElementById(id+'-url')||{}).value;
  if(cur) _fosBLoad(id,cur);
}
function _fosBLoad(id,url){
  var ur=document.getElementById(id+'-url'); if(ur) ur.value=url;
  var fr=document.getElementById(id+'-fr');
  var cv=document.getElementById(id+'-chromeView');
  var info=document.getElementById(id+'-info');
  if(_fosBChromiumMode[id]){
    if(fr) fr.style.display='none';
    if(cv) cv.style.display='block';
    if(info){ info.style.display='block'; info.textContent='🖥 Rendering '+url+' with real headless Chromium…'; }
    var img=document.getElementById(id+'-chromeImg'); var msg=document.getElementById(id+'-chromeMsg');
    if(img) img.style.display='none'; if(msg) msg.textContent='';
    apiFetch('/api/fos/chromium',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({url:url,mode:'screenshot'})})
      .then(function(r){return r.json();})
      .then(function(d){
        if(info) info.style.display='none';
        if(d.ok && img){ img.src='data:'+d.content_type+';base64,'+d.data; img.style.display='inline-block'; }
        else if(msg){ msg.textContent='⚠️ '+(d.data||'Chromium render failed.'); }
      }).catch(function(e){ if(info) info.style.display='none'; if(msg) msg.textContent='⚠️ '+e.message; });
    return;
  }
  if(!fr) return;
  if(cv) cv.style.display='none';
  fr.style.display='block';
  if(info){ info.style.display='block'; info.textContent='🔄 Loading '+url+'…'; }
  // Route through server-side proxy to bypass X-Frame-Options
  var proxyUrl='/api/fos/proxy?url='+encodeURIComponent(url)+'&_auth='+encodeURIComponent(localStorage.getItem('token')||'');
  fr.src=proxyUrl;
  fr.onload=function(){ if(info) info.style.display='none'; };
  fr.onerror=function(){ if(info){ info.style.display='block'; info.textContent='Failed — try opening in a new tab'; info.innerHTML=info.textContent+' <a href="'+url+'" target="_blank" style="color:#06b6d4">↗</a>'; } };
}
function _fosBRld(id){ var fr=document.getElementById(id+'-fr'); if(fr) fr.src=fr.src; }
function _fosBNav(id,dir){ var fr=document.getElementById(id+'-fr'); if(fr){try{dir==='back'?fr.contentWindow.history.back():fr.contentWindow.history.forward();}catch(e){}} }

// ── Agent ─────────────────────────────────────────────────────────────────────
function _fosAgent(body,id){
  body.innerHTML='<div class="fos-agent"><div class="fos-alog" id="'+id+'-al"><div style="padding:16px;text-align:center;color:rgba(255,255,255,.28);font-size:12px">🤖 Give me a task and I\'ll plan &amp; run real commands autonomously.</div></div><div class="fos-afooter"><div style="display:flex;gap:6px"><select id="'+id+'-sp" style="flex:1;background:rgba(255,255,255,.045);border:1px solid rgba(255,255,255,.09);border-radius:7px;padding:5px 6px;color:#b8cce8;font-size:10.5px"><option value="fast">⚡ Fast (quick tasks)</option><option value="balanced" selected>⚖️ Balanced</option><option value="thorough">🧠 Thorough (complex)</option></select><select id="'+id+'-ms" style="background:rgba(255,255,255,.045);border:1px solid rgba(255,255,255,.09);border-radius:7px;padding:5px 6px;color:#b8cce8;font-size:10.5px"><option value="6">6 steps</option><option value="10" selected>10 steps</option><option value="14">14 steps</option><option value="20">20 steps</option></select></div><textarea id="'+id+'-g" rows="2" placeholder="e.g. generate 10 passwords, save to passwords.txt and show them" style="background:rgba(255,255,255,.045);border:1px solid rgba(255,255,255,.09);border-radius:8px;padding:8px 10px;color:#c8ddf0;font-size:12px;resize:none;outline:none"></textarea><button id="'+id+'-rb" onclick="_fosAgentRun(\''+id+'\')" style="background:linear-gradient(135deg,#4a9eff,#a855f7);border:none;border-radius:8px;padding:9px;color:#fff;font-size:12px;font-weight:700;cursor:pointer">✨ Run Agent</button></div></div>';
}
async function _fosAgentRun(id,resumeTr,resumeStep){
  var goalEl=document.getElementById(id+'-g'); var goal=goalEl.value.trim(); if(!goal){showToast('Describe a task first');return;}
  var log=document.getElementById(id+'-al'); var btn=document.getElementById(id+'-rb');
  var speed=(document.getElementById(id+'-sp')||{value:'balanced'}).value;
  var maxS=parseInt((document.getElementById(id+'-ms')||{value:'10'}).value,10);
  if(!resumeTr){ log.innerHTML='<div class="fos-astep" style="border-color:rgba(168,85,247,.28)"><b style="color:#c084fc">🎯 Goal:</b> '+escHtml(goal)+'</div>'; }
  btn.disabled=true; btn.textContent='⏳ Working…'; goalEl.disabled=true;
  var tr=resumeTr||[]; var sn=resumeStep||1;
  try{
    while(sn<=maxS){
      var sd=document.createElement('div'); sd.className='fos-astep'; sd.innerHTML='<span class="fos-dim">🧠 Step '+sn+'…</span>'; log.appendChild(sd); log.scrollTop=log.scrollHeight;
      var d=await _fosPost('/api/fos/agent',{goal:goal,transcript:tr,step:sn,speed:speed,max_steps:maxS});
      if(d.action==='run'){
        sd.classList.add(d.code===0?'fos-ok':'fos-err');
        sd.innerHTML=(d.explain?'<div style="color:rgba(255,255,255,.5);margin-bottom:3px;font-size:10.5px">💭 '+escHtml(d.explain)+'</div>':'')
          +'<div class="fos-acmd">$ '+escHtml(d.cmd)+'</div>'
          +(d.stdout?'<div class="fos-aout">'+escHtml(d.stdout)+'</div>':'')
          +(d.stderr?'<div class="fos-aout" style="color:#ff8484">'+escHtml(d.stderr)+'</div>':'');
        tr.push({cmd:d.cmd,stdout:d.stdout,stderr:d.stderr,code:d.code}); sn++;
        log.scrollTop=log.scrollHeight; await new Promise(function(r){setTimeout(r,180);});
      } else if(d.action==='done'){
        sd.classList.add('fos-ok'); sd.innerHTML='<div style="color:#7ee8a2;font-weight:600;margin-bottom:3px">✅ Done</div><div style="color:rgba(255,255,255,.72)">'+escHtml(d.summary||'Completed.')+'</div>';
        log.scrollTop=log.scrollHeight; break;
      } else {
        sd.classList.add('fos-err'); var snc=sn,trc=JSON.stringify(tr);
        sd.dataset.tr=trc;
        sd.innerHTML='<div style="color:#ff8484;font-weight:600;margin-bottom:3px">⚠️ Stopped</div><div style="color:rgba(255,255,255,.6);margin-bottom:6px">'+escHtml(d.reason||'Could not continue.')+'</div><button onclick="_fosAgentResume(\''+id+'\','+snc+',this)" style="background:rgba(74,158,255,.16);border:1px solid rgba(74,158,255,.32);border-radius:6px;padding:4px 10px;color:#7ec3ff;font-size:10.5px;cursor:pointer">↻ Retry step</button>';
        log.scrollTop=log.scrollHeight; break;
      }
    }
  }catch(e){ var ed=document.createElement('div'); ed.className='fos-astep fos-err'; ed.textContent='❌ '+e.message; log.appendChild(ed); }
  btn.disabled=false; btn.textContent='✨ Run Agent'; goalEl.disabled=false;
}
function _fosAgentResume(id,step,btn){ var sd=btn.closest('.fos-astep'); var tr=[]; try{tr=JSON.parse(sd.dataset.tr||'[]');}catch(e){} _fosAgentRun(id,tr,step); }

// ── Monitor ───────────────────────────────────────────────────────────────────
function _fosMon(body,id){
  body.innerHTML='<div class="fos-mon" id="'+id+'-mn"><div style="text-align:center;color:rgba(255,255,255,.28);padding:26px;font-size:11.5px">Loading…</div></div>';
  _fosMonRef(id); _fos._monTimer=setInterval(function(){ if(!_fos.wins[id]){clearInterval(_fos._monTimer);return;} _fosMonRef(id); },5000);
}
async function _fosMonRef(id){
  var el=document.getElementById(id+'-mn'); if(!el)return;
  try{
    var d=await _fosPost('/api/fos/stats',{});
    el.innerHTML='<div style="display:grid;grid-template-columns:1fr 1fr;gap:8px">'
      +'<div class="fos-mcard"><div class="fos-mlbl">Sandbox Files</div><div class="fos-mval">'+d.files+' files / '+d.dirs+' dirs</div></div>'
      +'<div class="fos-mcard"><div class="fos-mlbl">Sandbox Size</div><div class="fos-mval">'+d.size_kb+' KB</div></div>'
      +'<div class="fos-mcard"><div class="fos-mlbl">Python</div><div class="fos-mval">v'+d.python+'</div></div>'
      +'<div class="fos-mcard"><div class="fos-mlbl">Server Uptime</div><div class="fos-mval">'+Math.floor(d.uptime_s/60)+'m '+Math.floor(d.uptime_s%60)+'s</div></div>'
      +'</div>'
      +'<div class="fos-mcard"><div class="fos-mlbl">Disk</div><div class="fos-mval">'+d.disk_used_gb+'GB used / '+d.disk_total_gb+'GB total ('+d.disk_free_gb+'GB free)</div></div>'
      +'<div class="fos-mcard"><div class="fos-mlbl">AI Models</div>'
      +'<div style="display:flex;gap:12px;flex-wrap:wrap;margin-top:4px">'
      +'<span style="font-size:11px;color:'+(d.groq?'#7ee8a2':'#ff8484')+'">●  Groq '+(d.groq?'online':'offline')+'</span>'
      +'<span style="font-size:11px;color:'+(d.openrouter?'#7ee8a2':'#ff8484')+'">●  OpenRouter '+(d.openrouter?'online':'offline')+'</span>'
      +'<span style="font-size:11px;color:'+(d.chromium?'#7ee8a2':'#ff8484')+'">●  Chromium '+(d.chromium?'ready':'not installed')+'</span>'
      +'</div></div>'
      +'<div class="fos-mcard"><div class="fos-mlbl">Platform</div><div class="fos-mval" style="font-size:9.5px">'+escHtml(d.platform||'')+'</div></div>';
  }catch(e){}
}

// ═══════════════════════════════════════════════════════════════════════════
// AI Computer
// ═══════════════════════════════════════════════════════════════════════════
var _compLastResult='',_compLastSources=[];
function openComputer(){ document.getElementById('computerOverlay').classList.remove('hidden'); setTimeout(function(){document.getElementById('compQuery').focus();},150); }
function closeComputer(){ document.getElementById('computerOverlay').classList.add('hidden'); }
function setComp(q){ document.getElementById('compQuery').value=q; runComputer(); }

async function runComputer(){
  var query=document.getElementById('compQuery').value.trim(); if(!query){showToast('Ask a question first');return;}
  var depth=document.getElementById('compDepth').value;
  document.getElementById('compPlaceholder').style.display='none';
  document.getElementById('compSendBtn').style.display='none';
  var res=document.getElementById('compResult'); res.style.display='block';
  res.innerHTML='<div style="display:flex;align-items:center;gap:10px;padding:20px"><div class="thinking-dots"><span></span><span></span><span></span></div><span style="color:var(--tx2);font-size:13px" id="compStep">🔍 Searching…</span></div>';
  document.getElementById('compStatus').textContent='';
  var prog; function setStep(t){var el=document.getElementById('compStep');if(el)el.textContent=t;}
  try{
    var r=await apiFetch('/api/computer',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({query:query,depth:depth,num_sources:{'fast':3,'normal':6,'deep':10}[depth]||6})});
    if(!r.ok){var e=await r.json();throw new Error(e.error||'Failed');}
    setStep('📖 Reading sources…');
    var d=await r.json();
    _compLastResult=d.answer||''; _compLastSources=d.sources||[];
    var srcHtml='<div style="margin-bottom:14px"><div style="font-size:9.5px;font-weight:700;letter-spacing:.8px;color:var(--tx3);text-transform:uppercase;margin-bottom:7px">📚 Sources ('+d.sources.length+')</div><div style="display:grid;grid-template-columns:repeat(auto-fill,minmax(200px,1fr));gap:5px">';
    d.sources.forEach(function(s,i){ srcHtml+='<div class="comp-src-card" onclick="window.open(\''+escHtml(s.url)+'\',\'_blank\')"><div style="display:flex;align-items:center;gap:6px"><span class="comp-cite">'+(i+1)+'</span><div style="min-width:0"><div style="font-size:11px;font-weight:600;color:var(--tx);white-space:nowrap;overflow:hidden;text-overflow:ellipsis">'+escHtml(s.title||s.url)+'</div><div style="font-size:9px;color:var(--tx3);white-space:nowrap;overflow:hidden;text-overflow:ellipsis">'+escHtml(s.url)+'</div></div></div></div>'; });
    srcHtml+='</div></div>';
    var ansHtml='<div style="font-size:13.5px;line-height:1.75;color:var(--tx);margin-bottom:16px">'+fmt(d.answer||'No answer generated.')+'</div>';
    var relHtml='';
    if(d.related&&d.related.length){ relHtml='<div><div style="font-size:9.5px;font-weight:700;letter-spacing:.8px;color:var(--tx3);text-transform:uppercase;margin-bottom:7px">💡 Related</div><div style="display:flex;flex-direction:column;gap:5px">'; d.related.forEach(function(r2){ relHtml+='<button onclick="setComp('+JSON.stringify(r2)+')" style="background:rgba(255,255,255,.038);border:1px solid var(--glass-bdr);border-radius:7px;padding:6px 11px;font-size:11.5px;color:var(--tx2);cursor:pointer;text-align:left;transition:all .12s">→ '+escHtml(r2)+'</button>'; }); relHtml+='</div></div>'; }
    res.innerHTML=srcHtml+ansHtml+relHtml;
    document.getElementById('compStatus').textContent=d.sources.length+' sources · '+(d.elapsed_ms?Math.round(d.elapsed_ms/100)/10+'s':'');
    document.getElementById('compSendBtn').style.display='inline-block';
  }catch(e){ res.innerHTML='<div style="color:var(--red);padding:20px">❌ '+escHtml(e.message)+'</div>'; }
}
function sendCompToChat(){ var q=document.getElementById('compQuery').value.trim(); closeComputer(); addMsg('ai','**Computer Search: '+q+'**\n\n'+_compLastResult,null,null); }

// ═══════════════════════════════════════════════════════════════════════════
// SVG Free Art Gen
// ═══════════════════════════════════════════════════════════════════════════
var _svgB64='';
function openSvgGen(){ document.getElementById('svgOverlay').classList.remove('hidden'); document.getElementById('svgPreviewBox').style.display='none'; ['svgDlBtn','svgChatBtn'].forEach(function(x){document.getElementById(x).style.display='none';}); setTimeout(function(){document.getElementById('svgPrompt').focus();},150); }
function closeSvgGen(){ document.getElementById('svgOverlay').classList.add('hidden'); }
async function runSvgGen(){
  var prompt=document.getElementById('svgPrompt').value.trim(); if(!prompt){showToast('Describe the art');return;}
  var style=document.getElementById('svgStyle').value; var sz=document.getElementById('svgSize').value.split(' ');
  var W=parseInt(sz[0])||800,H=parseInt(sz[1])||600;
  var btn=document.getElementById('svgGenBtn'); btn.textContent='⏳ Generating…'; btn.disabled=true;
  document.getElementById('svgPreviewBox').style.display='none';
  try{
    var r=await apiFetch('/api/gen_svg_img',{method:'POST',headers:{'Content-Type':'application/json'},body:JSON.stringify({prompt:prompt,style:style,width:W,height:H})});
    var d=await r.json(); if(d.error)throw new Error(d.error);
    var canvas=document.getElementById('svgCanvas'); canvas.width=W; canvas.height=H;
    var ctx=canvas.getContext('2d'); var blob=new Blob([d.svg],{type:'image/svg+xml'});
    var url=URL.createObjectURL(blob); var img=new Image();
    img.onload=function(){ ctx.drawImage(img,0,0,W,H); URL.revokeObjectURL(url); _svgB64=canvas.toDataURL('image/png').split(',')[1]; document.getElementById('svgPreviewBox').style.display='block'; ['svgDlBtn','svgChatBtn'].forEach(function(x){document.getElementById(x).style.display='inline-block';}); showToast('🎨 Generated!'); };
    img.onerror=function(){ document.getElementById('svgPreviewBox').innerHTML='<div style="padding:10px">'+d.svg+'</div>'; document.getElementById('svgPreviewBox').style.display='block'; };
    img.src=url;
  }catch(e){showToast('❌ '+e.message);}
  btn.textContent='✨ Generate Art'; btn.disabled=false;
}
function svgDownload(){ if(!_svgB64){showToast('Generate first');return;} var a=document.createElement('a'); a.href='data:image/png;base64,'+_svgB64; a.download='fusion-art.png'; a.click(); }
function svgToChat(){ if(!_svgB64){showToast('Generate first');return;} closeSvgGen(); var bbl=addMsg('ai','',null,null); bbl.innerHTML='<img src="data:image/png;base64,'+_svgB64+'" style="max-width:100%;border-radius:10px;border:1px solid var(--glass-bdr2)"/>'; }
</script>
</body></html>"""



@app.post("/api/update_avatar")
async def update_avatar(request:Request):
    u=auth_user(request); d=await request.json()
    emoji=d.get("emoji","")[:4]
    if not emoji: return err("emoji required")
    with db() as c: c.execute("UPDATE users SET avatar_emoji=? WHERE id=?",(emoji,u["id"]))
    return J({"ok":True,"emoji":emoji})

@app.get("/api/profile/{username}")
async def get_profile(username:str,request:Request):
    with db() as c:
        row=c.execute("SELECT username,display_name,avatar_emoji,guest_id,created FROM users WHERE username=?",(username,)).fetchone()
    if not row: return err("User not found",404)
    return J({"username":row["username"],"display_name":row["display_name"],"emoji":row["avatar_emoji"],"guest_id":row["guest_id"],"joined":row["created"][:10]})

@app.post("/api/ai_tools/translate")
async def ai_translate(request:Request):
    u=auth_user(request); d=await request.json()
    text=d.get("text","").strip()[:4000]; target=d.get("target","English"); source=d.get("source","auto")
    if not text: return err("text required")
    avail=get_available(u["id"],u["salt"])
    fast=[k for k in ["groq_llama31_8b","or_lfm_instruct","groq_llama33_70b"] if MODELS[k]["provider"] in avail]
    if not fast: return err("No model available")
    key=fast[0]; m=MODELS[key]; prov=m["provider"]
    sys_p=f"You are a professional translator. Translate the user text to {target}. Source language: {source if source!='auto' else 'detect automatically'}. Return ONLY the translated text, no explanations, no quotes."
    msgs=[{"role":"system","content":sys_p},{"role":"user","content":text}]
    try:
        hdrs={"Authorization":f"Bearer {avail[prov]}","Content-Type":"application/json"}
        body={"messages":msgs,"max_tokens":2000,"stream":False,"model":m["model"]}
        r=req.post(ENDPOINTS.get(prov,""),headers=hdrs,json=body,timeout=20)
        if not r.ok: return err("Translation failed")
        out=r.json()["choices"][0]["message"]["content"].strip()
        return J({"ok":True,"translation":out,"model":m["label"],"target":target})
    except Exception as ex: return err(str(ex))

@app.post("/api/ai_tools/summarise")
async def ai_summarise(request:Request):
    u=auth_user(request); d=await request.json()
    text=d.get("text","").strip()[:8000]; style=d.get("style","bullet")  # bullet | paragraph | eli5 | tldr
    if not text: return err("text required")
    avail=get_available(u["id"],u["salt"])
    fast=[k for k in ["groq_llama33_70b","gh_phi4","or_deepseek_v3"] if MODELS[k]["provider"] in avail]
    if not fast: return err("No model")
    key=fast[0]; m=MODELS[key]; prov=m["provider"]
    styles={"bullet":"as clear bullet points (5-8 key points)","paragraph":"as 2-3 concise paragraphs","eli5":"as if explaining to a 10-year-old in simple language","tldr":"as a single TL;DR sentence followed by 3 key points"}
    sys_p=f"Summarise the following text {styles.get(style,'concisely')}. Be accurate and comprehensive."
    msgs=[{"role":"system","content":sys_p},{"role":"user","content":text}]
    try:
        hdrs={"Authorization":f"Bearer {avail[prov]}","Content-Type":"application/json"}
        body={"messages":msgs,"max_tokens":800,"stream":False,"model":m["model"]}
        r=req.post(ENDPOINTS.get(prov,""),headers=hdrs,json=body,timeout=20)
        if not r.ok: return err("Summarise failed")
        out=r.json()["choices"][0]["message"]["content"].strip()
        return J({"ok":True,"summary":out,"model":m["label"],"style":style})
    except Exception as ex: return err(str(ex))

@app.post("/api/ai_tools/rewrite")
async def ai_rewrite(request:Request):
    u=auth_user(request); d=await request.json()
    text=d.get("text","").strip()[:4000]; tone=d.get("tone","professional")
    if not text: return err("text required")
    avail=get_available(u["id"],u["salt"])
    fast=[k for k in ["groq_llama33_70b","gh_phi4","or_deepseek_v3"] if MODELS[k]["provider"] in avail]
    if not fast: return err("No model")
    key=fast[0]; m=MODELS[key]; prov=m["provider"]
    tones={"professional":"professional and formal","casual":"casual and friendly","concise":"concise and direct (remove fluff)","creative":"creative and engaging","academic":"academic and scholarly","persuasive":"persuasive and compelling"}
    sys_p=f"Rewrite the following text to be {tones.get(tone,'professional')}. Keep the core meaning. Return ONLY the rewritten text."
    msgs=[{"role":"system","content":sys_p},{"role":"user","content":text}]
    try:
        hdrs={"Authorization":f"Bearer {avail[prov]}","Content-Type":"application/json"}
        body={"messages":msgs,"max_tokens":2000,"stream":False,"model":m["model"]}
        r=req.post(ENDPOINTS.get(prov,""),headers=hdrs,json=body,timeout=20)
        if not r.ok: return err("Rewrite failed")
        out=r.json()["choices"][0]["message"]["content"].strip()
        return J({"ok":True,"result":out,"model":m["label"]})
    except Exception as ex: return err(str(ex))

@app.post("/api/ai_tools/quiz")
async def ai_quiz(request:Request):
    u=auth_user(request); d=await request.json()
    topic=d.get("topic","General Knowledge").strip()[:300]; count=min(int(d.get("count",5)),10)
    avail=get_available(u["id"],u["salt"])
    fast=[k for k in ["groq_llama33_70b","gh_phi4","or_deepseek_v3"] if MODELS[k]["provider"] in avail]
    if not fast: return err("No model")
    key=fast[0]; m=MODELS[key]; prov=m["provider"]
    sys_p='You are a quiz master. Generate exactly '+str(count)+' multiple-choice questions about: '+topic+'. Return ONLY valid JSON array: [{"q":"question","opts":["A","B","C","D"],"ans":0}] where ans is the index (0-3) of the correct answer. No extra text.'
    msgs=[{"role":"system","content":sys_p},{"role":"user","content":"Generate the quiz now."}]
    try:
        hdrs={"Authorization":f"Bearer {avail[prov]}","Content-Type":"application/json"}
        body={"messages":msgs,"max_tokens":1200,"stream":False,"model":m["model"]}
        r=req.post(ENDPOINTS.get(prov,""),headers=hdrs,json=body,timeout=25)
        if not r.ok: return err("Quiz failed")
        import json as _json2, re as _re5
        raw=r.json()["choices"][0]["message"]["content"]
        m2=_re5.search(r'\[.*\]',raw,_re5.DOTALL)
        if not m2: return err("Bad quiz format")
        qs=_json2.loads(m2.group())
        return J({"ok":True,"questions":qs,"topic":topic})
    except Exception as ex: return err(str(ex))

# ══ Binary file generation endpoints ══════════════════════════════════════════

# ══ Multi-Model Arena ═══════════════════════════════════════════════════════
@app.post("/api/arena")
async def api_arena(request: Request):
    """Run N models in parallel, return all answers for user to pick."""
    import time as _ta, re as _rea
    t0=_ta.time()
    u=auth_user(request); d=await request.json()
    prompt=d.get("prompt","").strip()
    chosen_keys=d.get("models",[])   # list of model keys, empty = auto-pick 5
    if not prompt: return err("Prompt required")
    avail=get_available(u["id"],u["salt"])

    CHAT_MODELS=[k for k,m in MODELS.items()
        if m.get("type","chat")=="chat"
        and m["provider"] in avail
        and m["provider"] in ("groq","openrouter","github","cloudflare")]

    if chosen_keys:
        run_keys=[k for k in chosen_keys if k in MODELS and MODELS[k]["provider"] in avail]
    else:
        # Auto-pick 5: diverse providers, best quality
        PREFER=["gh_gpt4o","or_deepseek_v3","groq_llama33_70b","cf_qwen3_30b","gh_phi4",
                "or_qwen3_235b","groq_compound","gh_deepseek_v3","cf_nemotron","or_nemotron_super"]
        run_keys=[k for k in PREFER if k in CHAT_MODELS][:5]
        if len(run_keys)<5:
            extras=[k for k in CHAT_MODELS if k not in run_keys]
            run_keys+=(extras[:5-len(run_keys)])

    sys_p=f"You are a brilliant AI. Answer concisely, clearly, and completely. Today: {datetime.now().strftime('%B %d, %Y')}."
    msgs=[{"role":"system","content":sys_p},{"role":"user","content":prompt}]

    def _fetch(key):
        m=MODELS[key]; prov=m["provider"]; t1=_ta.time()
        try:
            hdrs={"Authorization":f"Bearer {avail[prov]}","Content-Type":"application/json"}
            if prov=="openrouter": hdrs.update({"X-Title":"Fusion.AI","HTTP-Referer":"https://fusionai.space"})
            body={"messages":msgs,"max_tokens":800,"stream":False}
            if prov=="cloudflare":
                r=req.post(cf_ep(m["model"]),headers={"Authorization":f"Bearer {avail['cloudflare']}","Content-Type":"application/json"},json=body,timeout=25)
                dj=r.json(); txt=(dj.get("result",{})or{}).get("response","") or (dj.get("choices",[{}])[0].get("message",{})or{}).get("content","")
            else:
                body["model"]=m["model"]
                r=req.post(ENDPOINTS[prov],headers=hdrs,json=body,timeout=25)
                if not r.ok: return None
                txt=r.json()["choices"][0]["message"]["content"]
            txt=_rea.sub(r"<think>[\s\S]*?</think>","",txt).strip()
            return {"key":key,"model":m["label"],"company":m.get("company",""),"emoji":m.get("emoji","🤖"),
                    "text":txt,"ms":round((_ta.time()-t1)*1000)} if txt else None
        except: return None

    with ThreadPoolExecutor(max_workers=10) as ex:
        futs=[ex.submit(_fetch,k) for k in run_keys]
        raw=[f.result() for f in as_completed(futs,timeout=30)]
    results=[r for r in raw if r]
    results.sort(key=lambda x:x["ms"])
    return J({"ok":True,"results":results,"total_ms":round((_ta.time()-t0)*1000),"prompt":prompt})

# ══ Extreme Deep Think ════════════════════════════════════════════════════════
@app.post("/api/extreme-think")
async def api_extreme_think(request: Request):
    """10+ web searches + all models thinking 5 rounds + final synthesis."""
    import time as _te, urllib.parse as _up3, re as _re4, random as _rnd
    t0=_te.time()
    u=auth_user(request); d=await request.json()
    prompt=d.get("prompt","").strip()
    if not prompt: return err("Prompt required")
    avail=get_available(u["id"],u["salt"])

    CHAT_MODELS=[k for k,m in MODELS.items()
        if m.get("type","chat")=="chat"
        and m["provider"] in avail
        and m["provider"] in ("groq","openrouter","github","cloudflare")]

    steps=[]; web_ctx_parts=[]

    # ── Phase 1: 12 web searches (varied angles) ─────────────────────────
    def _bing_search(q):
        try:
            import urllib.parse as _up4
            r=req.get(f"http://openserp.alwaysdata.net/bing/search?text={_up4.quote(q)}",
                      headers={"User-Agent":"FusionAI/2.0","Accept":"application/json"},timeout=8)
            if r.ok:
                items=r.json() if isinstance(r.json(),list) else []
                snippets=" | ".join((item.get("description","") or item.get("snippet",""))[:180] for item in items[:4] if item.get("description") or item.get("snippet"))
                if snippets: return f"[{q}]: {snippets}"
        except: pass
        return ""

    search_angles=[
        prompt,
        prompt+" latest research 2024 2025",
        prompt+" key concepts explained",
        prompt+" expert opinion analysis",
        prompt+" statistics data evidence",
        prompt+" pros cons tradeoffs",
        prompt+" practical applications examples",
        prompt+" historical context background",
        prompt+" future implications",
        prompt+" common misconceptions debunked",
        prompt+" step by step guide",
        prompt+" case studies real world",
    ]
    steps.append({"phase":"search","status":"Performing 12 web searches…"})
    with ThreadPoolExecutor(max_workers=12) as ex:
        sr=[ex.submit(_bing_search,q) for q in search_angles]
        raw_searches=[f.result() for f in as_completed(sr,timeout=18)]
    web_ctx="\n".join(r for r in raw_searches if r)
    steps.append({"phase":"search","status":f"Web search complete: {len([r for r in raw_searches if r])}/12 returned data","ctx_chars":len(web_ctx)})

    # ── Phase 2: All models think Round 1 (understand the problem) ───────
    THINK_MODELS=[k for k in CHAT_MODELS if MODELS[k]["provider"] in ("groq","github","openrouter")]
    # Limit to 12 best
    PREFER_T=["gh_gpt4o","gh_o4_mini","or_deepseek_r1","groq_compound","groq_qwen3_32b",
              "or_qwen3_235b","gh_deepseek_r1","cf_deepseek_r1","or_nemotron_super",
              "groq_llama33_70b","gh_phi4","or_mistral_small"]
    use_models=[k for k in PREFER_T if k in CHAT_MODELS][:12]
    if len(use_models)<6: use_models+=([k for k in CHAT_MODELS if k not in use_models][:6-len(use_models)])

    thoughts={}  # key → list of 5 round texts

    def _think_round(key,prev_thoughts,round_num,sys_prompt):
        m=MODELS[key]; prov=m["provider"]
        try:
            hdrs={"Authorization":f"Bearer {avail[prov]}","Content-Type":"application/json"}
            if prov=="openrouter": hdrs.update({"X-Title":"Fusion.AI","HTTP-Referer":"https://fusionai.space"})
            body={"messages":sys_prompt,"max_tokens":600,"stream":False}
            if prov=="cloudflare":
                r=req.post(cf_ep(m["model"]),headers={"Authorization":f"Bearer {avail['cloudflare']}","Content-Type":"application/json"},json=body,timeout=22)
                dj=r.json(); txt=(dj.get("result",{})or{}).get("response","") or (dj.get("choices",[{}])[0].get("message",{})or{}).get("content","")
            else:
                body["model"]=m["model"]
                r=req.post(ENDPOINTS.get(prov,""),headers=hdrs,json=body,timeout=22)
                if not r.ok: return ""
                txt=r.json()["choices"][0]["message"]["content"]
            return _re4.sub(r"<think>[\s\S]*?</think>","",txt or "").strip()
        except: return ""

    ROUND_SYSPS=[
        "You are a deep analytical thinker. Given the question and web research context, identify the 3 most important sub-problems or aspects that need addressing. Be specific and concrete.",
        "You are a critical analyst. Review the sub-problems identified in round 1. Now generate 2-3 concrete hypotheses or potential solutions/answers. What evidence supports each?",
        "You are a devil's advocate. Challenge the hypotheses from round 2. What are the strongest counterarguments, edge cases, or missing information? What might be wrong?",
        "You are a synthesiser. Given the analysis so far, what is the most accurate, nuanced answer? Integrate the counterarguments. State your confidence level and why.",
        "You are a final editor. Write the definitive, comprehensive answer to the original question. Be clear, precise, and complete. Use evidence from web research where relevant.",
    ]

    for round_num in range(5):
        steps.append({"phase":f"think_r{round_num+1}","status":f"Round {round_num+1}/5: {ROUND_SYSPS[round_num][:60]}…"})
        def _do_round(key, r=round_num):
            prev="\n".join(thoughts.get(key,[]))
            web_snippet=web_ctx[:1200] if r==0 else web_ctx[:600]
            sys_content=f"Web Research Context:\n{web_snippet}\n\nPrevious thinking:\n{prev}\n\nTask: {ROUND_SYSPS[r]}"
            msgs_r=[{"role":"system","content":sys_content},{"role":"user","content":prompt}]
            txt=_think_round(key,prev,r,msgs_r)
            if key not in thoughts: thoughts[key]=[]
            thoughts[key].append(f"[Round {r+1}] {txt}" if txt else "")
            return key,txt

        with ThreadPoolExecutor(max_workers=12) as ex:
            futs=[ex.submit(_do_round,k) for k in use_models]
            for f in as_completed(futs,timeout=35): f.result()

    # ── Phase 3: Collect all final (round 5) thoughts ────────────────────
    final_thoughts=[]
    for k in use_models:
        rds=thoughts.get(k,[])
        if rds and rds[-1]:
            txt=rds[-1].replace(f"[Round 5] ","")
            final_thoughts.append({"model":MODELS[k]["label"],"emoji":MODELS[k].get("emoji","🤖"),"text":txt[:800],"rounds":len(rds)})

    # ── Phase 4: Master synthesis ─────────────────────────────────────────
    steps.append({"phase":"synthesis","status":"Creating master synthesis…"})
    synth_key=None
    for k in ["gh_gpt4o","groq_compound","or_deepseek_v3","gh_deepseek_v3","groq_llama33_70b"]:
        if k in CHAT_MODELS: synth_key=k; break
    synthesis=""
    if synth_key:
        m=MODELS[synth_key]; prov=m["provider"]
        combined="".join(f"[{ft['model']}]: {ft['text']}" for ft in final_thoughts[:10])
        synth_sys=f"""You are the world's best research synthesiser. You have:
1. Performed 12 web searches on the topic
2. Consulted {len(use_models)} AI models, each thinking 5 deep rounds
3. Collected their final conclusions

Web research summary (top):
{web_ctx[:1500]}

Model final conclusions:
{combined[:3000]}

Task: Write the single most complete, accurate, nuanced answer to the question. 
- Use ## headers for sections
- Cite specific insights from models where useful
- Reference web research data where relevant  
- End with a "Key Takeaway" section
- Be comprehensive but clear"""
        try:
            hdrs={"Authorization":f"Bearer {avail[prov]}","Content-Type":"application/json"}
            if prov=="openrouter": hdrs.update({"X-Title":"Fusion.AI","HTTP-Referer":"https://fusionai.space"})
            body={"messages":[{"role":"system","content":synth_sys},{"role":"user","content":prompt}],"max_tokens":2000,"stream":False}
            if prov!="cloudflare": body["model"]=m["model"]
            ep=cf_ep(m["model"]) if prov=="cloudflare" else ENDPOINTS.get(prov,"")
            r=req.post(ep,headers=hdrs,json=body,timeout=45)
            if r.ok:
                rj=r.json()
                synthesis=rj.get("choices",[{}])[0].get("message",{}).get("content","") or (rj.get("result",{})or{}).get("response","")
                synthesis=_re4.sub(r"<think>[\s\S]*?</think>","",synthesis).strip()
        except: pass

    elapsed=round(_te.time()-t0,1)
    return J({"ok":True,"synthesis":synthesis,"perspectives":final_thoughts,
              "web_searches":len([r for r in raw_searches if r]),
              "models_used":len(use_models),"rounds":5,
              "steps":steps,"total_s":elapsed,"prompt":prompt})

@app.post("/api/gen_pptx")
async def api_gen_pptx(request: Request):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from lxml import etree
        import re as _re3

        data   = await request.json()
        md     = data.get("content", "")
        theme  = data.get("theme", "dark")
        fname  = data.get("filename", "slides.pptx")

        # ── Themes ────────────────────────────────────────────────────────
        THEMES = {
            "dark":      dict(bg1=(0x04,0x06,0x14), bg2=(0x06,0x0c,0x24),
                              a1=(0x1a,0x6e,0xf5), a2=(0x7c,0x3a,0xed),
                              tx=(0xff,0xff,0xff), tx2=(0xb0,0xcc,0xee), tx3=(0x60,0x90,0xb8)),
            "light":     dict(bg1=(0xf5,0xf7,0xff), bg2=(0xe8,0xed,0xff),
                              a1=(0x1a,0x6e,0xf5), a2=(0x7c,0x3a,0xed),
                              tx=(0x08,0x0c,0x28), tx2=(0x28,0x40,0x80), tx3=(0x60,0x80,0xb8)),
            "ocean":     dict(bg1=(0x02,0x0e,0x1c), bg2=(0x04,0x18,0x2e),
                              a1=(0x00,0xb4,0xd8), a2=(0x00,0x77,0xb6),
                              tx=(0xe8,0xf4,0xff), tx2=(0x88,0xcc,0xee), tx3=(0x40,0x88,0xaa)),
            "forest":    dict(bg1=(0x04,0x12,0x06), bg2=(0x06,0x1e,0x08),
                              a1=(0x2d,0xa4,0x4e), a2=(0x14,0x74,0x2e),
                              tx=(0xe8,0xff,0xec), tx2=(0x7c,0xcc,0x90), tx3=(0x3a,0x88,0x4e)),
            "sunset":    dict(bg1=(0x16,0x05,0x03), bg2=(0x26,0x08,0x00),
                              a1=(0xf5,0x6e,0x1a), a2=(0xed,0x3a,0x7c),
                              tx=(0xff,0xf0,0xe8), tx2=(0xee,0xcc,0xaa), tx3=(0xb0,0x78,0x58)),
            "corporate": dict(bg1=(0x04,0x04,0x0e), bg2=(0x08,0x08,0x1c),
                              a1=(0x2a,0x7a,0xff), a2=(0x5a,0x2a,0xd0),
                              tx=(0xff,0xff,0xff), tx2=(0xcc,0xdd,0xff), tx3=(0x80,0xa0,0xcc)),
        }
        C  = THEMES.get(theme, THEMES["dark"])
        def rgb(*t): return RGBColor(*t)
        BG1=rgb(*C["bg1"]); BG2=rgb(*C["bg2"])
        A1=rgb(*C["a1"]);   A2=rgb(*C["a2"])
        TX=rgb(*C["tx"]);   TX2=rgb(*C["tx2"]); TX3=rgb(*C["tx3"])

        # ── Parse slides ──────────────────────────────────────────────────
        slides_raw = [s.strip() for s in _re3.split(r'(?m)\n---+\n|^---+$', md) if s.strip()]
        if not slides_raw:
            slides_raw = [md.strip() or "Slide 1"]
        total = len(slides_raw)

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]   # completely blank

        # ── Shared helpers ────────────────────────────────────────────────
        def _grad(fill, c1, c2, angle=270):
            fill.gradient()
            fill.gradient_stops[0].position = 0.0; fill.gradient_stops[0].color.rgb = c1
            fill.gradient_stops[1].position = 1.0; fill.gradient_stops[1].color.rgb = c2
            fill.gradient_angle = angle

        def _bg(slide):
            _grad(slide.background.fill, BG1, BG2, 225)

        def _rect(slide, x, y, w_in, h_emu, *, solid=None, g1=None, g2=None, angle=90):
            s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w_in), Emu(h_emu))
            if solid: s.fill.solid(); s.fill.fore_color.rgb = solid
            else:     _grad(s.fill, g1 or A1, g2 or A2, angle)
            s.line.fill.background()
            return s

        def _oval(slide, x, y, sz, col, alpha):
            s = slide.shapes.add_shape(9, Inches(x), Inches(y), Inches(sz), Inches(sz))
            s.fill.solid(); s.fill.fore_color.rgb = col
            s.fill.fore_color.transparency = 1.0 - alpha
            s.line.fill.background()
            return s

        def _tb(slide, x, y, w, h, txt, sz, *, bold=False, italic=False,
                col=None, align=PP_ALIGN.LEFT):
            bx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
            tf = bx.text_frame; tf.word_wrap = True
            par = tf.paragraphs[0]; par.alignment = align
            run = par.add_run(); run.text = txt
            run.font.size = Pt(sz); run.font.bold = bold; run.font.italic = italic
            run.font.color.rgb = col or TX
            return bx

        def _footer(slide, num, tot):
            # dark strip at bottom
            strip = _rect(slide, 0, 7.1, 13.33, 55000,
                          solid=rgb(*[max(0, v-10) for v in C["bg1"]]))
            # slide number (right)
            nb = slide.shapes.add_textbox(Inches(11.8), Inches(7.15), Inches(1.4), Inches(0.28))
            nf = nb.text_frame; np = nf.paragraphs[0]; np.alignment = PP_ALIGN.RIGHT
            nr = np.add_run(); nr.text = f"{num} / {tot}"
            nr.font.size = Pt(9); nr.font.color.rgb = TX3
            # logo (left)  — plain ASCII to avoid emoji issues
            lb = slide.shapes.add_textbox(Inches(0.4), Inches(7.15), Inches(3), Inches(0.28))
            lf = lb.text_frame; lp = lf.paragraphs[0]
            lr = lp.add_run(); lr.text = "Fusion.AI"
            lr.font.size = Pt(9); lr.font.color.rgb = TX3

        def _parse(md_txt):
            title = ""; body = []
            for ln in md_txt.split("\n"):
                s = ln.strip()
                if not s: continue
                if s.startswith("# ") and not title:   title = s[2:].strip()
                elif s.startswith("## ") and not title: title = s[3:].strip()
                else: body.append(s)
            return title, body

        # ── Animation (single global counter per presentation) ─────────────
        _id = [200]
        def nid(): _id[0] += 1; return _id[0]

        def _anim(slide, pairs):
            """pairs = list of (shape, delay_ms). Single timing block per slide."""
            if not pairs: return
            seq_xml = ""
            bld_xml = ""
            for gi, (shp, dms) in enumerate(pairs):
                sid = shp.shape_id
                i1=nid(); i2=nid(); i3=nid(); i4=nid(); i5=nid()
                nt = "clickEffect" if gi == 0 else "afterEffect"
                seq_xml += (
                    f'<p:par><p:cTn id="{i1}" fill="hold">'
                    f'<p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>'
                    f'<p:childTnLst><p:par><p:cTn id="{i2}" fill="hold">'
                    f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
                    f'<p:childTnLst><p:par>'
                    f'<p:cTn id="{i3}" presetID="10" presetClass="entr" presetSubtype="0"'
                    f' fill="hold" grpId="{gi}" nodeType="{nt}">'
                    f'<p:stCondLst><p:cond delay="{dms}"/></p:stCondLst>'
                    f'<p:childTnLst>'
                    f'<p:set><p:cBhvr>'
                    f'<p:cTn id="{i4}" dur="1"/>'
                    f'<p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl>'
                    f'<p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>'
                    f'</p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set>'
                    f'<p:animEffect transition="in" filter="fade">'
                    f'<p:cBhvr><p:cTn id="{i5}" dur="500"/>'
                    f'<p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl>'
                    f'</p:cBhvr></p:animEffect>'
                    f'</p:childTnLst></p:cTn>'
                    f'</p:par></p:childTnLst></p:cTn></p:par>'
                    f'</p:childTnLst></p:cTn></p:par>'
                )
                bld_xml += f'<p:bldP spid="{sid}" grpId="{gi}" uiExpand="1" build="p"/>'

            r1 = nid(); r2 = nid()
            timing = (
                f'<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
                f'<p:tnLst><p:par>'
                f'<p:cTn id="{r1}" dur="indefinite" restart="whenNotActive" nodeType="tmRoot">'
                f'<p:childTnLst><p:seq concurrent="1" nextAc="seek">'
                f'<p:cTn id="{r2}" dur="indefinite" nodeType="mainSeq">'
                f'<p:childTnLst>{seq_xml}</p:childTnLst>'
                f'</p:cTn></p:seq></p:childTnLst>'
                f'</p:cTn></p:par></p:tnLst>'
                f'<p:bldLst>{bld_xml}</p:bldLst>'
                f'</p:timing>'
            )
            try:
                slide._element.append(etree.fromstring(timing))
            except Exception as xe:
                print(f"[pptx anim] {xe}", flush=True)

        def _trans(slide, fx="fade"):
            ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
            xml = f'<p:transition xmlns:p="{ns}" spd="med"><p:{fx}/></p:transition>'
            slide._element.append(etree.fromstring(xml))

        # ── Render slides ─────────────────────────────────────────────────
        for si, slide_md in enumerate(slides_raw):
            slide = prs.slides.add_slide(blank)
            _bg(slide)
            anims = []

            # Background orbs — decorative, no animation (keeps XML simple)
            _oval(slide, 9.0,  -1.5, 7, A1, 0.09)
            _oval(slide, -3.0,  3.5, 6, A2, 0.07)

            title_txt, body = _parse(slide_md)
            title_txt = title_txt or f"Slide {si+1}"

            # Detect layout
            is_title   = (si == 0)
            is_closing = (si == total-1 and total > 1)
            has_table  = "| " in slide_md and _re3.search(r'\|.+\|', slide_md)
            is_section = (not is_title and not is_closing
                          and len(slide_md) < 200
                          and slide_md.count("\n") < 5
                          and not body)
            layout = ("title"   if is_title   else
                      "closing" if is_closing  else
                      "section" if is_section  else
                      "data"    if has_table   else
                      "content")

            # Left accent bar (all slides)
            bar = _rect(slide, 0, 0, 0.07, int(7.5*914400), g1=A1, g2=A2, angle=90)
            anims.append((bar, 0))

            # ── TITLE layout ──────────────────────────────────────────────
            if layout == "title":
                # Gradient overlay (bottom half)
                ov = slide.shapes.add_shape(1, Inches(0), Inches(3.6), Inches(13.33), Inches(3.9))
                _grad(ov.fill, rgb(*[max(0,v-6) for v in C["bg2"]]), BG1, 270)
                ov.line.fill.background()

                t = _tb(slide, 1.2, 1.3, 11.0, 2.0, title_txt, 48, bold=True)
                anims.append((t, 0))

                div = _rect(slide, 1.2, 3.2, 4.5, 38000, g1=A1, g2=A2, angle=0)
                anims.append((div, 160))

                if body:
                    sub = "  |  ".join(b.lstrip("-*• ") for b in body[:2] if b.lstrip("-*• "))
                    if sub:
                        s2 = _tb(slide, 1.2, 3.5, 10.0, 0.8, sub, 18, col=TX2)
                        anims.append((s2, 320))

                _trans(slide, "fade")

            # ── SECTION layout ────────────────────────────────────────────
            elif layout == "section":
                cent = _tb(slide, 0.8, 2.4, 11.7, 1.8, title_txt, 42, bold=True,
                            align=PP_ALIGN.CENTER)
                anims.append((cent, 0))
                div = _rect(slide, 4.2, 4.1, 5.0, 36000, g1=A1, g2=A2, angle=0)
                anims.append((div, 160))
                _trans(slide, "push")

            # ── CLOSING layout ────────────────────────────────────────────
            elif layout == "closing":
                ov2 = slide.shapes.add_shape(1, Inches(0), Inches(2.5), Inches(13.33), Inches(5.0))
                _grad(ov2.fill, rgb(*[max(0,v-4) for v in C["bg2"]]), BG1, 270)
                ov2.line.fill.background()

                t2 = _tb(slide, 1.2, 1.5, 11.0, 1.8, title_txt, 44, bold=True,
                          align=PP_ALIGN.CENTER)
                anims.append((t2, 0))
                div2 = _rect(slide, 3.5, 3.35, 6.5, 36000, g1=A1, g2=A2, angle=0)
                anims.append((div2, 160))
                if body:
                    sub2 = " ".join(b.lstrip("-*• ") for b in body[:2] if b.lstrip("-*• "))
                    if sub2:
                        sb = _tb(slide, 1.5, 3.75, 10.3, 0.7, sub2, 17,
                                  col=TX2, align=PP_ALIGN.CENTER)
                        anims.append((sb, 300))
                _trans(slide, "fade")

            # ── DATA (table) layout ───────────────────────────────────────
            elif layout == "data":
                ht = _tb(slide, 0.9, 0.3, 11.5, 0.9, title_txt, 28, bold=True)
                anims.append((ht, 0))
                hd = _rect(slide, 0.9, 1.15, 5.0, 28000, solid=A1)
                anims.append((hd, 80))

                tbl_rows = [ln for ln in body
                            if "|" in ln and not _re3.match(r'^\|[-\s|]+\|$', ln)]
                y = 1.42
                for ri, row in enumerate(tbl_rows[:9]):
                    cells = [c.strip() for c in row.strip("|").split("|")]
                    row_str = "    ".join(cells)
                    is_h = (ri == 0)
                    rc = rgb(*[min(255,v+14) for v in C["bg2"]]) if ri%2==0 else BG2
                    rb = _rect(slide, 0.9, y, 11.5, int(0.42*914400), solid=rc)
                    anims.append((rb, 90+ri*55))
                    rt = _tb(slide, 1.0, y+0.03, 11.3, 0.37, row_str, 13,
                              bold=is_h, col=TX if is_h else TX2)
                    anims.append((rt, 110+ri*55))
                    y += 0.42
                    if y > 6.6: break
                _trans(slide, "fade")

            # ── CONTENT (default) layout ──────────────────────────────────
            else:
                ht = _tb(slide, 0.9, 0.3, 11.5, 0.95, title_txt, 30, bold=True)
                anims.append((ht, 0))
                hd = _rect(slide, 0.9, 1.2, 6.0, 30000, g1=A1, g2=A2, angle=0)
                anims.append((hd, 80))

                y = 1.62; delay = 200
                for bi, bl in enumerate(body[:8]):
                    clean = _re3.sub(r'\*\*(.+?)\*\*', r'\1', bl)
                    clean = _re3.sub(r'\*(.+?)\*',     r'\1', clean)
                    is_bul = clean[:1] in "-*•·>"
                    txt = clean.lstrip("-*•·> ").strip()
                    if not txt: continue
                    if is_bul:
                        dot = _rect(slide, 0.9, y+0.12, 0.05, int(0.09*914400), solid=A1)
                        anims.append((dot, delay-60))
                        bt = _tb(slide, 1.08, y, 11.1, 0.55, txt, 17, col=TX2)
                    else:
                        bt = _tb(slide, 0.9, y, 11.3, 0.55, txt, 17,
                                  bold=(bi==0), col=TX2)
                    anims.append((bt, delay))
                    y += 0.58; delay += 170
                    if y > 6.7: break
                _trans(slide, "push")

            _footer(slide, si+1, total)
            _anim(slide, anims)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return StreamingResponse(
            buf,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f'attachment; filename="{fname}"'}
        )
    except Exception as ex:
        import traceback as _tb2; _tb2.print_exc()
        return JSONResponse({"error": str(ex)}, status_code=500)

@app.post("/api/gen_xlsx")
async def api_gen_xlsx(request: Request):
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        from openpyxl.utils import get_column_letter
        data=await request.json(); csv_text=data.get("content","")
        rows=[]
        for line in csv_text.strip().split("\n"):
            cells=[]; inq=False; cur=""
            for ch in line:
                if ch=='"': inq=not inq
                elif ch==',' and not inq: cells.append(cur.strip().strip('"')); cur=""
                else: cur+=ch
            cells.append(cur.strip().strip('"'))
            if any(c for c in cells): rows.append(cells)
        wb=openpyxl.Workbook(); ws=wb.active; ws.title=data.get("sheet","Data")
        HDR=PatternFill("solid",fgColor="1A6EF5"); ALT=PatternFill("solid",fgColor="0A1228"); ALT2=PatternFill("solid",fgColor="060A1A")
        HF=Font(bold=True,color="FFFFFF",name="Calibri",size=11); BF=Font(color="C8E0FF",name="Calibri",size=11)
        TH=Side(style="thin",color="1A3060"); BD=Border(left=TH,right=TH,top=TH,bottom=TH)
        for ri,row in enumerate(rows):
            for ci,val in enumerate(row):
                cell=ws.cell(row=ri+1,column=ci+1)
                try: cell.value=float(val) if val.replace(".","",1).replace("-","",1).isdigit() else val
                except: cell.value=val
                cell.border=BD; cell.alignment=Alignment(horizontal="left",vertical="center",wrap_text=True)
                if ri==0: cell.font=HF; cell.fill=HDR; cell.alignment=Alignment(horizontal="center",vertical="center")
                else: cell.font=BF; cell.fill=ALT if ri%2==0 else ALT2
        for col in ws.columns:
            mx=max((len(str(c.value or "")) for c in col),default=8)
            ws.column_dimensions[get_column_letter(col[0].column)].width=min(mx+4,45)
        ws.row_dimensions[1].height=24; ws.freeze_panes="A2"
        buf=io.BytesIO(); wb.save(buf); buf.seek(0)
        return StreamingResponse(buf,media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            headers={"Content-Disposition":f'attachment; filename="{data.get("filename","data.xlsx")}"'})
    except Exception as ex:
        return JSONResponse({"error":str(ex)},status_code=500)

@app.post("/api/gen_docx")
async def api_gen_docx(request: Request):
    try:
        from docx import Document
        from docx.shared import Pt
        data=await request.json(); md=data.get("content","")
        doc=Document()
        for line in md.split("\n"):
            raw=line.rstrip()
            if raw.startswith("# "): doc.add_heading(raw[2:].strip(),level=1)
            elif raw.startswith("## "): doc.add_heading(raw[3:].strip(),level=2)
            elif raw.startswith("### "): doc.add_heading(raw[4:].strip(),level=3)
            elif raw.startswith("- ") or raw.startswith("* "): doc.add_paragraph(raw[2:].strip(),style="List Bullet")
            elif _re.match(r'^\d+\.\s',raw): doc.add_paragraph(_re.sub(r'^\d+\.\s','',raw).strip(),style="List Number")
            elif raw=="": doc.add_paragraph("")
            else:
                p_=doc.add_paragraph()
                for part in _re.split(r'(\*\*[^*]+\*\*|\*[^*]+\*|`[^`]+`)',raw):
                    if part.startswith("**") and part.endswith("**"): r_=p_.add_run(part[2:-2]); r_.bold=True
                    elif part.startswith("*") and part.endswith("*"): r_=p_.add_run(part[1:-1]); r_.italic=True
                    elif part.startswith("`") and part.endswith("`"): r_=p_.add_run(part[1:-1]); r_.font.name="Courier New"; r_.font.size=Pt(10)
                    else: p_.add_run(part)
        buf=io.BytesIO(); doc.save(buf); buf.seek(0)
        return StreamingResponse(buf,media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={"Content-Disposition":f'attachment; filename="{data.get("filename","document.docx")}"'})
    except Exception as ex:
        return JSONResponse({"error":str(ex)},status_code=500)


# ════════════════════════════════════════════════════════════════════════════
# FusionOS — sandboxed AI desktop (Terminal, Files, Editor, Agent, Browser, Monitor)
# ════════════════════════════════════════════════════════════════════════════
import tempfile as _tf, shutil as _sh, platform as _plat, concurrent.futures as _cf2

_FOS_ROOT = os.path.join(_tf.gettempdir(), "fusionos_sandboxes")
os.makedirs(_FOS_ROOT, exist_ok=True)
_FOS_BLOCKED = ["rm -rf /","rm -rf /*",":(){ :|:& };:","mkfs","dd if=/dev/zero",
                 "dd if=/dev/random","> /dev/sda","shutdown","reboot","poweroff",
                 "chmod -R 777 /","wget http://169.254","curl http://169.254"]
_SERVER_START = _time.time()

def _fos_sandbox(uid):
    d = os.path.join(_FOS_ROOT, str(uid))
    os.makedirs(d, exist_ok=True)
    wf = os.path.join(d, "README.txt")
    if not os.path.exists(wf):
        open(wf,"w").write(
            "Welcome to FusionOS!\n\n"
            "This is your personal sandboxed Linux environment.\n\n"
            "Available: python3, pip, node, npm, curl, git, gcc, g++, bash\n\n"
            "Try in Terminal:\n"
            "  python3 -c \"print(2**32)\"\n"
            "  curl -s https://api.github.com/zen\n"
            "  node -e \"console.log('hello')\"\n"
            "  chromium https://example.com   (real headless-Chromium render)\n\n"
            "Or use the AI Agent — give it a goal in plain English!\n"
        )
    return d

def _fos_safe_path(sandbox, rel):
    rel = (rel or "").strip().lstrip("/")
    t = os.path.normpath(os.path.join(sandbox, rel))
    if not (t == sandbox or t.startswith(sandbox + os.sep)): t = sandbox
    return t

def _fos_blocked(cmd):
    low = cmd.lower()
    for b in _FOS_BLOCKED:
        if b in low: return b
    return None

def _fos_run(cmd, cwd, timeout=30):
    env = os.environ.copy(); env["HOME"] = cwd; env["TERM"] = "xterm-256color"
    try:
        r = subprocess.run(cmd, shell=True, cwd=cwd, capture_output=True,
                            text=True, timeout=timeout, env=env)
        return r.stdout[-6000:], r.stderr[-3000:], r.returncode
    except subprocess.TimeoutExpired:
        return "", f"⏱ Timed out ({timeout}s limit — try a shorter command)", 124
    except Exception as ex:
        return "", str(ex), 1

# ── Real headless Chromium — used by the FusionOS Browser app ────────────────
_CHROMIUM_BIN_CACHE = {"path": None, "checked": False}
def _find_chromium():
    """Locate a real Chromium/Chrome binary on the host, cached after first check."""
    if _CHROMIUM_BIN_CACHE["checked"]: return _CHROMIUM_BIN_CACHE["path"]
    for name in ("chromium","chromium-browser","google-chrome","google-chrome-stable","chrome"):
        p = _sh.which(name)
        if p: _CHROMIUM_BIN_CACHE.update(path=p, checked=True); return p
    _CHROMIUM_BIN_CACHE["checked"] = True
    return None

def _run_chromium(url, mode="screenshot", timeout=25):
    """Run real headless Chromium against a URL. mode: screenshot | pdf | dom | text.
    Returns (ok, data_or_error, content_type)."""
    bin_path = _find_chromium()
    if not bin_path:
        return False, ("No real Chromium binary found on this host. Install it with "
                        "'apt-get install -y chromium' (Debian/Ubuntu) or add it to your "
                        "Dockerfile, then restart the server."), "text/plain"
    tmp_dir = _tf.mkdtemp(prefix="fos_chromium_")
    try:
        base_flags = [bin_path,"--headless=new","--disable-gpu","--no-sandbox",
                      "--disable-dev-shm-usage","--hide-scrollbars","--window-size=1280,1600",
                      "--virtual-time-budget=6000","--run-all-compositor-stages-before-draw"]
        if mode == "screenshot":
            out_path = os.path.join(tmp_dir,"shot.png")
            cmd = base_flags + [f"--screenshot={out_path}", url]
            r = subprocess.run(cmd, capture_output=True, text=True, timeout=timeout)
            if r.returncode != 0 or not os.path.exists(out_path):
                return False, (r.stderr or "Chromium failed to render the page")[-1500:], "text/plain"
            with open(out_path,"rb") as f: data = f.read()
            return True, base64.b64encode(data).decode(), "image/png"
        elif mode == "pdf":
            out_path = os.path.join(tmp_dir,"page.pdf")
            cmd = base_flags + [f"--print-to-pdf={out_path}", url]
            r = subprocess.run(cmd, capture_output=True, text=True, timeout=timeout)
            if r.returncode != 0 or not os.path.exists(out_path):
                return False, (r.stderr or "Chromium failed to print the page")[-1500:], "text/plain"
            with open(out_path,"rb") as f: data = f.read()
            return True, base64.b64encode(data).decode(), "application/pdf"
        elif mode == "dom":
            cmd = base_flags + ["--dump-dom", url]
            r = subprocess.run(cmd, capture_output=True, text=True, timeout=timeout)
            if r.returncode != 0:
                return False, (r.stderr or "Chromium failed to load the page")[-1500:], "text/plain"
            return True, r.stdout[-500000:], "text/html"
        else:  # "text" — DOM stripped of tags, cheap readable-page mode
            cmd = base_flags + ["--dump-dom", url]
            r = subprocess.run(cmd, capture_output=True, text=True, timeout=timeout)
            if r.returncode != 0:
                return False, (r.stderr or "Chromium failed to load the page")[-1500:], "text/plain"
            html = r.stdout
            txt = _re.sub(r"<script[\s\S]*?</script>","",html,flags=_re.I)
            txt = _re.sub(r"<style[\s\S]*?</style>","",txt,flags=_re.I)
            txt = _re.sub(r"<[^>]+>"," ",txt)
            txt = _re.sub(r"\s+"," ",txt).strip()
            return True, txt[:20000], "text/plain"
    except subprocess.TimeoutExpired:
        return False, f"Chromium timed out after {timeout}s — the page may be too heavy or slow.", "text/plain"
    except Exception as ex:
        return False, str(ex), "text/plain"
    finally:
        try: _sh.rmtree(tmp_dir, ignore_errors=True)
        except: pass

@app.post("/api/fos/chromium")
async def api_fos_chromium(request: Request):
    """Real headless Chromium — screenshot, PDF, DOM, or extracted text of any URL."""
    auth_user(request); d = await request.json()
    url = (d.get("url") or "").strip()
    mode = (d.get("mode") or "screenshot").strip().lower()
    if not url: return err("url required")
    if not (url.startswith("http://") or url.startswith("https://")): url = "https://" + url
    ok, data, ctype = _run_chromium(url, mode=mode)
    return J({"ok": ok, "mode": mode, "content_type": ctype, "data": data, "url": url})

@app.get("/api/fos/chromium/status")
async def api_fos_chromium_status(request: Request):
    """Report whether a real Chromium binary is available on this host."""
    auth_user(request)
    p = _find_chromium()
    return J({"available": bool(p), "path": p or ""})

@app.post("/api/fos/exec")
async def api_fos_exec(request: Request):
    u = auth_user(request); d = await request.json()
    cmd = (d.get("cmd") or "").strip()
    cwd_rel = (d.get("cwd") or "").strip().lstrip("/")
    sandbox = _fos_sandbox(u["id"])
    if not cmd: return J({"stdout":"","stderr":"","code":0,"cwd":cwd_rel})
    blocked = _fos_blocked(cmd)
    if blocked: return J({"stdout":"","stderr":f"⛔ Blocked: '{blocked}'","code":1,"cwd":cwd_rel})
    workdir = _fos_safe_path(sandbox, cwd_rel); os.makedirs(workdir, exist_ok=True)
    low = cmd.strip().lower()
    if low in ("clear","cls"): return J({"stdout":"__CLEAR__","stderr":"","code":0,"cwd":cwd_rel})
    if low == "cd" or low.startswith("cd "):
        target = cmd[2:].strip() or "~"
        if target in ("~","","/"): return J({"stdout":"","stderr":"","code":0,"cwd":""})
        newpath = os.path.normpath(os.path.join(workdir, target))
        if not (newpath == sandbox or newpath.startswith(sandbox+os.sep)): newpath = sandbox
        if not os.path.isdir(newpath): return J({"stdout":"","stderr":f"cd: {target}: No such directory","code":1,"cwd":cwd_rel})
        newrel = os.path.relpath(newpath, sandbox)
        return J({"stdout":"","stderr":"","code":0,"cwd":"" if newrel=="." else newrel})
    if low.startswith("chromium ") or low == "chromium":
        arg = cmd[len("chromium"):].strip()
        if not arg:
            return J({"stdout":"Usage: chromium <url>  — renders a real headless-Chromium screenshot.\nTip: open the FusionOS Browser app and click '🖥 Chromium' for an inline view.","stderr":"","code":0,"cwd":cwd_rel})
        url = arg if arg.startswith("http") else "https://"+arg
        bin_path = _find_chromium()
        if not bin_path:
            return J({"stdout":"","stderr":("No real Chromium binary found on this host.\nInstall with: apt-get install -y chromium   (Debian/Ubuntu)\nor add 'chromium' to your Dockerfile, then restart the server."),"code":1,"cwd":cwd_rel})
        ok, data, ctype = _run_chromium(url, mode="text")
        if not ok: return J({"stdout":"","stderr":str(data),"code":1,"cwd":cwd_rel})
        return J({"stdout":f"✅ Rendered {url} with real Chromium ({bin_path}):\n\n{data}","stderr":"","code":0,"cwd":cwd_rel})
    stdout, stderr, code = _fos_run(cmd, workdir)
    return J({"stdout": stdout, "stderr": stderr, "code": code, "cwd": cwd_rel})

@app.post("/api/fos/files")
async def api_fos_files(request: Request):
    u = auth_user(request); d = await request.json()
    action = d.get("action","list"); rel = d.get("path","")
    sandbox = _fos_sandbox(u["id"]); target = _fos_safe_path(sandbox, rel)
    try:
        if action == "list":
            if not os.path.isdir(target): target = sandbox; rel = ""
            items = []
            for name in sorted(os.listdir(target), key=lambda n:(not os.path.isdir(os.path.join(target,n)),n.lower())):
                fp = os.path.join(target, name); is_dir = os.path.isdir(fp)
                items.append({"name":name,"is_dir":is_dir,"size":0 if is_dir else os.path.getsize(fp),"mtime":int(os.path.getmtime(fp))})
            return J({"items":items,"path":rel.strip("/")})
        elif action == "read":
            if not os.path.isfile(target): return J({"error":"File not found"})
            if os.path.getsize(target) > 2_000_000: return J({"error":"File too large (>2MB)"})
            return J({"content": open(target,"r",errors="replace").read(), "path": rel})
        elif action == "write":
            os.makedirs(os.path.dirname(target), exist_ok=True)
            open(target,"w").write(d.get("content",""))
            return J({"ok":True})
        elif action in ("mkdir","new_file"):
            if action=="mkdir": os.makedirs(target, exist_ok=True)
            else:
                os.makedirs(os.path.dirname(target) or sandbox, exist_ok=True)
                if not os.path.exists(target): open(target,"w").write("")
            return J({"ok":True})
        elif action == "delete":
            if os.path.isdir(target): _sh.rmtree(target, ignore_errors=True)
            elif os.path.isfile(target): os.remove(target)
            return J({"ok":True})
        elif action == "rename":
            new_t = _fos_safe_path(sandbox, d.get("new_path",""))
            os.makedirs(os.path.dirname(new_t) or sandbox, exist_ok=True)
            os.rename(target, new_t); return J({"ok":True})
    except Exception as ex: return J({"error": str(ex)})
    return J({"error":"Unknown action"})

@app.post("/api/fos/agent")
async def api_fos_agent(request: Request):
    u = auth_user(request); d = await request.json()
    goal = (d.get("goal") or "").strip()
    transcript = d.get("transcript", []) or []
    step_num = int(d.get("step", 1))
    speed = d.get("speed", "balanced")
    max_steps = max(3, min(20, int(d.get("max_steps", 10))))
    if not goal: return J({"action":"fail","reason":"No goal provided"})
    if step_num > max_steps:
        return J({"action":"done","summary":f"Reached the step limit ({max_steps}). Check the output above."})
    SCFG = {"fast":{"max_tokens":500,"timeout":20,"cheap":True},
             "balanced":{"max_tokens":900,"timeout":35,"cheap":False},
             "thorough":{"max_tokens":1400,"timeout":50,"cheap":False}}.get(speed,{"max_tokens":900,"timeout":35,"cheap":False})
    sandbox = _fos_sandbox(u["id"])
    hist_txt = ""
    for i, t in enumerate(transcript[-6:], 1):
        hist_txt += f"\nStep {i}: ran `{t.get('cmd','')}`\n  stdout: {(t.get('stdout','') or '')[:350]}\n  exit: {t.get('code')}\n"
    sys_p = f"""You are the FusionOS AI Agent. You run real shell commands in a sandboxed Linux VM.
Available: python3, pip, node, npm, curl, git, gcc, g++, bash, and standard unix tools. No sudo.
CWD: {sandbox}

Respond with ONLY a single JSON object — no markdown, no explanation outside it:
{{"action":"run","cmd":"shell command here","explain":"one sentence why"}}
or when done: {{"action":"done","summary":"what was accomplished"}}
or if stuck: {{"action":"fail","reason":"why"}}

CRITICAL: Escape ALL newlines as \\n inside the "cmd" string — never put raw line breaks in JSON.
To write a file: cmd = "printf 'line1\\nline2\\n' > file.py"  (or use printf, not heredoc)
Verify work — run scripts after writing them. Finish in 2-8 steps."""
    user_p = f"GOAL: {goal}\n\nTRANSCRIPT:{hist_txt or ' (first step)'}\n\nStep {step_num}:"
    raw = _call_overseer(sys_p, user_p, max_tokens=SCFG["max_tokens"], timeout=SCFG["timeout"], cheap=SCFG["cheap"])
    if not raw: return J({"action":"fail","reason":f"AI model timed out or is unavailable. Try 'Fast' mode or try again."})
    clean = _re.sub(r'```\w*','',raw).strip().strip('`').strip()
    m = _re.search(r'\{[\s\S]*\}', clean); clean = m.group(0) if m else clean
    plan = None
    try: plan = json.loads(clean, strict=False)
    except:
        am = _re.search(r'"action"\s*:\s*"(run|done|fail)"', clean)
        if am:
            act2 = am.group(1)
            if act2 == "run":
                cm = _re.search(r'"cmd"\s*:\s*"((?:[^"\\]|\\.)*)"', clean)
                em = _re.search(r'"explain"\s*:\s*"((?:[^"\\]|\\.)*)"', clean)
                if cm: plan = {"action":"run","cmd":cm.group(1).replace('\\n','\n'),"explain": em.group(1) if em else ""}
            elif act2 == "done":
                sm = _re.search(r'"summary"\s*:\s*"((?:[^"\\]|\\.)*)"', clean)
                plan = {"action":"done","summary": sm.group(1) if sm else "Task completed."}
            else:
                rm = _re.search(r'"reason"\s*:\s*"((?:[^"\\]|\\.)*)"', clean)
                plan = {"action":"fail","reason": rm.group(1) if rm else "Agent could not continue."}
    if plan is None: return J({"action":"fail","reason":"Malformed response from AI — try again or switch speed mode."})
    act = plan.get("action")
    if act == "run":
        cmd = (plan.get("cmd") or "").strip()
        if not cmd: return J({"action":"fail","reason":"Empty command"})
        b = _fos_blocked(cmd)
        if b: return J({"action":"fail","reason":f"Blocked: '{b}'"})
        stdout, stderr, code = _fos_run(cmd, sandbox, timeout=30)
        return J({"action":"run","cmd":cmd,"explain":plan.get("explain",""),"stdout":stdout,"stderr":stderr,"code":code})
    elif act == "done": return J({"action":"done","summary":plan.get("summary","Task completed.")})
    else: return J({"action":"fail","reason":plan.get("reason","Agent stopped.")})

@app.post("/api/fos/stats")
async def api_fos_stats(request: Request):
    u = auth_user(request); sandbox = _fos_sandbox(u["id"])
    total_size=0; fc=0; dc=0
    for root,dirs,files in os.walk(sandbox):
        dc+=len(dirs)
        for f in files:
            try: total_size+=os.path.getsize(os.path.join(root,f)); fc+=1
            except: pass
    du = _sh.disk_usage(sandbox)
    return J({"files":fc,"dirs":dc,"size_kb":round(total_size/1024,1),
               "python":_plat.python_version(),"platform":_plat.platform(),
               "uptime_s":round(_time.time()-_SERVER_START,1),
               "disk_total_gb":round(du.total/1e9,1),"disk_used_gb":round(du.used/1e9,1),"disk_free_gb":round(du.free/1e9,1),
               "groq":bool(GROQ_KEY.strip()),"openrouter":bool(OPENROUTER_KEY.strip()),
               "chromium":bool(_find_chromium())})

# ── AI Computer (Perplexity-style) with parallel page fetch ──────────────────
@app.post("/api/computer")
async def api_computer(request: Request):
    import urllib.parse as _up, re as _re5, time as _t5
    auth_user(request); d=await request.json()
    query=d.get("query","").strip(); depth=d.get("depth","normal")
    num_src=int(d.get("num_sources",6)); t0=_t5.time(); sources=[]
    try:
        br=req.get(f"http://openserp.alwaysdata.net/bing/search?text={_up.quote(query)}",
                    headers={"User-Agent":"FusionAI/2.0","Accept":"application/json"},timeout=9)
        if br.ok:
            items=br.json() if isinstance(br.json(),list) else []
            for item in items[:num_src+3]:
                t2=item.get("title",""); ru=item.get("url",""); sn=item.get("description","") or item.get("snippet","")
                if t2 and ru and ru.startswith("http") and not any(s["url"]==ru for s in sources):
                    sources.append({"title":t2,"url":ru,"snippet":sn[:400],"full":sn})
    except: pass
    fetch_n={"fast":2,"normal":3,"deep":5}.get(depth,3); ptmo=6
    skip=["youtube.com","reddit.com","facebook.com","twitter.com","instagram.com","tiktok.com"]
    cands=[s for s in sources if not any(sd in s["url"] for sd in skip)][:fetch_n]
    def _fp(src):
        try:
            pr=req.get(src["url"],headers={"User-Agent":"Mozilla/5.0 (Windows NT 10.0)","Accept":"text/html"},timeout=ptmo,allow_redirects=True)
            if pr.ok and "text/html" in pr.headers.get("Content-Type",""):
                txt=_re5.sub(r'<script[\s\S]*?</script>','',pr.text,flags=_re5.IGNORECASE)
                txt=_re5.sub(r'<style[\s\S]*?</style>','',txt,flags=_re5.IGNORECASE)
                txt=_re5.sub(r'<[^>]+>','',txt); txt=_re5.sub(r'\s+',' ',txt).strip()
                if len(txt)>200: return src["url"],txt[:3000]
        except: pass
        return src["url"],None
    if cands:
        with _cf2.ThreadPoolExecutor(max_workers=min(5,len(cands))) as pool:
            futs=[pool.submit(_fp,s) for s in cands]
            for fut in _cf2.as_completed(futs,timeout=ptmo+2):
                try:
                    url,txt=fut.result()
                    if txt:
                        for s in sources:
                            if s["url"]==url: s["full"]=txt; break
                except: pass
    sources=sources[:num_src]
    ctx="\n\n---\n\n".join(f"[{i}] {s['title']} ({s['url']})\n{s.get('full',s.get('snippet',''))[:1200]}" for i,s in enumerate(sources,1))
    sys_p=f"You are an AI research assistant. Today is {datetime.utcnow().strftime('%Y-%m-%d')} UTC.\nAnswer the question using the web sources. Use markdown. Cite sources inline as [1],[2] etc. Be thorough but concise."
    ans=_call_overseer(sys_p,f"Question: {query}\n\nSources:\n{ctx}\n\nAnswer:",max_tokens=1400,timeout=40) or ""
    if not ans and sources: ans="\n\n".join(f"**[{i}] {s['title']}**\n{s.get('snippet','')}" for i,s in enumerate(sources,1))
    related=[]
    try:
        rr=_call_overseer("Return ONLY a JSON array of 3 short follow-up question strings. No markdown.",f"Query: {query}",max_tokens=120,cheap=True,timeout=15)
        if rr:
            import json as _jj; clean2=_re.sub(r'```\w*','',rr).strip().strip('`')
            m2=_re.search(r'\[[\s\S]*\]',clean2)
            if m2: related=_jj.loads(m2.group(0))
    except: pass
    return J({"answer":ans,"sources":sources,"related":related,"elapsed_ms":int((_t5.time()-t0)*1000)})

# ── Free SVG image gen ───────────────────────────────────────────────────────
@app.post("/api/gen_svg_img")
async def api_gen_svg_img(request: Request):
    auth_user(request); d=await request.json()
    prompt=d.get("prompt","").strip(); style=d.get("style","detailed illustration")
    w=int(d.get("width",800)); h=int(d.get("height",600))
    if not prompt: return J({"error":"prompt required"})
    sys_p=f"""You are an expert SVG artist. Output ONLY raw SVG starting with <svg and ending with </svg>.
Use viewBox="0 0 {w} {h}" width="{w}" height="{h}". Create rich {style} art with gradients, shapes and vivid color. No text unless asked."""
    raw=_call_overseer(sys_p,f"Create: {prompt}",max_tokens=4096,timeout=45)
    if not raw: return J({"error":"LLM unavailable"})
    m=_re.search(r'(<svg[\s\S]*?</svg>)',raw,_re.IGNORECASE)
    svg=m.group(1) if m else raw
    if 'viewBox' not in svg: svg=svg.replace('<svg',f'<svg viewBox="0 0 {w} {h}"',1)
    return J({"svg":svg})


# ── FusionOS browser proxy — strips X-Frame-Options so pages load in iframe ──
@app.get("/api/fos/proxy")
async def api_fos_proxy(request: Request, url: str = ""):
    auth_user(request)
    if not url: return HTMLResponse("<p>No URL provided</p>", status_code=400)
    if not url.startswith("http"): url = "https://" + url
    from urllib.parse import urlparse; import re as _rfp
    host = urlparse(url).hostname or ""
    if any(b in host for b in ["localhost","127.","0.0.0.0","169.254","::1"]):
        return HTMLResponse("<p>Blocked</p>", status_code=403)
    try:
        r = req.get(url, headers={
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 Chrome/124",
            "Accept": "text/html,application/xhtml+xml,*/*",
            "Accept-Language": "en-US,en;q=0.9",
        }, timeout=12, allow_redirects=True)
        ct = r.headers.get("Content-Type","text/html")
        if "text/html" not in ct:
            return HTMLResponse(f"<p style='font-family:sans-serif;padding:20px'>Cannot display non-HTML content ({ct}).<br><br><a href='{url}' target='_blank'>Open in new tab ↗</a></p>")
        html = r.text
        base = f'<base href="{url}" target="_blank">'
        html = _rfp.sub(r'(?i)<head>', '<head>' + base, html, count=1)
        if '<base' not in html: html = base + html
        # Strip frame-buster scripts
        html = _rfp.sub(r'(?is)<script[^>]*>.*?</script>', '', html)
        return HTMLResponse(html, headers={
            "Content-Type": "text/html; charset=utf-8",
            "X-Frame-Options": "ALLOWALL",
            "Content-Security-Policy": "default-src * 'unsafe-inline' 'unsafe-eval' data: blob:",
        })
    except Exception as ex:
        return HTMLResponse(f"<div style='font-family:sans-serif;padding:24px;color:#e88;background:#0a0c14'><b>Failed to load:</b><br>{url}<br><br>{ex}<br><br><a href='{url}' target='_blank' style='color:#4a9eff'>Open in new tab ↗</a></div>")


@app.get("/")
async def index(): return HTMLResponse(_HTML.replace("__HSITE__",HCAPTCHA_SITE_KEY))

@app.exception_handler(404)
async def not_found(r,e): return JSONResponse({"error":"Not found"},404)
@app.exception_handler(405)
async def method_not_allowed(r,e): return JSONResponse({"error":"Method not allowed"},405)
@app.exception_handler(Exception)
async def server_error(r,e):
    import traceback; traceback.print_exc()
    return JSONResponse({"error":"Internal server error","dev_error":str(e)},500)

def _startup_banner(port):
    sep = "="*60
    print(f"\n{sep}",flush=True)
    print("⚡  Fusion.AI v14 — FastAPI Edition",flush=True)
    print(sep,flush=True)
    print(f"📦  Database      : {DB}",flush=True)
    print(f"🔑  GROQ_KEY      : {'✅' if GROQ_KEY else '❌'}",flush=True)
    print(f"🔑  OPENROUTER    : {'✅' if OPENROUTER_KEY else '❌'}",flush=True)
    print(f"☁️   CF Workers AI : {'✅' if CF_ACCOUNT_ID and CF_KEY else '❌'}",flush=True)
    print(f"☁️   CF Account 2  : {'✅' if CF_ACCOUNT_ID2 and CF_KEY2 else 'not set'}",flush=True)
    print(f"💻  GitHub AI     : {'✅' if GITHUB_TOKEN else '❌'}",flush=True)
    print(f"🔒  SECRET_KEY    : {'✅' if _raw_secret else '⚠️  not set'}",flush=True)
    print(f"🔐  HCAPTCHA      : {'✅' if HCAPTCHA_SECRET else '⚠️  not set'}",flush=True)
    print(f"🌐  GOOGLE_OAUTH  : {'✅' if GOOGLE_CLIENT_ID else '❌'}",flush=True)
    print(f"👤  Dev user      : {TP_USERNAME or 'not set'}",flush=True)
    print(f"🤖  Models        : {len(MODELS)}",flush=True)
    print(f"👉  Open: http://localhost:{port}",flush=True)
    print(f"{sep}\n",flush=True)

if __name__ == "__main__":
    import uvicorn
    port = int(os.environ.get("PORT", 7860))
    _startup_banner(port)
    uvicorn.run(app, host="0.0.0.0", port=port, log_level="warning")
